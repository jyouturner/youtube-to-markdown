#!/usr/bin/env python3
"""
youtube-to-markdown — Turn a YouTube video (or local MP4 + SRT) into a readable
Markdown digest with embedded frame images. Optional PowerPoint export.

Default flow (digest):
  yt2md "https://youtu.be/..."
  yt2md input.mp4 transcript.srt

Also build a deck:
  yt2md "https://youtu.be/..." --deck

Just the deck (no API key needed):
  yt2md input.mp4 transcript.srt --deck-only

Requirements:
  System:  ffmpeg, ffprobe
  API:     ANTHROPIC_API_KEY (prompted on first run unless --deck-only)
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple


# ---------- Step 1: Frame extraction (scene detection + periodic sampling) ----------

# Hard cap on scene frames a single video can yield. Slide-heavy talks
# with subtle reveal animations can blow past 500-1000 candidates at the
# default threshold; beyond a few hundred there's negligible additional
# visual signal but real hashing/IO overhead. Truncating keeps the dedupe
# and vision-pick stages bounded.
SCENE_FRAME_HARD_CAP = 500


def extract_scene_frames(video: Path, out_dir: Path, threshold: float = 0.2) -> List[Tuple[Path, float]]:
    """Run ffmpeg scene detection. Returns list of (frame_path, timestamp_seconds)."""
    out_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"select='eq(n,0)+gt(scene,{threshold})',showinfo",
        "-vsync", "vfr",
        "-q:v", "3",
        str(out_dir / "scene_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg failed:\n{proc.stderr}")

    timestamps = [float(m) for m in re.findall(r"pts_time:([\d.]+)", proc.stderr)]
    frame_files = sorted(out_dir.glob("scene_*.jpg"))
    if len(frame_files) != len(timestamps):
        n = min(len(frame_files), len(timestamps))
        frame_files, timestamps = frame_files[:n], timestamps[:n]

    if len(frame_files) > SCENE_FRAME_HARD_CAP:
        # Evenly-spaced subsample so we keep visual diversity across the
        # whole video, not just the first chunk.
        step = len(frame_files) / SCENE_FRAME_HARD_CAP
        keep = [int(i * step) for i in range(SCENE_FRAME_HARD_CAP)]
        frame_files = [frame_files[i] for i in keep]
        timestamps = [timestamps[i] for i in keep]

    return list(zip(frame_files, timestamps))


def extract_interval_frames(video: Path, out_dir: Path, interval: float, duration: float) -> List[Tuple[Path, float]]:
    """Sample one frame every `interval` seconds. Returns (frame_path, timestamp) list."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if interval <= 0:
        return []

    fps = 1.0 / interval
    cmd = [
        "ffmpeg", "-hide_banner", "-y",
        "-i", str(video),
        "-vf", f"fps={fps}",
        "-q:v", "3",
        str(out_dir / "interval_%04d.jpg"),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg interval sampling failed:\n{proc.stderr}")

    frame_files = sorted(out_dir.glob("interval_*.jpg"))
    # ffmpeg's fps filter places the first frame at t≈interval/2; reconstruct timestamps.
    timestamps = [min(duration, interval * (i + 0.5)) for i in range(len(frame_files))]
    return list(zip(frame_files, timestamps))


def extract_scene_and_interval_frames(
    video: Path,
    scene_dir: Path,
    interval_dir: Path,
    scene_threshold: float,
    interval: float,
    duration: float,
) -> Tuple[List[Tuple[Path, float]], List[Tuple[Path, float]]]:
    """Single-pass replacement for the parallel scene+interval pair. Uses
    ffmpeg's `filter_complex` with `split` to fan the decoded video into
    two filter chains and produce both frame sets from ONE decode of the
    source MP4.

    Why this matters: the previous design ran two ffmpeg subprocesses in
    a thread pool, each doing its own full decode of the input. On a slide-
    heavy talk that meant 2× the I/O + CPU contending on the same file —
    typical wall time was ~5–7 minutes for a 25-minute video. Folding both
    extractions into one decode roughly halves that.
    """
    scene_dir.mkdir(parents=True, exist_ok=True)
    interval_dir.mkdir(parents=True, exist_ok=True)

    # Build the filter graph. The `split=2` filter duplicates the decoded
    # video into two virtual streams; each chain then does its own thing.
    # `showinfo` writes per-frame metadata to stderr — we parse `pts_time`
    # from there to get scene timestamps.
    fps = 1.0 / interval if interval > 0 else 0.0
    if fps > 0:
        filter_complex = (
            "[0:v]split=2[s][i];"
            f"[s]select='eq(n,0)+gt(scene,{scene_threshold})',showinfo[scene_out];"
            f"[i]fps={fps}[interval_out]"
        )
        cmd = [
            "ffmpeg", "-hide_banner", "-y",
            "-i", str(video),
            "-filter_complex", filter_complex,
            "-fps_mode", "vfr", "-q:v", "3",
            "-map", "[scene_out]", str(scene_dir / "scene_%04d.jpg"),
            "-q:v", "3",
            "-map", "[interval_out]", str(interval_dir / "interval_%04d.jpg"),
        ]
    else:
        # Interval sampling disabled — single chain for scene only.
        cmd = [
            "ffmpeg", "-hide_banner", "-y",
            "-i", str(video),
            "-vf", f"select='eq(n,0)+gt(scene,{scene_threshold})',showinfo",
            "-fps_mode", "vfr", "-q:v", "3",
            str(scene_dir / "scene_%04d.jpg"),
        ]

    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg single-pass extraction failed:\n{proc.stderr}")

    # Scene timestamps come from showinfo's pts_time in stderr.
    scene_ts = [float(m) for m in re.findall(r"pts_time:([\d.]+)", proc.stderr)]
    scene_files = sorted(scene_dir.glob("scene_*.jpg"))
    if len(scene_files) != len(scene_ts):
        n = min(len(scene_files), len(scene_ts))
        scene_files, scene_ts = scene_files[:n], scene_ts[:n]
    if len(scene_files) > SCENE_FRAME_HARD_CAP:
        step = len(scene_files) / SCENE_FRAME_HARD_CAP
        keep = [int(i * step) for i in range(SCENE_FRAME_HARD_CAP)]
        scene_files = [scene_files[i] for i in keep]
        scene_ts = [scene_ts[i] for i in keep]
    scene_frames = list(zip(scene_files, scene_ts))

    # Interval timestamps reconstructed from frame index (ffmpeg's fps filter
    # places frame i at t = interval * (i + 0.5)).
    if fps > 0:
        interval_files = sorted(interval_dir.glob("interval_*.jpg"))
        interval_ts = [
            min(duration, interval * (i + 0.5)) for i in range(len(interval_files))
        ]
        interval_frames = list(zip(interval_files, interval_ts))
    else:
        interval_frames = []

    return scene_frames, interval_frames


def merge_frames(*frame_lists: List[Tuple[Path, float]]) -> List[Tuple[Path, float]]:
    """Merge multiple frame lists, sort by timestamp."""
    combined: List[Tuple[Path, float]] = []
    for fl in frame_lists:
        combined.extend(fl)
    combined.sort(key=lambda x: x[1])
    return combined


# ---------- Step 2: Perceptual-hash dedup (consecutive only) ----------

def dedupe_frames(frames: List[Tuple[Path, float]], hash_distance: int = 4) -> List[Tuple[Path, float]]:
    """Cluster runs of near-identical consecutive frames; keep the LAST of each cluster.

    Animated slide reveals (bullets / diagram elements appearing over time) and slow
    pans both hit this case: each intermediate frame is similar to its neighbor, but
    only the final frame shows the fully-revealed / settled state. Keeping the last
    frame of the cluster preserves that state rather than the partial opening one.

    Discrete scene changes (dist > threshold between neighbors) break the cluster,
    so recurring views — e.g. switching to an editor and back to slides — still
    survive as separate kept frames.
    """
    if not frames:
        return []

    import imagehash
    from PIL import Image

    hashed: List[Tuple[Path, float, "imagehash.ImageHash"]] = []
    for path, ts in frames:
        with Image.open(path) as im:
            hashed.append((path, ts, imagehash.phash(im)))

    clusters: List[List[Tuple[Path, float, "imagehash.ImageHash"]]] = [[hashed[0]]]
    for i in range(1, len(hashed)):
        if (hashed[i][2] - hashed[i - 1][2]) <= hash_distance:
            clusters[-1].append(hashed[i])
        else:
            clusters.append([hashed[i]])

    kept: List[Tuple[Path, float]] = []
    for cluster in clusters:
        for path, _, _ in cluster[:-1]:
            path.unlink(missing_ok=True)
        kept.append((cluster[-1][0], cluster[-1][1]))
    return kept


def global_phash_cluster(
    frames: List[Tuple[Path, float]],
    distance: int = 4,
    time_window_secs: float = 90.0,
) -> List[Tuple[Path, float]]:
    """Drop frames whose pHash is close to an earlier kept frame WITHIN a
    bounded time window.

    Consecutive dedup (dedupe_frames) catches near-identical adjacent
    frames but misses the talk-deck pattern where the camera cuts
    speaker → slide → speaker → SAME slide: the speaker cuts break the
    consecutive-similarity chain so each return to the slide opens a new
    cluster, and we end up with N copies.

    The time-window matters: pHash compares low-frequency components,
    which is fine for "same slide that just came back" but routinely
    false-positives on template-similar slides from anywhere in the talk
    (e.g. two different slides that share the same title-bar + two-card
    layout). Without the window, a slide at 21:10 can be falsely merged
    with a different-content-but-same-template slide from 14:12. Bounding
    the comparison to ~90s keeps the "speaker-cut returns to same slide"
    case (happens within ~30s) while preserving distinct slides that
    happen to share a visual template.

    The default distance (4) matches the consecutive threshold — we only
    drop frames that look genuinely identical, not just structurally
    similar. The vision-LLM classifier downstream catches anything that
    slips through.
    """
    if not frames:
        return []
    import imagehash
    from PIL import Image

    kept: List[Tuple[Path, float]] = []
    kept_hashes: List[Tuple["imagehash.ImageHash", float]] = []
    for path, ts in frames:
        with Image.open(path) as im:
            h = imagehash.phash(im)
        drop = False
        # Iterate kept frames in reverse — most matches are recent, so
        # we exit early. Also lets us bail as soon as we see a kept frame
        # older than the time window (since the list is in chronological
        # order, anything older would also fail the time check).
        for kh, kts in reversed(kept_hashes):
            if (ts - kts) > time_window_secs:
                break
            if (h - kh) <= distance:
                drop = True
                break
        if not drop:
            kept.append((path, ts))
            kept_hashes.append((h, ts))
    return kept


# ---------- Vision-grid slide classifier ----------
#
# Tile candidate frames into 3×3 grid images and ask Claude (Haiku, by
# default) to classify each cell as NEW_SLIDE / SAME_AS_PREVIOUS_CELL /
# TALKING_HEAD / TRANSITION. Sending grids instead of single frames is
# ~10× cheaper in input tokens and lets the LLM compare adjacent cells
# directly — easier than reasoning across separate API calls.

_GRID_COLS = 3
_GRID_ROWS = 3
_GRID_CELLS = _GRID_COLS * _GRID_ROWS  # 9
_GRID_CELL_W = 400  # 16:9 thumbnails
_GRID_CELL_H = 225


def _render_classification_grids(
    frames: List[Tuple[Path, float]],
    out_dir: Path,
) -> List[Tuple[Path, List[int]]]:
    """Render frames into one or more numbered 3×3 grid images.

    Returns a list of (grid_path, [original_frame_index_for_each_cell]). The
    last grid may be padded with blank cells; the index list reflects only
    real frames (length ≤ 9 per grid). Cells are numbered 1..9 with a small
    overlay so the LLM can address them unambiguously.

    Each grid carries a 1-cell overlap from the previous grid (the last real
    frame of grid N-1 is duplicated as cell 1 of grid N) so the LLM has
    context for the SAME-AS-PREVIOUS judgement at boundaries.
    """
    from PIL import Image, ImageDraw, ImageFont

    out_dir.mkdir(parents=True, exist_ok=True)
    if not frames:
        return []

    # Try platform-appropriate system fonts in order; fall back to PIL's
    # tiny built-in default if none of them are reachable.
    font = None
    for font_path in (
        "/System/Library/Fonts/Helvetica.ttc",  # macOS
        "C:\\Windows\\Fonts\\arial.ttf",        # Windows
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",  # common Linux
        "/usr/share/fonts/TTF/DejaVuSans-Bold.ttf",              # other Linux
    ):
        try:
            font = ImageFont.truetype(font_path, 32)
            break
        except (OSError, IOError):
            continue
    if font is None:
        font = ImageFont.load_default()

    grids: List[Tuple[Path, List[int]]] = []
    grid_w = _GRID_CELL_W * _GRID_COLS
    grid_h = _GRID_CELL_H * _GRID_ROWS

    # Walk the frame list in chunks of 8 (leaving cell 0 of each grid for
    # the overlap from the previous grid). The first grid has no overlap so
    # it gets all 9 cells.
    i = 0
    grid_idx = 0
    prev_last_idx: Optional[int] = None
    while i < len(frames):
        canvas = Image.new("RGB", (grid_w, grid_h), color="black")
        draw = ImageDraw.Draw(canvas)
        cell_indices: List[int] = []

        # Position 0: overlap cell from previous grid (None on the first grid).
        # We embed the overlap so the model can judge "is cell 1 the same as
        # cell 0?" right at the grid boundary.
        if prev_last_idx is not None and grid_idx > 0:
            cells_to_fill = [prev_last_idx]
            slots_for_new = _GRID_CELLS - 1
        else:
            cells_to_fill = []
            slots_for_new = _GRID_CELLS

        # Fill the remaining cells with new frames.
        end = min(i + slots_for_new, len(frames))
        cells_to_fill.extend(range(i, end))

        for cell_pos, frame_idx in enumerate(cells_to_fill):
            row, col = divmod(cell_pos, _GRID_COLS)
            x0 = col * _GRID_CELL_W
            y0 = row * _GRID_CELL_H
            with Image.open(frames[frame_idx][0]) as im:
                im.thumbnail((_GRID_CELL_W, _GRID_CELL_H))
                # Center the thumbnail in its cell (frames may not be exact 16:9).
                tw, th = im.size
                px = x0 + (_GRID_CELL_W - tw) // 2
                py = y0 + (_GRID_CELL_H - th) // 2
                canvas.paste(im, (px, py))
            # Big numbered label, white-on-black with a small offset, so it
            # stands out against either light or dark cell backgrounds.
            label = str(cell_pos + 1)
            draw.rectangle(
                [(x0 + 4, y0 + 4), (x0 + 44, y0 + 44)],
                fill="black", outline="white",
            )
            draw.text((x0 + 12, y0 + 4), label, fill="white", font=font)
            cell_indices.append(frame_idx)

        grid_path = out_dir / f"grid_{grid_idx:03d}.jpg"
        canvas.save(grid_path, "JPEG", quality=85)
        # Subtract 1 from cell_indices to get the "real new frames" range
        # (cell 0 may be the overlap, which the caller should skip when
        # mapping LLM responses back).
        grids.append((grid_path, cell_indices))

        prev_last_idx = cells_to_fill[-1]
        i = end
        grid_idx += 1

    return grids


def classify_slides_via_grids(
    frames: List[Tuple[Path, float]],
    *,
    backend,
    model: str,
    workdir: Path,
    log_video_id: Optional[str] = None,
) -> List[Tuple[Path, float]]:
    """Use vision LLM (typically Haiku via grids) to filter frames down to
    distinct deck slides. Returns the kept frames in chronological order.

    On any failure (LLM error, malformed response, or implausibly small
    output), returns the input frames unchanged so the caller still gets a
    deck — just a less-pruned one.

    log_video_id (optional): when provided, each per-grid LLM call is
    recorded to the cost-audit log under kind='slide_classifier'.
    """
    from pydantic import BaseModel
    from typing import List as TList, Literal

    class CellLabel(BaseModel):
        cell: int  # 1..9
        label: Literal["NEW_SLIDE", "SAME_AS_PREVIOUS_CELL",
                       "TALKING_HEAD", "TRANSITION"]

    class GridLabels(BaseModel):
        labels: TList[CellLabel]

    if not frames:
        return frames
    if not getattr(backend, "vision_supported", False):
        return frames  # backend without vision (e.g. Claude Code w/o opt-in)

    grid_dir = workdir / "slide_classifier_grids"
    grids = _render_classification_grids(frames, grid_dir)
    if not grids:
        return frames

    system_prompt = (
        "You classify frames extracted from a video that contains a slide "
        "deck. Each grid image you receive is a 3×3 layout (cells numbered "
        "1–9, top-left to bottom-right). For each cell, decide what the "
        "frame is and whether it shows a slide we haven't already seen.\n\n"
        "Labels:\n"
        "- NEW_SLIDE — a slide whose content (text, layout, charts) is "
        "different from cell N-1 in the same grid AND from any slide you've "
        "already labeled NEW_SLIDE in earlier grids.\n"
        "- SAME_AS_PREVIOUS_CELL — visually the same slide as cell N-1 (or, "
        "for cell 1 of grids after the first, the same as the overlap cell). "
        "Animations and reveals count as same.\n"
        "- TALKING_HEAD — frame is dominated by the speaker, no slide visible.\n"
        "- TRANSITION — fade, blur, mid-cut, or otherwise not a clean slide.\n\n"
        "Important:\n"
        "- Cell 1 of grids 2+ duplicates the LAST real cell of the previous "
        "grid as boundary context. Use it to compare cell 2 against the "
        "previous grid's last frame.\n"
        "- Output exactly one label per cell present in the grid, in cell-"
        "number order. If the grid has fewer than 9 real frames, only "
        "output labels for the cells that contain images (skip blank ones)."
    )

    kept_global: set = set()  # original frame indices we've decided to keep
    for grid_idx, (grid_path, cell_indices) in enumerate(grids):
        # Build the message: grid image + a short reminder of cell count.
        with open(grid_path, "rb") as f:
            grid_bytes = f.read()
        import base64 as _b64
        b64 = _b64.b64encode(grid_bytes).decode("ascii")
        content_blocks = [
            {"type": "text", "text":
             f"Grid {grid_idx + 1} of {len(grids)} — "
             f"{len(cell_indices)} cell(s) populated."},
            {"type": "image", "source": {
                "type": "base64", "media_type": "image/jpeg", "data": b64,
            }},
            {"type": "text", "text":
             "Classify each cell. Return one label per real cell, in order."},
        ]
        try:
            parsed, grid_usage = backend.vision_parse(
                system=system_prompt,
                content_blocks=content_blocks,
                model=model,
                max_tokens=600,
                schema=GridLabels,
            )
        except Exception:
            # Defensive: any LLM failure leaves the candidate set intact.
            return frames
        if log_video_id is not None:
            record_llm_usage(
                video_id=log_video_id, kind="slide_classifier",
                model=model, backend_name=backend.name, usage=grid_usage,
            )

        for cell_label in parsed.labels:
            cell_pos = cell_label.cell - 1  # back to 0-indexed
            if cell_pos < 0 or cell_pos >= len(cell_indices):
                continue
            frame_idx = cell_indices[cell_pos]
            # Skip the overlap cell on grids 2+ (cell 0) — already considered.
            if grid_idx > 0 and cell_pos == 0:
                continue
            if cell_label.label == "NEW_SLIDE":
                kept_global.add(frame_idx)

    if not kept_global:
        # Implausible result (LLM said no NEW_SLIDE anywhere) — fall back.
        return frames
    return [frames[i] for i in sorted(kept_global)]


# ---------- Step 3: SRT parser ----------

@dataclass
class TranscriptSegment:
    start: float
    end: float
    text: str


def _srt_time_to_seconds(t: str) -> float:
    """'00:01:23,456' -> 83.456"""
    h, m, rest = t.split(":")
    s, ms = rest.split(",")
    return int(h) * 3600 + int(m) * 60 + int(s) + int(ms) / 1000


def parse_srt(srt_path: Path) -> List[TranscriptSegment]:
    """Parse a .srt file. Tolerant of BOM, CRLF, dot-vs-comma ms separator, simple markup."""
    text = srt_path.read_text(encoding="utf-8-sig")
    blocks = re.split(r"\r?\n\r?\n", text.strip())

    segments: List[TranscriptSegment] = []
    time_re = re.compile(r"(\d{2}:\d{2}:\d{2}[,.]\d{3})\s*-->\s*(\d{2}:\d{2}:\d{2}[,.]\d{3})")

    for block in blocks:
        lines = [ln for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        time_line_idx = next((i for i, ln in enumerate(lines) if time_re.search(ln)), None)
        if time_line_idx is None:
            continue
        m = time_re.search(lines[time_line_idx])
        start = _srt_time_to_seconds(m.group(1).replace(".", ","))
        end = _srt_time_to_seconds(m.group(2).replace(".", ","))
        body = " ".join(lines[time_line_idx + 1:])
        body = re.sub(r"<[^>]+>", "", body)         # <i>, <b>, <font>
        body = re.sub(r"\{[^}]+\}", "", body).strip()  # {\an8}, etc.
        if body:
            segments.append(TranscriptSegment(start=start, end=end, text=body))

    return segments


# ---------- Step 4: Align transcript to frames ----------

def assign_transcript_to_frames(
    frames: List[Tuple[Path, float]],
    segments: List[TranscriptSegment],
    video_duration: float,
) -> List[Tuple[Path, float, float, str]]:
    """
    For each frame at time t_i, the slide covers [t_i, t_{i+1}).
    Collect transcript segments whose midpoint falls in that window.
    Returns: list of (frame_path, slide_start, slide_end, transcript_text).
    """
    results = []
    for i, (path, start) in enumerate(frames):
        end = frames[i + 1][1] if i + 1 < len(frames) else video_duration
        chunk = " ".join(
            seg.text for seg in segments
            if start <= (seg.start + seg.end) / 2 < end
        )
        results.append((path, start, end, chunk))
    return results


def get_video_duration(video: Path) -> float:
    out = subprocess.check_output([
        "ffprobe", "-v", "error", "-show_entries", "format=duration",
        "-of", "json", str(video),
    ])
    return float(json.loads(out)["format"]["duration"])


# ---------- Step 5: Build the .pptx ----------

def format_timestamp(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    return f"{h:d}:{m:02d}:{s:02d}" if h else f"{m:d}:{s:02d}"


def build_deck(slides_data, output: Path, video_name: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = video_name
    p.font.size = Pt(40)
    p.font.bold = True
    sub = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.5))
    sub.text_frame.paragraphs[0].text = f"{len(slides_data)} slides extracted from video"
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Layout: image fills most of the slide (image-first, deck-replication
    # framing). The full transcript chunk lives in PowerPoint's speaker
    # notes so anyone who wants the narration can open it via View → Notes
    # Page; the slide itself stays clean.
    img_top_in = 0.3
    img_max_h_in = 6.8
    img_max_w_in = 12.33
    footer_top_in = 7.15

    for idx, (frame_path, start, end, transcript) in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        with Image.open(frame_path) as im:
            iw, ih = im.size
        scale = min(img_max_w_in / (iw / 96), img_max_h_in / (ih / 96))
        disp_w = (iw / 96) * scale
        disp_h = (ih / 96) * scale
        left = (13.33 - disp_w) / 2

        slide.shapes.add_picture(
            str(frame_path),
            Inches(left), Inches(img_top_in),
            width=Inches(disp_w), height=Inches(disp_h),
        )

        # Footer with time range + slide number — useful provenance, small
        # enough not to compete with the image.
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(footer_top_in), Inches(12.7), Inches(0.3))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f"{format_timestamp(start)} – {format_timestamp(end)}   |   Slide {idx}"
        fp.font.size = Pt(10)

        if transcript:
            slide.notes_slide.notes_text_frame.text = transcript

    prs.save(str(output))


# ---------- Config / API key handling ----------

# Everything lives in one visible directory under HOME. Override with YT2MD_DATA.
# Layout under that directory:
#   .env, channels.txt, state.json, digests/, meta/, downloads/, logs/

DEFAULT_DATA_DIR = Path.home() / "yt2md"


# Cross-platform "detach this child from the parent" Popen kwargs. On POSIX,
# start_new_session=True calls setsid() so the child becomes its own session
# leader and survives even if the parent dies. On Windows, DETACHED_PROCESS
# + CREATE_NEW_PROCESS_GROUP achieve the same effect (the child is detached
# from the console and won't receive the parent's Ctrl-C).
if sys.platform == "win32":
    _DETACH_KWARGS = {
        "creationflags": (
            subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP
        ),
    }
else:
    _DETACH_KWARGS = {"start_new_session": True}


def get_data_dir() -> Path:
    return Path(os.environ.get("YT2MD_DATA", str(DEFAULT_DATA_DIR))).expanduser()


def env_file() -> Path:
    return get_data_dir() / ".env"


def load_env_files() -> None:
    """Populate os.environ from .env files. Real env vars always win.

    Order (lowest priority first; later loads do NOT override earlier-set keys):
      1. Real env vars (from the shell)
      2. CWD/.env (project-local override)
      3. <data dir>/.env (default: ~/yt2md/.env)
    """
    from dotenv import load_dotenv

    load_dotenv()  # CWD/.env, only fills in missing
    e = env_file()
    if e.exists():
        load_dotenv(e)


def set_env_var(name: str, value: str) -> Path:
    """Persist <name>=<value> to ~/yt2md/.env, preserving other entries.

    Uses dotenv.set_key for round-trip safe updates (vs. naive overwrite, which
    would clobber co-resident keys). Also updates os.environ so the running
    process sees the new value immediately. Returns the .env path.
    """
    from dotenv import set_key

    e = env_file()
    e.parent.mkdir(parents=True, exist_ok=True)
    if not e.exists():
        e.touch()
    try:
        os.chmod(e, 0o600)
    except OSError:
        pass
    set_key(str(e), name, value, quote_mode="never")
    os.environ[name] = value
    return e


API_KEY_COST_NOTE = (
    "Anthropic bills your API key per request (separate from any Claude.ai "
    "subscription). Rough costs: a 30-min digest is ~$0.03 with "
    "<code>claude-sonnet-4-6</code> (default), ~$0.15 with "
    "<code>claude-opus-4-7</code>. The panel discussion adds one Opus call "
    "(~$0.10). Add a payment method at "
    '<a href="https://console.anthropic.com/settings/billing" target="_blank" '
    'rel="noopener">console.anthropic.com/settings/billing</a>.'
)


def validate_api_key(key: str) -> Optional[str]:
    """Send a 1-token request to the cheapest model to verify auth. Returns
    None on success, otherwise a short human-readable error string.
    """
    import anthropic

    try:
        client = anthropic.Anthropic(api_key=key)
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1,
            messages=[{"role": "user", "content": "ok"}],
        )
        # Even validation pings cost a few tokens — record so the audit
        # log is complete.
        record_llm_usage(
            video_id=None, kind="validation",
            model="claude-haiku-4-5-20251001",
            backend_name="api", usage=resp.usage,
        )
        return None
    except anthropic.AuthenticationError:
        return "key rejected by Anthropic (authentication failed)."
    except anthropic.PermissionDeniedError as e:
        return f"key authenticated but lacks permission: {e}"
    except anthropic.APIConnectionError as e:
        return f"could not reach Anthropic: {e}"
    except Exception as e:
        return f"unexpected error: {type(e).__name__}: {e}"


def ensure_api_key() -> None:
    """Make sure ANTHROPIC_API_KEY is set, prompting + saving on first run if interactive."""
    if os.environ.get("ANTHROPIC_API_KEY"):
        return

    e = env_file()
    msg = (
        "ANTHROPIC_API_KEY is not set.\n"
        "Get a key from: https://console.anthropic.com/settings/keys"
    )
    if not sys.stdin.isatty():
        sys.exit(
            f"{msg}\n"
            "Then either export it (`export ANTHROPIC_API_KEY=...`) or save it via:\n"
            f"  mkdir -p {e.parent} && echo 'ANTHROPIC_API_KEY=sk-ant-...' >> {e}"
        )

    print(msg)
    key = input("Paste your API key (or press Enter to abort): ").strip()
    if not key:
        sys.exit("Aborted.")

    save = input(
        f"Save it to {e} so future runs find it automatically? [Y/n] "
    ).strip().lower()
    if save in ("", "y", "yes"):
        set_env_var("ANTHROPIC_API_KEY", key)
        print(f"      saved to {e}")
    else:
        os.environ["ANTHROPIC_API_KEY"] = key


# ---------- Claude Code sandbox (alternative auth path) ----------
#
# yt2md ships a private Claude Code install under <data dir>/claude-code/.
# Driving the official `claude` binary as a subprocess is the supported way
# for a third-party tool to leverage a user's Claude.ai subscription auth
# without violating ToS (vs. extracting OAuth tokens, which is forbidden).
#
# Sandbox layout:
#   <data>/claude-code/node_modules/.bin/claude   <- the binary we invoke
#   <data>/claude-code/config/                     <- CLAUDE_CONFIG_DIR target
#       settings.json, projects/, plugins/, .credentials.json (Linux/Win)
# On macOS, CLAUDE_CONFIG_DIR isolates settings/projects but credentials still
# go to the system Keychain — this is a documented Claude Code limitation.

CLAUDE_CODE_NPM_PACKAGE = "@anthropic-ai/claude-code"  # always pulls latest
MIN_NODE_MAJOR = 18


def claude_sandbox_dir() -> Path:
    return get_data_dir() / "claude-code"


def claude_config_dir() -> Path:
    return claude_sandbox_dir() / "config"


def claude_binary_path() -> Path:
    """Path to the sandboxed claude executable. npm installs both a Unix
    shell shim (`claude`) and a Windows batch shim (`claude.cmd`) into
    `node_modules/.bin/`. Python's subprocess on Windows can't execute the
    shell shim directly, so we point at the .cmd shim there."""
    bin_dir = claude_sandbox_dir() / "node_modules" / ".bin"
    if sys.platform == "win32":
        return bin_dir / "claude.cmd"
    return bin_dir / "claude"


def claude_subprocess_env() -> dict:
    """Env dict for invoking the sandboxed claude binary. Pins CLAUDE_CONFIG_DIR
    into our sandbox so settings/projects/plugins don't collide with any
    system-wide Claude Code install. (macOS Keychain owns credentials regardless.)
    """
    env = os.environ.copy()
    env["CLAUDE_CONFIG_DIR"] = str(claude_config_dir())
    return env


def detect_node() -> Optional[Tuple[int, str]]:
    """Returns (major_version, full_version_string) if Node ≥ 18 is available,
    else None. We need npm to install the sandbox; npm requires Node.
    """
    node = shutil.which("node")
    if not node:
        return None
    try:
        out = subprocess.check_output([node, "-v"], text=True, timeout=5).strip()
    except (subprocess.SubprocessError, OSError):
        return None
    # `node -v` prints e.g. "v18.19.0"
    m = re.match(r"^v(\d+)\.", out)
    if not m:
        return None
    major = int(m.group(1))
    if major < MIN_NODE_MAJOR:
        return None
    return major, out


def claude_code_installed() -> bool:
    """True if our sandbox has a working claude binary."""
    return claude_binary_path().exists()


def install_claude_code(stream_to: Optional[Path] = None) -> Tuple[int, str]:
    """Run `npm install --prefix <sandbox> <package>` to materialize a private
    Claude Code install. Returns (returncode, combined_output). If stream_to
    is given, output is appended there as it's produced (for live UI polling).
    """
    sandbox = claude_sandbox_dir()
    sandbox.mkdir(parents=True, exist_ok=True)
    claude_config_dir().mkdir(parents=True, exist_ok=True)

    npm = shutil.which("npm")
    if npm is None:
        return 127, "npm not found on PATH (install Node.js 18+ first)."

    cmd = [npm, "install", "--prefix", str(sandbox),
           "--no-fund", "--no-audit", "--silent",
           CLAUDE_CODE_NPM_PACKAGE]
    proc = subprocess.Popen(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True,
    )
    chunks: list = []
    assert proc.stdout is not None
    for line in proc.stdout:
        chunks.append(line)
        if stream_to is not None:
            try:
                with open(stream_to, "a") as f:
                    f.write(line)
            except OSError:
                pass
    proc.wait()
    return proc.returncode, "".join(chunks)


def claude_code_logged_in() -> bool:
    """Heuristic: if there's a working binary AND `claude /status` (or a tiny
    -p call) succeeds without auth error, we're logged in. Cheap probe via the
    presence of credentials state."""
    if not claude_code_installed():
        return False
    # The cheapest probe is a 1-token --print call; cache the result for the
    # lifetime of the process (set when we explicitly log in/out).
    return _claude_code_session_state.get("logged_in", False)


# Module-level cache of login probe state. Populated by validate_claude_code()
# after install/login; checked by claude_code_logged_in() to avoid spawning a
# subprocess on every page load.
_claude_code_session_state: dict = {}


def _claude_login_sentinel() -> Path:
    """Touched after a successful validation; lets us assume logged-in across
    server restarts without burning a token-cost probe per boot. Cleared by
    claude_logout(). Real auth failures still reset the session state when
    encountered."""
    return claude_config_dir() / ".yt2md-logged-in"


def claude_probe_login_state() -> None:
    """Cheap startup-time probe: read the sentinel file and populate the
    session-state cache. No subprocess call. Real auth-failure reset happens
    in validate_claude_code() and on first failed LLM call.
    """
    if claude_code_installed() and _claude_login_sentinel().exists():
        _claude_code_session_state["logged_in"] = True
    else:
        _claude_code_session_state["logged_in"] = False


def validate_claude_code() -> Optional[str]:
    """Run a 1-token call through the sandboxed claude binary. Returns None on
    success, otherwise a short error string. Updates the session-state cache
    and the on-disk sentinel.
    """
    if not claude_code_installed():
        return "Claude Code is not installed in the sandbox."
    cmd = [
        str(claude_binary_path()), "-p", "ok",
        "--model", "claude-haiku-4-5-20251001",
        "--output-format", "json",
    ]
    try:
        proc = subprocess.run(
            cmd, env=claude_subprocess_env(),
            capture_output=True, text=True, timeout=60,
        )
    except subprocess.TimeoutExpired:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        return "validation timed out (60s) — the OAuth flow may not have completed."
    except OSError as e:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        return f"could not invoke claude binary: {e}"
    if proc.returncode != 0:
        _claude_code_session_state["logged_in"] = False
        _claude_login_sentinel().unlink(missing_ok=True)
        msg = (proc.stderr or proc.stdout or "").strip().splitlines()
        last = msg[-1] if msg else f"exit code {proc.returncode}"
        return f"claude returned an error: {last}"
    _claude_code_session_state["logged_in"] = True
    try:
        _claude_login_sentinel().parent.mkdir(parents=True, exist_ok=True)
        _claude_login_sentinel().touch()
    except OSError:
        pass
    return None


# In-memory tracker for async setup jobs (install, login). Keyed by job name
# ("install" | "login"). Each entry: {"proc": Popen, "log": Path, "started":
# epoch, "error": Optional[str]}. The web reader is single-process, so a plain
# dict suffices. State resets on server restart, which is fine — finished jobs
# leave their result on disk (sandbox dir exists; credentials persist).

_claude_setup_jobs: dict = {}


def _claude_setup_log(name: str) -> Path:
    p = get_data_dir() / "logs" / f"claude-setup-{name}.log"
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def claude_setup_job_running(name: str) -> bool:
    job = _claude_setup_jobs.get(name)
    if job is None:
        return False
    proc = job.get("proc")
    return proc is not None and proc.poll() is None


def start_install_job() -> Optional[str]:
    """Spawn `npm install` for Claude Code in the sandbox. Idempotent: returns
    None if it's already running, or an error string if Node is missing.
    """
    if claude_setup_job_running("install"):
        return None
    if detect_node() is None:
        return (
            f"Node.js {MIN_NODE_MAJOR}+ is required. Install with "
            "`brew install node` (macOS) / `winget install OpenJS.NodeJS` (Windows), then retry."
        )
    sandbox = claude_sandbox_dir()
    sandbox.mkdir(parents=True, exist_ok=True)
    claude_config_dir().mkdir(parents=True, exist_ok=True)
    log_path = _claude_setup_log("install")
    log_path.write_text("")  # truncate prior run
    npm = shutil.which("npm") or "npm"
    cmd = [npm, "install", "--prefix", str(sandbox),
           "--no-fund", "--no-audit",
           CLAUDE_CODE_NPM_PACKAGE]
    log_f = open(log_path, "a")
    proc = subprocess.Popen(
        cmd, stdout=log_f, stderr=subprocess.STDOUT, text=True,
        **_DETACH_KWARGS,
    )
    import time as _t
    _claude_setup_jobs["install"] = {
        "proc": proc, "log": log_path, "started": _t.time(),
        "log_f": log_f, "error": None,
    }
    return None


def start_login_job() -> Optional[str]:
    """Spawn `claude /login`. The CLI auto-opens the user's default browser
    via `open` for OAuth; the random localhost callback is captured by the
    subprocess itself. We wait for exit-0 then validate.
    """
    if claude_setup_job_running("login"):
        return None
    if not claude_code_installed():
        return "Claude Code is not installed yet. Install it first."
    log_path = _claude_setup_log("login")
    log_path.write_text("")
    cmd = [str(claude_binary_path()), "/login"]
    log_f = open(log_path, "a")
    proc = subprocess.Popen(
        cmd, env=claude_subprocess_env(),
        stdout=log_f, stderr=subprocess.STDOUT,
        stdin=subprocess.DEVNULL,  # no terminal stdin → CLI relies on browser callback
        text=True,
        **_DETACH_KWARGS,
    )
    import time as _t
    _claude_setup_jobs["login"] = {
        "proc": proc, "log": log_path, "started": _t.time(),
        "log_f": log_f, "error": None, "validated": False,
    }
    return None


def claude_logout() -> Tuple[int, str]:
    """Run `claude /logout` to clear stored credentials. Resets session cache."""
    if not claude_code_installed():
        return 0, "(not installed)"
    cmd = [str(claude_binary_path()), "/logout"]
    proc = subprocess.run(
        cmd, env=claude_subprocess_env(),
        capture_output=True, text=True, timeout=30,
    )
    _claude_code_session_state["logged_in"] = False
    _claude_login_sentinel().unlink(missing_ok=True)
    return proc.returncode, (proc.stdout or "") + (proc.stderr or "")


def claude_setup_snapshot() -> dict:
    """Status blob consumed by the /setup polling JS. Reaps finished jobs and
    runs validation after a successful login.
    """
    install_job = _claude_setup_jobs.get("install")
    login_job = _claude_setup_jobs.get("login")

    # Reap install if it finished — surface the exit code as an error if non-zero.
    if install_job and install_job["proc"].poll() is not None:
        rc = install_job["proc"].returncode
        if rc != 0 and not install_job.get("error"):
            install_job["error"] = f"npm install failed (exit {rc}). See log."
        try:
            install_job["log_f"].close()
        except Exception:
            pass

    # Reap login. On success, validate with a 1-token call (cached).
    if login_job and login_job["proc"].poll() is not None:
        try:
            login_job["log_f"].close()
        except Exception:
            pass
        if not login_job.get("validated"):
            rc = login_job["proc"].returncode
            if rc == 0:
                err = validate_claude_code()
                if err:
                    login_job["error"] = err
            else:
                login_job["error"] = f"login subprocess exited with code {rc}. See log."
            login_job["validated"] = True

    def _tail(path: Path, n: int = 30) -> str:
        try:
            lines = path.read_text(errors="replace").splitlines()
            return "\n".join(lines[-n:])
        except OSError:
            return ""

    return {
        "node_ok": detect_node() is not None,
        "installed": claude_code_installed(),
        "logged_in": _claude_code_session_state.get("logged_in", False),
        "install_running": claude_setup_job_running("install"),
        "login_running": claude_setup_job_running("login"),
        "install_log_tail": _tail(install_job["log"]) if install_job else "",
        "login_log_tail": _tail(login_job["log"]) if login_job else "",
        "install_error": install_job.get("error") if install_job else None,
        "login_error": login_job.get("error") if login_job else None,
    }


# ---------- YouTube fetch ----------

URL_RE = re.compile(r"^https?://", re.IGNORECASE)


def is_url(s: str) -> bool:
    return bool(URL_RE.match(s))


DEFAULT_WHISPER_MODEL = "medium"


def _whisper_secs_to_srt(secs: float) -> str:
    """Convert seconds to SRT-style HH:MM:SS,mmm timestamp."""
    if secs < 0:
        secs = 0.0
    h = int(secs // 3600)
    m = int((secs % 3600) // 60)
    s = int(secs % 60)
    ms = int(round((secs - int(secs)) * 1000))
    if ms == 1000:  # rounding can push us a full ms over
        s += 1
        ms = 0
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def _transcribe_with_whisper(
    media_path: Path,
    out_dir: Path,
    video_id: str,
    model_name: str = DEFAULT_WHISPER_MODEL,
) -> Tuple[Path, str]:
    """Transcribe a media file with faster-whisper, write an SRT, return (srt_path, lang).

    Model weights are downloaded on first use to ~/.cache/huggingface and reused
    afterwards. Detected language is used as the SRT filename suffix so the
    existing cache-by-glob logic in fetch_youtube picks it up on re-runs.
    """
    try:
        from faster_whisper import WhisperModel
    except ImportError as e:
        raise RuntimeError(
            "Whisper transcription needed but faster-whisper is not installed. "
            "Run `uv sync` (or `pip install faster-whisper`) and try again."
        ) from e

    print(f"      loading whisper model '{model_name}' (first run downloads weights)")
    # int8 keeps memory low and runs well on CPU; faster-whisper picks Metal/CUDA
    # automatically when device='auto'.
    model = WhisperModel(model_name, device="auto", compute_type="int8")

    print(f"      transcribing audio with whisper ({media_path.name})...")
    segments_iter, info = model.transcribe(
        str(media_path),
        beam_size=5,
        vad_filter=True,  # cuts long silences so the transcript stays tight
    )

    lang = info.language or "und"
    srt_path = out_dir / f"{video_id}.{lang}.srt"

    n = 0
    with srt_path.open("w") as fh:
        for seg in segments_iter:
            n += 1
            text = seg.text.strip().replace("\n", " ")
            if not text:
                continue
            fh.write(
                f"{n}\n"
                f"{_whisper_secs_to_srt(seg.start)} --> {_whisper_secs_to_srt(seg.end)}\n"
                f"{text}\n\n"
            )

    print(
        f"      whisper: {n} segments, language='{lang}' "
        f"(prob={info.language_probability:.2f})"
    )
    return srt_path, lang


def _ensure_js_runtime_available() -> Optional[str]:
    """Find a JS runtime usable for yt-dlp's n-challenge solver.

    yt-dlp accepts deno / node / bun. As of 2026, it marks Node <20 as
    'unsupported' — silently failing the n-challenge and producing only
    storyboard formats. So for Node we collect all candidates (PATH match +
    common version-manager locations: nvm, fnm, asdf, volta), version-rank
    them, and prepend the dir of the highest version to os.environ['PATH']
    so yt-dlp's internal lookups pick the right one.
    """
    # deno / bun: trust the first PATH match (no version-rank needed).
    for rt in ("deno", "bun"):
        if shutil.which(rt):
            return rt

    home = Path.home()
    seen: set = set()
    node_candidates: List[Path] = []

    def _add(p: Optional[Path]):
        if p and p.is_file() and str(p) not in seen:
            seen.add(str(p))
            node_candidates.append(p)

    path_node = shutil.which("node")
    if path_node:
        _add(Path(path_node))
    for p in sorted((home / ".nvm" / "versions" / "node").glob("*/bin/node"), reverse=True):
        _add(p)
    for p in sorted(
        (home / ".local" / "share" / "fnm" / "node-versions").glob("*/installation/bin/node"),
        reverse=True,
    ):
        _add(p)
    for p in sorted((home / ".asdf" / "installs" / "nodejs").glob("*/bin/node"), reverse=True):
        _add(p)
    _add(home / ".volta" / "bin" / "node")

    best_path: Optional[Path] = None
    best_version: Tuple[int, ...] = (0, 0, 0)
    for c in node_candidates:
        try:
            r = subprocess.run(
                [str(c), "--version"], capture_output=True, text=True, timeout=5
            )
            v = tuple(int(p) for p in r.stdout.strip().lstrip("v").split(".")[:3])
        except Exception:
            continue
        if v > best_version:
            best_version = v
            best_path = c

    if best_path is None:
        return None

    # Prepend its dir so yt-dlp's shutil.which lookups pick this one.
    bin_dir = str(best_path.parent)
    current = shutil.which("node")
    if current != str(best_path):
        os.environ["PATH"] = f"{bin_dir}{os.pathsep}{os.environ.get('PATH', '')}"
        print(
            f"[yt2md] using node v{'.'.join(map(str, best_version))} at "
            f"{best_path} for yt-dlp"
        )

    if best_version < (20, 0, 0):
        print(
            f"[yt2md] warning: node v{'.'.join(map(str, best_version))} is below "
            "yt-dlp's required v20+ for YouTube JS challenges. Install a newer "
            "node (`nvm install 20`) or deno (`brew install deno`).",
            file=sys.stderr,
        )

    return "node"


def _pick_caption_lang(info: dict) -> Optional[Tuple[str, bool]]:
    """Choose best caption track from a yt-dlp info dict.

    Returns (lang_code, is_manual) or None if no captions exist at all.

    Priority:
      1. Manual English (any en-* code)
      2. Manual matching the audio language (info['language'])
      3. Any manual track
      4. Auto matching the audio language
      5. Auto English
      6. Any auto track

    Manual is preferred over auto everywhere. Within auto, we prefer the
    original audio language over English: YouTube's auto-EN on a non-English
    video is translation-of-auto-caption (double degradation), whereas
    Claude translating the auto-caption-of-original-audio is single degradation
    and produces better digests.
    """
    manual = info.get("subtitles") or {}
    auto = info.get("automatic_captions") or {}
    audio_lang = (info.get("language") or "").lower()

    def first_starts_with(d, prefix):
        if not prefix:
            return None
        for k in d:
            if k.lower().startswith(prefix):
                return k
        return None

    def first_any(d):
        return next(iter(d), None)

    for picker, source, is_manual in (
        (lambda d: first_starts_with(d, "en"), manual, True),
        (lambda d: first_starts_with(d, audio_lang), manual, True),
        (first_any, manual, True),
        (lambda d: first_starts_with(d, audio_lang), auto, False),
        (lambda d: first_starts_with(d, "en"), auto, False),
        (first_any, auto, False),
    ):
        k = picker(source)
        if k:
            return (k, is_manual)
    return None


def strip_markdown_for_tts(text: str) -> str:
    """Convert markdown to plain text suitable for TTS narration.

    Removes images, HTML tags, and markdown formatting characters while
    preserving the paragraph structure (so the synthesizer gets natural
    pauses). Links collapse to their display text.

    Keeping this minimal on purpose — we'd rather narrate a slightly-
    awkward line than swallow content trying to be clever.
    """
    # Images (markdown + raw HTML img tags)
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"<img\s[^>]*>", "", text)
    # Auto-links < https://... > → drop entirely (URLs read awfully aloud)
    text = re.sub(r"<https?://[^>]+>", "", text)
    # Bracketed timestamp links — these are the inline [3:15](url) anchors
    # the takeaway prompt emits to ground specific claims. They read as
    # "Goldman 4:00 and OpenAI" mid-sentence, which is just confusing
    # aloud. Drop them entirely (text + URL).
    text = re.sub(r"\[\d+:\d+(?::\d+)?\]\([^)]+\)", "", text)
    # Other markdown links → keep just the display text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    # Strip headings markers but keep the heading text on its own line
    text = re.sub(r"^#+\s+", "", text, flags=re.MULTILINE)
    # Strip emphasis markers
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"__([^_]+)__", r"\1", text)
    text = re.sub(r"(?<!\w)_([^_]+)_(?!\w)", r"\1", text)
    # Strip list markers
    text = re.sub(r"^[-*+]\s+", "", text, flags=re.MULTILINE)
    # Keep inner text of inline HTML wrappers (<sub>, <em>, <code>, <strong>…)
    text = re.sub(r"<(\w+)(\s[^>]*)?>(.*?)</\1>", r"\3", text, flags=re.DOTALL)
    # Drop any remaining standalone HTML
    text = re.sub(r"<[^>]+>", "", text)
    # Collapse runs of blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _elevenlabs_api_key() -> Optional[str]:
    """Look up the ElevenLabs API key from the environment. The canonical
    name is ELEVENLABS_API_KEY, but we accept any case variant since
    .env files quietly tolerate typos and humans don't notice them."""
    direct = os.environ.get("ELEVENLABS_API_KEY") or os.environ.get("ELEVEN_API_KEY")
    if direct:
        return direct
    for k, v in os.environ.items():
        if k.upper() in ("ELEVENLABS_API_KEY", "ELEVEN_API_KEY") and v:
            return v
    return None


def _tts_macos(text: str, mp3_path: Path,
               *, voice: Optional[str], rate: Optional[int]) -> None:
    """Render text → MP3 via macOS `say` + ffmpeg transcode."""
    if sys.platform != "darwin":
        raise RuntimeError(
            "macOS `say` is only available on macOS. Switch your "
            "TTS provider to 'elevenlabs' in Settings."
        )
    if shutil.which("say") is None:
        raise RuntimeError("`say` not on PATH (expected on macOS).")
    if shutil.which("ffmpeg") is None:
        raise RuntimeError("ffmpeg not on PATH — install with `brew install ffmpeg`.")

    workdir = Path(tempfile.mkdtemp(prefix="yt2md_audio_"))
    try:
        aiff_path = workdir / "out.aiff"
        say_cmd: List[str] = ["say", "-o", str(aiff_path)]
        if voice:
            say_cmd += ["-v", str(voice)]
        if rate is not None:
            try:
                rate_int = int(str(rate).strip())
                say_cmd += ["-r", str(rate_int)]
            except (ValueError, TypeError):
                pass
        # -f - reads text from stdin (avoids argv length + shell quoting).
        say_cmd += ["-f", "-"]
        proc = subprocess.run(
            say_cmd, input=text, text=True, capture_output=True,
        )
        if proc.returncode != 0:
            raise RuntimeError(
                f"`say` failed (exit {proc.returncode}): "
                f"{(proc.stderr or '').strip()[-500:]}"
            )

        mp3_path.parent.mkdir(parents=True, exist_ok=True)
        mp3_tmp = mp3_path.with_suffix(".mp3.tmp")
        proc = subprocess.run(
            ["ffmpeg", "-hide_banner", "-y",
             "-i", str(aiff_path),
             "-codec:a", "libmp3lame", "-b:a", "64k",
             "-f", "mp3", str(mp3_tmp)],
            capture_output=True, text=True,
        )
        if proc.returncode != 0:
            mp3_tmp.unlink(missing_ok=True)
            raise RuntimeError(
                f"ffmpeg transcode failed (exit {proc.returncode}): "
                f"{(proc.stderr or '').strip()[-500:]}"
            )
        mp3_tmp.replace(mp3_path)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def _chunk_text_for_tts(text: str, max_chars: int) -> List[str]:
    """Split text into chunks of at most max_chars, preferring paragraph
    boundaries (split on blank lines), then sentence boundaries, then
    hard chops as a last resort.
    """
    if len(text) <= max_chars:
        return [text]
    out: List[str] = []
    cur = ""
    for para in text.split("\n\n"):
        if len(cur) + len(para) + 2 <= max_chars:
            cur = (cur + "\n\n" + para).strip() if cur else para
            continue
        if cur:
            out.append(cur)
            cur = ""
        # Paragraph itself bigger than the budget — split on sentences.
        if len(para) <= max_chars:
            cur = para
            continue
        sentences = re.split(r"(?<=[.!?])\s+", para)
        for sent in sentences:
            if len(cur) + len(sent) + 1 <= max_chars:
                cur = (cur + " " + sent).strip() if cur else sent
            else:
                if cur:
                    out.append(cur)
                # Sentence too long even alone — hard chop.
                while len(sent) > max_chars:
                    out.append(sent[:max_chars])
                    sent = sent[max_chars:]
                cur = sent
    if cur.strip():
        out.append(cur.strip())
    return out


def _tts_elevenlabs(text: str, mp3_path: Path,
                    *, voice_id: str, model_id: str) -> None:
    """Render text → MP3 via the ElevenLabs API.

    Chunks text on paragraph/sentence boundaries to stay under the
    per-request character limit (~5k chars). MP3 frames are
    self-contained so concatenating chunk bytes produces a valid file
    with negligible audible artifacts at paragraph boundaries.

    Network failures, auth errors, and rate limits all surface as
    RuntimeError with the response body included so the user can act on
    them.
    """
    import urllib.request
    import urllib.error

    api_key = _elevenlabs_api_key()
    if not api_key:
        raise RuntimeError(
            "ELEVENLABS_API_KEY not set. Add it to ~/yt2md/.env and "
            "restart yt2md serve."
        )

    chunks = _chunk_text_for_tts(text, max_chars=4500)
    audio_bytes = bytearray()
    for i, chunk in enumerate(chunks):
        req = urllib.request.Request(
            f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}",
            method="POST",
            data=json.dumps({
                "text": chunk,
                "model_id": model_id,
                "voice_settings": {
                    "stability": 0.5,
                    "similarity_boost": 0.75,
                },
            }).encode("utf-8"),
            headers={
                "xi-api-key": api_key,
                "Content-Type": "application/json",
                "Accept": "audio/mpeg",
                "User-Agent": "yt2md/1.0 (ElevenLabs TTS)",
            },
        )
        try:
            with urllib.request.urlopen(req, timeout=180) as resp:
                audio_bytes.extend(resp.read())
        except urllib.error.HTTPError as e:
            body = e.read().decode("utf-8", "replace")[:500]
            raise RuntimeError(
                f"ElevenLabs API error {e.code} on chunk {i + 1}/"
                f"{len(chunks)}: {body}"
            )
        except urllib.error.URLError as e:
            raise RuntimeError(
                f"ElevenLabs API request failed on chunk {i + 1}/"
                f"{len(chunks)}: {e.reason}"
            )

    if not audio_bytes:
        raise RuntimeError("ElevenLabs returned no audio bytes.")

    mp3_path.parent.mkdir(parents=True, exist_ok=True)
    mp3_tmp = mp3_path.with_suffix(".mp3.tmp")
    mp3_tmp.write_bytes(bytes(audio_bytes))
    mp3_tmp.replace(mp3_path)


def generate_audio_from_markdown(
    md_path: Path,
    mp3_path: Path,
    *,
    provider: Optional[str] = None,
    voice: Optional[str] = None,
    rate: Optional[int] = None,
    elevenlabs_voice_id: Optional[str] = None,
    elevenlabs_model: Optional[str] = None,
) -> None:
    """Strip markdown to plain text and synthesize an MP3 via the
    configured TTS provider.

    provider: "macos" (default, free, lower quality) or "elevenlabs"
    (high quality, costs against your ElevenLabs plan credits).
    """
    text = strip_markdown_for_tts(md_path.read_text())
    if not text:
        raise RuntimeError("Markdown produced no narrate-able text.")

    chosen = (provider or "macos").lower()
    if chosen == "elevenlabs":
        _tts_elevenlabs(
            text, mp3_path,
            voice_id=elevenlabs_voice_id or "nPczCjzI2devNBz1zQrb",
            model_id=elevenlabs_model or "eleven_multilingual_v2",
        )
    elif chosen == "macos":
        _tts_macos(text, mp3_path, voice=voice, rate=rate)
    else:
        raise RuntimeError(f"Unknown tts_provider: {chosen!r}")


# Maps the audio "kind" string used in URLs / job keys to the source
# markdown filename. Kept in one place so the routes, viewer UI, and
# background worker all agree on what the supported kinds are.
AUDIO_SOURCE_BY_KIND = {
    "digest":   "digest.md",
    "panel":    "panel.md",
    "takeaway": "takeaway.md",
}


def download_image(url: str, dest: Path, *, timeout: float = 15.0) -> bool:
    """Save a remote image to dest. Returns True on success. Atomic write
    via .tmp suffix → rename so a partial download never replaces an
    existing file. Best-effort: any failure is swallowed and returned as
    False, since thumbnails are nice-to-have, not load-bearing.
    """
    if not url:
        return False
    tmp: Optional[Path] = None
    try:
        import urllib.request
        dest.parent.mkdir(parents=True, exist_ok=True)
        tmp = dest.with_suffix(dest.suffix + ".tmp")
        req = urllib.request.Request(
            url, headers={"User-Agent": "yt2md/1.0 (thumbnail fetch)"},
        )
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            tmp.write_bytes(resp.read())
        tmp.replace(dest)
        return True
    except Exception:
        if tmp is not None:
            try:
                tmp.unlink(missing_ok=True)
            except Exception:
                pass
        return False


def probe_channel_thumbnail_url(
    channel_url: str,
    cookies_from_browser: Optional[str] = None,
) -> Optional[str]:
    """yt-dlp probe for the channel's avatar URL. Channel avatars aren't
    in the standard video info, so this is a separate extract. Best-effort:
    returns None on any failure (private channel, anti-bot wall, no
    channel_url). The caller falls back to displaying just the channel
    name without an avatar.
    """
    if not channel_url:
        return None
    try:
        import yt_dlp
        opts: dict = {
            "quiet": True, "no_warnings": True, "skip_download": True,
            "extract_flat": False, "playlist_items": "0",
        }
        if cookies_from_browser:
            opts["cookiesfrombrowser"] = (cookies_from_browser,)
        with yt_dlp.YoutubeDL(opts) as ydl:
            ch_info = ydl.extract_info(channel_url, download=False)
    except Exception:
        return None
    if not isinstance(ch_info, dict):
        return None
    # Channel info shapes vary by extractor. Try the standard list-of-
    # thumbnails first, picking the largest by pixel area, then fall back
    # to a single `thumbnail` field.
    thumbs = ch_info.get("thumbnails") or []
    sized = [
        t for t in thumbs
        if isinstance(t, dict) and t.get("url")
    ]
    if sized:
        sized.sort(key=lambda t: (t.get("height") or 0) * (t.get("width") or 0))
        return sized[-1].get("url")
    return ch_info.get("thumbnail")


def fetch_youtube(
    url: str,
    cache_root: Path,
    whisper_model: str = DEFAULT_WHISPER_MODEL,
    allow_whisper: bool = True,
    cookies_from_browser: Optional[str] = None,
) -> dict:
    """Download mp4 + best-available SRT from YouTube. Cached by video ID under cache_root.

    Returns a dict with:
      mp4 (Path), srt (Path), lang (str),
      title (str), webpage_url (str),
      download_secs (float), whisper_secs (float),
      used_whisper (bool), whisper_model (Optional[str]).

    Falls back to local Whisper transcription when YouTube has no captions.
    Set allow_whisper=False to fail fast instead of falling back.
    """
    import yt_dlp
    import time as _time

    cache_root.mkdir(parents=True, exist_ok=True)

    # YouTube increasingly requires logged-in cookies to bypass the bot challenge.
    # yt-dlp accepts `cookiesfrombrowser` as a tuple — single-item is the simplest
    # form (no profile / domain filter).
    cookie_opt: dict = {}
    if cookies_from_browser:
        cookie_opt["cookiesfrombrowser"] = (cookies_from_browser,)

    # YouTube's "n challenge" obfuscates real format URLs behind a JavaScript
    # function that yt-dlp must execute to deobfuscate. Without a JS runtime +
    # the challenge solver scripts, only thumbnail storyboards come back.
    rt = _ensure_js_runtime_available()
    yt_dlp_runtime_opt: dict = {}
    if rt is not None:
        yt_dlp_runtime_opt["js_runtimes"] = {rt: {}}
        yt_dlp_runtime_opt["remote_components"] = ["ejs:github"]

    base_opts = {**cookie_opt, **yt_dlp_runtime_opt}

    # Probe first to get the video ID for stable cache layout. We use the same
    # format selector as the download below so probe doesn't reject videos whose
    # default yt-dlp selector ('bestvideo*+bestaudio') happens to match nothing
    # (some YouTube videos return a format pool that the default doesn't span).
    probe_opts = {
        "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
        "quiet": True,
        "no_warnings": True,
        **base_opts,
    }
    with yt_dlp.YoutubeDL(probe_opts) as ydl:
        info = ydl.extract_info(url, download=False)
    video_id = info["id"]
    title = info.get("title") or video_id
    webpage_url = info.get("webpage_url") or url
    upload_date = info.get("upload_date")  # YYYYMMDD per yt-dlp; None if missing
    # Thumbnail + channel info — used by the sidebar to render small visual
    # anchors next to each digest. Best-effort; any field may be missing
    # (especially for non-YouTube sources via yt-dlp).
    thumbnail_url = info.get("thumbnail")
    channel_id = info.get("channel_id") or info.get("uploader_id") or ""
    channel_name = info.get("channel") or info.get("uploader") or ""
    channel_url = info.get("channel_url") or info.get("uploader_url") or ""
    # yt-dlp video info doesn't carry a channel-avatar URL in the standard
    # fields; we keep the channel_url and let the caller probe separately
    # if it wants a channel thumbnail.
    out_dir = cache_root / video_id
    out_dir.mkdir(parents=True, exist_ok=True)

    mp4_path = out_dir / f"{video_id}.mp4"

    # Cache hit: lang is in the filename (works for legacy *.en.srt too).
    existing_srt = next(iter(out_dir.glob(f"{video_id}.*.srt")), None)
    if mp4_path.exists() and existing_srt is not None:
        lang = existing_srt.stem[len(video_id) + 1:]
        print(f"      using cached {out_dir}/ (lang: {lang})")
        return {
            "mp4": mp4_path, "srt": existing_srt, "lang": lang,
            "title": title, "webpage_url": webpage_url, "upload_date": upload_date, "thumbnail_url": thumbnail_url, "channel_id": channel_id, "channel_name": channel_name, "channel_url": channel_url,
            "download_secs": 0.0, "whisper_secs": 0.0,
            "used_whisper": False, "whisper_model": None,
        }

    picked = _pick_caption_lang(info)

    if picked is not None:
        picked_lang, _is_manual = picked
        ydl_opts = {
            "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
            "merge_output_format": "mp4",
            "outtmpl": str(out_dir / f"{video_id}.%(ext)s"),
            "writesubtitles": True,
            "writeautomaticsub": True,
            "subtitleslangs": [picked_lang],
            "subtitlesformat": "srt/vtt/best",
            "postprocessors": [{"key": "FFmpegSubtitlesConvertor", "format": "srt"}],
            "quiet": True,
            "no_warnings": True,
            **base_opts,
        }
    else:
        # No captions of any kind. Download just the mp4 and transcribe locally.
        if not allow_whisper:
            raise RuntimeError(
                f"No subtitles available for {url} in any language and "
                "Whisper fallback is disabled."
            )
        print("      no captions on YouTube; will transcribe with Whisper after download")
        ydl_opts = {
            "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
            "merge_output_format": "mp4",
            "outtmpl": str(out_dir / f"{video_id}.%(ext)s"),
            "quiet": True,
            "no_warnings": True,
            **base_opts,
        }

    download_t0 = _time.monotonic()
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    download_secs = _time.monotonic() - download_t0

    candidates_mp4 = list(out_dir.glob(f"{video_id}.*"))
    mp4_found = next((p for p in candidates_mp4 if p.suffix in (".mp4", ".mkv", ".webm")), None)
    if mp4_found and mp4_found != mp4_path:
        mp4_found.rename(mp4_path)
    elif not mp4_path.exists():
        raise RuntimeError(f"yt-dlp finished but no video file found in {out_dir}")

    if picked is not None:
        # Re-derive lang from the filename yt-dlp actually wrote (it normalizes
        # codes like en-US -> en, so the picked code may differ from the result).
        srt_found = next(iter(out_dir.glob(f"{video_id}.*.srt")), None)
        if srt_found is None:
            raise RuntimeError(
                f"yt-dlp claimed '{picked[0]}' subtitles for {url} but produced no SRT."
            )
        lang = srt_found.stem[len(video_id) + 1:]
        return {
            "mp4": mp4_path, "srt": srt_found, "lang": lang,
            "title": title, "webpage_url": webpage_url, "upload_date": upload_date, "thumbnail_url": thumbnail_url, "channel_id": channel_id, "channel_name": channel_name, "channel_url": channel_url,
            "download_secs": download_secs, "whisper_secs": 0.0,
            "used_whisper": False, "whisper_model": None,
        }

    whisper_t0 = _time.monotonic()
    srt_path, lang = _transcribe_with_whisper(mp4_path, out_dir, video_id, model_name=whisper_model)
    whisper_secs = _time.monotonic() - whisper_t0
    return {
        "mp4": mp4_path, "srt": srt_path, "lang": lang,
        "title": title, "webpage_url": webpage_url, "upload_date": upload_date, "thumbnail_url": thumbnail_url, "channel_id": channel_id, "channel_name": channel_name, "channel_url": channel_url,
        "download_secs": download_secs, "whisper_secs": whisper_secs,
        "used_whisper": True, "whisper_model": whisper_model,
    }


# ---------- LLM backend abstraction ----------
#
# Two backends:
#   AnthropicAPIBackend   — direct anthropic.Anthropic() SDK calls; uses
#                           messages.parse for structured output and
#                           cache_control for prompt caching. Requires
#                           ANTHROPIC_API_KEY.
#   ClaudeCodeBackend     — shells out to the sandboxed `claude` binary with
#                           --output-format json (and --json-schema for parse).
#                           Uses the user's Claude.ai subscription auth as
#                           configured in our sandbox. No prompt caching
#                           (each subprocess starts cold) and vision support
#                           is opt-in (off by default per user preference).
#
# Both expose: text(...), parse(schema=...), vision_parse(content_blocks=...).
# Returns (response, usage_namespace) where usage has .input_tokens,
# .output_tokens, .cache_read_input_tokens, .cache_creation_input_tokens.
#
# Backends raise VisionUnsupported when vision is requested but unavailable;
# callers fall back to non-vision paths.

from types import SimpleNamespace as _SN


class VisionUnsupported(Exception):
    """Raised when a backend cannot process image inputs in the current config."""


def _zero_usage() -> _SN:
    return _SN(input_tokens=0, output_tokens=0,
               cache_read_input_tokens=0, cache_creation_input_tokens=0)


# ---------- Cost audit / pricing ----------
#
# Per-million-token rates in USD. Used for the cost-transparency layer
# so the user can see per-call and aggregate spend. Update when Anthropic
# adjusts pricing; users can override via the `model_pricing` setting.
# Ranges as of May 2026; treat as estimates.
#
# Cache-read is the discounted price for tokens the API serves from prompt
# cache; cache-creation is the surcharge for the FIRST time the cached
# block is seen. Anthropic publishes both; both matter for our usage shape
# (digest + panel + takeaway re-cite the transcript).

DEFAULT_MODEL_PRICING: dict = {
    # Sonnet 4.6 — digest, takeaway, vision-pick, on-demand panel re-runs
    "claude-sonnet-4-6": {
        "input": 3.0, "output": 15.0,
        "cache_read": 0.30, "cache_creation": 3.75,
    },
    # Opus 4.7 — panel discussion (highest-quality multi-perspective).
    # Rates per platform.claude.com/docs/.../pricing (Opus 4.5+ tier):
    # $5 in / $25 out, 5m cache-write 1.25x ($6.25), cache-read 0.1x ($0.50).
    "claude-opus-4-7": {
        "input": 5.0, "output": 25.0,
        "cache_read": 0.50, "cache_creation": 6.25,
    },
    # Opus 4.8 — same pricing tier as 4.5/4.6/4.7.
    "claude-opus-4-8": {
        "input": 5.0, "output": 25.0,
        "cache_read": 0.50, "cache_creation": 6.25,
    },
    # Haiku 4.5 — slide classifier, validation pings. $1 in / $5 out,
    # 5m cache-write $1.25, cache-read $0.10.
    "claude-haiku-4-5-20251001": {
        "input": 1.0, "output": 5.0,
        "cache_read": 0.10, "cache_creation": 1.25,
    },
    # Older / aliases that may show up in saved settings.json:
    "claude-haiku-4-5": {
        "input": 1.0, "output": 5.0,
        "cache_read": 0.10, "cache_creation": 1.25,
    },
}


def _model_pricing(model: str) -> Optional[dict]:
    """Resolve pricing for a model name. Resolution order:
      1. billing-calibrated cache (~/yt2md/pricing_cache.json, when fresh) —
         merged over the hardcoded fields so a partial calibration keeps the
         rest of the table.
      2. settings['model_pricing'][model] override.
      3. DEFAULT_MODEL_PRICING hardcoded fallback.
    Returns None for a fully-unknown model so estimate_cost_usd surfaces 'n/a'
    rather than a silent $0.
    """
    try:
        cache = load_pricing_cache()
        if cache and model in cache:
            merged = dict(DEFAULT_MODEL_PRICING.get(model) or {})
            merged.update(cache[model])
            if merged:
                return merged
    except Exception:
        pass
    try:
        s = load_settings()
        override = (s.get("model_pricing") or {}).get(model)
        if override:
            return override
    except Exception:
        pass
    return DEFAULT_MODEL_PRICING.get(model)


def estimate_cost_usd(model: str, usage) -> float:
    """Dollar estimate for a single LLM call. Returns 0.0 when pricing is
    unknown for the model — caller can choose to surface that as 'n/a'.
    Token attributes default to 0 if missing (Claude Code backend).
    """
    rates = _model_pricing(model)
    if not rates:
        return 0.0

    def _tok(name: str) -> int:
        return int(getattr(usage, name, 0) or 0)

    return (
        _tok("input_tokens") / 1_000_000 * rates.get("input", 0.0)
        + _tok("output_tokens") / 1_000_000 * rates.get("output", 0.0)
        + _tok("cache_read_input_tokens") / 1_000_000 * rates.get("cache_read", 0.0)
        + _tok("cache_creation_input_tokens") / 1_000_000 * rates.get("cache_creation", 0.0)
    )


def _llm_usage_log_path() -> Path:
    p = get_data_dir() / "logs" / "llm_usage.jsonl"
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def record_llm_usage(
    *,
    video_id: Optional[str],
    kind: str,
    model: str,
    backend_name: str,
    usage,
) -> dict:
    """Append a single usage record to ~/yt2md/logs/llm_usage.jsonl and
    return the recorded dict. Cost is set to 0.0 for the claude-code backend
    (subscription bills the user via their Anthropic plan, not per-call) so
    the audit log stays consistent — token counts still recorded for
    rate-limit awareness.
    """
    import time as _t

    cost = 0.0 if backend_name == "claude-code" else estimate_cost_usd(model, usage)
    entry = {
        "ts": _t.time(),
        "video_id": video_id or "",
        "kind": kind,
        "model": model,
        "backend": backend_name,
        "input_tokens": int(getattr(usage, "input_tokens", 0) or 0),
        "output_tokens": int(getattr(usage, "output_tokens", 0) or 0),
        "cache_read_input_tokens": int(getattr(usage, "cache_read_input_tokens", 0) or 0),
        "cache_creation_input_tokens": int(getattr(usage, "cache_creation_input_tokens", 0) or 0),
        "cost_usd": round(cost, 6),
    }
    try:
        with open(_llm_usage_log_path(), "a") as f:
            f.write(json.dumps(entry) + "\n")
    except OSError:
        pass  # Audit log must never break the digest pipeline.
    return entry


def read_llm_usage_log() -> List[dict]:
    """Load the full usage log into memory. The file stays small in
    practice (a few hundred bytes per entry, one entry per LLM call).
    """
    p = _llm_usage_log_path()
    if not p.exists():
        return []
    rows: List[dict] = []
    try:
        with open(p) as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rows.append(json.loads(line))
                except json.JSONDecodeError:
                    continue
    except OSError:
        pass
    return rows


# ====================================================================
# Admin API: self-calibrating pricing + workspace budget gate
#
# yt2md can read its own billing via the Anthropic Admin API (Usage &
# Cost endpoints; requires an org Admin key, sk-ant-admin...). Two uses:
#   1. Self-calibrating pricing — derive effective $/Mtok from real
#      billing so DEFAULT_MODEL_PRICING can't silently go stale.
#   2. Budget gate — refuse to START a new digest once this workspace's
#      month-to-date spend crosses a threshold (under any Console cap).
# Everything degrades gracefully (returns None / allows the action) when
# no admin key is present, so non-org users are unaffected.
# ====================================================================

ADMIN_API_BASE = "https://api.anthropic.com"
PRICING_CACHE_MAX_AGE_DAYS = 30
_PRICING_DIVERGENCE_GUARD = 3.0    # reject a calibrated rate >3x off the fallback
_PRICING_MIN_TOKENS = 50_000       # too little volume to trust a rate divide
_WS_COST_CACHE_TTL_SECS = 300      # cost data is fresh ~5 min; don't poll faster


def _admin_api_key() -> Optional[str]:
    """Resolve the Admin API key. load_env_files() already fills os.environ
    from repo ./.env and ~/yt2md/.env, so a plain environ read covers the IDE,
    installed, and shell-export workflows. Returns None when absent."""
    key = os.environ.get("ANTHROPIC_ADMIN_KEY")
    if not key:
        try:
            load_env_files()
            key = os.environ.get("ANTHROPIC_ADMIN_KEY")
        except Exception:
            key = None
    return key or None


def _admin_get(path: str, params: Optional[dict] = None) -> Optional[dict]:
    """GET an Admin API endpoint. Returns parsed JSON, or None on any failure
    (no key, HTTP error, network). Never raises — cost features must not break
    the pipeline. params values may be lists (repeated query keys, e.g.
    group_by[])."""
    import urllib.request, urllib.parse
    key = _admin_api_key()
    if not key:
        return None
    url = ADMIN_API_BASE + path
    if params:
        flat = []
        for k, v in params.items():
            if isinstance(v, (list, tuple)):
                flat.extend((k, str(x)) for x in v)
            else:
                flat.append((k, str(v)))
        url += "?" + urllib.parse.urlencode(flat)
    req = urllib.request.Request(
        url, headers={"anthropic-version": "2023-06-01", "x-api-key": key})
    try:
        with urllib.request.urlopen(req, timeout=30) as r:
            return json.load(r)
    except Exception:
        return None


def _admin_get_pages(path: str, params: dict) -> Optional[List[dict]]:
    """Fetch all pages of a list endpoint. Returns the concatenated `data`
    buckets, [] when reachable-but-empty, or None when the API is unreachable
    (so callers can distinguish 'zero spend' from 'no admin key')."""
    first = _admin_get(path, params)
    if first is None:
        return None
    out: List[dict] = list(first.get("data", []) or [])
    page = first
    p = dict(params)
    for _ in range(50):  # defensive hard cap
        if not page.get("has_more"):
            break
        nxt = page.get("next_page")
        if not nxt:
            break
        p["page"] = nxt
        page = _admin_get(path, p)
        if page is None:
            break
        out.extend(page.get("data", []) or [])
    return out


def _utc_iso(dt) -> str:
    return dt.strftime("%Y-%m-%dT%H:%M:%SZ")


def _utc_now_iso() -> str:
    from datetime import datetime, timezone
    return _utc_iso(datetime.now(timezone.utc))


def _utc_month_start_iso() -> str:
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc)
    return _utc_iso(now.replace(day=1, hour=0, minute=0, second=0, microsecond=0))


def _utc_days_ago_iso(days: int) -> str:
    from datetime import datetime, timezone, timedelta
    return _utc_iso(datetime.now(timezone.utc) - timedelta(days=days))


def _norm_model(s: str) -> str:
    """Normalize a model id ('claude-opus-4-7', 'claude-haiku-4-5-20251001')
    or a billing display name ('Claude Opus 4.7') to a join key. Strips a
    trailing -YYYYMMDD date and all non-alphanumerics: both forms collapse to
    e.g. 'claudeopus47'."""
    s = re.sub(r"-20\d{6}$", "", s)
    return re.sub(r"[^a-z0-9]", "", s.lower())


# ---- Self-calibrating pricing ----------------------------------------

def fetch_usage_by_model(start_iso: str, end_iso: str) -> Optional[dict]:
    """Token counts grouped by model over a window:
    {model_id: {input, output, cache_read, cache_creation}}. None if API down."""
    pages = _admin_get_pages(
        "/v1/organizations/usage_report/messages",
        {"starting_at": start_iso, "ending_at": end_iso,
         "group_by[]": ["model"], "bucket_width": "1d"})
    if pages is None:
        return None
    tok: dict = {}
    for bucket in pages:
        for r in bucket.get("results", []):
            mdl = r.get("model") or "unknown"
            d = tok.setdefault(mdl, {"input": 0, "output": 0,
                                     "cache_read": 0, "cache_creation": 0})
            d["input"] += int(r.get("uncached_input_tokens", 0) or 0)
            d["output"] += int(r.get("output_tokens", 0) or 0)
            d["cache_read"] += int(r.get("cache_read_input_tokens", 0) or 0)
            d["cache_creation"] += int(r.get("cache_creation_input_tokens", 0) or 0)
    return tok


def fetch_cost_by_description(start_iso: str, end_iso: str) -> Optional[dict]:
    """Billed USD grouped by description (model + token type) over a window:
    {description: usd}. None if API down."""
    pages = _admin_get_pages(
        "/v1/organizations/cost_report",
        {"starting_at": start_iso, "ending_at": end_iso,
         "group_by[]": ["description"], "bucket_width": "1d"})
    if pages is None:
        return None
    out: dict = {}
    for bucket in pages:
        for item in bucket.get("results", []):
            desc = item.get("description") or item.get("model") or "unknown"
            out[desc] = out.get(desc, 0.0) + float(item.get("amount", "0") or 0) / 100.0
    return out


# Billing-description substrings -> pricing field. Order matters: the
# cache-write line reads "Input Tokens, Cache Write", so "cache write" must
# be tested before the bare "input tokens".
_COST_FIELD_NEEDLES = [
    ("cache write", "cache_creation"),
    ("cache read", "cache_read"),
    ("cache hit", "cache_read"),
    ("output tokens", "output"),
    ("input tokens", "input"),
]


def calibrate_pricing_from_billing(lookback_days: int = 14) -> dict:
    """Derive effective $/Mtok per (model, field) from real billing and write
    ~/yt2md/pricing_cache.json. Joins usage tokens to cost dollars by a
    normalized model key. Skips lines with too little volume or a rate that
    diverges >3x from the hardcoded fallback (guards a mapping bug). Returns a
    summary dict with ok / derived / warnings."""
    start, end = _utc_days_ago_iso(lookback_days), _utc_now_iso()
    usage = fetch_usage_by_model(start, end)
    cost = fetch_cost_by_description(start, end)
    if usage is None or cost is None:
        return {"ok": False, "reason": "admin API unavailable (no key or error)"}

    usage_norm = {_norm_model(mid): {**d, "_id": mid} for mid, d in usage.items()}
    derived: dict = {}
    warnings: List[str] = []
    for desc, usd in cost.items():
        dl = desc.lower()
        field = next((f for needle, f in _COST_FIELD_NEEDLES if needle in dl), None)
        if not field:
            continue
        disp = re.split(r"\s+usage", dl)[0]
        u = usage_norm.get(_norm_model(disp))
        if not u:
            continue
        toks = u.get(field, 0)
        if toks < _PRICING_MIN_TOKENS:
            continue
        rate = usd / (toks / 1_000_000)
        mid = u["_id"]
        fb = (DEFAULT_MODEL_PRICING.get(mid) or {}).get(field)
        if fb and (rate > fb * _PRICING_DIVERGENCE_GUARD
                   or rate < fb / _PRICING_DIVERGENCE_GUARD):
            warnings.append(
                f"{mid}.{field}: billing-implied ${rate:.2f}/M diverges >3x "
                f"from fallback ${fb}/M — skipped")
            continue
        derived.setdefault(mid, {})[field] = round(rate, 4)

    # Merge over any previous cache so a partial pull never wipes known rates.
    prev: dict = {}
    p = _pricing_cache_path()
    if p.exists():
        try:
            prev = (json.loads(p.read_text()).get("rates")) or {}
        except Exception:
            prev = {}
    rates = dict(prev)
    for mid, fields in derived.items():
        rates.setdefault(mid, {}).update(fields)

    import time as _t
    payload = {"_source": "anthropic-billing", "_calibrated_at": _utc_now_iso(),
               "_calibrated_at_epoch": _t.time(), "_lookback_days": lookback_days,
               "rates": rates}
    try:
        get_data_dir().mkdir(parents=True, exist_ok=True)
        p.write_text(json.dumps(payload, indent=2) + "\n")
    except OSError:
        pass
    return {"ok": True, "derived": derived, "warnings": warnings,
            "models": list(derived)}


def _pricing_cache_path() -> Path:
    return get_data_dir() / "pricing_cache.json"


def load_pricing_cache() -> Optional[dict]:
    """{model: {field: rate}} from the billing-calibrated cache when present
    and fresh (< PRICING_CACHE_MAX_AGE_DAYS), else None."""
    p = _pricing_cache_path()
    if not p.exists():
        return None
    try:
        data = json.loads(p.read_text())
    except Exception:
        return None
    import time as _t
    ts = data.get("_calibrated_at_epoch", 0)
    if not ts or (_t.time() - ts) > PRICING_CACHE_MAX_AGE_DAYS * 86400:
        return None
    return data.get("rates") or None


# ---- Workspace budget gate -------------------------------------------

def detect_runtime_workspace_id() -> Optional[str]:
    """Match the runtime ANTHROPIC_API_KEY to its workspace via the Admin API
    (by partial_key_hint suffix). Returns the workspace id, or None."""
    rk = os.environ.get("ANTHROPIC_API_KEY") or ""
    if len(rk) < 8:
        return None
    tail = rk[-4:]
    keys = _admin_get_pages("/v1/organizations/api_keys", {"limit": 100})
    if not keys:
        return None
    for k in keys:
        if (k.get("partial_key_hint") or "").endswith(tail):
            return k.get("workspace_id")
    return None


def fetch_workspace_month_to_date_usd(workspace_id: str) -> Optional[float]:
    """Month-to-date billed USD for one workspace (group_by workspace_id, sum
    the matching rows). None if the API is unreachable; 0.0 if reachable with
    no spend."""
    pages = _admin_get_pages(
        "/v1/organizations/cost_report",
        {"starting_at": _utc_month_start_iso(), "ending_at": _utc_now_iso(),
         "group_by[]": ["workspace_id"], "bucket_width": "1d"})
    if pages is None:
        return None
    total = 0.0
    for bucket in pages:
        for item in bucket.get("results", []):
            if item.get("workspace_id") == workspace_id:
                total += float(item.get("amount", "0") or 0) / 100.0
    return total


_ws_cost_cache = {"ts": 0.0, "usd": None, "ws": None}


def _local_month_to_date_cost() -> float:
    """Fallback: sum the local usage log for the current UTC month."""
    from datetime import datetime, timezone
    month_start = datetime.now(timezone.utc).replace(
        day=1, hour=0, minute=0, second=0, microsecond=0).timestamp()
    return sum(float(r.get("cost_usd") or 0)
               for r in read_llm_usage_log() if (r.get("ts") or 0) >= month_start)


def workspace_month_to_date_cost(workspace_id: Optional[str]) -> Optional[float]:
    """Authoritative month-to-date USD for a workspace (Admin Cost API, cached
    ~5 min). Falls back to the local usage log when the API is unreachable.
    Returns None only if there is no workspace and no local data."""
    import time as _t
    if workspace_id:
        c = _ws_cost_cache
        if (c["ws"] == workspace_id and c["usd"] is not None
                and (_t.time() - c["ts"]) < _WS_COST_CACHE_TTL_SECS):
            return c["usd"]
        usd = fetch_workspace_month_to_date_usd(workspace_id)
        if usd is not None:
            _ws_cost_cache.update(ts=_t.time(), usd=usd, ws=workspace_id)
            return usd
    return _local_month_to_date_cost()


def budget_status() -> dict:
    """Resolve workspace, month-to-date spend, and the warn/block thresholds.
    Auto-detects + persists budget_workspace_id on first call when possible."""
    s = load_settings()
    ws = s.get("budget_workspace_id") or None
    if not ws:
        ws = detect_runtime_workspace_id()
        if ws:
            try:
                s["budget_workspace_id"] = ws
                save_settings(s)
            except Exception:
                pass
    return {
        "workspace_id": ws,
        "month_to_date_usd": workspace_month_to_date_cost(ws),
        "warn_usd": float(s.get("budget_warn_usd") or 0),
        "block_usd": float(s.get("budget_block_usd") or 0),
        "source": "billing" if ws and _admin_api_key() else "local",
    }


def check_budget(action: str = "start a new digest") -> Optional[str]:
    """Return a human message if the action should be BLOCKED (mtd >= block),
    else None. Warns to stderr in the warn..block band. Always allows when
    budget is unconfigured or month-to-date cost is unknown."""
    try:
        st = budget_status()
    except Exception:
        return None
    mtd, warn, block = st.get("month_to_date_usd"), st.get("warn_usd") or 0, st.get("block_usd") or 0
    if mtd is None:
        return None
    if block and mtd >= block:
        return (f"yt2md budget gate: workspace month-to-date spend is "
                f"${mtd:.2f} >= block threshold ${block:.2f}. Refusing to {action}. "
                f"Raise budget_block_usd in settings.json or wait until next month "
                f"(the Console hard cap still applies as a backstop).")
    if warn and mtd >= warn:
        sys.stderr.write(f"[yt2md] budget warning: workspace month-to-date "
                         f"${mtd:.2f} >= warn ${warn:.2f} (block at ${block:.2f}).\n")
    return None


def cmd_refresh_pricing(args) -> int:
    """yt2md refresh-pricing — recompute the model price table from your actual
    Anthropic billing (Admin API) and cache it to ~/yt2md/pricing_cache.json."""
    res = calibrate_pricing_from_billing(lookback_days=args.lookback)
    if not res.get("ok"):
        print(f"could not calibrate: {res.get('reason')}")
        print("(set ANTHROPIC_ADMIN_KEY in the repo .env or ~/yt2md/.env)")
        return 1
    if not res["derived"]:
        print("no model lines had enough billing volume to calibrate "
              f"(last {args.lookback}d); table unchanged.")
    else:
        print(f"calibrated {len(res['derived'])} model(s) from billing "
              f"(last {args.lookback}d) -> {_pricing_cache_path()}")
        for mid, fields in res["derived"].items():
            print(f"  {mid}: " + ", ".join(f"{k}=${v}/M" for k, v in fields.items()))
    for w in res.get("warnings", []):
        print(f"  ! {w}")
    return 0


class AnthropicAPIBackend:
    name = "api"
    vision_supported = True

    def __init__(self):
        import anthropic
        self._client = anthropic.Anthropic()

    def text(self, *, system: str, user_text: str, model: str,
             max_tokens: int, cache: bool = False):
        block = {"type": "text", "text": user_text}
        if cache:
            block["cache_control"] = {"type": "ephemeral"}
        response = self._client.messages.create(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": [block]}],
        )
        text = next(b.text for b in response.content if b.type == "text")
        return text, response.usage

    def parse(self, *, system: str, user_text: str, model: str,
              max_tokens: int, schema, cache: bool = False):
        block = {"type": "text", "text": user_text}
        if cache:
            block["cache_control"] = {"type": "ephemeral"}
        response = self._client.messages.parse(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": [block]}],
            output_format=schema,
        )
        return response.parsed_output, response.usage

    def vision_parse(self, *, system: str, content_blocks: list,
                     model: str, max_tokens: int, schema):
        response = self._client.messages.parse(
            model=model, max_tokens=max_tokens, system=system,
            messages=[{"role": "user", "content": content_blocks}],
            output_format=schema,
        )
        return response.parsed_output, response.usage


class ClaudeCodeBackend:
    name = "claude-code"

    def __init__(self, *, vision_enabled: bool = False):
        if not claude_code_installed():
            raise RuntimeError("Claude Code is not installed in the sandbox.")
        self._binary = str(claude_binary_path())
        self._env = claude_subprocess_env()
        self._vision_enabled = vision_enabled

    @property
    def vision_supported(self) -> bool:
        return self._vision_enabled

    def _run(self, *, prompt: str, model: str, schema=None,
             timeout: float = 600.0):
        cmd = [self._binary, "-p", prompt, "--model", model,
               "--output-format", "json"]
        if schema is not None:
            cmd += ["--json-schema", json.dumps(schema.model_json_schema())]
        proc = subprocess.run(
            cmd, env=self._env, capture_output=True, text=True, timeout=timeout,
        )
        if proc.returncode != 0:
            raise RuntimeError(
                f"claude -p failed (exit {proc.returncode}): "
                f"{(proc.stderr or proc.stdout or '').strip()[:500]}"
            )
        return self._parse_output(proc.stdout)

    @staticmethod
    def _parse_output(stdout: str):
        """`claude -p --output-format json` emits a JSON envelope. Extract the
        text/result and a usage namespace. Tolerant of envelope shape changes.
        Raises RuntimeError when the envelope's `is_error` flag is set (which
        Claude Code uses for auth + tool failures even when the process exits 0).
        """
        try:
            payload = json.loads(stdout)
        except json.JSONDecodeError:
            # Defensive fallback: treat raw stdout as the result.
            return stdout.strip(), _zero_usage()
        result_text = (
            payload.get("result")
            or payload.get("text")
            or payload.get("response")
            or ""
        )
        if payload.get("is_error"):
            raise RuntimeError(f"claude reported error: {result_text or '(no message)'}")
        usage_dict = payload.get("usage") or {}
        usage = _SN(
            input_tokens=int(usage_dict.get("input_tokens") or 0),
            output_tokens=int(usage_dict.get("output_tokens") or 0),
            cache_read_input_tokens=int(usage_dict.get("cache_read_input_tokens") or 0),
            cache_creation_input_tokens=int(
                usage_dict.get("cache_creation_input_tokens") or 0
            ),
        )
        return result_text, usage

    def text(self, *, system: str, user_text: str, model: str,
             max_tokens: int, cache: bool = False):
        # No system-prompt CLI flag we rely on — combine system + user into
        # one prompt with explicit delineation. cache is ignored (subprocess
        # invocations don't share Anthropic's prompt cache).
        prompt = f"<system>\n{system}\n</system>\n\n{user_text}"
        return self._run(prompt=prompt, model=model)

    def parse(self, *, system: str, user_text: str, model: str,
              max_tokens: int, schema, cache: bool = False):
        prompt = f"<system>\n{system}\n</system>\n\n{user_text}"
        result_text, usage = self._run(prompt=prompt, model=model, schema=schema)
        # --json-schema constrains the output to the schema; the result string
        # IS the JSON we need to parse into the Pydantic model.
        try:
            data = json.loads(result_text) if isinstance(result_text, str) else result_text
        except json.JSONDecodeError as e:
            raise RuntimeError(
                f"claude returned non-JSON for a schema-constrained call: {e}\n"
                f"Output: {result_text[:500]}"
            )
        return schema.model_validate(data), usage

    def vision_parse(self, *, system: str, content_blocks: list,
                     model: str, max_tokens: int, schema):
        if not self._vision_enabled:
            raise VisionUnsupported(
                "Claude Code vision is disabled. Enable claude_code_vision in "
                "Settings to base64-embed images in prompts (token-heavy)."
            )
        # Opt-in path: serialize image blocks as base64 markers in the prompt.
        # This is expensive — each image expands ~33% over its byte size and
        # there's no native multipart in the CLI -p mode.
        text_parts: list = []
        for block in content_blocks:
            if block.get("type") == "text":
                text_parts.append(block["text"])
            elif block.get("type") == "image":
                src = block.get("source") or {}
                if src.get("type") == "base64":
                    text_parts.append(
                        f"[image media_type={src.get('media_type','image/jpeg')} "
                        f"data:base64]\n{src.get('data','')}\n[/image]"
                    )
        prompt = (
            f"<system>\n{system}\n</system>\n\n" + "\n".join(text_parts)
        )
        result_text, usage = self._run(prompt=prompt, model=model, schema=schema)
        try:
            data = json.loads(result_text) if isinstance(result_text, str) else result_text
        except json.JSONDecodeError as e:
            raise RuntimeError(
                f"claude returned non-JSON for vision call: {e}\nOutput: {result_text[:500]}"
            )
        return schema.model_validate(data), usage


def _resolve_claude_binary_for_pty() -> Tuple[Optional[str], Optional[str]]:
    """Find a `claude` on PATH whose --help advertises the flags the PTY
    backend needs (--tools, --system-prompt, --disable-slash-commands). Older
    npm-installed 1.x lacks --tools entirely; modern 2.x has it. Returns
    (binary_path, None) on success, or (None, error_message) on failure."""
    seen: list = []
    for entry in (os.environ.get("PATH") or "").split(os.pathsep):
        if not entry:
            continue
        candidate = os.path.join(entry, "claude")
        if not os.path.isfile(candidate) or not os.access(candidate, os.X_OK):
            continue
        real = os.path.realpath(candidate)
        if real in seen:
            continue
        seen.append(real)
        try:
            help_out = subprocess.run(
                [candidate, "--help"],
                capture_output=True, text=True, timeout=15,
            ).stdout
        except (OSError, subprocess.TimeoutExpired):
            continue
        if "--tools" in help_out and "--system-prompt" in help_out:
            return candidate, None
    if not seen:
        return None, (
            "claude binary not found on PATH. Install Claude Code "
            "(https://docs.claude.com/en/docs/claude-code/quickstart) "
            "and sign in to your Pro/Max plan."
        )
    return None, (
        f"None of the {len(seen)} `claude` install(s) on PATH support the "
        "flags the PTY backend needs (--tools, --system-prompt). Candidates: "
        + ", ".join(seen)
        + ". Upgrade via `claude install` so the modern 2.x version is found "
        "first on PATH."
    )


class ClaudeCodePtyBackend:
    """Drives the sandboxed `claude` binary as an interactive REPL via a PTY,
    instead of `claude -p`. Why: starting June 15 2026 Anthropic bills `-p` and
    Agent SDK usage from a separate metered "Agent SDK credit" pool, while
    *interactive* session usage continues to count against the user's Pro/Max
    plan. Driving the REPL programmatically keeps the bill on the subscription.

    Mechanics: spawn `claude` (no `-p`) under a PTY, wait for the input box to
    settle, send the prompt via bracketed-paste so the REPL treats long inputs
    as one atomic paste, then \\r to submit. Read stdout through a pyte virtual
    terminal so the TUI's cursor-positioning ANSI is rendered to a clean
    scrollable screen. Extract the response between the `⏺ ` marker and the
    `✻ Brewed`/next-input-prompt footer.

    Caveats:
    - No token-usage envelope in interactive mode; usage counts come back zero.
      That's accurate for cost-audit purposes — subscription billing has no
      per-call dollar amount.
    - Vision is not supported (interactive mode has no clean image-input path).
    - The TUI is a moving target; pin the claude binary version and re-verify
      the screen markers on upgrade. End-of-response detection is via "✻ Brewed"
      / "✻ Cooking" / input-box redraw — Anthropic could change these.
    - One PTY = one inflight request. Caller must serialize.
    """
    name = "claude-code-pty"
    vision_supported = False

    # Screen geometry. Wide because pyte preserves line wraps as real
    # newlines — for parse() that breaks JSON strings whose values exceed the
    # terminal width. Tall + history covers long-form responses (panel ≈ 2k
    # words). We read both display + history when extracting.
    _COLS = 1200
    _ROWS = 500
    _HISTORY = 20000

    def __init__(self, *, vision_enabled: bool = False):
        # Resolve the user's PRIMARY claude install (logged into claude.ai
        # Pro/Max), not the sandboxed one. The sandbox runs -p with an API
        # key fallback; the PTY backend's whole point is to hit the user's
        # subscription, which is OAuth-authenticated in their primary config.
        # Probe candidates because users can have multiple `claude` installs
        # on PATH (e.g. an old npm-installed 1.x via nvm AND a current 2.x
        # via `claude install`) — pick the first one whose --help mentions
        # the flags we depend on.
        binary, why = _resolve_claude_binary_for_pty()
        if not binary:
            raise RuntimeError(why)
        try:
            import pyte  # noqa: F401
        except ImportError as e:
            raise RuntimeError(
                "pyte is required for the claude-code-pty backend. Run `uv sync`."
            ) from e
        self._binary = binary
        # Inherit the user's environment so OAuth credentials and default
        # CLAUDE_CONFIG_DIR (~/.claude) are visible. Explicitly strip
        # ANTHROPIC_API_KEY so the REPL can't silently fall back to API
        # billing if OAuth login is missing.
        env = os.environ.copy()
        env.pop("ANTHROPIC_API_KEY", None)
        self._env = env
        # CWD: a yt2md-controlled directory under the data dir so the REPL's
        # workspace-trust state is isolated from the user's real projects.
        workdir = get_data_dir() / "claude-pty-workdir"
        workdir.mkdir(parents=True, exist_ok=True)
        self._cwd = str(workdir)

    def text(self, *, system: str, user_text: str, model: str,
             max_tokens: int, cache: bool = False):
        # max_tokens, cache: ignored — no equivalents in interactive mode.
        text, usage = self._run(system=system, user_text=user_text, model=model)
        return text, usage

    def parse(self, *, system: str, user_text: str, model: str,
              max_tokens: int, schema, cache: bool = False):
        # No --json-schema flag in interactive mode; instruct via system prompt
        # and validate on the way out.
        schema_json = json.dumps(schema.model_json_schema())
        sys_with_schema = (
            system
            + "\n\nReply with ONLY a single JSON object matching this schema. "
            "Output raw JSON only — no markdown fences, no prose before or after.\n"
            + schema_json
        )
        text, usage = self._run(
            system=sys_with_schema, user_text=user_text, model=model,
        )
        m = re.search(r"\{.*\}", text, re.DOTALL)
        if not m:
            raise RuntimeError(
                f"claude PTY returned non-JSON for a schema-constrained call. "
                f"Output: {text[:500]}"
            )
        raw = m.group(0)
        try:
            return schema.model_validate(json.loads(raw)), usage
        except json.JSONDecodeError:
            # pyte hard-wraps long lines into separate screen rows, which
            # surfaces as raw newlines inside JSON string values — invalid by
            # spec. Fall back to collapsing all whitespace runs to a single
            # space; JSON parsers treat whitespace as insignificant outside
            # strings, and yt2md schemas use short string values where
            # collapsing inside-string whitespace is acceptable.
            collapsed = re.sub(r"\s+", " ", raw)
            try:
                return schema.model_validate(json.loads(collapsed)), usage
            except json.JSONDecodeError as e:
                raise RuntimeError(
                    f"claude PTY returned invalid JSON even after wrap-aware "
                    f"collapse: {e}\nOutput: {raw[:800]}"
                )

    def vision_parse(self, *, system: str, content_blocks: list,
                     model: str, max_tokens: int, schema):
        raise VisionUnsupported(
            "claude-code-pty backend does not support vision input. Use the api "
            "backend, or switch to claude-code with claude_code_vision enabled."
        )

    def _run(self, *, system: str, user_text: str, model: str,
             timeout: float = 600.0):
        import pty
        import select
        import fcntl
        import termios
        import struct
        import time
        import pyte

        # Use --tools="" (equals form) instead of --tools "" — commander.js's
        # variadic <tools...> parser intermittently rejects the empty positional
        # form with `error: unknown option '--tools'`, which then fails our boot
        # detection. The equals form binds the empty value directly.
        cmd = [
            self._binary,
            "--model", model,
            "--system-prompt", system,
            "--tools=",
            "--disable-slash-commands",
        ]

        screen = pyte.HistoryScreen(self._COLS, self._ROWS, history=self._HISTORY)
        stream = pyte.ByteStream(screen)

        master, slave = pty.openpty()
        fcntl.ioctl(
            master, termios.TIOCSWINSZ,
            struct.pack("HHHH", self._ROWS, self._COLS, 0, 0),
        )
        env = {
            **self._env,
            "TERM": "xterm-256color",
            "COLUMNS": str(self._COLS),
            "LINES": str(self._ROWS),
        }
        proc = subprocess.Popen(
            cmd, stdin=slave, stdout=slave, stderr=slave,
            close_fds=True, cwd=self._cwd, env=env,
        )
        os.close(slave)

        def drain():
            while True:
                r, _, _ = select.select([master], [], [], 0.1)
                if not r:
                    return
                try:
                    chunk = os.read(master, 8192)
                except OSError:
                    return
                if not chunk:
                    return
                stream.feed(chunk)

        def screen_text() -> str:
            buf: list = []
            top = getattr(screen.history, "top", None)
            if top is not None:
                for line in top:
                    buf.append("".join(c.data for c in line).rstrip())
            for line in screen.display:
                buf.append(line.rstrip())
            return "\n".join(buf).rstrip()

        try:
            # Boot phase: wait for the input prompt to settle. The sandboxed
            # claude may show first-run onboarding modals (theme picker,
            # workspace trust, welcome splash) if it's never been driven
            # interactively before — dismiss each by pressing Enter (the
            # default-selected option is always the safe choice). Loop until
            # we see the input-box footer.
            boot_deadline = time.time() + 45.0
            last_dismiss_at = 0.0
            booted = False
            while time.time() < boot_deadline:
                drain()
                t = screen_text()
                if "for shortcuts" in t or "for agents" in t:
                    booted = True
                    break
                # Modal heuristic: a numbered-options dialog with a "❯" pointer
                # awaiting selection. The default highlighted option is what
                # we'd choose interactively (trust=Yes, theme=Dark mode, etc.).
                # Rate-limit our Enter presses so we don't double-submit.
                looks_like_modal = (
                    "❯" in t
                    and re.search(r"\b1\.\s", t) is not None
                    and ("confirm" in t.lower() or "select" in t.lower()
                         or "choose" in t.lower() or "trust" in t.lower()
                         or "let's get started" in t.lower())
                )
                if looks_like_modal and time.time() - last_dismiss_at > 1.5:
                    os.write(master, b"\r")
                    last_dismiss_at = time.time()
                    time.sleep(0.8)
                    continue
                time.sleep(0.25)
            if not booted:
                raise RuntimeError(
                    f"claude PTY did not reach an interactive prompt within 45s. "
                    f"Last screen tail:\n{screen_text()[-800:]}"
                )
            time.sleep(0.4)
            drain()

            # Send the prompt via bracketed paste, then \r to submit. Without
            # bracketed-paste markers, long inputs get chunked by kernel writes
            # and the REPL treats each chunk as a separate paste.
            BPS, BPE = b"\x1b[200~", b"\x1b[201~"
            payload = BPS + user_text.encode("utf-8") + BPE
            for i in range(0, len(payload), 1024):
                os.write(master, payload[i:i + 1024])
                time.sleep(0.005)
            time.sleep(0.3)
            os.write(master, b"\r")

            # Wait for the response to settle. End-of-response signal is the
            # `✻ <verb> for <duration>` footer Claude Code emits AFTER a
            # response finishes (verb rotates: "Brewed", "Cogitated",
            # "Cooking", "Baked", "Pondered" — match on shape, not the word).
            # Duration switches format past 60s: "9s" / "1m 0s" / "1h 5m 30s"
            # — match one-or-more "<digits><unit-letter>" chunks separated by
            # whitespace. Earlier verb-only match missed everything > 60s,
            # which is every digest/panel/takeaway call on a real video.
            done_re = re.compile(r"✻\s+\w+\s+for\s+\d+\w(?:\s+\d+\w)*")
            deadline = time.time() + timeout
            seen_marker = False
            while time.time() < deadline:
                drain()
                t = screen_text()
                if "⏺" in t:
                    seen_marker = True
                lower = t.lower()
                if "rate limit" in lower or "usage limit" in lower:
                    raise RuntimeError(
                        f"Claude Code reported a usage/rate limit. Tail:\n"
                        f"{t[-500:]}"
                    )
                if seen_marker and done_re.search(t):
                    time.sleep(0.5)
                    drain()
                    break
                time.sleep(0.25)

            response = self._extract_response(screen_text())
            if not response:
                raise RuntimeError(
                    f"claude PTY produced no response within {timeout:.0f}s. "
                    f"Final screen tail:\n{screen_text()[-600:]}"
                )
            return response, _zero_usage()
        finally:
            # Ctrl-C twice is Claude Code's escape to exit.
            try:
                os.write(master, b"\x03")
                time.sleep(0.2)
                os.write(master, b"\x03")
                time.sleep(0.2)
            except OSError:
                pass
            proc.terminate()
            try:
                proc.wait(timeout=3)
            except subprocess.TimeoutExpired:
                proc.kill()
            try:
                os.close(master)
            except OSError:
                pass

    @staticmethod
    def _extract_response(text: str) -> str:
        """Pull the assistant's text out of the rendered TUI. Response lives
        between the LAST `⏺ ` marker and the subsequent `✻` footer / next `❯`
        input box / horizontal rule.
        """
        lines = text.splitlines()
        start: Optional[int] = None
        for i in range(len(lines) - 1, -1, -1):
            if lines[i].lstrip().startswith("⏺"):
                start = i
                break
        if start is None:
            return ""
        out: list = []
        for line in lines[start:]:
            stripped = line.lstrip()
            if stripped.startswith("✻") or stripped.startswith("❯"):
                break
            if stripped.startswith("─") and out:
                break
            if stripped.startswith("⏺"):
                line = stripped[1:].lstrip()
            out.append(line)
        return "\n".join(l for l in out if l.strip()).strip()


def select_backend(*, vision_enabled: Optional[bool] = None,
                   for_vision: bool = False):
    """Resolve the active LLM backend from settings + environment.

    Honors settings["llm_backend"] in {"auto", "api", "claude-code", "claude-code-pty"}:
      - "auto":            prefer "api" when ANTHROPIC_API_KEY is set, else
                           "claude-code" when sandboxed claude is installed and
                           logged in, else raises RuntimeError.
      - "api":             requires ANTHROPIC_API_KEY.
      - "claude-code":     requires sandbox install + login. Uses `claude -p`;
                           after Jun 15 2026 this bills against the Agent SDK
                           credit pool, not the Pro/Max subscription.
      - "claude-code-pty": requires `claude` on PATH + OAuth Pro/Max login.
                           Drives the REPL via a PTY so text/parse calls stay
                           on subscription billing. No native vision; vision
                           callers should pass for_vision=True to opportunistically
                           route those calls to the API backend (preserves frame
                           quality while keeping text/panel/etc. on subscription).

    for_vision=True: caller is about to make a vision_parse() call. If the
    primary backend can't do vision and ANTHROPIC_API_KEY is set, transparently
    return AnthropicAPIBackend() instead — implements the hybrid PTY-for-text +
    API-for-vision routing without exposing a new user-facing setting. If no
    API key is available, returns the primary backend anyway and lets the
    caller's VisionUnsupported fallback handle it.
    """
    s = load_settings()
    choice = (s.get("llm_backend") or "auto").lower()
    if vision_enabled is None:
        vision_enabled = bool(s.get("claude_code_vision", False))

    # Hybrid short-circuit: PTY mode has no vision path; route image calls
    # through the API when a key is available. The text/parse split keeps the
    # vast majority of token volume on the subscription while preserving the
    # frame-picking quality that timestamp-based fallback can't match.
    if for_vision and choice == "claude-code-pty" and os.environ.get("ANTHROPIC_API_KEY"):
        return AnthropicAPIBackend()

    if choice == "api":
        if not os.environ.get("ANTHROPIC_API_KEY"):
            raise RuntimeError(
                "llm_backend=api but ANTHROPIC_API_KEY is not set. "
                "Configure it in Settings or switch backend to claude-code."
            )
        return AnthropicAPIBackend()

    if choice == "claude-code":
        if not claude_code_installed():
            raise RuntimeError(
                "llm_backend=claude-code but Claude Code is not installed. "
                "Run /setup to install it."
            )
        return ClaudeCodeBackend(vision_enabled=vision_enabled)

    if choice == "claude-code-pty":
        # Doesn't use the yt2md sandbox — uses the user's primary `claude` so
        # OAuth (Pro/Max) auth applies. ClaudeCodePtyBackend.__init__ raises
        # a descriptive error if `claude` is not on PATH.
        return ClaudeCodePtyBackend(vision_enabled=vision_enabled)

    # auto: prefer API when key is set (keeps prompt caching + native vision);
    # else fall through to Claude Code only when both installed AND we have a
    # cached login signal (sentinel). Otherwise raise so the caller redirects
    # to /setup rather than burning a doomed call.
    if os.environ.get("ANTHROPIC_API_KEY"):
        return AnthropicAPIBackend()
    if (claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)):
        return ClaudeCodeBackend(vision_enabled=vision_enabled)
    raise RuntimeError(
        "No LLM backend configured. Set ANTHROPIC_API_KEY or install + sign "
        "in to Claude Code via /setup."
    )


# ---------- Markdown digest (transcript-primary, LLM-summarized) ----------

DIGEST_SYSTEM_PROMPT = (
    "You distill video transcripts into readable digests so a reader can grasp "
    "the essence of a video without watching it.\n\n"
    "Given a timestamped transcript, identify the natural topic segments. "
    "Return 5–12 sections depending on length and content density (favor fewer "
    "for short videos, more for long talks).\n\n"
    "For each topic:\n"
    "- title: a short, descriptive heading (not a question, not a teaser)\n"
    "- start_time: the timestamp in SECONDS where the topic begins, taken from "
    "the transcript's bracketed timestamps\n"
    "- summary: 2–4 sentences of informative prose distilling what is actually "
    "said. Concrete claims, names, numbers, and reasoning — not vague paraphrases.\n"
    "- key_points: 3–6 short bullet points capturing the most useful takeaways. "
    "Bullets should add detail beyond the summary, not restate it.\n\n"
    "Also produce a 2–3 sentence overview of the entire video.\n\n"
    "Write for a reader, not a viewer. Skip filler ('in this video I'll show you'); "
    "go straight to the substance."
)


def generate_digest(
    segments: List[TranscriptSegment],
    video_title: str,
    model: str,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Call Claude to segment the transcript into topics. Returns a parsed VideoDigest.

    source_lang is the BCP-47 language code of the transcript (e.g. 'en', 'zh-Hans').
    output_language: 'auto' (write in source language) or 'en' (force English).
    backend: an LLMBackend; defaults to select_backend() (auto-resolved).
    """
    from pydantic import BaseModel
    from typing import List as TList

    class Topic(BaseModel):
        title: str
        start_time: float
        summary: str
        key_points: TList[str]

    class VideoDigest(BaseModel):
        title: str
        overview: str
        topics: TList[Topic]

    transcript = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_note = ""
    if not is_english_source:
        if output_language == "en":
            lang_note = (
                f"NOTE: This transcript is in language code '{source_lang}', not English. "
                "Translate to English while distilling. Title, overview, topic titles, "
                "summaries, and key points must all be written in English regardless of "
                "the source language. Preserve proper nouns (people, places, products) in "
                "their original form when there is no established English rendering.\n\n"
            )
        else:  # "auto" — match the source language
            lang_note = (
                f"NOTE: This transcript is in language code '{source_lang}'. Write the "
                "digest in the SAME language as the transcript. Title, overview, topic "
                "titles, summaries, and key points must all be in the source language. "
                "Preserve proper nouns and established technical terms in their original "
                "form (including English technical jargon when the field uses it that way).\n\n"
            )

    user_text = (
        f"{lang_note}"
        f"Video title: {video_title}\n\n"
        f"Total duration: {format_timestamp(segments[-1].end if segments else 0)}\n\n"
        f"Timestamped transcript:\n\n{transcript}"
    )

    if backend is None:
        backend = select_backend()
    return backend.parse(
        system=DIGEST_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=16000, schema=VideoDigest,
    )


DEFAULT_PANEL_MODEL = "claude-opus-4-7"


PANEL_SYSTEM_PROMPT = (
    "You facilitate a panel of domain experts critically analyzing video content. "
    "Read the digest and transcript carefully, then:\n\n"
    "1. Infer 3–5 experts whose perspectives would best illuminate this material. "
    "Choose them from the actual domain of the video — a neuroscientist for a brain "
    "talk, a hardware engineer + an ML practitioner for a chip-design talk, a historian "
    "of science + a contemporary researcher for a science-history piece. Avoid generic "
    "labels (\"a thoughtful generalist\"); make each expert's specialty concrete enough "
    "that their angle on this material is distinct.\n\n"
    "2. Run a 1500–2500 word panel discussion in markdown. Open with one short paragraph "
    "introducing each panelist (name, role, one credential or claim-to-relevance). Then "
    "the discussion proper, with each turn labeled by the speaker's name.\n\n"
    "Goals for the discussion:\n"
    "- Surface what the speaker glossed over, hand-waved, or assumed without arguing.\n"
    "- Bring contrary readings — where would a competing school of thought disagree?\n"
    "- Connect to adjacent domains the speaker didn't mention.\n"
    "- Examine concrete claims (numbers, names, mechanisms) for how robust they actually are.\n"
    "- Synthesize, but don't paper over disagreements: if two panelists land in "
    "different places, leave them there.\n\n"
    "Style: skip restating the digest — the reader already read it. Open directly with "
    "the moderator framing the first question. No conclusion-summary at the end; let "
    "the discussion close naturally."
)


def generate_panel_discussion(
    digest_md_text: str,
    segments: List["TranscriptSegment"],
    model: str,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Call Claude to simulate a panel of domain-relevant experts discussing a video.
    Returns (markdown_text, usage).

    source_lang / output_language follow the same convention as generate_digest:
    'auto' writes the panel in the transcript's language; 'en' forces English.

    Costs ~one Opus call per click (≈ 4–8k input + 2–4k output tokens). Output is
    one markdown document the caller writes to digests/<id>/panel.md.
    """
    transcript_str = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_directive = ""
    if not is_english_source and output_language == "auto":
        lang_directive = (
            f"\n\nIMPORTANT: The transcript is in language code '{source_lang}'. "
            "Write the entire panel discussion in the SAME language — expert names "
            "(transliterated when appropriate), credentials, the moderator's "
            "questions, every speaker's turns. Preserve proper nouns and technical "
            "terms in their original form when the field uses them that way."
        )
    # output_language == "en" with non-English source: rely on the existing
    # system prompt (no explicit translate directive needed; English is the
    # default Claude output style for this prompt).

    user_text = (
        "## Existing digest (the reader has already seen this)\n\n"
        f"{digest_md_text}\n\n"
        "## Full timestamped transcript\n\n"
        f"{transcript_str}"
        f"{lang_directive}\n\n"
        "Now: introduce the panelists, then run the discussion."
    )

    if backend is None:
        backend = select_backend()
    return backend.text(
        system=PANEL_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=8000,
    )


DEFAULT_TAKEAWAY_MODEL = "claude-sonnet-4-6"


TAKEAWAY_SYSTEM_PROMPT = (
    "You write the audience-facing takeaway for a video the reader has just "
    "finished. They've read the digest and (often) the panel discussion. "
    "Now they want a friend's-eye-view: 'Here's what I got out of this; "
    "here's what to walk away with; here's what's contested.'\n\n"
    "Format: 1–3 short paragraphs of plain prose. Open with the single most "
    "important thing the reader should leave with — the bottom line, "
    "stated directly. Then weave in supporting context: where the speaker's "
    "framing is solid, where the panel pushed back, what's still open. "
    "Close with the implications for the reader (so what / why does this "
    "matter), one or two sentences.\n\n"
    "Genre awareness — first identify (silently, internally) what kind of "
    "video this is, and shape the takeaway accordingly:\n"
    "- Tech talk / explainer → the working position on the frameworks "
    "presented, with critical pushback woven in.\n"
    "- Market / finance → the trade thesis, what's priced in, the catalysts "
    "to watch — with an explicit 'as of <date>' anchor.\n"
    "- News / current events → what changed and what it means.\n"
    "- How-to / tutorial → what they teach you to do, and the gotcha "
    "experienced practitioners flag.\n"
    "- Interview → where the speakers' positions actually differ.\n"
    "- Product launch → what's genuinely new, what's hype, what to "
    "actually use.\n"
    "Do NOT label the genre in the output. Just let it shape the writing.\n\n"
    "Grounding claims in the video — when you state a specific fact, claim, "
    "or quote that lives at a particular moment, mark it with bracketed "
    "timestamps using the original video's M:SS or H:MM:SS format, e.g. "
    "[3:15] or [1:02:48]. Use the bracketed form exactly — no parentheses, "
    "no markdown link syntax. The renderer will turn these into clickable "
    "links to the source video. Use them sparingly (3–8 across the whole "
    "takeaway) and only on substantive points worth verifying — not on "
    "every sentence.\n\n"
    "Time-sensitive content: if the video discusses dated material — market "
    "state, recent product launches, current numbers — mention the publish "
    "date inline ('As of <date>...') so a future reader knows the freshness "
    "window. Use the publish date provided in the user message. For "
    "evergreen content (general knowledge, frameworks, well-established "
    "claims), don't add a date.\n\n"
    "Style:\n"
    "- Conversational, not academic. A friend telling you what they got out "
    "of the talk — not a research abstract.\n"
    "- Concrete > abstract. Use real names, numbers, and frameworks from "
    "the video.\n"
    "- Honest about contested points. If the panel disagreed, say so "
    "(\"though the panel pushed back on the file-system framing as a long-"
    "term abstraction\"). Don't paper it over.\n"
    "- Don't restate the digest. The reader just read it. Synthesize and "
    "go beyond.\n"
    "- No headings, no bullet lists, no preamble like 'Here's my takeaway'. "
    "Just the prose."
)


def generate_takeaway_prose(
    digest_md_text: str,
    panel_md_text: Optional[str],
    segments: List["TranscriptSegment"],
    model: str,
    *,
    publish_date: Optional[str] = None,
    source_lang: str = "en",
    output_language: str = "auto",
    backend=None,
):
    """Final pipeline step: write the audience-facing takeaway as 1-3 short
    paragraphs of prose. Synthesizes digest + panel into a personal
    'what to walk away with' read.

    publish_date: YYYYMMDD or YYYY-MM-DD as returned by yt-dlp's
    `info["upload_date"]`; threaded into the prompt so time-sensitive
    takeaways can anchor with 'as of <date>'.

    Returns (takeaway_text: str, usage). The text contains [M:SS] bracket
    markers that the renderer converts into clickable timestamp links.
    """
    transcript_str = "\n".join(
        f"[{format_timestamp(seg.start)}] {seg.text}" for seg in segments
    )

    # Normalize publish_date to YYYY-MM-DD for the prompt (yt-dlp returns
    # YYYYMMDD by default).
    pub_str = ""
    if publish_date:
        pd = publish_date.replace("-", "")
        if len(pd) == 8 and pd.isdigit():
            pub_str = f"{pd[0:4]}-{pd[4:6]}-{pd[6:8]}"
        else:
            pub_str = publish_date

    is_english_source = (source_lang or "").lower().startswith("en")
    lang_directive = ""
    if not is_english_source and output_language == "auto":
        lang_directive = (
            f"\n\nIMPORTANT: The transcript is in language code '{source_lang}'. "
            "Write the takeaway in the SAME language as the transcript. "
            "Preserve proper nouns and technical terms in their original form."
        )

    panel_section = (
        "## Panel discussion\n\n" + panel_md_text + "\n\n"
        if panel_md_text else
        "## Panel discussion\n\n(none generated)\n\n"
    )

    user_text = (
        f"Video publish date: {pub_str or '(unknown)'}\n\n"
        "## Digest (the reader has seen this)\n\n"
        f"{digest_md_text}\n\n"
        f"{panel_section}"
        "## Full timestamped transcript\n\n"
        f"{transcript_str}"
        f"{lang_directive}\n\n"
        "Now: write the takeaway."
    )

    if backend is None:
        backend = select_backend()
    return backend.text(
        system=TAKEAWAY_SYSTEM_PROMPT, user_text=user_text,
        model=model, max_tokens=2000,
    )



def _transcript_slice(
    segments: List["TranscriptSegment"],
    topic_start: float,
    topic_end: float,
) -> str:
    """Render the transcript segments inside [topic_start, topic_end) as one timestamped string.

    Used by vision_pick_frames to ground picks against what the narrator is saying
    at a candidate frame's timestamp. Truncates very long topics to keep token cost
    bounded — first 30 + last 30 segments with a marker between, which captures the
    narrator's framing at start and conclusion at end.
    """
    in_window = [s for s in segments if topic_start <= s.start < topic_end]
    if len(in_window) > 70:
        head = in_window[:30]
        tail = in_window[-30:]
        in_window = head + [None] + tail  # type: ignore[list-item]
    lines: List[str] = []
    for seg in in_window:
        if seg is None:
            lines.append("[…]")
            continue
        lines.append(f"[{format_timestamp(seg.start)}] {seg.text}")
    return "\n".join(lines)


def _candidates_for_topic(
    topic_start: float,
    topic_end: float,
    frames: List[Tuple[Path, float]],
    max_per_topic: int = 5,
    overlap_pre: float = 5.0,
) -> List[Tuple[Path, float]]:
    """Frames whose timestamp falls in [start - overlap, end). Downsample to max_per_topic."""
    in_window = [
        (p, t) for p, t in frames if (topic_start - overlap_pre) <= t < topic_end
    ]
    if len(in_window) <= max_per_topic:
        return in_window
    # Even spacing across the window
    step = len(in_window) / max_per_topic
    return [in_window[int(i * step)] for i in range(max_per_topic)]


def _encode_frame_for_vision(path: Path, max_long_edge: int = 1024) -> str:
    """Resize a frame to <= max_long_edge on the long side and return base64-encoded JPEG bytes."""
    import base64
    import io
    from PIL import Image

    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        scale = max_long_edge / max(w, h)
        if scale < 1:
            im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=82)
    return base64.standard_b64encode(buf.getvalue()).decode("ascii")


def vision_pick_frames(
    digest,
    frames: List[Tuple[Path, float]],
    video_duration: float,
    model: str,
    segments: Optional[List["TranscriptSegment"]] = None,
    backend=None,
):
    """Use Claude's vision to pick the best frame per topic from in-window candidates.

    Returns a dict {topic_index -> chosen_frame_path}, plus the API usage object.
    Topics with no in-window candidates are omitted (caller falls back to timestamp-based pick).

    If `segments` is provided, the per-topic transcript slice is included so vision
    can ground picks on what the narrator is saying at each candidate's timestamp
    (e.g. "speaker says 'as you can see in this diagram' at 04:23 → frame at 04:23").

    Raises VisionUnsupported when the active backend can't process images
    (e.g. Claude Code with vision opt-out). Caller should fall back.
    """
    from pydantic import BaseModel
    from typing import List as TList

    class TopicChoice(BaseModel):
        topic_index: int
        candidate_index: int
        rationale: str

    class FrameChoices(BaseModel):
        choices: TList[TopicChoice]

    topics = digest.topics

    # Build per-topic candidate lists and a flat list of (topic_idx, cand_idx, path, ts)
    per_topic: List[List[Tuple[Path, float]]] = []
    per_topic_transcript: List[str] = []
    for i, topic in enumerate(topics):
        end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
        per_topic.append(_candidates_for_topic(topic.start_time, end, frames))
        if segments:
            per_topic_transcript.append(
                _transcript_slice(segments, topic.start_time, end)
            )
        else:
            per_topic_transcript.append("")

    # Build the message: text intro -> for each topic, label + summary + numbered candidate images
    content: list = []
    intro = (
        "For each topic below, pick the candidate frame that best illustrates what the "
        "narrator is discussing. Prefer frames showing the most informative visual content "
        "(diagrams, code, distinctive UI) over generic framing or talking-head shots.\n\n"
        "When multiple candidates show the same scene at different stages of an animation "
        "or progressive reveal — bullets appearing one at a time, diagram elements being "
        "added, code typed line by line — prefer the LATEST candidate in the sequence. "
        "The final state shows the most complete information; partial/early states omit "
        "content the narrator goes on to add.\n\n"
        "Use the per-topic transcript to ground your pick: when the narrator says things "
        "like \"as you can see here\" or refers to a specific element at a specific "
        "moment, prefer the candidate whose timestamp is closest to that mention.\n\n"
        "Return one choice per topic that has candidates.\n\n"
        f"Total topics: {len(topics)}\n"
    )
    content.append({"type": "text", "text": intro})

    for ti, topic in enumerate(topics):
        cands = per_topic[ti]
        if not cands:
            content.append({
                "type": "text",
                "text": f"\n--- Topic {ti} (no candidates available — skip) ---\n"
                        f"Title: {topic.title}\n",
            })
            continue
        header_parts = [
            f"\n--- Topic {ti} ---",
            f"Title: {topic.title}",
            f"Summary: {topic.summary}",
        ]
        if per_topic_transcript[ti]:
            header_parts.append("Transcript:\n" + per_topic_transcript[ti])
        header_parts.append(f"Candidates ({len(cands)} frames):\n")
        content.append({"type": "text", "text": "\n".join(header_parts)})
        for ci, (path, ts) in enumerate(cands):
            content.append({
                "type": "text",
                "text": f"Candidate {ci} (at {format_timestamp(ts)}):",
            })
            content.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/jpeg",
                    "data": _encode_frame_for_vision(path),
                },
            })

    system = (
        "You select illustrative frames for a video digest. You will be shown a list of "
        "topics, each with the topic's title, summary, the transcript spoken during that "
        "topic (timestamped), and a small set of candidate frames (also timestamped). "
        "For each topic with candidates, return the (topic_index, candidate_index) pair "
        "that best illustrates the topic, with a one-sentence rationale. Use the "
        "transcript to ground your choice on what the narrator is saying when each "
        "candidate frame was captured. Skip topics that say 'no candidates available'."
    )

    if backend is None:
        backend = select_backend()
    if not getattr(backend, "vision_supported", False):
        raise VisionUnsupported(
            f"Backend {backend.name!r} does not support vision in the current "
            "configuration."
        )
    parsed, usage = backend.vision_parse(
        system=system, content_blocks=content,
        model=model, max_tokens=4000, schema=FrameChoices,
    )

    chosen: dict = {}
    for choice in parsed.choices:
        if 0 <= choice.topic_index < len(topics):
            cands = per_topic[choice.topic_index]
            if 0 <= choice.candidate_index < len(cands):
                chosen[choice.topic_index] = cands[choice.candidate_index][0]
    return chosen, usage


def _pick_topic_frame(
    topic_start: float,
    topic_end: float,
    candidates: List[Tuple[Path, float]],
    used: set,
) -> Optional[Tuple[Path, float]]:
    """Pick the best frame for a topic: prefer a frame inside [start, end), else closest."""
    in_window = [(p, t) for p, t in candidates if topic_start <= t < topic_end and p not in used]
    if in_window:
        midpoint = (topic_start + min(topic_end, in_window[-1][1] + 1)) / 2
        return min(in_window, key=lambda x: abs(x[1] - midpoint))
    available = [(p, t) for p, t in candidates if p not in used]
    if not available:
        return None
    return min(available, key=lambda x: abs(x[1] - topic_start))


def write_markdown_digest(
    digest,
    candidate_frames: List[Tuple[Path, float]],
    video_duration: float,
    output_md: Path,
    images_dir: Path,
    vision_picks: Optional[dict] = None,
    video_title: Optional[str] = None,
    video_url: Optional[str] = None,
) -> None:
    """Write the digest to disk. Picks frames, copies them into
    images_dir, builds the structured JSON shape from the in-memory
    Digest object, and writes BOTH digest.md AND digest.json atomically.

    Phase B: JSON is the source of truth — `digest.md` is rendered from
    the JSON via `render_digest_md`, so a future read via the agent API
    skips the markdown parser entirely. (Legacy digests written before
    this change still work via the parse-on-read fallback in
    `load_digest_json`.)

    If vision_picks is provided, prefer those mappings; fall back to
    timestamp-based picks for any topic not covered.

    video_title (when provided) overrides the LLM-generated title so the
    digest's heading matches the original YouTube title — readers can
    map it back to the source. video_url renders a "Watch on YouTube"
    link directly under the title.
    """
    images_dir.mkdir(parents=True, exist_ok=True)

    used: set = set()
    topic_images: List[Optional[Path]] = []
    topics = digest.topics
    for i, topic in enumerate(topics):
        # Vision pick takes priority
        if vision_picks and i in vision_picks:
            src_path = vision_picks[i]
            used.add(src_path)
        else:
            end = topics[i + 1].start_time if i + 1 < len(topics) else video_duration
            pick = _pick_topic_frame(topic.start_time, end, candidate_frames, used)
            if pick is None:
                topic_images.append(None)
                continue
            src_path, _ = pick
            used.add(src_path)
        dest = images_dir / f"topic_{i + 1:02d}.jpg"
        shutil.copy(src_path, dest)
        topic_images.append(dest)

    # Build the canonical JSON shape directly from the Digest object —
    # no round-trip through markdown. Sidecar metadata.json (if present)
    # fills in channel info and upload date.
    json_dict = _digest_json_from_object(
        digest,
        topic_images=topic_images,
        images_rel_dir=images_dir.name,
        video_title=video_title,
        video_url=video_url,
        video_id=output_md.parent.name,
        metadata_path=output_md.parent / "metadata.json",
    )

    # Write .md first so its mtime is older than .json; the mtime check
    # in load_digest_json then sees the .json cache as fresh on read.
    _atomic_write_text(output_md, render_digest_md(json_dict))
    _atomic_write_json(output_md.with_suffix(".json"), json_dict)


def render_takeaway_markdown(
    takeaway_text: str,
    *,
    video_url: Optional[str] = None,
) -> str:
    """Post-process the LLM's takeaway prose for writing to takeaway.md.
    Converts the LLM's bracketed [M:SS] / [H:MM:SS] markers into clickable
    markdown links to the source video when video_url is known. Returns the
    body text alone — the file's heading is rendered by the viewer chrome,
    not embedded in the markdown.
    """
    body = takeaway_text.strip()

    if not video_url:
        return body + "\n"

    sep = "&" if "?" in video_url else "?"

    def _ts_to_seconds(ts: str) -> int:
        parts = [int(p) for p in ts.split(":")]
        if len(parts) == 2:
            return parts[0] * 60 + parts[1]
        if len(parts) == 3:
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        return 0

    # Match bare bracketed timestamps like [3:15] or [1:02:48], but not ones
    # that are already a markdown link's display text — skip patterns
    # immediately followed by '(' which indicates the LLM already wrote a
    # full [text](url) link.
    def _link(m):
        ts = m.group(1)
        sec = _ts_to_seconds(ts)
        return f"[{ts}]({video_url}{sep}t={sec}s)"

    return re.sub(r"\[(\d+:\d+(?::\d+)?)\](?!\()", _link, body) + "\n"


# ---------- Subcommands: watch / meta / serve ----------

LATEST_LIMIT = 10
MAX_NEW_PER_RUN = 3


def channels_file() -> Path:
    return get_data_dir() / "channels.txt"


def state_file() -> Path:
    return get_data_dir() / "state.json"


def read_channels() -> List[str]:
    p = channels_file()
    if not p.exists():
        return []
    return [
        line.strip()
        for line in p.read_text().splitlines()
        if line.strip() and not line.lstrip().startswith("#")
    ]


def write_channels(channels: List[str]) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    body = "# YouTube channels to watch. One URL per line. Lines starting with # are ignored.\n"
    body += "\n".join(channels) + ("\n" if channels else "")
    channels_file().write_text(body)


def load_state() -> dict:
    p = state_file()
    if not p.exists():
        return {"channels": {}}
    return json.loads(p.read_text())


def save_state(state: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    state_file().write_text(json.dumps(state, indent=2, sort_keys=True) + "\n")


# ---- watch subcommands ----

def normalize_channel_url(url: str) -> str:
    """Light normalization for YouTube channel URLs.

    Accepts: '@handle', 'youtube.com/@handle', 'https://www.youtube.com/@handle/videos'.
    Always returns a fully-qualified URL. Adds '/videos' to bare @handle URLs so
    yt-dlp targets the videos tab specifically.
    """
    url = url.strip()
    if url.startswith("@"):
        url = f"https://www.youtube.com/{url}"
    if not url.startswith(("http://", "https://")):
        url = f"https://{url}"
    m = re.match(r"^(https?://(?:www\.)?youtube\.com/@[^/]+)/?$", url)
    if m:
        url = m.group(1) + "/videos"
    return url


def cmd_watch_add(args) -> int:
    url = normalize_channel_url(args.url)
    if not is_url(url):
        sys.exit(f"Not a URL: {url}")
    channels = read_channels()
    if url in channels:
        print(f"Already watching: {url}")
        return 0
    channels.append(url)
    write_channels(channels)
    print(f"Added: {url}")
    return 0


def cmd_watch_list(args) -> int:
    channels = read_channels()
    if not channels:
        print("No channels configured. Add one with: yt2md watch add <URL>")
        return 0
    print(f"Watching {len(channels)} channel(s):")
    for ch in channels:
        print(f"  {ch}")
    print(f"\nConfig: {channels_file()}")
    print(f"State:  {state_file()}")
    print(f"Data:   {get_data_dir()}")
    return 0


def cmd_watch_remove(args) -> int:
    url = args.url.strip()
    channels = read_channels()
    if url not in channels:
        sys.exit(f"Not in list: {url}")
    channels = [c for c in channels if c != url]
    write_channels(channels)
    print(f"Removed: {url}")
    return 0


def _list_channel_videos(url: str, limit: int = LATEST_LIMIT) -> List[str]:
    out = subprocess.check_output(
        ["yt-dlp", "--flat-playlist", "--playlist-end", str(limit), "--print", "%(id)s", url],
        text=True,
    )
    return [line.strip() for line in out.splitlines() if line.strip()]


def _digest_video(
    video_id: str, output_dir: Path, *, source: str = "subscription",
) -> Tuple[int, str]:
    """Run yt2md on a video. Returns (exit_code, combined_stdout_stderr).

    Streams output to the parent's stdout in real time (so poll.log captures
    it as it happens) AND collects it into a buffer the caller can scan for
    permanent-failure patterns.

    `source` is stamped into digest_meta as the ingestion provenance.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    digest_path = output_dir / "digest.md"
    yt2md = shutil.which("yt2md") or sys.argv[0]
    proc = subprocess.Popen(
        [yt2md, f"https://youtu.be/{video_id}", "-o", str(digest_path),
         "--source", source],
        cwd=output_dir,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
    )
    chunks: list = []
    assert proc.stdout is not None
    for line in proc.stdout:
        sys.stdout.write(line)
        sys.stdout.flush()
        chunks.append(line)
    proc.wait()
    return proc.returncode, "".join(chunks)


# Substrings (case-insensitive) that indicate a video can never be digested by
# this account — gating it forever, deleted, etc. Matches mark the video
# as seen so polling stops re-trying it. Network/transient errors don't match
# and continue to retry next cycle.
_PERMANENT_FAILURE_PATTERNS = (
    "members-only",
    "join this channel",
    "private video",
    "video unavailable",
    "this video is no longer available",
    "removed by the uploader",
    "removed for violating",
    "sign in to confirm your age",
    "video has been removed",
)


def _is_permanent_failure(output: str) -> bool:
    low = output.lower()
    return any(p in low for p in _PERMANENT_FAILURE_PATTERNS)


def cmd_watch_run(args) -> int:
    ensure_api_key()
    channels = read_channels()
    if not channels:
        print("No channels configured. Add one with: yt2md watch add <URL>")
        return 0

    data_dir = get_data_dir()
    digests_dir = data_dir / "digests"
    state = load_state()
    any_failures = False

    for channel_url in channels:
        print(f"--- {channel_url}")
        seen = set(state["channels"].get(channel_url, {}).get("seen", []))
        latest_ids = _list_channel_videos(channel_url)

        if not seen:
            print(f"  first run, seeding state with {len(latest_ids)} videos (no backfill)")
            state["channels"][channel_url] = {"seen": sorted(latest_ids)}
            continue

        new_ids = [vid for vid in latest_ids if vid not in seen][:MAX_NEW_PER_RUN]
        if not new_ids:
            print("  no new videos")
            continue

        print(f"  {len(new_ids)} new: {new_ids}")
        for vid in reversed(new_ids):
            print(f"  processing {vid}...")
            rc, output = _digest_video(vid, digests_dir / vid)
            if rc == 0:
                seen.add(vid)
                state["channels"][channel_url] = {"seen": sorted(seen)}
                save_state(state)
            elif _is_permanent_failure(output):
                # Permanent: mark seen so polling stops cycling on it.
                # Wipe the partial dir (mp4 download, empty digest, etc.) to
                # keep digests/ tidy.
                print(f"  PERMANENTLY UNAVAILABLE: {vid} — marking seen and wiping partial dir")
                shutil.rmtree(digests_dir / vid, ignore_errors=True)
                seen.add(vid)
                state["channels"][channel_url] = {"seen": sorted(seen)}
                save_state(state)
            else:
                print(f"  FAILED on {vid} (transient — will retry next poll)", file=sys.stderr)
                any_failures = True

    save_state(state)
    return 1 if any_failures else 0


# ---- in-process scheduler ----
#
# Cadenced background runner for `yt2md watch run` (subscription poll). Lives
# inside `yt2md serve`: a daemon thread ticks every ~30s, fires due jobs as
# detached subprocesses, tracks pid + exit code in schedule_state.json for
# the /schedule UI.
#
# Tradeoff: scheduling pauses while serve is down — for an interactively-used
# reader this is fine; missed slots fire on next start (catch-up semantics).

DEFAULT_SCHEDULE_CONFIG = {
    "poll_interval_hours": 6,
}

_SCHED_TICK_SECS = 30
_scheduler_jobs: dict = {}
_scheduler_thread = None
_scheduler_lock_obj = None


def _schedule_lock():
    global _scheduler_lock_obj
    if _scheduler_lock_obj is None:
        import threading
        _scheduler_lock_obj = threading.Lock()
    return _scheduler_lock_obj


def _schedule_config_file() -> Path:
    return get_data_dir() / "schedule.json"


def _schedule_state_file() -> Path:
    return get_data_dir() / "schedule_state.json"


def load_schedule_config() -> dict:
    p = _schedule_config_file()
    if not p.exists():
        return dict(DEFAULT_SCHEDULE_CONFIG)
    try:
        cfg = json.loads(p.read_text())
        merged = dict(DEFAULT_SCHEDULE_CONFIG)
        merged.update({k: v for k, v in cfg.items() if k in DEFAULT_SCHEDULE_CONFIG})
        return merged
    except Exception:
        return dict(DEFAULT_SCHEDULE_CONFIG)


def save_schedule_config(cfg: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    _schedule_config_file().write_text(json.dumps(cfg, indent=2) + "\n")


def _load_schedule_state() -> dict:
    p = _schedule_state_file()
    default = {"poll": {}}
    if not p.exists():
        return default
    try:
        s = json.loads(p.read_text())
        return {"poll": s.get("poll") or {}}
    except Exception:
        return default


def _save_schedule_state(state: dict) -> None:
    p = _schedule_state_file()
    p.parent.mkdir(parents=True, exist_ok=True)
    tmp = p.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(state, indent=2) + "\n")
    tmp.replace(p)


# ---- user-tunable settings (model + tooling choices) ----
#
# Lives at ~/yt2md/settings.json. Editable via the /settings page; flows into
# every subprocess spawn (one-off, scheduled poll) as YT2MD_* env vars, so the
# digest CLI's argparse defaults pick them up. The /digests/<id>/discuss route
# reads settings directly (it runs in-process).

DEFAULT_SETTINGS = {
    "digest_model": "claude-sonnet-4-6",
    "panel_model": "claude-opus-4-7",
    "whisper_model": DEFAULT_WHISPER_MODEL,
    "cookies_from_browser": "",
    # "auto" = write the digest in the same language as the transcript;
    # "en" = always translate to English. Applies to both the per-video digest
    # and the panel discussion.
    "digest_language": "auto",
    # Which auth path to use for LLM calls.
    #   "auto":        prefer ANTHROPIC_API_KEY when set; else use the bundled
    #                  Claude Code sandbox if installed + logged in.
    #   "api":         force direct Anthropic API (requires ANTHROPIC_API_KEY).
    #   "claude-code": force the bundled Claude Code subprocess backend.
    "llm_backend": "auto",
    # Vision frame-picking is automatic for the API backend (cheap & native)
    # but disabled by default for Claude Code (no -p image flag; we'd have to
    # base64-embed → token-heavy). Toggle on if you want vision via Claude
    # Code despite the cost.
    "claude_code_vision": False,
    # When generating slides, use a vision-LLM (typically Haiku) to filter
    # raw extracted frames down to actual deck slides. Trades a small LLM
    # cost (~$0.005 per video via 3×3 grid batching with Haiku) for a much
    # cleaner deck. Set False to use pure pHash dedup only.
    "slide_classification": True,
    "slide_classifier_model": "claude-haiku-4-5-20251001",
    # Which TTS backend the 🎧 Listen button uses.
    #   "macos":      built-in `say` + ffmpeg. Free, offline, lower quality
    #                 unless you set Siri Voice 1 as your system voice.
    #   "elevenlabs": ElevenLabs API. Much higher quality, costs against
    #                 your plan credits. Requires ELEVENLABS_API_KEY in
    #                 ~/yt2md/.env.
    "tts_provider": "macos",
    # macOS-only: voice name (blank → system default). For best quality
    # leave this blank AND set Siri Voice 1/2/3 as your System Voice in
    # System Settings → Accessibility → Spoken Content. Otherwise try
    # "Fiona" or "Samantha (Enhanced)". `say -v ?` lists all voices.
    "tts_voice": "",
    # macOS-only: speaking rate in words/min. Blank → system default (~175).
    "tts_rate": "",
    # ElevenLabs voice ID. Default = "Brian" — a calm US narration voice
    # in the free-tier "Default voices" set, so it works on every plan.
    # (The older "Rachel" / `21m00Tcm4TlvDq8ikWAM` is now treated as a
    # library voice and requires a paid plan via API.)
    "elevenlabs_voice_id": "nPczCjzI2devNBz1zQrb",
    # ElevenLabs model. eleven_multilingual_v2 is the high-quality default;
    # eleven_turbo_v2_5 is faster + cheaper; eleven_flash_v2_5 is fastest.
    "elevenlabs_model": "eleven_multilingual_v2",
    # --- Budget gate (requires an Admin key for authoritative billing; falls
    # back to the local usage log otherwise). Refuses to START a new digest
    # once this workspace's month-to-date spend crosses block_usd; warns past
    # warn_usd. A running digest is never interrupted. Any Console spend cap is
    # an independent hard backstop above these.
    # budget_workspace_id: auto-detected from the runtime API key on first use;
    # leave blank to auto-fill, or set explicitly to pin attribution.
    "budget_workspace_id": "",
    "budget_warn_usd": 15.0,
    "budget_block_usd": 18.0,
    # Per-model price overrides, e.g. {"claude-opus-4-7": {"input": 5.0, ...}}.
    # Lowest precedence after the billing-calibrated cache; mainly an escape
    # hatch. Registered here so load_settings() doesn't filter it out.
    "model_pricing": {},
}


def _settings_file() -> Path:
    return get_data_dir() / "settings.json"


def load_settings() -> dict:
    p = _settings_file()
    if not p.exists():
        return dict(DEFAULT_SETTINGS)
    try:
        s = json.loads(p.read_text())
        merged = dict(DEFAULT_SETTINGS)
        merged.update({k: v for k, v in s.items() if k in DEFAULT_SETTINGS})
        return merged
    except Exception:
        return dict(DEFAULT_SETTINGS)


def save_settings(s: dict) -> None:
    get_data_dir().mkdir(parents=True, exist_ok=True)
    tmp = _settings_file().with_suffix(".json.tmp")
    tmp.write_text(json.dumps(s, indent=2) + "\n")
    tmp.replace(_settings_file())


def _settings_to_env(settings: dict) -> dict:
    """Map settings to YT2MD_* env vars for subprocess invocation. Empty values
    are dropped so the subprocess sees only meaningful overrides."""
    out: dict = {}
    if settings.get("digest_model"):
        out["YT2MD_DIGEST_MODEL"] = settings["digest_model"]
    if settings.get("whisper_model"):
        out["YT2MD_WHISPER_MODEL"] = settings["whisper_model"]
    if settings.get("cookies_from_browser"):
        out["YT2MD_COOKIES_FROM_BROWSER"] = settings["cookies_from_browser"]
    if settings.get("panel_model"):
        out["YT2MD_PANEL_MODEL"] = settings["panel_model"]
    if settings.get("digest_language"):
        out["YT2MD_DIGEST_LANGUAGE"] = settings["digest_language"]
    # Spawned subprocesses re-resolve llm_backend / claude_code_vision via
    # load_settings() in the child — no env-var passthrough needed for those.
    # CLAUDE_CONFIG_DIR is required so the bundled `claude` binary in the
    # child finds our sandboxed credentials/settings instead of any system
    # install's defaults.
    out["CLAUDE_CONFIG_DIR"] = str(claude_config_dir())
    return out


def _format_schedule_summary(cfg: dict) -> str:
    """Human-readable description of the schedule config."""
    poll = cfg["poll_interval_hours"]
    poll_str = f"every {poll} hour{'s' if poll != 1 else ''}"
    return f"polling {poll_str}"


def _compute_next_poll(cfg: dict, last_started_at: Optional[float]) -> float:
    """Timestamp of the next scheduled poll. First-run convention: fire ASAP."""
    interval = max(60.0, float(cfg.get("poll_interval_hours", 6)) * 3600.0)
    if last_started_at is None:
        import time as _t
        return _t.time()
    return last_started_at + interval


def _fire_scheduled_job(kind: str) -> Optional[subprocess.Popen]:
    """Spawn yt2md as a subprocess for the given kind. Returns the running
    Popen or None if already running / yt2md not on PATH. Caller holds lock."""
    if kind != "poll":
        return None
    existing = _scheduler_jobs.get(kind)
    if existing is not None and existing.poll() is None:
        return existing
    yt2md_path = shutil.which("yt2md")
    if not yt2md_path:
        print(f"[scheduler] yt2md not on PATH; cannot fire {kind}", file=sys.stderr)
        return None
    args = [yt2md_path, "watch", "run"]
    log_path = get_data_dir() / "logs" / f"{kind}.log"
    log_path.parent.mkdir(parents=True, exist_ok=True)
    import time as _t
    log_fd = open(log_path, "a")
    log_fd.write(f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} {kind} run =====\n")
    log_fd.flush()
    proc = subprocess.Popen(
        args,
        stdout=log_fd,
        stderr=subprocess.STDOUT,
        env={**os.environ, **_settings_to_env(load_settings()), "PYTHONUNBUFFERED": "1"},
        **_DETACH_KWARGS,
    )
    log_fd.close()
    _scheduler_jobs[kind] = proc

    state = _load_schedule_state()
    state[kind] = {
        **(state.get(kind) or {}),
        "last_started_at": _t.time(),
        "last_pid": proc.pid,
        "last_exit_code": None,
        "last_finished_at": None,
    }
    _save_schedule_state(state)
    return proc


def _reap_scheduled_job(kind: str) -> None:
    """Capture exit code if the kind's subprocess has exited. Caller holds lock."""
    proc = _scheduler_jobs.get(kind)
    if proc is None:
        return
    rc = proc.poll()
    if rc is None:
        return
    import time as _t
    state = _load_schedule_state()
    state[kind] = {
        **(state.get(kind) or {}),
        "last_finished_at": _t.time(),
        "last_exit_code": rc,
    }
    _save_schedule_state(state)
    del _scheduler_jobs[kind]


def _scheduler_tick() -> None:
    with _schedule_lock():
        _reap_scheduled_job("poll")
        cfg = load_schedule_config()
        state = _load_schedule_state()
        import time as _t
        now = _t.time()
        if "poll" not in _scheduler_jobs:
            if now >= _compute_next_poll(cfg, (state.get("poll") or {}).get("last_started_at")):
                _fire_scheduled_job("poll")


def _scheduler_loop() -> None:
    import time as _t
    while True:
        try:
            _scheduler_tick()
        except Exception as e:
            print(f"[scheduler] tick error: {e}", file=sys.stderr)
        _t.sleep(_SCHED_TICK_SECS)


def start_scheduler() -> None:
    """Start the daemon scheduler thread. Idempotent."""
    global _scheduler_thread
    if _scheduler_thread is not None and _scheduler_thread.is_alive():
        return
    import threading
    _scheduler_thread = threading.Thread(target=_scheduler_loop, daemon=True)
    _scheduler_thread.start()
    print("[scheduler] started (in-process; ticks every "
          f"{_SCHED_TICK_SECS}s)")


def _cleanup_legacy_launchd() -> None:
    """Best-effort one-time removal of the old launchctl plists from
    ~/Library/LaunchAgents — so the user doesn't end up with duplicate
    scheduling after this migration. Silent if nothing is present.
    macOS-only — launchctl doesn't exist anywhere else."""
    if sys.platform != "darwin":
        return
    launchd_dir = Path.home() / "Library" / "LaunchAgents"
    removed = []
    for label in ("com.youtube-to-markdown.poll", "com.youtube-to-markdown.meta"):
        plist = launchd_dir / f"{label}.plist"
        if plist.exists():
            subprocess.run(
                ["launchctl", "bootout", f"gui/{os.getuid()}", str(plist)],
                check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
            try:
                plist.unlink()
                removed.append(label)
            except OSError:
                pass
    if removed:
        print(f"[scheduler] removed legacy launchd plists: {', '.join(removed)}")


def _scheduler_status_summary(kind: str, state: dict) -> Tuple[str, str]:
    """English summary + dot class for /schedule and /channels surfaces."""
    s = state.get(kind) or {}
    if kind in _scheduler_jobs:
        return ("Running now.", "dot-on")
    last_started = s.get("last_started_at")
    last_exit = s.get("last_exit_code")
    if last_started is None:
        return ("Set up — first run hasn't happened yet.", "dot-warn")
    if last_exit not in (0, None):
        return (f"Last run failed (exit code {last_exit}). Check the log.", "dot-warn")
    import datetime as dt
    age = dt.datetime.now() - dt.datetime.fromtimestamp(last_started)
    age_secs = int(age.total_seconds())
    if age_secs < 60:
        age_str = f"{age_secs}s ago"
    elif age_secs < 3600:
        age_str = f"{age_secs // 60}m ago"
    elif age_secs < 86400:
        age_str = f"{age_secs // 3600}h ago"
    else:
        age_str = f"{age_secs // 86400}d ago"
    return (f"Healthy — last run {age_str} (clean).", "dot-on")


def _format_next_run(ts: float) -> str:
    """Render an upcoming-run timestamp as 'in 2h 15m' / 'in 3 days' / 'overdue'."""
    import time as _t
    delta = ts - _t.time()
    if delta < 0:
        return "due now"
    if delta < 60:
        return f"in {int(delta)}s"
    if delta < 3600:
        return f"in {int(delta // 60)}m"
    if delta < 86400:
        h = int(delta // 3600)
        m = int((delta % 3600) // 60)
        return f"in {h}h {m}m"
    return f"in {int(delta // 86400)}d"


def _tail_log(path: Path, n: int = 20) -> str:
    if not path.exists():
        return "(no log yet)"
    try:
        text = path.read_text(errors="replace")
    except Exception as e:
        return f"(error reading log: {e})"
    lines = text.splitlines()
    return "\n".join(lines[-n:]) if lines else "(empty)"




# ---- one-off digest jobs (in-memory tracking; cheap) ----

# Module-level dict: PID -> {"video_id": str, "started": float, "url": str, "proc": Popen}.
# Lost on server restart, which is fine — the digest still completes (detached
# subprocess) and shows up in the sidebar when done.
_oneoff_jobs: dict = {}

# Recent failures, most-recent-first, capped at _ONEOFF_FAILURE_CAP.
# Lost on server restart (matches _oneoff_jobs lifecycle).
_oneoff_failures: list = []
_ONEOFF_FAILURE_CAP = 20


_VIDEO_ID_RE = re.compile(r"(?:v=|youtu\.be/|/shorts/|/embed/)([A-Za-z0-9_-]{11})")


def extract_video_id(url: str) -> str:
    """Pull a YouTube video ID from common URL forms. Returns '' if not found."""
    m = _VIDEO_ID_RE.search(url)
    return m.group(1) if m else ""


def _extract_last_error(log_path: Path, video_id: str) -> str:
    """Pull the most relevant error line from oneoff.log for a given video_id.

    The log uses '===== {ts} starting {video_id} ({url}) =====' as section
    markers. We bound the section, then prefer 'RuntimeError: ...' / similar
    summary lines over the bare 'Traceback' header.
    """
    try:
        text = log_path.read_text(errors="replace")
    except OSError:
        return ""
    marker = f"starting {video_id}"
    idx = text.rfind(marker)
    if idx < 0:
        return ""
    next_idx = text.find("\n===== ", idx + len(marker))
    section = text[idx:next_idx if next_idx > 0 else len(text)]
    candidates = [
        ln.strip() for ln in section.splitlines()
        if ln.strip()
        and not ln.lstrip().startswith("[download")
        and ("Error" in ln or "Traceback" in ln)
    ]
    for ln in reversed(candidates):
        if ":" in ln and not ln.startswith("Traceback"):
            return ln
    return candidates[-1] if candidates else ""


# Pipeline stage markers, in pipeline order. The status endpoint scans the log
# section for the LAST matching substring to determine the current stage.
# Reordering here changes which stage is reported when two markers happen to
# match — keep this list in pipeline order (later entries override earlier).
_ONEOFF_STAGE_MARKERS: list = [
    ("starting",            "starting "),  # section header is always present
    ("downloading",         "[0/5] Fetching YouTube video"),
    ("downloading",         "[download]"),
    ("loading whisper",     "loading whisper model"),
    ("transcribing",        "transcribing audio with whisper"),
    ("extracting frames",   "[1/5] Extracting frames"),
    ("deduping frames",     "[2/5] Deduping"),
    ("parsing transcript",  "[3/5] Parsing SRT"),
    ("aligning",            "[4/5] Aligning"),
    ("building deck",       "[5/5] Building deck"),
    ("digesting",           "[+] Generating digest"),
    ("vision pass",         "[+] Vision-picking"),
    ("writing digest",      "Digest written"),
]


def _describe_job_stage(log_text: str, video_id: str) -> str:
    """Return a short human-readable label for the latest stage of a job.

    Scans the log section bounded by the start marker for video_id (or EOF /
    next start marker) and returns the most pipeline-advanced stage whose
    substring marker appears in that section.
    """
    marker = f"starting {video_id}"
    idx = log_text.rfind(marker)
    if idx < 0:
        return "starting"
    next_idx = log_text.find("\n===== ", idx + len(marker))
    section = log_text[idx:next_idx if next_idx > 0 else len(log_text)]
    current = "starting"
    for label, needle in _ONEOFF_STAGE_MARKERS:
        if needle in section:
            current = label
    return current


def _extract_run_summary(log_text: str, video_id: str) -> Optional[dict]:
    """Pull the last `[summary] {...}` JSON line emitted by the pipeline for a job.

    The pipeline prints one line of the form `[summary] {...}` on successful
    completion. Returns the parsed dict, or None if absent / malformed.
    """
    import json as _json
    marker = f"starting {video_id}"
    idx = log_text.rfind(marker)
    if idx < 0:
        return None
    next_idx = log_text.find("\n===== ", idx + len(marker))
    section = log_text[idx:next_idx if next_idx > 0 else len(log_text)]
    last = None
    for ln in section.splitlines():
        s = ln.lstrip()
        if s.startswith("[summary] "):
            last = s[len("[summary] "):]
    if not last:
        return None
    try:
        return _json.loads(last)
    except Exception:
        return None


def _runs_jsonl_path() -> Path:
    return get_data_dir() / "logs" / "runs.jsonl"


def _record_run(row: dict) -> None:
    """Persist a run completion to library.db AND append a JSONL line.

    Two stores by design: SQLite for the activity UI's queries; JSONL for
    `tail -f` debugging and downstream scripts. Both reflect the same data.
    """
    import json as _json
    cols = (
        "video_id", "url", "source", "started_at", "ended_at", "duration_secs",
        "exit_code", "success", "stage_reached", "error", "source_lang",
        "used_whisper", "whisper_model",
        "download_secs", "whisper_secs", "frames_secs", "digest_secs", "vision_secs",
        "digest_input_tokens", "digest_output_tokens",
        "digest_cache_read_tokens", "digest_cache_creation_tokens",
        "digest_path",
    )
    placeholders = ", ".join("?" for _ in cols)
    values = tuple(row.get(c) for c in cols)
    try:
        with _library_connect() as conn:
            conn.execute(
                f"INSERT INTO runs ({', '.join(cols)}) VALUES ({placeholders})",
                values,
            )
    except Exception as e:
        # DB failure shouldn't bring down the reaper. JSONL still gets written.
        print(f"[runs] db insert failed: {e}", file=sys.stderr)

    jsonl = _runs_jsonl_path()
    try:
        jsonl.parent.mkdir(parents=True, exist_ok=True)
        with jsonl.open("a") as fh:
            fh.write(_json.dumps(row) + "\n")
    except OSError as e:
        print(f"[runs] jsonl append failed: {e}", file=sys.stderr)


def _recent_runs(limit: int = 100) -> list:
    """Read the last `limit` runs from SQLite, newest first."""
    try:
        with _library_connect() as conn:
            conn.row_factory = __import__("sqlite3").Row
            rows = conn.execute(
                "SELECT * FROM runs ORDER BY started_at DESC LIMIT ?", (limit,)
            ).fetchall()
            return [dict(r) for r in rows]
    except Exception:
        return []


def _build_run_row(info: dict, exit_code: int, log_text: str) -> dict:
    """Translate a finished job + its log into a row dict suitable for runs table."""
    import time as _t
    video_id = info["video_id"]
    started = info["started"]
    ended = _t.time()
    summary = _extract_run_summary(log_text, video_id)
    success = exit_code == 0 and summary is not None
    error = "" if success else _extract_last_error(_oneoff_log_path(), video_id)
    stage = _describe_job_stage(log_text, video_id)
    timings = (summary or {}).get("timings") or {}
    tokens = (summary or {}).get("tokens") or {}
    return {
        "video_id": video_id,
        "url": info.get("url"),
        "source": "oneoff",
        "started_at": started,
        "ended_at": ended,
        "duration_secs": ended - started,
        "exit_code": exit_code,
        "success": 1 if success else 0,
        "stage_reached": stage,
        "error": error or None,
        "source_lang": (summary or {}).get("source_lang"),
        "used_whisper": 1 if (summary or {}).get("used_whisper") else 0,
        "whisper_model": (summary or {}).get("whisper_model"),
        "download_secs": timings.get("download"),
        "whisper_secs": timings.get("whisper"),
        "frames_secs": timings.get("frames"),
        "digest_secs": timings.get("digest"),
        "vision_secs": timings.get("vision"),
        "digest_input_tokens": tokens.get("input"),
        "digest_output_tokens": tokens.get("output"),
        "digest_cache_read_tokens": tokens.get("cache_read"),
        "digest_cache_creation_tokens": tokens.get("cache_creation"),
        "digest_path": (summary or {}).get("digest_path"),
    }


def _oneoff_log_path() -> Path:
    return get_data_dir() / "logs" / "oneoff.log"


def _record_oneoff_failure(info: dict, exit_code: int) -> None:
    import time as _t
    log_path = _oneoff_log_path()
    last_error = _extract_last_error(log_path, info["video_id"]) if log_path.exists() else ""
    _oneoff_failures.insert(0, {
        "video_id": info["video_id"],
        "url": info["url"],
        "exit_code": exit_code,
        "started": info["started"],
        "ended": _t.time(),
        "error": last_error,
    })
    del _oneoff_failures[_ONEOFF_FAILURE_CAP:]


def _list_active_oneoff_jobs() -> list:
    """Return one-off jobs whose subprocesses are still alive.

    Side effect: jobs that have exited are removed from _oneoff_jobs;
    non-zero exits are appended to _oneoff_failures.
    """
    active = []
    for pid in list(_oneoff_jobs.keys()):
        info = _oneoff_jobs[pid]
        proc = info.get("proc")
        if proc is None:
            # legacy entries with no Popen handle — fall back to kill probe
            try:
                os.kill(pid, 0)
                active.append({"pid": pid, **{k: v for k, v in info.items() if k != "proc"}})
            except (ProcessLookupError, PermissionError):
                del _oneoff_jobs[pid]
            continue
        rc = proc.poll()
        if rc is None:
            active.append({"pid": pid, **{k: v for k, v in info.items() if k != "proc"}})
            continue
        del _oneoff_jobs[pid]
        # Treat any non-zero exit, or zero-exit-with-no-digest, as a failure.
        digest_path = get_data_dir() / "digests" / info["video_id"] / "digest.md"
        try:
            log_text = _oneoff_log_path().read_text(errors="replace")
        except OSError:
            log_text = ""
        if rc != 0 or not digest_path.exists():
            _record_oneoff_failure(info, rc)
        _record_run(_build_run_row(info, rc, log_text))
    return active


def _list_recent_oneoff_failures() -> list:
    """Return a copy of recent failures (most recent first)."""
    return list(_oneoff_failures)


# ---- read-state library (SQLite-backed) ----

def _library_path() -> Path:
    return get_data_dir() / "library.db"


def _library_connect():
    """Open (and lazily migrate) the read-state SQLite database."""
    import sqlite3
    get_data_dir().mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(_library_path())
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS digest_reads (
            digest_id TEXT PRIMARY KEY,
            opened_at INTEGER NOT NULL
        );
        -- LLM- or user-assigned topic tags. Composite PK lets the same
        -- tag come from both sources without collision; source filters
        -- on read distinguish LLM noise from user intent.
        CREATE TABLE IF NOT EXISTS digest_topics (
            digest_id TEXT NOT NULL,
            topic     TEXT NOT NULL,
            source    TEXT NOT NULL,   -- 'llm' | 'user'
            added_at  INTEGER NOT NULL,
            PRIMARY KEY (digest_id, topic, source)
        );
        CREATE INDEX IF NOT EXISTS idx_topics_topic
            ON digest_topics(topic);
        CREATE INDEX IF NOT EXISTS idx_topics_digest
            ON digest_topics(digest_id);
        -- Per-digest curation flags + source provenance. Mirrors fields
        -- in metadata.json — the JSON is ground truth, this row is the
        -- query index for list_digests filters.
        CREATE TABLE IF NOT EXISTS digest_meta (
            digest_id      TEXT PRIMARY KEY,
            source_kind    TEXT,           -- 'subscription' | 'oneoff' | 'meta'
            added_at       INTEGER,
            user_dismissed INTEGER NOT NULL DEFAULT 0,
            user_saved     INTEGER NOT NULL DEFAULT 0
        );
        CREATE TABLE IF NOT EXISTS runs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT NOT NULL,
            url TEXT,
            source TEXT NOT NULL,           -- 'oneoff' | 'poll' | 'meta'
            started_at REAL NOT NULL,
            ended_at REAL NOT NULL,
            duration_secs REAL NOT NULL,
            exit_code INTEGER NOT NULL,
            success INTEGER NOT NULL,
            stage_reached TEXT,
            error TEXT,
            source_lang TEXT,
            used_whisper INTEGER NOT NULL DEFAULT 0,
            whisper_model TEXT,
            download_secs REAL,
            whisper_secs REAL,
            frames_secs REAL,
            digest_secs REAL,
            vision_secs REAL,
            digest_input_tokens INTEGER,
            digest_output_tokens INTEGER,
            digest_cache_read_tokens INTEGER,
            digest_cache_creation_tokens INTEGER,
            digest_path TEXT
        );
        CREATE INDEX IF NOT EXISTS idx_runs_started ON runs(started_at DESC);
    """)
    return conn


def _mark_digest_read(digest_id: str) -> None:
    import time
    with _library_connect() as conn:
        conn.execute(
            "INSERT OR REPLACE INTO digest_reads(digest_id, opened_at) VALUES (?, ?)",
            (digest_id, int(time.time())),
        )


def _read_digest_ids() -> set:
    with _library_connect() as conn:
        rows = conn.execute("SELECT digest_id FROM digest_reads").fetchall()
        return {r[0] for r in rows}


# ----------------------------------------------------------------------
# Agent-facing data API (Phase A: parse-on-read from existing markdown,
# cache as JSON sidecar). The principle is that markdown is a render
# target — these functions yield structured data that an agent (or any
# non-HTML caller) can consume without parsing the markdown itself.
#
# Per-artifact JSON files (digest.json, panel.json, takeaway.json) live
# next to the corresponding .md and are regenerated whenever the .md is
# newer. Until Phase B, the .md is the source of truth and the .json is
# a cache; in Phase B that inverts.
# ----------------------------------------------------------------------

import json as _json
import re as _re


def _ts_label_to_seconds(label: str) -> Optional[int]:
    """Convert "M:SS" or "H:MM:SS" into integer seconds. Returns None for
    anything that doesn't fit either shape so callers can keep going."""
    try:
        parts = [int(p) for p in label.split(":")]
    except ValueError:
        return None
    if len(parts) == 2:
        return parts[0] * 60 + parts[1]
    if len(parts) == 3:
        return parts[0] * 3600 + parts[1] * 60 + parts[2]
    return None


def _upload_date_iso(yyyymmdd: Optional[str]) -> Optional[str]:
    """yt-dlp records upload_date as YYYYMMDD with no separator. Convert
    to ISO YYYY-MM-DD for consumers (or pass through if unparseable)."""
    if not yyyymmdd or len(yyyymmdd) != 8 or not yyyymmdd.isdigit():
        return yyyymmdd
    return f"{yyyymmdd[:4]}-{yyyymmdd[4:6]}-{yyyymmdd[6:8]}"


_DIGEST_TOPIC_HEADER_RE = _re.compile(
    # Matches:    "## Title  <sub>*[5:23](https://...)*</sub>"
    # And also:   "## Title  <sub>*5:23*</sub>"  (older format, no link)
    r"^## (?P<title>.+?)\s+<sub>\*"
    r"(?:\[(?P<ts_linked>\d+(?::\d+){1,2})\]\((?P<ts_url>[^)]+)\)"
    r"|(?P<ts_bare>\d+(?::\d+){1,2}))"
    r"\*</sub>\s*$",
    _re.M,
)

_DIGEST_IMG_RE = _re.compile(r'<img src="([^"]+)"[^>]*>\s*')


def digest_md_to_json(
    md_path: Path,
    metadata_path: Optional[Path] = None,
) -> dict:
    """Parse a digest.md file (the format emitted by write_markdown_digest)
    into structured JSON: {video, overview, topics: [...]}.

    Pulls video metadata (channel, upload date) from the sibling
    metadata.json when present — those fields aren't recoverable from
    the markdown alone. Sections that fail to parse fall back to a raw
    chunk so an agent can still see the source text.
    """
    text = md_path.read_text()

    video_meta = {}
    if metadata_path is None:
        metadata_path = md_path.parent / "metadata.json"
    if metadata_path.exists():
        try:
            video_meta = _json.loads(metadata_path.read_text())
        except Exception:
            video_meta = {}

    title_match = _re.match(r"^# (.+?)\n", text)
    title = (
        title_match.group(1).strip()
        if title_match
        else (video_meta.get("title") or "")
    )

    url_match = _re.search(
        r"^\*\*Watch on YouTube:\*\* <(.+?)>\s*$", text, _re.M
    )
    video_url = (
        url_match.group(1) if url_match else (video_meta.get("url") or "")
    )

    matches = list(_DIGEST_TOPIC_HEADER_RE.finditer(text))

    # Overview: text between the header block and the first topic. Strip
    # the title line and Watch-on-YouTube line so only the prose remains.
    if matches:
        pre = text[: matches[0].start()]
    else:
        pre = text
    pre = _re.sub(r"^# .+?\n", "", pre, count=1)
    pre = _re.sub(
        r"^\*\*Watch on YouTube:\*\* <.+?>\s*$", "", pre, count=1, flags=_re.M
    )
    overview = pre.strip()

    topics = []
    for i, m in enumerate(matches):
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[m.end():end].strip()

        ts_label = m.group("ts_linked") or m.group("ts_bare")
        ts_url = m.group("ts_url")

        img_match = _DIGEST_IMG_RE.match(block)
        image = img_match.group(1) if img_match else None
        if img_match:
            block = block[img_match.end():].strip()

        # Topic body is the paragraph(s) before the first `- ` bullet.
        # The current writer emits "body then bullets" with no trailing
        # paragraph, so this split is lossless against present output.
        lines = block.split("\n")
        first_bullet = next(
            (j for j, ln in enumerate(lines) if ln.startswith("- ")), None
        )
        if first_bullet is None:
            body_para = block.strip()
            bullets: list = []
        else:
            body_para = "\n".join(lines[:first_bullet]).strip()
            bullets = [
                ln[2:].strip() for ln in lines[first_bullet:] if ln.startswith("- ")
            ]

        topics.append({
            "index": i + 1,
            "title": m.group("title").strip(),
            "ts_start_s": _ts_label_to_seconds(ts_label) if ts_label else None,
            "ts_link": ts_url,
            "image": image,
            "body": body_para,
            "bullets": bullets,
        })

    return {
        "video": {
            "id": video_meta.get("video_id") or md_path.parent.name,
            "title": title,
            "url": video_url,
            "channel": video_meta.get("channel_name") or "",
            "channel_url": video_meta.get("channel_url") or "",
            "published_at": _upload_date_iso(video_meta.get("upload_date")),
        },
        "overview": overview,
        "topics": topics,
    }


def render_digest_md(d: dict) -> str:
    """Render a digest JSON shape back to the markdown the viewer + the
    existing parser expect. Inverse of `digest_md_to_json` — kept
    byte-compatible so a re-render produces no gratuitous diff.

    Used by `write_markdown_digest` (Phase B: JSON is the source of
    truth; markdown is a deterministic view of it)."""
    lines: list = []
    lines.append(f"# {d['video']['title']}")
    lines.append("")
    if d["video"].get("url"):
        lines.append(f"**Watch on YouTube:** <{d['video']['url']}>")
        lines.append("")
    if d.get("overview"):
        lines.append(d["overview"])
        lines.append("")
    for t in d.get("topics", []):
        ts_s = t.get("ts_start_s")
        ts_label = format_timestamp(ts_s) if ts_s is not None else ""
        if t.get("ts_link"):
            lines.append(f"## {t['title']}  <sub>*[{ts_label}]({t['ts_link']})*</sub>")
        elif ts_label:
            lines.append(f"## {t['title']}  <sub>*{ts_label}*</sub>")
        else:
            lines.append(f"## {t['title']}")
        lines.append("")
        if t.get("image"):
            alt = (t["title"] or "").replace('"', "'")
            lines.append(f'<img src="{t["image"]}" alt="{alt}" width="800">')
            lines.append("")
        if t.get("body"):
            lines.append(t["body"])
            lines.append("")
        for b in t.get("bullets") or []:
            lines.append(f"- {b}")
        if t.get("bullets"):
            lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def _digest_json_from_object(
    digest_obj,
    *,
    topic_images: list,
    images_rel_dir: str,
    video_title: Optional[str],
    video_url: Optional[str],
    video_id: Optional[str] = None,
    metadata_path: Optional[Path] = None,
) -> dict:
    """Build a digest.json dict directly from the in-memory `Digest`
    Pydantic object + frame picks. Bypasses the markdown parser — this
    is the Phase B path for new digests. Output shape is identical to
    what `digest_md_to_json` produces for the rendered markdown, so the
    two code paths are interchangeable from the read side."""
    heading = video_title if video_title else digest_obj.title

    topics_json = []
    for i, t in enumerate(digest_obj.topics):
        ts_start_s = int(t.start_time)
        ts_link = None
        if video_url:
            sep = "&" if "?" in video_url else "?"
            ts_link = f"{video_url}{sep}t={ts_start_s}s"
        img = None
        if i < len(topic_images) and topic_images[i] is not None:
            img = f"{images_rel_dir}/{topic_images[i].name}"
        topics_json.append({
            "index": i + 1,
            "title": t.title,
            "ts_start_s": ts_start_s,
            "ts_link": ts_link,
            "image": img,
            "body": t.summary,
            "bullets": list(t.key_points),
        })

    channel = ""
    channel_url = ""
    published_at = None
    if metadata_path and metadata_path.exists():
        try:
            meta = _json.loads(metadata_path.read_text())
            channel = meta.get("channel_name") or ""
            channel_url = meta.get("channel_url") or ""
            published_at = _upload_date_iso(meta.get("upload_date"))
        except Exception:
            pass

    return {
        "video": {
            "id": video_id or "",
            "title": heading,
            "url": video_url or "",
            "channel": channel,
            "channel_url": channel_url,
            "published_at": published_at,
        },
        "overview": digest_obj.overview,
        "topics": topics_json,
    }


def _atomic_write_text(path: Path, content: str) -> None:
    """Write `content` to `path` via a .tmp suffix + rename so a reader
    never sees a half-written file."""
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(content)
    tmp.replace(path)


def _atomic_write_json(path: Path, obj) -> None:
    """JSON-encode and atomic-write."""
    _atomic_write_text(path, _json.dumps(obj, indent=2, ensure_ascii=False))


def _load_artifact_json(
    digest_id: str,
    *,
    artifact: str,
    parser,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Shared cache logic for digest.json / panel.json / takeaway.json.
    Re-parses when the .md is newer than the .json sidecar (or the
    sidecar doesn't exist / is corrupt). Cache write is best-effort —
    a parse failure on disk doesn't block the in-memory return.
    """
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    dir_path = digests_dir / digest_id
    md_path = dir_path / f"{artifact}.md"
    json_path = dir_path / f"{artifact}.json"
    if not md_path.exists():
        raise FileNotFoundError(f"No {artifact}.md at {dir_path}")

    if json_path.exists():
        try:
            if json_path.stat().st_mtime >= md_path.stat().st_mtime:
                return _json.loads(json_path.read_text())
        except Exception:
            pass  # corrupt or unreadable — fall through to re-parse

    parsed = parser(md_path)
    try:
        json_path.write_text(_json.dumps(parsed, indent=2, ensure_ascii=False))
    except Exception:
        pass  # caching is best-effort
    return parsed


def load_digest_json(
    digest_id: str, *, digests_dir: Optional[Path] = None,
) -> dict:
    """Get the structured digest for a video. Parses digest.md on first
    call, caches the result to digest.json, and returns from cache on
    subsequent calls until digest.md changes."""
    return _load_artifact_json(
        digest_id,
        artifact="digest",
        parser=digest_md_to_json,
        digests_dir=digests_dir,
    )


_PANEL_TURN_RE = _re.compile(
    r"^\*\*(?P<speaker>[^*\n]+?):\*\*\s*(?P<text>.+?)(?=\n\s*\n|\Z)",
    _re.M | _re.S,
)
# Panelist intro: "**Name** — bio" or "**Name** is bio" (both shapes appear).
_PANELIST_INTRO_RE = _re.compile(
    r"^\*\*(?P<name>[^*\n]+?)\*\*\s*(?:[—–-]\s*|is\s+)(?P<bio>.+?)(?=\n\s*\n|\Z)",
    _re.M | _re.S,
)


def panel_text_to_json(text: str) -> dict:
    """Core panel parser; operates on a string so the generator can
    parse in-memory before writing both .md and .json atomically.
    `panel_md_to_json(path)` is a thin wrapper over this."""
    title_match = _re.search(r"^## (.+?)\s*$", text, _re.M)
    title = title_match.group(1).strip() if title_match else ""

    parts = _re.split(r"\n---+\s*\n", text, maxsplit=1)
    head = parts[0]
    body = parts[1] if len(parts) == 2 else ""

    panelists = []
    for m in _PANELIST_INTRO_RE.finditer(head):
        panelists.append({
            "name": m.group("name").strip(),
            "bio": " ".join(m.group("bio").split()).strip(),
        })

    turns = []
    for m in _PANEL_TURN_RE.finditer(body):
        turns.append({
            "speaker": m.group("speaker").strip(),
            "text": " ".join(m.group("text").split()).strip(),
        })

    return {"title": title, "panelists": panelists, "turns": turns}


def panel_md_to_json(md_path: Path) -> dict:
    """Parse panel.md (as produced by generate_panel_discussion) into
    {title, panelists: [{name, bio}], turns: [{speaker, text}]}.

    The "## The Panel" / "## Panel: ..." header is captured as `title`.
    Anything between the header and the first `---` line is the
    panelists block; anything after `---` is the discussion."""
    return panel_text_to_json(md_path.read_text())


def load_panel_json(
    digest_id: str, *, digests_dir: Optional[Path] = None,
) -> dict:
    """Get the structured panel discussion. Parses panel.md on first
    call, caches to panel.json, re-parses on mtime invalidation."""
    return _load_artifact_json(
        digest_id, artifact="panel",
        parser=panel_md_to_json, digests_dir=digests_dir,
    )


_TAKEAWAY_CITATION_RE = _re.compile(
    r"\[(?P<label>\d+(?::\d+){1,2})\]\((?P<url>[^)]+)\)"
)


def takeaway_text_to_json(text: str) -> dict:
    """Core takeaway parser; operates on a string so the generator can
    parse in-memory before writing both .md and .json atomically.
    `takeaway_md_to_json(path)` is a thin wrapper over this."""
    text = text.strip()
    paragraphs = []
    for block in _re.split(r"\n\s*\n+", text):
        block = block.strip()
        if not block:
            continue
        citations = []
        for m in _TAKEAWAY_CITATION_RE.finditer(block):
            citations.append({
                "label": m.group("label"),
                "ts_s": _ts_label_to_seconds(m.group("label")),
                "url": m.group("url"),
            })
        paragraphs.append({"text": block, "citations": citations})
    return {"paragraphs": paragraphs}


def takeaway_md_to_json(md_path: Path) -> dict:
    """Parse takeaway.md (a small block of plain prose with inline
    timestamp links) into {paragraphs: [{text, citations}]}.

    Citations are `[M:SS](url)` or `[H:MM:SS](url)` links — collected
    per paragraph so an agent can deep-link to the source claim without
    re-parsing markdown."""
    return takeaway_text_to_json(md_path.read_text())


def load_takeaway_json(
    digest_id: str, *, digests_dir: Optional[Path] = None,
) -> dict:
    """Get the structured takeaway. Parses takeaway.md on first call,
    caches to takeaway.json, re-parses on mtime invalidation."""
    return _load_artifact_json(
        digest_id, artifact="takeaway",
        parser=takeaway_md_to_json, digests_dir=digests_dir,
    )


def read_digest(
    digest_id: str,
    section: str = "full",
    *,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Agent-facing read for a single digest.

    section: one of "full", "meta", "overview", "topics",
             "topic:<N>" (1-indexed), "topic:<slug>" (substring match
             on topic title; first hit wins).
             "panel" / "takeaway" join after Phase A — until then they
             raise NotImplementedError.

    Every return includes a `video` block so the caller can ground the
    content (link back to YouTube, render the title, etc.) without
    juggling separate lookups.
    """
    d = load_digest_json(digest_id, digests_dir=digests_dir)
    video = d["video"]

    if section == "full":
        return {"section": section, "video": video, "content": d}
    if section == "meta":
        return {"section": section, "video": video, "content": video}
    if section == "overview":
        return {
            "section": section, "video": video,
            "content": {"overview": d.get("overview", "")},
        }
    if section == "topics":
        return {
            "section": section, "video": video,
            "content": {"topics": d.get("topics", [])},
        }
    if section.startswith("topic:"):
        ref = section.split(":", 1)[1]
        topics = d.get("topics") or []
        topic = None
        if ref.isdigit():
            n = int(ref)
            if 1 <= n <= len(topics):
                topic = topics[n - 1]
        else:
            needle = ref.lower()
            for t in topics:
                if needle in (t.get("title") or "").lower():
                    topic = t
                    break
        if topic is None:
            raise ValueError(
                f"No topic matching {ref!r} in {digest_id} "
                f"(have {len(topics)} topics)"
            )
        return {"section": section, "video": video, "content": topic}
    if section == "panel" or section.startswith("panel:"):
        try:
            panel = load_panel_json(digest_id, digests_dir=digests_dir)
        except FileNotFoundError:
            raise ValueError(f"No panel for {digest_id}")
        if section == "panel":
            return {"section": section, "video": video, "content": panel}
        ref = section.split(":", 1)[1]
        # panel:turn:<N>  → single turn (1-indexed)
        # panel:panelists → just the panelist list
        if ref == "panelists":
            return {
                "section": section, "video": video,
                "content": {"panelists": panel.get("panelists", [])},
            }
        if ref.startswith("turn:"):
            n_str = ref.split(":", 1)[1]
            turns = panel.get("turns") or []
            if n_str.isdigit():
                n = int(n_str)
                if 1 <= n <= len(turns):
                    return {
                        "section": section, "video": video,
                        "content": turns[n - 1],
                    }
            raise ValueError(
                f"No turn matching {n_str!r} (have {len(turns)} turns)"
            )
        raise ValueError(f"unknown panel sub-section: {ref!r}")
    if section == "takeaway":
        try:
            tk = load_takeaway_json(digest_id, digests_dir=digests_dir)
        except FileNotFoundError:
            raise ValueError(f"No takeaway for {digest_id}")
        return {"section": section, "video": video, "content": tk}
    raise ValueError(f"unknown section: {section!r}")


def list_digests(
    *,
    channel: Optional[str] = None,
    since: Optional[str] = None,
    unread: bool = False,
    q: Optional[str] = None,
    topic: Optional[str] = None,
    source: Optional[str] = None,
    saved: Optional[bool] = None,
    dismissed: Optional[bool] = None,
    limit: int = 20,
    digests_dir: Optional[Path] = None,
) -> list:
    """Agent-facing index of the local digest library. Returns lightweight
    entries (no topic body) so a caller can scan, then drill in with
    `read_digest(id, section=...)`.

    Filters compose (AND):
      channel   : substring match against channel name (case-insensitive)
      since     : ISO date "YYYY-MM-DD"; keeps digests whose published_at
                  is >= the given date. Digests with no published_at fall
                  back to the on-disk mtime as a coarse proxy.
      unread    : True → only digests not yet marked read.
      q         : substring match against title (case-insensitive).
      topic     : exact-match on a topic tag (LLM or user-assigned).
      source    : 'subscription' | 'oneoff' | 'meta' (provenance filter).
      saved     : True → only user-saved; False → only un-saved.
      dismissed : True → only user-dismissed; False → only not-dismissed.
                  Defaults None for both → don't filter on the flag.
      limit     : max results returned. The list is sorted most-recent
                  first; limit caps the head of that list.
    """
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if not digests_dir.exists():
        return []

    read_ids = _read_digest_ids() if unread else set()

    needle_channel = (channel or "").strip().lower() or None
    needle_q = (q or "").strip().lower() or None

    # Pre-fetch the digests matching topic / source / saved / dismissed
    # filters via the indexed tables, so we can short-circuit large
    # libraries without parsing every metadata.json.
    topic_match: Optional[set] = None
    if topic:
        with _library_connect() as conn:
            rows = conn.execute(
                "SELECT DISTINCT digest_id FROM digest_topics WHERE topic = ?",
                (topic.strip().lower(),),
            ).fetchall()
        topic_match = {r[0] for r in rows}

    # digest_meta join: we read all rows once and look up per-id since
    # the table is at most O(library size).
    meta_index = {}
    with _library_connect() as conn:
        for row in conn.execute(
            "SELECT digest_id, source_kind, user_saved, user_dismissed "
            "FROM digest_meta"
        ):
            meta_index[row[0]] = {
                "source_kind": row[1],
                "user_saved": bool(row[2]),
                "user_dismissed": bool(row[3]),
            }

    dirs = sorted(
        (d for d in digests_dir.iterdir()
         if d.is_dir() and (d / "digest.md").exists()),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )

    out = []
    for d in dirs:
        if unread and d.name in read_ids:
            continue
        if topic_match is not None and d.name not in topic_match:
            continue

        meta_row = meta_index.get(d.name, {})
        if source and meta_row.get("source_kind") != source:
            continue
        if saved is not None and meta_row.get("user_saved", False) != saved:
            continue
        if dismissed is not None and meta_row.get("user_dismissed", False) != dismissed:
            continue

        meta: dict = {}
        meta_path = d / "metadata.json"
        if meta_path.exists():
            try:
                meta = _json.loads(meta_path.read_text())
            except Exception:
                meta = {}

        # Cheap title pull (digest.json may not exist yet on first list)
        title = meta.get("title") or d.name
        if not meta.get("title"):
            try:
                for line in (d / "digest.md").read_text().splitlines():
                    if line.startswith("# "):
                        title = line[2:].strip()
                        break
            except Exception:
                pass

        ch_name = meta.get("channel_name") or ""
        if needle_channel and needle_channel not in ch_name.lower():
            continue
        if needle_q and needle_q not in title.lower():
            continue

        published = _upload_date_iso(meta.get("upload_date"))
        if since:
            ref = published or _iso_from_mtime(d.stat().st_mtime)
            if ref < since:
                continue

        # Topics: union of LLM + user, deduped, sorted for stable output.
        topics_split = _read_topics_for_digest(d.name)
        all_topics = sorted(set(topics_split["llm"] + topics_split["user"]))

        out.append({
            "id": d.name,
            "title": title,
            "url": meta.get("url") or f"https://www.youtube.com/watch?v={d.name}",
            "channel": ch_name,
            "channel_url": meta.get("channel_url") or "",
            "published_at": published,
            "mtime": d.stat().st_mtime,
            "read": d.name in read_ids if unread else (d.name in _read_digest_ids()),
            "topics": all_topics,
            "topics_split": topics_split,
            "source": meta.get("source") or {"kind": meta_row.get("source_kind")},
            "user_saved": meta_row.get("user_saved", False),
            "user_dismissed": meta_row.get("user_dismissed", False),
            "has_panel": (d / "panel.md").exists(),
            "has_takeaway": (d / "takeaway.md").exists(),
            "has_slides": (d / "slides.pptx").exists(),
            "has_audio": {
                "digest":   (d / "digest.mp3").exists(),
                "panel":    (d / "panel.mp3").exists(),
                "takeaway": (d / "takeaway.mp3").exists(),
            },
        })
        if len(out) >= limit:
            break
    return out


def _iso_from_mtime(mtime: float) -> str:
    """YYYY-MM-DD for a filesystem mtime — used as a fallback when a
    digest has no upload_date in metadata.json."""
    import datetime as _dt
    return _dt.date.fromtimestamp(mtime).isoformat()


# ----------------------------------------------------------------------
# Topic tagging — LLM-derived (Haiku) topic tags per digest, plus user-
# curation flags. The metadata layer that lets the agent answer
# questions like "what AI-policy stuff have I read this month?" without
# scanning every digest body.
# ----------------------------------------------------------------------

DEFAULT_TAGGING_MODEL = "claude-haiku-4-5-20251001"
TAXONOMY_TOP_N = 50         # how many existing tags to surface to the LLM
TAG_RE = _re.compile(r"^[a-z0-9]+(?:-[a-z0-9]+)*$")

_TAGGING_SYSTEM_PROMPT = """\
You tag a video digest for retrieval. Given the digest, return 2-5
topic tags as JSON.

Rules:
- Tags are lowercase, hyphen-separated, 1-3 words: "agent-pricing",
  "distributed-systems", "ai-policy", "vc-strategy".
- REUSE existing tags from the taxonomy when they fit. Do not create
  synonyms of existing tags (no "ai-regulation" if "ai-policy" exists).
- Only propose a NEW tag when no existing tag fits. Mark new tags with
  is_new: true so the system can flag vocabulary drift.
- Tags describe SUBJECT matter, not format. "interview" is wrong;
  "venture-capital" is right.
- Quality over quantity. 3 sharp tags > 5 fuzzy ones.

Output strict JSON only:
  {"tags": [{"tag": "<slug>", "is_new": <bool>}, ...], "reason": "<one sentence>"}
"""


def _current_taxonomy(limit: int = TAXONOMY_TOP_N) -> list:
    """Return [(tag, n_digests), ...] sorted by frequency desc. Drives
    the 'reuse existing tags' guidance in the tagging prompt."""
    with _library_connect() as conn:
        rows = conn.execute(
            "SELECT topic, COUNT(DISTINCT digest_id) AS n "
            "FROM digest_topics WHERE source = 'llm' "
            "GROUP BY topic ORDER BY n DESC, topic ASC LIMIT ?",
            (limit,),
        ).fetchall()
    return [(r[0], r[1]) for r in rows]


def _topic_tagging_user_prompt(d: dict, taxonomy: list) -> str:
    """Build the per-digest user message for the tagging call.
    `d` is a digest.json dict."""
    if taxonomy:
        tax_lines = "\n".join(f"  {tag} ({n})" for tag, n in taxonomy)
    else:
        tax_lines = "  (empty — this is the first digest being tagged)"
    topics_block = "\n".join(
        f"  {i + 1}. {t.get('title', '')}"
        for i, t in enumerate(d.get("topics") or [])
    )
    return (
        f"EXISTING TAXONOMY (most-used first; reuse where reasonable):\n"
        f"{tax_lines}\n\n"
        f"DIGEST:\n"
        f"Title: {d['video'].get('title', '')}\n\n"
        f"Overview: {d.get('overview', '')}\n\n"
        f"Topic titles:\n{topics_block}\n"
    )


def _validate_tag(tag: str) -> Optional[str]:
    """Return a normalized tag or None if it fails the format check.
    Guards against the LLM emitting CamelCase, spaces, punctuation, etc."""
    if not isinstance(tag, str):
        return None
    norm = tag.strip().lower().replace("_", "-").replace(" ", "-")
    norm = _re.sub(r"-+", "-", norm).strip("-")
    if not norm or len(norm) > 40:
        return None
    if not TAG_RE.match(norm):
        return None
    return norm


def tag_digest_via_llm(
    digest_id: str,
    *,
    digests_dir: Optional[Path] = None,
    model: Optional[str] = None,
    backend=None,
) -> dict:
    """Run the Haiku tagging step for a single digest. Persists tags to
    metadata.json (under `topics`) AND to digest_topics with source=llm.
    Returns {tags: [...], new_tags: [...], reason: str}."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if backend is None:
        backend = select_backend()
    if model is None:
        model = DEFAULT_TAGGING_MODEL

    d = load_digest_json(digest_id, digests_dir=digests_dir)
    taxonomy = _current_taxonomy()
    user_text = _topic_tagging_user_prompt(d, taxonomy)

    # Use the backend's structured-output path: a Pydantic model that
    # forces the LLM to emit the right shape. Same pattern as
    # generate_digest.
    from pydantic import BaseModel
    from typing import List as TList

    class _Tag(BaseModel):
        tag: str
        is_new: bool = False

    class _Response(BaseModel):
        tags: TList[_Tag]
        reason: str = ""

    parsed, _usage = backend.parse(
        system=_TAGGING_SYSTEM_PROMPT,
        user_text=user_text,
        schema=_Response,
        model=model,
        max_tokens=400,
    )
    record_llm_usage(
        video_id=digest_id, kind="tagging", model=model,
        backend_name=backend.name, usage=_usage,
    )

    cleaned: list = []
    new_tags: list = []
    taxonomy_set = {t for t, _ in taxonomy}
    for entry in parsed.tags:
        norm = _validate_tag(entry.tag)
        if norm is None:
            continue
        if norm in cleaned:
            continue
        cleaned.append(norm)
        if entry.is_new and norm not in taxonomy_set:
            new_tags.append(norm)
    cleaned = cleaned[:5]  # hard cap

    _write_topics_to_metadata(digest_id, cleaned, digests_dir=digests_dir)
    _write_topics_to_db(digest_id, cleaned, source="llm")
    return {"tags": cleaned, "new_tags": new_tags, "reason": parsed.reason}


def _write_topics_to_metadata(
    digest_id: str,
    topics: list,
    *,
    digests_dir: Optional[Path] = None,
) -> None:
    """Persist LLM tags into metadata.json. Preserves existing fields
    (channel info, user_tags, etc.) — merges rather than replaces."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    meta_path = digests_dir / digest_id / "metadata.json"
    meta = {}
    if meta_path.exists():
        try:
            meta = _json.loads(meta_path.read_text())
        except Exception:
            meta = {}
    meta["topics"] = list(topics)
    _atomic_write_json(meta_path, meta)


def _write_topics_to_db(
    digest_id: str,
    topics: list,
    *,
    source: str,
) -> None:
    """Insert into digest_topics (PK conflicts → no-op, idempotent).
    `source` is 'llm' or 'user' — kept separate so re-tagging from one
    side never wipes the other."""
    import time as _t
    now = int(_t.time())
    with _library_connect() as conn:
        # Replace LLM tags wholesale (re-tagging means new vocabulary).
        # User tags are append-only — handled by tag_digest separately.
        if source == "llm":
            conn.execute(
                "DELETE FROM digest_topics WHERE digest_id = ? AND source = 'llm'",
                (digest_id,),
            )
        for tag in topics:
            conn.execute(
                "INSERT OR IGNORE INTO digest_topics "
                "(digest_id, topic, source, added_at) VALUES (?, ?, ?, ?)",
                (digest_id, tag, source, now),
            )


def _read_topics_for_digest(digest_id: str) -> dict:
    """Return {llm: [...], user: [...]} for a digest. Sorted, unique."""
    with _library_connect() as conn:
        rows = conn.execute(
            "SELECT topic, source FROM digest_topics WHERE digest_id = ? "
            "ORDER BY source ASC, topic ASC",
            (digest_id,),
        ).fetchall()
    out = {"llm": [], "user": []}
    for tag, src in rows:
        out.setdefault(src, []).append(tag)
    return out


def _digest_meta_row(digest_id: str) -> dict:
    """Return the digest_meta row as a dict, or {} if absent."""
    with _library_connect() as conn:
        row = conn.execute(
            "SELECT source_kind, added_at, user_dismissed, user_saved "
            "FROM digest_meta WHERE digest_id = ?",
            (digest_id,),
        ).fetchone()
    if not row:
        return {}
    return {
        "source_kind": row[0],
        "added_at": row[1],
        "user_dismissed": bool(row[2]),
        "user_saved": bool(row[3]),
    }


def _record_digest_added(
    digest_id: str,
    *,
    source_kind: str = "oneoff",
) -> None:
    """Stamp source provenance into digest_meta on ingestion. Idempotent
    — re-ingest preserves the original added_at + source_kind."""
    import time as _t
    now = int(_t.time())
    with _library_connect() as conn:
        conn.execute(
            "INSERT OR IGNORE INTO digest_meta "
            "(digest_id, source_kind, added_at) VALUES (?, ?, ?)",
            (digest_id, source_kind, now),
        )


# Per-section weight when computing search relevance. Title hits matter
# most (closest to user intent), section bodies count progressively
# less. Tuned for "what would I want surfaced first if I searched for X
# across my library" — i.e. matches on the title of a video are almost
# always more relevant than matches deep in a panelist's seventh turn.
_SEARCH_WEIGHTS = {
    "video.title":   10,
    "overview":       5,
    "topic.title":    4,
    "takeaway":       3,
    "topic.body":     2,
    "topic.bullet":   2,
    "panel.turn":     1,
}


def _snippet(text: str, query: str, *, radius: int = 80) -> str:
    """Return a ~radius*2-char excerpt of `text` centered on the first
    case-insensitive occurrence of any token from `query`. Falls back
    to the head of the text when no token is found."""
    tokens = [t for t in query.lower().split() if t]
    low = text.lower()
    idx = -1
    for tok in tokens:
        i = low.find(tok)
        if i >= 0 and (idx < 0 or i < idx):
            idx = i
    if idx < 0:
        return text[: radius * 2].strip() + ("…" if len(text) > radius * 2 else "")
    start = max(0, idx - radius)
    end = min(len(text), idx + radius)
    out = text[start:end].strip()
    if start > 0:
        out = "…" + out
    if end < len(text):
        out = out + "…"
    return out


def _matches_all_tokens(text: str, tokens: list) -> bool:
    """True iff every token appears in text (case-insensitive)."""
    if not text:
        return False
    low = text.lower()
    return all(tok in low for tok in tokens)


def search_library(
    q: str,
    *,
    k: int = 10,
    digests_dir: Optional[Path] = None,
) -> list:
    """Substring search across the whole library (digest / panel /
    takeaway). v1 is case-insensitive whole-token AND: every space-
    separated token in `q` must appear in the same section to count as
    a hit. Embeddings are out of scope here — this is the cheap
    deterministic baseline an agent can rely on.

    Returns up to k hits, each:
      {digest_id, title, section, snippet, score, url, video}

    `section` is the same string `read_digest(section=...)` understands,
    so an agent can pipe a hit straight into a follow-up read for the
    full content."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if not digests_dir.exists():
        return []

    tokens = [t.lower() for t in (q or "").split() if t]
    if not tokens:
        return []

    hits = []
    for entry in list_digests(limit=10_000, digests_dir=digests_dir):
        did = entry["id"]
        try:
            d = load_digest_json(did, digests_dir=digests_dir)
        except Exception:
            continue
        video_meta = {"id": did, "title": entry["title"], "url": entry["url"]}

        def add(section, source_text, weight):
            if not _matches_all_tokens(source_text, tokens):
                return
            hits.append({
                "digest_id": did,
                "title": entry["title"],
                "section": section,
                "snippet": _snippet(source_text, q),
                "score": weight,
                "url": entry["url"],
                "video": video_meta,
            })

        add("meta",     d.get("video", {}).get("title", ""), _SEARCH_WEIGHTS["video.title"])
        add("overview", d.get("overview", ""),               _SEARCH_WEIGHTS["overview"])
        for t in d.get("topics", []):
            ref = f"topic:{t['index']}"
            add(ref, t.get("title", ""), _SEARCH_WEIGHTS["topic.title"])
            add(ref, t.get("body", ""),  _SEARCH_WEIGHTS["topic.body"])
            for b in t.get("bullets", []):
                add(ref, b, _SEARCH_WEIGHTS["topic.bullet"])

        # Panel + takeaway are optional artifacts — read defensively.
        if entry["has_panel"]:
            try:
                panel = load_panel_json(did, digests_dir=digests_dir)
                for i, turn in enumerate(panel.get("turns", []), 1):
                    add(f"panel:turn:{i}", turn.get("text", ""),
                        _SEARCH_WEIGHTS["panel.turn"])
            except Exception:
                pass
        if entry["has_takeaway"]:
            try:
                tk = load_takeaway_json(did, digests_dir=digests_dir)
                for para in tk.get("paragraphs", []):
                    add("takeaway", para.get("text", ""),
                        _SEARCH_WEIGHTS["takeaway"])
            except Exception:
                pass

    # Sort by score desc, then most-recent digest first (caller-supplied
    # mtime not in hits, so re-derive a tiebreak from digest_id order in
    # the already-mtime-sorted list_digests output).
    digest_order = {entry["id"]: i for i, entry in enumerate(
        list_digests(limit=10_000, digests_dir=digests_dir)
    )}
    hits.sort(key=lambda h: (-h["score"], digest_order.get(h["digest_id"], 0)))
    return hits[:k]


# ---- Module-level artifact builders ------------------------------------
#
# Originally these lived inside cmd_serve() and closed over `digests_dir`.
# Lifted to module level so the agent-facing API can call them without
# spinning up a Flask app. cmd_serve's inner wrappers (kept for back-compat
# with the existing call sites) delegate here.

def _resolve_cached_srt_for(
    video_id: str,
    digests_dir: Optional[Path] = None,
) -> Tuple[Path, str]:
    """Locate the cached SRT for a digest and pull the BCP-47 lang from
    its filename. Raises RuntimeError when the transcript isn't on disk
    (means the video needs to be re-digested first)."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    srt_dir = digests_dir / video_id / "downloads" / video_id
    srt_files = list(srt_dir.glob("*.srt")) if srt_dir.exists() else []
    if not srt_files:
        raise RuntimeError(
            "No cached transcript. Re-digest the video first."
        )
    srt_path = srt_files[0]
    lang = (
        srt_path.stem[len(video_id) + 1:]
        if srt_path.stem.startswith(video_id + ".") else "en"
    )
    return srt_path, lang


def build_panel_for_video(
    video_id: str,
    *,
    digests_dir: Optional[Path] = None,
) -> Path:
    """Run the panel-generation pipeline for a digested video. Synchronous;
    raises on failure. Atomic write via .tmp → rename. Returns the
    written panel.md path."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    digest_md = digests_dir / video_id / "digest.md"
    if not digest_md.exists():
        raise RuntimeError("digest.md missing — re-digest first.")
    srt_path, lang = _resolve_cached_srt_for(video_id, digests_dir)
    s = load_settings()
    backend = select_backend()
    panel_model = s.get("panel_model") or DEFAULT_PANEL_MODEL
    segments = parse_srt(srt_path)
    text, p_usage = generate_panel_discussion(
        digest_md.read_text(), segments,
        model=panel_model, source_lang=lang,
        output_language=s.get("digest_language") or "auto",
        backend=backend,
    )
    record_llm_usage(
        video_id=video_id, kind="panel", model=panel_model,
        backend_name=backend.name, usage=p_usage,
    )
    # Phase B: parse the LLM output to JSON once, write both files
    # atomically. Read path then never invokes the parser on this digest.
    panel_path = digests_dir / video_id / "panel.md"
    panel_json_path = digests_dir / video_id / "panel.json"
    _atomic_write_text(panel_path, text)
    _atomic_write_json(panel_json_path, panel_text_to_json(text))
    return panel_path


def build_takeaway_for_video(
    video_id: str,
    *,
    digests_dir: Optional[Path] = None,
) -> Path:
    """Run the takeaway-generation pipeline. Uses panel.md when present
    so the takeaway can integrate the panel's critique; copes without
    it. Synchronous; raises on failure. Returns takeaway.md path."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    digest_md = digests_dir / video_id / "digest.md"
    if not digest_md.exists():
        raise RuntimeError("digest.md missing — re-digest first.")
    srt_path, lang = _resolve_cached_srt_for(video_id, digests_dir)
    panel_md_path = digests_dir / video_id / "panel.md"
    panel_text = panel_md_path.read_text() if panel_md_path.exists() else None
    s = load_settings()
    backend = select_backend()
    takeaway_model = (os.environ.get("YT2MD_TAKEAWAY_MODEL")
                      or DEFAULT_TAKEAWAY_MODEL)
    segments = parse_srt(srt_path)
    takeaway_text, t_usage = generate_takeaway_prose(
        digest_md.read_text(), panel_text, segments,
        model=takeaway_model,
        publish_date=None,
        source_lang=lang,
        output_language=s.get("digest_language") or "auto",
        backend=backend,
    )
    record_llm_usage(
        video_id=video_id, kind="takeaway", model=takeaway_model,
        backend_name=backend.name, usage=t_usage,
    )
    body = render_takeaway_markdown(
        takeaway_text,
        video_url=f"https://www.youtube.com/watch?v={video_id}",
    )
    # Phase B: parse body to JSON once, write both files atomically.
    takeaway_path = digests_dir / video_id / "takeaway.md"
    takeaway_json_path = digests_dir / video_id / "takeaway.json"
    _atomic_write_text(takeaway_path, body)
    _atomic_write_json(takeaway_json_path, takeaway_text_to_json(body))
    return takeaway_path


def build_slides_for_video(
    video_id: str,
    *,
    digests_dir: Optional[Path] = None,
) -> Path:
    """Run the slides-only pipeline against the cached MP4 + SRT for
    a digested video. Atomic write. Returns the slides.pptx path."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    cache_dir = digests_dir / video_id / "downloads" / video_id
    if not cache_dir.exists():
        raise RuntimeError("No cached video. Re-digest the video first.")
    mp4_path = cache_dir / f"{video_id}.mp4"
    srt_files = list(cache_dir.glob("*.srt"))
    if not mp4_path.exists() or not srt_files:
        raise RuntimeError(
            "Cached video or SRT missing. Re-digest the video first."
        )
    srt_path = srt_files[0]
    workdir = Path(tempfile.mkdtemp(prefix="v2d_slides_"))
    scene_dir = workdir / "scene"
    interval_dir = workdir / "interval"
    try:
        duration = get_video_duration(mp4_path)
        scene_frames, interval_frames = extract_scene_and_interval_frames(
            mp4_path, scene_dir, interval_dir,
            scene_threshold=0.2, interval=20.0, duration=duration,
        )
        frames = merge_frames(scene_frames, interval_frames)
        frames = dedupe_frames(frames, 4)
        deck_frames = global_phash_cluster(frames)
        settings = load_settings()
        if (settings.get("slide_classification", True)
                and len(deck_frames) > _GRID_CELLS):
            try:
                backend = select_backend(for_vision=True)
                if getattr(backend, "vision_supported", False):
                    deck_frames = classify_slides_via_grids(
                        deck_frames, backend=backend,
                        model=settings.get("slide_classifier_model")
                              or "claude-haiku-4-5-20251001",
                        workdir=workdir,
                        log_video_id=video_id,
                    )
            except Exception:
                pass
        segments = parse_srt(srt_path)
        slides_data = assign_transcript_to_frames(deck_frames, segments, duration)
        digest_md_path = digests_dir / video_id / "digest.md"
        title = video_id
        if digest_md_path.exists():
            for line in digest_md_path.read_text().splitlines():
                if line.startswith("# "):
                    title = line[2:].strip()
                    break
        out_path = digests_dir / video_id / "slides.pptx"
        tmp_out = digests_dir / video_id / "slides.pptx.tmp"
        build_deck(slides_data, tmp_out, title)
        tmp_out.replace(out_path)
        return out_path
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def build_audio_for_artifact(
    video_id: str,
    kind: str,
    *,
    digests_dir: Optional[Path] = None,
) -> Path:
    """Render one artifact's markdown to MP3 using the configured TTS
    provider (macOS `say` or ElevenLabs). kind ∈ {"digest","panel","takeaway"}.
    Returns the written .mp3 path; raises on failure."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    src_name = AUDIO_SOURCE_BY_KIND.get(kind)
    if src_name is None:
        raise RuntimeError(f"unknown audio kind: {kind!r}")
    md_path = digests_dir / video_id / src_name
    if not md_path.exists():
        raise RuntimeError(
            f"{src_name} not found — generate it first before requesting audio."
        )
    mp3_path = digests_dir / video_id / f"{kind}.mp3"
    s = load_settings()
    generate_audio_from_markdown(
        md_path, mp3_path,
        provider=s.get("tts_provider"),
        voice=s.get("tts_voice") or None,
        rate=s.get("tts_rate") or None,
        elevenlabs_voice_id=s.get("elevenlabs_voice_id") or None,
        elevenlabs_model=s.get("elevenlabs_model") or None,
    )
    return mp3_path


# ---- Agent-facing ingestion + generation + subscriptions ---------------

def digest_video(
    url: str,
    *,
    blocking: bool = False,
    source: str = "oneoff",
    digests_dir: Optional[Path] = None,
) -> dict:
    """Kick off the full digest pipeline for a YouTube URL.

    blocking=False (default): spawns `yt2md <url>` as a detached child
        and returns {"video_id", "job_id" (pid), "status": "started",
        "log_path"} immediately. Mirrors the /one-off route's pattern
        so it survives a parent crash.

    blocking=True: runs the pipeline synchronously, returns
        {"video_id", "status": "done", "digest": <full digest JSON>}
        on success — or raises RuntimeError on a non-zero exit.

    Skip-if-exists: if the digest is already in the library, returns
        {"video_id", "status": "exists", "digest": <full digest JSON>}
        without re-running the pipeline.
    """
    import time as _t
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    video_id = extract_video_id(url)
    if not video_id:
        raise ValueError(f"Couldn't extract a YouTube video ID from: {url!r}")

    existing = digests_dir / video_id / "digest.md"
    if existing.exists():
        return {
            "video_id": video_id, "status": "exists",
            "digest": load_digest_json(video_id, digests_dir=digests_dir),
        }

    # Budget gate: refuse to start a new digest once month-to-date workspace
    # spend crosses the block threshold. Returns a status dict (rather than
    # raising) so programmatic callers — CLI, MCP tool, web — handle it
    # uniformly. None == allowed.
    _blocked = check_budget(action=f"digest {video_id}")
    if _blocked:
        return {"video_id": video_id, "status": "blocked", "reason": _blocked}

    digest_path = digests_dir / video_id / "digest.md"
    digest_path.parent.mkdir(parents=True, exist_ok=True)
    log_path = get_data_dir() / "logs" / "oneoff.log"
    log_path.parent.mkdir(parents=True, exist_ok=True)

    yt2md_path = shutil.which("yt2md")
    if not yt2md_path:
        raise RuntimeError("yt2md not on PATH — install with `uv tool install`.")

    if blocking:
        with open(log_path, "a") as log_fd:
            log_fd.write(
                f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} (blocking) "
                f"starting {video_id} ({url}) =====\n"
            )
            result = subprocess.run(
                [yt2md_path, url, "-o", str(digest_path),
                 "--source", source],
                cwd=digest_path.parent,
                stdout=log_fd,
                stderr=subprocess.STDOUT,
                env={**os.environ, **_settings_to_env(load_settings())},
            )
        if result.returncode != 0:
            raise RuntimeError(
                f"yt2md exited {result.returncode}; see {log_path}"
            )
        return {
            "video_id": video_id, "status": "done",
            "digest": load_digest_json(video_id, digests_dir=digests_dir),
        }

    log_fd = open(log_path, "a")
    log_fd.write(
        f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} starting {video_id} "
        f"({url}) =====\n"
    )
    log_fd.flush()
    try:
        proc = subprocess.Popen(
            [yt2md_path, url, "-o", str(digest_path),
             "--source", source],
            cwd=digest_path.parent,
            stdout=log_fd,
            stderr=subprocess.STDOUT,
            env={**os.environ, **_settings_to_env(load_settings())},
            **_DETACH_KWARGS,
        )
    finally:
        log_fd.close()
    return {
        "video_id": video_id, "job_id": proc.pid, "status": "started",
        "log_path": str(log_path),
    }


def _start_generation(video_id: str, kind: str, build_fn) -> dict:
    """Shared shape for generate_{panel,takeaway,slides}. Uses
    start_local_job for in-process daemon-thread execution so a poll
    against local_job_status(f"{id}:{kind}") sees the running phase."""
    started = start_local_job(f"{video_id}:{kind}", build_fn, video_id)
    return {
        "video_id": video_id, "kind": kind,
        "status": "started" if started else "already_running",
        "job_key": f"{video_id}:{kind}",
    }


def generate_panel(
    video_id: str, *, blocking: bool = False,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Generate the panel discussion for an already-digested video. If
    blocking=True, runs synchronously and returns the panel JSON on
    completion; otherwise dispatches to a background thread and returns
    a job key the caller can poll via local_job_status()."""
    if blocking:
        build_panel_for_video(video_id, digests_dir=digests_dir)
        return {
            "video_id": video_id, "kind": "panel", "status": "done",
            "panel": load_panel_json(video_id, digests_dir=digests_dir),
        }
    return _start_generation(
        video_id, "panel",
        lambda vid: build_panel_for_video(vid, digests_dir=digests_dir),
    )


def generate_takeaway(
    video_id: str, *, blocking: bool = False,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Generate the takeaway for an already-digested video. See
    generate_panel for blocking semantics."""
    if blocking:
        build_takeaway_for_video(video_id, digests_dir=digests_dir)
        return {
            "video_id": video_id, "kind": "takeaway", "status": "done",
            "takeaway": load_takeaway_json(video_id, digests_dir=digests_dir),
        }
    return _start_generation(
        video_id, "takeaway",
        lambda vid: build_takeaway_for_video(vid, digests_dir=digests_dir),
    )


def generate_slides(
    video_id: str, *, blocking: bool = False,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Generate the slides deck for an already-digested video. See
    generate_panel for blocking semantics. Returns the .pptx path
    rather than parsed content (no JSON view for binary decks)."""
    if blocking:
        out = build_slides_for_video(video_id, digests_dir=digests_dir)
        return {
            "video_id": video_id, "kind": "slides", "status": "done",
            "slides_path": str(out),
        }
    return _start_generation(
        video_id, "slides",
        lambda vid: build_slides_for_video(vid, digests_dir=digests_dir),
    )


def generate_audio(
    video_id: str,
    kind: str,
    *,
    blocking: bool = False,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Render a digest/panel/takeaway markdown file to MP3 using the
    configured TTS provider. kind ∈ {"digest","panel","takeaway"}.

    blocking=False (default): dispatches to a background thread; poll
        via local_job_status(f"{id}:audio_{kind}").
    blocking=True: runs synchronously, returns the mp3 path.

    Skip-if-exists: if the mp3 already exists, returns immediately."""
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if kind not in AUDIO_SOURCE_BY_KIND:
        raise ValueError(
            f"kind must be one of {list(AUDIO_SOURCE_BY_KIND)!r}; got {kind!r}"
        )
    mp3_path = digests_dir / video_id / f"{kind}.mp3"
    if mp3_path.exists():
        return {
            "video_id": video_id, "kind": kind, "status": "exists",
            "mp3_path": str(mp3_path),
        }
    if blocking:
        out = build_audio_for_artifact(video_id, kind, digests_dir=digests_dir)
        return {
            "video_id": video_id, "kind": kind, "status": "done",
            "mp3_path": str(out),
        }
    return _start_generation(
        video_id, f"audio_{kind}",
        lambda vid: build_audio_for_artifact(vid, kind, digests_dir=digests_dir),
    )


def mark_digest_read(digest_id: str) -> dict:
    """Mark a digest as read. Idempotent — re-marking just updates the
    opened_at timestamp."""
    _mark_digest_read(digest_id)
    return {"digest_id": digest_id, "status": "read"}


def mark_digest_unread(digest_id: str) -> dict:
    """Mark a digest as unread (removes the read record). Idempotent."""
    with _library_connect() as conn:
        conn.execute("DELETE FROM digest_reads WHERE digest_id = ?", (digest_id,))
    return {"digest_id": digest_id, "status": "unread"}


# ---- Topic taxonomy + curation ----------------------------------------

def list_topics(
    *, min_digests: int = 1, limit: int = 50, source: Optional[str] = None,
) -> list:
    """Return the topic taxonomy: [{topic, n_digests, last_seen,
    sources: {llm, user}}] sorted by n_digests desc.

    Use for "what have I been reading about?" and as input to
    `list_digests(topic=...)`. `source` filters to 'llm' or 'user' tags
    only; None returns the union."""
    src_clause = ""
    params: list = [min_digests]
    if source:
        src_clause = " WHERE source = ?"
        params = [source, min_digests]
    with _library_connect() as conn:
        rows = conn.execute(
            "SELECT topic, "
            "       COUNT(DISTINCT digest_id) AS n, "
            "       MAX(added_at) AS last_seen, "
            "       GROUP_CONCAT(DISTINCT source) AS sources "
            "FROM digest_topics" + src_clause + " "
            "GROUP BY topic HAVING n >= ? "
            "ORDER BY n DESC, topic ASC LIMIT ?",
            params + [limit],
        ).fetchall()
    out = []
    for tag, n, last_seen, sources in rows:
        srcs = set((sources or "").split(","))
        out.append({
            "topic": tag, "n_digests": n,
            "last_seen": last_seen,
            "sources": {"llm": "llm" in srcs, "user": "user" in srcs},
        })
    return out


def tag_digest(digest_id: str, tags: list) -> dict:
    """Add user-curated tags to a digest. Normalizes + dedupes input
    against existing user tags. Idempotent (re-adding is a no-op).

    Tags are kept separate from LLM-assigned tags so the LLM can be
    wrong without overwriting your intent — list_digests filters union
    both sets, but `topics_split` in the response shows which is which."""
    cleaned: list = []
    for t in tags:
        norm = _validate_tag(t)
        if norm and norm not in cleaned:
            cleaned.append(norm)
    if not cleaned:
        raise ValueError(f"No valid tags in {tags!r}")
    _write_topics_to_db(digest_id, cleaned, source="user")
    # Mirror to metadata.json user_tags (union with existing, deduped).
    digests_dir = get_data_dir() / "digests"
    meta_path = digests_dir / digest_id / "metadata.json"
    meta = {}
    if meta_path.exists():
        try:
            meta = _json.loads(meta_path.read_text())
        except Exception:
            meta = {}
    existing = list(meta.get("user_tags") or [])
    for t in cleaned:
        if t not in existing:
            existing.append(t)
    meta["user_tags"] = existing
    _atomic_write_json(meta_path, meta)
    return {"digest_id": digest_id, "user_tags": existing}


def untag_digest(digest_id: str, tags: list) -> dict:
    """Remove user tags from a digest. Idempotent (removing a missing
    tag is fine). LLM-assigned tags are never touched."""
    norm = {_validate_tag(t) for t in tags}
    norm.discard(None)
    if not norm:
        return {"digest_id": digest_id, "user_tags": []}
    with _library_connect() as conn:
        for t in norm:
            conn.execute(
                "DELETE FROM digest_topics WHERE digest_id = ? "
                "AND topic = ? AND source = 'user'",
                (digest_id, t),
            )
    digests_dir = get_data_dir() / "digests"
    meta_path = digests_dir / digest_id / "metadata.json"
    meta = {}
    if meta_path.exists():
        try:
            meta = _json.loads(meta_path.read_text())
        except Exception:
            meta = {}
    meta["user_tags"] = [t for t in (meta.get("user_tags") or []) if t not in norm]
    _atomic_write_json(meta_path, meta)
    return {"digest_id": digest_id, "user_tags": meta["user_tags"]}


def _set_curation_flag(digest_id: str, field: str, value: bool) -> dict:
    """Shared shape for save/unsave/dismiss/undismiss. Updates both the
    DB row and metadata.json."""
    import time as _t
    assert field in ("user_saved", "user_dismissed")
    with _library_connect() as conn:
        conn.execute(
            "INSERT INTO digest_meta (digest_id, added_at) "
            "VALUES (?, ?) ON CONFLICT(digest_id) DO NOTHING",
            (digest_id, int(_t.time())),
        )
        conn.execute(
            f"UPDATE digest_meta SET {field} = ? WHERE digest_id = ?",
            (1 if value else 0, digest_id),
        )
    digests_dir = get_data_dir() / "digests"
    meta_path = digests_dir / digest_id / "metadata.json"
    meta = {}
    if meta_path.exists():
        try:
            meta = _json.loads(meta_path.read_text())
        except Exception:
            meta = {}
    meta[field] = value
    _atomic_write_json(meta_path, meta)
    return {"digest_id": digest_id, field: value}


def save_digest(digest_id: str) -> dict:
    """Mark a digest as 'saved' (worth keeping / returning to)."""
    return _set_curation_flag(digest_id, "user_saved", True)


def unsave_digest(digest_id: str) -> dict:
    """Remove the 'saved' mark."""
    return _set_curation_flag(digest_id, "user_saved", False)


def dismiss_digest(digest_id: str) -> dict:
    """Mark a digest as 'dismissed' (not worth surfacing in briefings).
    list_digests(dismissed=False) will hide it; list_digests(dismissed=True)
    can recall the dismissed set."""
    return _set_curation_flag(digest_id, "user_dismissed", True)


def undismiss_digest(digest_id: str) -> dict:
    """Remove the 'dismissed' mark."""
    return _set_curation_flag(digest_id, "user_dismissed", False)


def retag_digest(digest_id: str) -> dict:
    """Re-run the LLM tagging step against the current taxonomy. Useful
    after the taxonomy has grown (so a digest tagged early can pick up
    tags introduced later) or after editing the prompt."""
    return tag_digest_via_llm(digest_id)


def retrofit_topics(
    *,
    since: Optional[str] = None,
    limit: Optional[int] = None,
    dry_run: bool = False,
    digests_dir: Optional[Path] = None,
) -> dict:
    """Tag any digest that doesn't already have LLM topics. Resumable —
    re-running picks up where it left off.

    since   : ISO date "YYYY-MM-DD"; skip digests older than this.
    limit   : max digests to tag in this run (None = no cap).
    dry_run : print what would be tagged without calling the LLM.
    """
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if not digests_dir.exists():
        return {"tagged": 0, "skipped": 0, "errors": []}

    candidates: list = []
    for d in sorted(digests_dir.iterdir(),
                    key=lambda p: p.stat().st_mtime, reverse=True):
        if not d.is_dir() or not (d / "digest.md").exists():
            continue
        if _read_topics_for_digest(d.name)["llm"]:
            continue  # already tagged
        if since:
            meta = {}
            if (d / "metadata.json").exists():
                try:
                    meta = _json.loads((d / "metadata.json").read_text())
                except Exception:
                    pass
            published = _upload_date_iso(meta.get("upload_date")) or \
                _iso_from_mtime(d.stat().st_mtime)
            if published < since:
                continue
        candidates.append(d.name)
        if limit and len(candidates) >= limit:
            break

    if dry_run:
        return {"would_tag": candidates, "count": len(candidates)}

    tagged = 0
    errors = []
    for vid in candidates:
        try:
            tag_digest_via_llm(vid)
            tagged += 1
            print(f"  ✓ tagged {vid}")
        except Exception as e:
            errors.append({"digest_id": vid, "error": f"{type(e).__name__}: {e}"})
            print(f"  ✗ {vid}: {e}")
    return {"tagged": tagged, "skipped": 0, "errors": errors,
            "candidates": len(candidates)}


def list_subscriptions() -> list:
    """Return the watched channels as structured entries (vs the raw
    URL strings in channels.txt). One per subscribed channel."""
    return [
        {"url": ch, "added_via": "channels.txt"}
        for ch in read_channels()
    ]


def add_subscription(channel_url: str) -> dict:
    """Subscribe to a YouTube channel. Idempotent — adding an
    already-subscribed channel returns status='exists' rather than
    erroring. Handles bare @handles and partial URLs the same way
    the CLI does."""
    url = normalize_channel_url(channel_url)
    if not is_url(url):
        raise ValueError(f"Not a URL: {channel_url!r}")
    channels = read_channels()
    if url in channels:
        return {"url": url, "status": "exists"}
    channels.append(url)
    write_channels(channels)
    return {"url": url, "status": "added"}


def remove_subscription(channel_url: str) -> dict:
    """Unsubscribe from a YouTube channel. Returns status='missing'
    if the channel wasn't subscribed (idempotent removal)."""
    url = normalize_channel_url(channel_url)
    channels = read_channels()
    if url not in channels:
        return {"url": url, "status": "missing"}
    write_channels([c for c in channels if c != url])
    return {"url": url, "status": "removed"}


# ---- serve subcommand (local reader UI) ----

# Shared between digest + panel + takeaway viewers. The button has
# data-copy-target pointing at a hidden <textarea> elsewhere on the page; on
# click we read its raw value and copy to the clipboard. Falls back to a
# select+execCommand path for browsers / contexts (eg http on a non-localhost)
# where the modern clipboard API is unavailable.
#
# Optional attribute data-then-open="<url>" navigates to the URL in a new tab
# after a successful copy — used by the "Continue in chat" handoff to land
# the user in claude.ai with the context already in their clipboard.
_COPY_BUTTON_JS = """
<script>
(function () {
  function flash(btn, text) {
    const orig = btn.getAttribute('data-orig-text') || btn.textContent;
    btn.setAttribute('data-orig-text', orig);
    btn.textContent = text;
    setTimeout(() => { btn.textContent = orig; }, 1500);
  }
  function maybeOpen(btn) {
    const url = btn.getAttribute('data-then-open');
    if (!url) return;
    // Brief delay so the user sees the "Copied!" feedback before the new tab.
    setTimeout(() => window.open(url, '_blank', 'noopener'), 600);
  }
  document.addEventListener('click', async (ev) => {
    const btn = ev.target.closest('[data-copy-target]');
    if (!btn) return;
    const id = btn.getAttribute('data-copy-target');
    const src = document.getElementById(id);
    if (!src) return;
    const text = src.value;
    try {
      if (navigator.clipboard && navigator.clipboard.writeText) {
        await navigator.clipboard.writeText(text);
        flash(btn, 'Copied!');
        maybeOpen(btn);
        return;
      }
    } catch (e) { /* fall through to legacy path */ }
    try {
      src.removeAttribute('hidden');
      src.style.position = 'absolute';
      src.style.left = '-9999px';
      src.select();
      document.execCommand('copy');
      src.setAttribute('hidden', '');
      flash(btn, 'Copied!');
      maybeOpen(btn);
    } catch (e) {
      flash(btn, 'Copy failed');
    }
  });
})();
</script>
"""


# ---- shared async job tracker for in-process work surfaced via the web UI ----
#
# In-process daemon threads keyed by "<video_id>:<kind>" — currently used for
# slides generation; designed to also serve panel + takeaway retrofits later.
# Threads die with the Flask process; the work is idempotent so a lost job
# just means the user clicks Generate again. State lives in a module dict
# (single-process Flask, single-writer assumption holds).

_local_jobs: dict = {}


def start_local_job(key: str, fn, *args, **kwargs) -> bool:
    """Start a daemon-thread job under `key` if one isn't already running.
    Returns True if a new job was started, False if one was already in flight.
    Idempotent: a second click while running is a no-op (not an error).
    """
    import threading
    import time as _t

    existing = _local_jobs.get(key)
    if existing and existing["thread"].is_alive():
        return False

    job = {
        "started": _t.time(),
        "kind": key.split(":", 1)[-1],
        "error": None,
        "thread": None,
    }

    def _wrapper():
        try:
            fn(*args, **kwargs)
        except Exception as e:
            job["error"] = f"{type(e).__name__}: {e}"

    t = threading.Thread(target=_wrapper, daemon=True)
    job["thread"] = t
    _local_jobs[key] = job
    t.start()
    return True


def local_job_status(key: str) -> dict:
    """Snapshot of a job's state. Includes a simple 'phase' string the UI can
    render directly: 'idle' (no job ever ran or completed long ago) /
    'running' / 'done' / 'error'.
    """
    import time as _t
    job = _local_jobs.get(key)
    if not job:
        return {"phase": "idle"}
    running = job["thread"].is_alive()
    elapsed = int(_t.time() - job["started"])
    if running:
        return {"phase": "running", "elapsed": elapsed,
                "started": job["started"]}
    if job.get("error"):
        return {"phase": "error", "elapsed": elapsed,
                "error": job["error"]}
    return {"phase": "done", "elapsed": elapsed}


# Job-status polling JS shared by any toolbar element with
# data-poll-url=<url> — polls every 2s per element, updates a `.elapsed`
# child with "{n}s", and reloads the page when the artifact appears (on
# success) or surfaces the error inline (on failure). Handles multiple
# concurrent jobs (e.g. panel + slides running for the same digest) by
# giving each [data-poll-url] element its own poll loop.
_JOB_POLL_JS = """
<script>
(function () {
  const els = document.querySelectorAll('[data-poll-url]');
  if (els.length === 0) return;
  els.forEach((el) => {
    const url = el.getAttribute('data-poll-url');
    const elapsedSpan = el.querySelector('.elapsed');
    let timer = null;
    async function tick() {
      try {
        const res = await fetch(url);
        if (!res.ok) return;
        const s = await res.json();
        if (s.phase === 'running') {
          if (elapsedSpan) elapsedSpan.textContent = s.elapsed + 's';
          return;
        }
        if (timer) { clearInterval(timer); timer = null; }
        if (s.phase === 'done' && s.artifact_exists) {
          // Reload so the toolbar swaps to "View X" / "Download X".
          window.location.reload();
          return;
        }
        if (s.phase === 'error') {
          el.innerHTML = '<span style="color: #c00;">Generation failed: ' +
            (s.error || '(unknown error)') + '</span>';
          return;
        }
        // Done but artifact missing — likely a write race. Reload anyway.
        window.location.reload();
      } catch (e) { /* network blip — try next tick */ }
    }
    tick();
    timer = setInterval(tick, 2000);
  });
})();
</script>
"""


def build_chat_handoff_prompt(video_id: str, digests_dir: Path) -> str:
    """Assemble a prompt the user can paste into claude.ai (or any chat) to
    continue thinking about a video. Includes whichever artifacts exist on
    disk — digest, panel, takeaway — so the chat has full context without
    a manual paste of each.
    """
    digest_md = digests_dir / video_id / "digest.md"
    panel_md = digests_dir / video_id / "panel.md"
    takeaway_md = digests_dir / video_id / "takeaway.md"

    have = []
    if digest_md.exists():
        have.append("digest")
    if panel_md.exists():
        have.append("panel-of-experts critique")
    if takeaway_md.exists():
        have.append("bottom-line takeaway")
    if not have:
        artifacts_phrase = "summary"
    elif len(have) == 1:
        artifacts_phrase = have[0]
    elif len(have) == 2:
        artifacts_phrase = f"{have[0]} and {have[1]}"
    else:
        artifacts_phrase = ", ".join(have[:-1]) + ", and " + have[-1]

    parts: List[str] = [
        f"I just read a distilled summary of a YouTube video. Below is the "
        f"{artifacts_phrase}. I have a follow-up question.",
        "",
        f"Source: https://www.youtube.com/watch?v={video_id}",
        "",
    ]
    if digest_md.exists():
        parts += ["# Digest", "", digest_md.read_text().rstrip(), ""]
    if panel_md.exists():
        parts += ["# Panel discussion", "", panel_md.read_text().rstrip(), ""]
    if takeaway_md.exists():
        parts += ["# Takeaway", "", takeaway_md.read_text().rstrip(), ""]
    parts += ["---", "", "My question: "]
    return "\n".join(parts)


SERVE_PAGE_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{{ title }} — yt2md</title>
{% if base_href %}<base href="{{ base_href }}">{% endif %}
<script>
(function() {
  const stored = localStorage.getItem('yt2md-theme');
  if (stored && stored !== 'auto') document.documentElement.setAttribute('data-theme', stored);
})();
</script>
<style>
:root {
  --bg: #fafaf7;
  --fg: #1a1a1a;
  --muted: #6b6b6b;
  --accent: #b65a2c;
  --unread: #2563eb;
  --border: #e5e3dc;
  --sidebar-bg: #f0eee6;
  --code-bg: #ececea;
}
@media (prefers-color-scheme: dark) {
  :root:not([data-theme="light"]) {
    --bg: #1a1a1a;
    --fg: #e8e8e8;
    --muted: #999;
    --accent: #d97a4d;
    --unread: #60a5fa;
    --border: #2e2e2e;
    --sidebar-bg: #141414;
    --code-bg: #232323;
  }
}
:root[data-theme="dark"] {
  --bg: #1a1a1a;
  --fg: #e8e8e8;
  --muted: #999;
  --accent: #d97a4d;
  --unread: #60a5fa;
  --border: #2e2e2e;
  --sidebar-bg: #141414;
  --code-bg: #232323;
}
* { box-sizing: border-box; }
html, body { margin: 0; padding: 0; height: 100%; }
body {
  display: flex;
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Helvetica Neue", sans-serif;
  background: var(--bg);
  color: var(--fg);
  line-height: 1.6;
}
aside {
  width: 320px;
  flex-shrink: 0;
  height: 100vh;
  /* Flex column so the Manage nav can stay pinned at the bottom while the
     Digests list grows + scrolls within its own pane. Without this, a
     long library pushes Manage below the viewport fold. */
  display: flex;
  flex-direction: column;
  padding: 20px 16px;
  background: var(--sidebar-bg);
  border-right: 1px solid var(--border);
}
aside .sidebar-header { flex: 0 0 auto; }
aside nav[aria-label="Per-video digests"] {
  flex: 1 1 auto;
  /* min-height: 0 lets a flex child shrink below its intrinsic content
     size — required for overflow-y: auto to actually scroll inside the
     flex column instead of forcing the whole sidebar taller. */
  min-height: 0;
  overflow-y: auto;
}
aside nav[aria-label="Manage"] {
  flex: 0 0 auto;
  border-top: 1px solid var(--border);
  margin-top: 12px;
  padding-top: 12px;
}
aside h1 { margin: 0 0 16px; font-size: 18px; }
aside h1 a { color: var(--fg); text-decoration: none; }
aside h2 {
  font-size: 11px;
  text-transform: uppercase;
  letter-spacing: 0.08em;
  color: var(--muted);
  margin: 20px 0 8px;
  font-weight: 600;
}
aside ul { list-style: none; padding: 0; margin: 0; }
aside li { margin: 0 0 3px; }
aside li a {
  display: flex;
  gap: 8px;
  align-items: flex-start;
  padding: 6px 8px;
  border-radius: 4px;
  color: var(--fg);
  text-decoration: none;
  font-size: 13px;
  line-height: 1.35;
}
/* Sidebar item layout: video thumbnail + textual body (title + channel
   avatar/name). Thumbnails are 16:9 mini-posters at 56px wide; the
   channel avatar is a small circle. When a digest has no thumbnails
   (legacy or fetch failure) the items lay out as text only. */
aside li .digest-thumb {
  flex: 0 0 56px;
  width: 56px;
  height: 32px;
  border-radius: 3px;
  background: var(--border);
  background-size: cover;
  background-position: center;
}
aside li .digest-body {
  flex: 1 1 auto;
  min-width: 0;
  display: -webkit-box;
  -webkit-line-clamp: 3;
  -webkit-box-orient: vertical;
  overflow: hidden;
}
aside li .digest-channel {
  display: flex; align-items: center; gap: 4px;
  margin-top: 3px;
  font-size: 11px;
  color: var(--muted);
  font-weight: normal;
}
aside li .channel-avatar {
  flex: 0 0 14px;
  width: 14px;
  height: 14px;
  border-radius: 50%;
  background: var(--border);
  background-size: cover;
  background-position: center;
}
aside li.active .digest-channel { color: rgba(255,255,255,0.85); }
aside li a:hover { background: rgba(0,0,0,0.05); }
@media (prefers-color-scheme: dark) {
  aside li a:hover { background: rgba(255,255,255,0.05); }
}
aside li.active a { background: var(--accent); color: white; }
aside li.unread a { font-weight: 600; }
aside .empty { color: var(--muted); font-size: 13px; padding: 6px 8px; }
aside .unread-count {
  display: inline-block; padding: 2px 7px; background: var(--unread); color: white;
  border-radius: 10px; font-size: 11px; font-weight: 600;
  text-transform: none; letter-spacing: 0; margin-left: 4px;
}
/* Banner shown by the sidebar poll JS when the library has changed
   since the page was rendered (new digest landed via one-off or
   scheduled poll, or a digest was deleted). Click reloads. */
#new-digest-banner {
  display: block;
  padding: 8px 12px;
  margin: 12px 0;
  background: var(--accent);
  color: white;
  border-radius: 4px;
  font-size: 13px;
  font-weight: 600;
  text-decoration: none;
  text-align: center;
  flex: 0 0 auto;
}
#new-digest-banner:hover { opacity: 0.9; }
aside .unread-dot {
  display: inline-block; width: 6px; height: 6px; border-radius: 50%;
  background: var(--unread); margin-right: 6px; vertical-align: middle;
}
aside .meta-card.unread { border-color: var(--unread); }
aside .meta-card.unread .week { font-weight: 700; }
/* When an item is the currently-viewed one, suppress the unread signals
   to avoid two competing color cues. */
aside li.active .unread-dot,
aside .meta-card.active .unread-dot { display: none; }
aside .meta-card.active.unread { border-color: var(--accent); }
aside .meta-card {
  display: block; padding: 10px 12px; border-radius: 4px;
  border: 1px solid var(--border); margin-bottom: 6px;
  text-decoration: none; color: var(--fg);
}
aside .meta-card:hover { border-color: var(--accent); }
aside .meta-card.active { background: var(--accent); color: white; border-color: var(--accent); }
aside .meta-card .week { font-weight: 600; font-size: 13px; line-height: 1.2; }
aside .meta-card .count {
  color: var(--muted); font-size: 11px; margin-top: 2px;
  text-transform: uppercase; letter-spacing: 0.04em;
}
aside .meta-card.active .count { color: white; opacity: 0.85; }
main {
  flex: 1;
  height: 100vh;
  overflow-y: auto;
  padding: 32px 48px 80px;
}
main .reader {
  max-width: 720px;
  margin: 0 auto;
}
main h1 { font-size: 28px; line-height: 1.25; margin-top: 0; }
main h2 { font-size: 22px; line-height: 1.3; margin-top: 32px; border-bottom: 1px solid var(--border); padding-bottom: 4px; }
main h3 { font-size: 17px; }
main img { max-width: 100%; height: auto; border: 1px solid var(--border); border-radius: 4px; }
main a { color: var(--accent); }
main code { background: var(--code-bg); padding: 2px 5px; border-radius: 3px; font-size: 0.9em; }
main pre { background: var(--code-bg); padding: 12px; border-radius: 4px; overflow-x: auto; }
main pre code { background: none; padding: 0; }
main blockquote { border-left: 3px solid var(--border); padding-left: 16px; color: var(--muted); margin-left: 0; }
main ul, main ol { padding-left: 24px; }
main hr { border: none; border-top: 1px solid var(--border); margin: 24px 0; }
main sub { color: var(--muted); font-size: 0.85em; }
.empty-state { color: var(--muted); margin-top: 80px; text-align: center; }
.meta-info { color: var(--muted); font-size: 13px; margin-top: -12px; margin-bottom: 24px; }
.featured-eyebrow {
  color: var(--muted); font-size: 12px; text-transform: uppercase;
  letter-spacing: 0.06em; margin-bottom: 24px;
}
.featured-eyebrow a { color: var(--accent); text-decoration: none; font-weight: 600; }
.cta {
  display: inline-block; background: var(--accent); color: white;
  padding: 12px 20px; border-radius: 4px; text-decoration: none;
  font-weight: 500; margin-top: 12px;
}
.cta:hover { opacity: 0.9; }
.recent-list { list-style: none; padding: 0; margin: 12px 0; }
.recent-list li { padding: 6px 0; border-bottom: 1px solid var(--border); }
.recent-list a { color: var(--fg); text-decoration: none; }
.recent-list a:hover { color: var(--accent); }
.add-form { display: flex; gap: 8px; margin: 24px 0; }
.add-form input[type="text"] {
  flex: 1; padding: 10px 12px; font-size: 14px;
  border: 1px solid var(--border); border-radius: 4px;
  background: var(--bg); color: var(--fg);
}
.add-form button, .channel-list button {
  padding: 10px 16px; font-size: 14px;
  border: 1px solid var(--accent); border-radius: 4px;
  background: var(--accent); color: white; cursor: pointer;
}
.channel-list { list-style: none; padding: 0; }
.channel-list li {
  display: flex; align-items: center; justify-content: space-between;
  padding: 10px 12px; border: 1px solid var(--border); border-radius: 4px;
  margin-bottom: 8px; gap: 12px;
}
.channel-list .url { flex: 1; word-break: break-all; font-size: 14px; }
.channel-list button {
  background: transparent; color: var(--muted);
  border-color: var(--border); padding: 6px 12px; font-size: 12px;
}
.channel-list button:hover { color: var(--accent); border-color: var(--accent); }
.flash { padding: 10px 14px; border-radius: 4px; margin: 16px 0;
  background: var(--code-bg); border-left: 3px solid var(--accent); }
.next-step {
  padding: 14px 18px; border-radius: 6px; margin: 16px 0 24px;
  background: var(--sidebar-bg); border: 1px solid var(--border);
  border-left: 3px solid var(--accent); font-size: 14px; line-height: 1.5;
}
.next-step strong { color: var(--accent); }
.next-step a { color: var(--accent); }
.schedule-form {
  background: var(--sidebar-bg); border: 1px solid var(--border);
  border-radius: 6px; padding: 16px 20px; margin: 16px 0;
}
.schedule-fields {
  display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
  gap: 12px 20px; margin-bottom: 16px;
}
.schedule-fields label {
  display: flex; flex-direction: column; gap: 4px;
  font-size: 12px; color: var(--muted); text-transform: uppercase;
  letter-spacing: 0.04em;
}
.schedule-fields input, .schedule-fields select {
  padding: 8px 10px; font-size: 14px; border: 1px solid var(--border);
  border-radius: 4px; background: var(--bg); color: var(--fg);
  font-family: inherit; text-transform: none; letter-spacing: normal;
}
.schedule-fields .suffix {
  position: absolute; right: 10px; top: 50%; transform: translateY(-50%);
  color: var(--muted); font-size: 12px; pointer-events: none;
}
.status-table { width: 100%; border-collapse: collapse; font-size: 13px;
  margin: 12px 0 16px; }
.status-table td { padding: 6px 12px; border-bottom: 1px solid var(--border); }
.status-table td:first-child { color: var(--muted); width: 140px; }
.status-table td:last-child { font-family: ui-monospace, "SF Mono", Menlo, monospace; }
.activity-table { width: 100%; border-collapse: collapse; font-size: 13px; margin-top: 12px; }
.activity-table th, .activity-table td {
  text-align: left; padding: 10px 8px; border-bottom: 1px solid var(--border);
  vertical-align: top;
}
.activity-table th { color: var(--muted); font-weight: 600; font-size: 11px;
  text-transform: uppercase; letter-spacing: 0.04em; }
.activity-table td a { color: var(--accent); text-decoration: none; }
.activity-table td a:hover { text-decoration: underline; }
.activity-table .ok { color: #4a9f56; font-weight: 600; }
.activity-table .fail { color: #d04545; font-weight: 600; }
.activity-meta { color: var(--muted); font-size: 12px; margin-top: 3px; }
.activity-error { font-family: ui-monospace, "SF Mono", Menlo, monospace; word-break: break-word; }
.activity-stages, .activity-tokens { font-family: ui-monospace, "SF Mono", Menlo, monospace;
  font-size: 12px; white-space: nowrap; }
.filter-row { display: flex; gap: 8px; margin: 12px 0 4px; flex-wrap: wrap; }
.filter-chip {
  padding: 5px 12px; border: 1px solid var(--border); border-radius: 999px;
  font-size: 12px; color: var(--fg); text-decoration: none; background: var(--bg);
}
.filter-chip:hover { border-color: var(--accent); }
.filter-chip.active { background: var(--accent); color: white; border-color: var(--accent); }
.filter-chip-count { opacity: 0.7; }
.delete-btn {
  padding: 6px 14px; font-size: 13px; cursor: pointer;
  background: transparent; color: #b13030;
  border: 1px solid #b13030; border-radius: 4px;
}
.delete-btn:hover { background: #b13030; color: white; }
.digest-actions { display: flex; gap: 12px; align-items: center; flex-wrap: wrap; }
.digest-toolbar { margin-top: -8px; margin-bottom: 8px; }
.digest-toolbar .delete-btn { margin-left: auto; }  /* push delete to the far right */
.discuss-btn {
  padding: 8px 16px; font-size: 13px; cursor: pointer;
  background: var(--accent); color: white;
  border: 1px solid var(--accent); border-radius: 4px;
  text-decoration: none; display: inline-block;
}
.discuss-btn:hover { opacity: 0.9; }
.discuss-btn-secondary {
  padding: 6px 12px; font-size: 12px; cursor: pointer;
  background: transparent; color: var(--fg);
  border: 1px solid var(--border); border-radius: 4px;
}
.discuss-btn-secondary:hover { border-color: var(--accent); }
.job-block { border: 1px solid var(--border); border-radius: 4px; padding: 16px 20px;
  margin: 16px 0; background: var(--sidebar-bg); }
.job-block h3 { margin: 0 0 8px; font-size: 15px; }
.job-actions { display: flex; gap: 8px; margin: 8px 0 16px; flex-wrap: wrap; }
.job-actions button {
  padding: 8px 14px; font-size: 13px; cursor: pointer;
  border: 1px solid var(--border); border-radius: 4px;
  background: var(--bg); color: var(--fg);
}
.job-actions button.primary { background: var(--accent); color: white; border-color: var(--accent); }
.job-actions button:hover { border-color: var(--accent); }
.job-summary { font-size: 15px; margin: 0 0 12px; }
details summary {
  cursor: pointer; color: var(--muted); font-size: 12px;
  text-transform: uppercase; letter-spacing: 0.04em;
  padding: 4px 0; user-select: none;
}
details summary:hover { color: var(--accent); }
details[open] summary { margin-bottom: 8px; }
.log-block {
  background: var(--code-bg); padding: 12px; border-radius: 4px;
  font-family: ui-monospace, "SF Mono", Menlo, monospace;
  font-size: 12px; line-height: 1.4; max-height: 240px;
  overflow: auto; white-space: pre-wrap; word-break: break-all;
}
.dot { display: inline-block; width: 8px; height: 8px; border-radius: 50%;
  margin-right: 6px; vertical-align: middle; }
.dot-on { background: #4caf50; }
.dot-off { background: #999; }
.dot-warn { background: #d97a4d; }

/* Accessibility: skip to content for keyboard users */
.skip-link {
  position: absolute; left: -1000px; top: 0; padding: 8px 12px;
  background: var(--accent); color: white; text-decoration: none;
  border-radius: 0 0 4px 0; z-index: 100;
}
.skip-link:focus { left: 0; }

/* Form labels (visually hidden, exposed to screen readers) */
.sr-only {
  position: absolute; width: 1px; height: 1px; padding: 0; margin: -1px;
  overflow: hidden; clip: rect(0,0,0,0); white-space: nowrap; border: 0;
}

/* Sidebar header (yt2md title + theme toggle inline) */
.sidebar-header {
  display: flex; justify-content: space-between; align-items: center;
  margin-bottom: 16px;
}
.sidebar-header h1 { margin: 0; }
.theme-toggle {
  background: transparent; border: 1px solid var(--border); border-radius: 4px;
  padding: 4px 8px; font-size: 14px; cursor: pointer; color: var(--fg);
  font-family: inherit; line-height: 1;
}
.theme-toggle:hover { border-color: var(--accent); }

/* Mobile: stack sidebar above main, slimmer padding. The flex-column
   layout still applies — header at top, scrollable digests in the
   middle of the constrained 40vh, manage pinned at the bottom. */
@media (max-width: 720px) {
  body { flex-direction: column; }
  aside { width: 100%; height: auto; max-height: 40vh; border-right: none; border-bottom: 1px solid var(--border); }
  main { height: auto; padding: 24px 20px 60px; }
  main .reader { max-width: 100%; }
}
</style>
</head>
<body data-digests-count="{{ digests_count }}" data-digests-mtime="{{ digests_max_mtime }}">
<a class="skip-link" href="#main-content">Skip to main content</a>
<aside>
  <div class="sidebar-header">
    <h1><a href="/">yt2md</a></h1>
    <button class="theme-toggle" type="button" onclick="cycleTheme()" aria-label="Cycle theme: auto / light / dark">🌓</button>
  </div>

  <a id="new-digest-banner" href="" onclick="window.location.reload(); return false;" style="display:none;">✨ New digest — refresh</a>

  <nav aria-label="Per-video digests">
  <h2>Digests {% if unread_digest_count %}<span class="unread-count">{{ unread_digest_count }} new</span>{% else %}({{ digests|length }}){% endif %}</h2>
  <ul>
    {% for d in digests %}
    <li{% if current == 'digest:' + d.id %} class="active"{% endif %}{% if d.unread %} class="unread"{% endif %}>
      <a href="/digests/{{ d.id }}/">
        {% if d.has_thumbnail %}<span class="digest-thumb" style="background-image: url('/digests/{{ d.id }}/thumbnail.jpg');" aria-hidden="true"></span>{% else %}<span class="digest-thumb" aria-hidden="true"></span>{% endif %}
        <span class="digest-body">{% if d.unread %}<span class="unread-dot" aria-label="unread"></span>{% endif %}{{ d.title }}{% if d.channel_name %}<span class="digest-channel">{% if d.has_channel_thumbnail %}<span class="channel-avatar" style="background-image: url('/channel-thumbnails/{{ d.channel_id }}.jpg');" aria-hidden="true"></span>{% endif %}{{ d.channel_name }}</span>{% endif %}</span>
      </a>
    </li>
    {% else %}
    <li class="empty">none yet</li>
    {% endfor %}
  </ul>
  </nav>

  <nav aria-label="Manage">
  <h2>Manage</h2>
  <ul>
    <li{% if current == 'channels' %} class="active"{% endif %}><a href="/channels">Subscriptions ({{ channel_count }})</a></li>
    <li{% if current == 'one-off' %} class="active"{% endif %}><a href="/one-off">One-off digest</a></li>
    <li{% if current == 'listen' %} class="active"{% endif %}><a href="/listen">📱 Listen on phone</a></li>
    <li{% if current == 'schedule' %} class="active"{% endif %}><a href="/schedule">Schedule</a></li>
    <li{% if current == 'activity' %} class="active"{% endif %}><a href="/activity">Activity</a></li>
    <li{% if current == 'settings' %} class="active"{% endif %}><a href="/settings">Settings</a></li>
  </ul>
  </nav>
</aside>
<main id="main-content" tabindex="-1">
  <div class="reader">
    {{ body|safe }}
  </div>
</main>
<script>
function applyTheme() {
  const stored = localStorage.getItem('yt2md-theme') || 'auto';
  const root = document.documentElement;
  if (stored === 'auto') root.removeAttribute('data-theme');
  else root.setAttribute('data-theme', stored);
  const btn = document.querySelector('.theme-toggle');
  if (btn) {
    const icons = {auto: '🌓', light: '☀️', dark: '🌙'};
    btn.textContent = icons[stored];
    btn.title = 'Theme: ' + stored + ' (click to cycle)';
  }
}
function cycleTheme() {
  const cur = localStorage.getItem('yt2md-theme') || 'auto';
  const next = {auto: 'light', light: 'dark', dark: 'auto'}[cur];
  localStorage.setItem('yt2md-theme', next);
  applyTheme();
}
applyTheme();

// Sidebar freshness check: poll /sidebar-status every 10s to detect
// new digests landing (auto-pipeline from one-off, scheduled poll, etc.).
// When the library has drifted, surface a "✨ New digest — refresh"
// banner above the digest list. Click reloads. We don't auto-reload
// because the user is usually mid-read; agency over interruption.
(function () {
  const banner = document.getElementById('new-digest-banner');
  if (!banner) return;
  const initialCount = parseInt(document.body.dataset.digestsCount || '0', 10);
  const initialMtime = parseFloat(document.body.dataset.digestsMtime || '0');
  async function check() {
    try {
      const res = await fetch('/sidebar-status', { cache: 'no-store' });
      if (!res.ok) return;
      const s = await res.json();
      const newCount = s.digests_count - initialCount;
      const drifted = (s.max_mtime > initialMtime + 0.5) || (s.digests_count !== initialCount);
      if (drifted) {
        banner.textContent = newCount > 0
          ? `✨ ${newCount} new digest${newCount > 1 ? 's' : ''} — refresh`
          : '✨ Library updated — refresh';
        banner.style.display = '';
      }
    } catch (e) { /* network blip — next tick */ }
  }
  setInterval(check, 10000);
})();
</script>
</body>
</html>
"""


def _list_digests(digests_dir: Path) -> List[dict]:
    if not digests_dir.exists():
        return []
    results = []
    for d in sorted(digests_dir.iterdir(), key=lambda p: p.stat().st_mtime, reverse=True):
        digest_md = d / "digest.md"
        if not d.is_dir() or not digest_md.exists():
            continue
        title = d.name
        try:
            for line in digest_md.read_text().splitlines():
                if line.startswith("# "):
                    title = line[2:].strip()
                    break
        except Exception:
            pass
        # Sidebar enrichment: try to load metadata.json + thumbnail presence
        # so the sidebar can render small visual anchors. All optional;
        # legacy digests without metadata.json fall back to text-only.
        entry: dict = {
            "id": d.name, "title": title, "mtime": d.stat().st_mtime,
            "has_thumbnail": (d / "thumbnail.jpg").exists(),
            "channel_id": None, "channel_name": None,
            "has_channel_thumbnail": False,
        }
        meta_path = d / "metadata.json"
        if meta_path.exists():
            try:
                meta = json.loads(meta_path.read_text())
                entry["channel_id"] = meta.get("channel_id")
                entry["channel_name"] = meta.get("channel_name")
                if entry["channel_id"]:
                    ch_thumb = (
                        get_data_dir() / "channel_thumbnails"
                        / f"{entry['channel_id']}.jpg"
                    )
                    entry["has_channel_thumbnail"] = ch_thumb.exists()
            except Exception:
                pass
        results.append(entry)
    return results


def _render_markdown(text: str) -> str:
    import markdown as md_lib
    html = md_lib.markdown(text, extensions=["fenced_code", "tables", "sane_lists"])
    # Rewrite cross-references to other digests (e.g. ../digests/X/digest.md) into view URLs.
    html = re.sub(r'href="[^"]*digests/([^"/]+)/digest\.md"', r'href="/digests/\1/"', html)
    # Open external links in a new tab so the reader doesn't lose their place
    # when clicking a YouTube timestamp / "Watch on YouTube" link. rel=noopener
    # blocks the new tab from manipulating window.opener (web-security best
    # practice). Skip anchors that already have a target= attribute.
    html = re.sub(
        r'<a (href="https?://[^"]+")(?![^>]*\btarget=)',
        r'<a \1 target="_blank" rel="noopener"',
        html,
    )
    return html


def _viewer_nav(
    video_id: str,
    current_view: str,
    digests_dir: Path,
) -> str:
    """Tab-bar style navigation shared across the digest / takeaway / panel
    viewer pages. The current_view ("digest" | "takeaway" | "panel") gets
    the primary highlighted styling; the other links are secondary outlines.
    Generation buttons replace "View X" links when an artifact is missing.

    Returns a single <div class='digest-actions digest-toolbar'>...</div>
    block ready to inject at the top of a viewer page body.
    """
    from html import escape as h

    takeaway_md = digests_dir / video_id / "takeaway.md"
    panel_md = digests_dir / video_id / "panel.md"
    slides_path = digests_dir / video_id / "slides.pptx"
    panel_job = local_job_status(f"{video_id}:panel")
    takeaway_job = local_job_status(f"{video_id}:takeaway")
    slides_job = local_job_status(f"{video_id}:slides")

    def cls(view: str) -> str:
        return "discuss-btn" if view == current_view else "discuss-btn-secondary"

    def running_placeholder(label: str, kind: str, elapsed: int) -> str:
        """Inline placeholder shown while a background job is in flight.
        The data-poll-url drives _JOB_POLL_JS which updates `.elapsed` and
        reloads the page when the artifact appears."""
        return (
            "<span class='discuss-btn-secondary' style='cursor: default;' "
            f"data-poll-url='/digests/{h(video_id)}/job-status?kind={kind}'>"
            f"Generating {label}… <span class='elapsed'>{elapsed}s</span>"
            "</span>"
        )

    parts: List[str] = ["<div class='digest-actions digest-toolbar'>"]

    # Always-visible: View digest. The page itself exists by definition (the
    # caller would have 404'd before reaching this helper).
    parts.append(
        f"<a class='{cls('digest')}' style='text-decoration:none;' "
        f"href='/digests/{h(video_id)}/'>View digest</a>"
    )

    # Panel — view link / running placeholder / generate button.
    if panel_md.exists():
        parts.append(
            f"<a class='{cls('panel')}' style='text-decoration:none;' "
            f"href='/digests/{h(video_id)}/panel/'>View panel discussion</a>"
        )
    elif panel_job.get("phase") == "running":
        parts.append(running_placeholder(
            "panel", "panel", panel_job.get("elapsed", 0),
        ))
    else:
        parts.append(
            f"<form method='post' action='/digests/{h(video_id)}/discuss' "
            f"style='display:inline;'>"
            "<button type='submit' class='discuss-btn-secondary' "
            "title='Generates a panel-of-experts discussion "
            "(~60–120s, one Opus call). Runs in the background; the page "
            "will refresh when ready.'>"
            "Generate panel</button></form>"
        )

    # Takeaway — view link / running placeholder / generate button.
    if takeaway_md.exists():
        parts.append(
            f"<a class='{cls('takeaway')}' style='text-decoration:none;' "
            f"href='/digests/{h(video_id)}/takeaway/'>View takeaway</a>"
        )
    elif takeaway_job.get("phase") == "running":
        parts.append(running_placeholder(
            "takeaway", "takeaway", takeaway_job.get("elapsed", 0),
        ))
    else:
        parts.append(
            f"<form method='post' action='/digests/{h(video_id)}/takeaway' "
            f"style='display:inline;'>"
            "<button type='submit' class='discuss-btn-secondary' "
            "title='Writes a 1–3 paragraph takeaway (~30s, one Sonnet call). "
            "Runs in the background; the page will refresh when ready.'>"
            "Generate takeaway</button></form>"
        )

    # Slides — download link / running placeholder / generate button.
    if slides_path.exists():
        parts.append(
            f"<a class='discuss-btn-secondary' style='text-decoration:none;' "
            f"href='/digests/{h(video_id)}/slides.pptx' "
            "title='Download the auto-generated PowerPoint deck "
            "(one slide per topic, with intelligently-picked frames).'>"
            "Download slides</a>"
        )
    elif slides_job.get("phase") == "running":
        parts.append(running_placeholder(
            "slides", "slides", slides_job.get("elapsed", 0),
        ))
    else:
        parts.append(
            f"<form method='post' action='/digests/{h(video_id)}/slides' "
            f"style='display:inline;'>"
            "<button type='submit' class='discuss-btn-secondary' "
            "title='Builds slides.pptx from the cached video — local frame "
            "extraction + alignment + PowerPoint assembly. No LLM call. "
            "Runs in the background; the page will refresh when ready.'>"
            "Generate slides</button></form>"
        )

    parts.append("</div>")
    return "".join(parts)


def _audio_section(
    video_id: str,
    kind: str,
    digests_dir: Path,
) -> str:
    """HTML block placed above the rendered markdown on a viewer page.
    Four states: inline <audio> player when the MP3 exists, a running
    placeholder while the background job is in flight, an error banner
    with a retry button when the previous attempt failed, or a Generate
    audio button when no MP3 has been requested yet.
    """
    from html import escape as h

    mp3_path = digests_dir / video_id / f"{kind}.mp3"
    job = local_job_status(f"{video_id}:audio_{kind}")
    phase = job.get("phase")

    if mp3_path.exists():
        return (
            "<div class='audio-row' style='margin: 16px 0 24px; "
            "display: flex; gap: 12px; align-items: center; flex-wrap: wrap;'>"
            f"<audio controls preload='none' "
            f"src='/digests/{h(video_id)}/audio/{h(kind)}.mp3' "
            "style='width: 100%; max-width: 520px; height: 36px;'></audio>"
            f"<a href='/digests/{h(video_id)}/audio/{h(kind)}.mp3' download "
            "style='font-size: 12px; color: var(--muted); "
            "text-decoration: none;'>Download MP3</a>"
            "</div>"
        )
    if phase == "running":
        # No live counter — the poll JS still reloads the page when the
        # MP3 lands (or surfaces an error), but a fixed message avoids
        # the misleading "stuck at 0s" appearance if a browser tab loses
        # focus or the poll lags. Duration hint tells the user it's
        # normal to wait rather than wondering if it hung.
        return (
            "<div class='audio-row' style='margin: 16px 0 24px;'>"
            "<span class='discuss-btn-secondary' style='cursor: default;' "
            f"data-poll-url='/digests/{h(video_id)}/"
            f"job-status?kind=audio_{h(kind)}'>"
            "🎧 Generating audio…"
            "</span>"
            "<span style='margin-left: 10px; font-size: 12px; "
            "color: var(--muted);'>"
            "Usually 30s for a takeaway, 1–3 min for digest / panel. "
            "The page will refresh when it's ready."
            "</span>"
            "</div>"
        )
    if phase == "error":
        # Last attempt failed — surface the error message inline so the
        # user can see what went wrong (e.g. "ElevenLabs library voice
        # requires paid plan" → switch voice in Settings) and offer a
        # one-click retry without making them dig into logs.
        err = (job.get("error") or "(unknown error)")
        return (
            "<div class='audio-row' style='margin: 16px 0 24px; "
            "padding: 12px 14px; border-radius: 4px; "
            "background: var(--code-bg); "
            "border-left: 3px solid #c00;'>"
            "<div style='font-weight: 600; margin-bottom: 6px;'>"
            "🎧 Audio generation failed</div>"
            f"<div style='font-size: 13px; color: var(--muted); "
            f"margin-bottom: 8px; word-break: break-word;'>{h(err)}</div>"
            f"<form method='post' action='/digests/{h(video_id)}/audio/{h(kind)}' "
            "style='display:inline;'>"
            "<button type='submit' class='discuss-btn-secondary' "
            "title='Retry generating audio. If this keeps failing, "
            "check your TTS provider settings (voice ID, API key).'>"
            "Retry</button></form>"
            "</div>"
        )
    return (
        "<div class='audio-row' style='margin: 16px 0 24px;'>"
        f"<form method='post' action='/digests/{h(video_id)}/audio/{h(kind)}' "
        "style='display:inline;'>"
        "<button type='submit' class='discuss-btn-secondary' "
        "title='Render this artifact to an MP3 using your configured TTS "
        "provider. Takes ~30s for a takeaway, 1-3 min for digest/panel.'>"
        "🎧 Listen (generate audio)</button></form></div>"
    )


def _any_local_job_running(video_id: str) -> bool:
    """True if any background job for this video is currently in flight.
    Drives whether the viewer page injects the polling JS — only emitted
    when there's actually something to poll."""
    kinds = (
        "panel", "takeaway", "slides",
        "audio_digest", "audio_panel", "audio_takeaway",
    )
    for kind in kinds:
        if local_job_status(f"{video_id}:{kind}").get("phase") == "running":
            return True
    return False


def cmd_serve(args) -> int:
    try:
        from flask import Flask, render_template_string, abort, send_from_directory
    except ImportError:
        sys.exit(
            "Flask is required for the reader. Reinstall with:\n"
            "  uv tool install --reinstall git+https://github.com/jyouturner/youtube-to-markdown"
        )

    data_dir = get_data_dir()
    digests_dir = data_dir / "digests"

    # Opportunistic, non-fatal: refresh billing-calibrated pricing if the cache
    # is missing or older than 7 days. No-op without an Admin key.
    try:
        _pc = _pricing_cache_path()
        import time as _t
        _stale = (not _pc.exists()
                  or (_t.time() - (json.loads(_pc.read_text()).get("_calibrated_at_epoch", 0))) > 7 * 86400)
        if _stale and _admin_api_key():
            _res = calibrate_pricing_from_billing()
            if _res.get("ok") and _res.get("derived"):
                print(f"[yt2md] refreshed pricing from billing: {', '.join(_res['models'])}")
    except Exception:
        pass

    app = Flask(__name__)
    # Disable Flask's default request logging — keep stdout clean.
    import logging
    logging.getLogger("werkzeug").setLevel(logging.WARNING)

    def page(body: str, *, title: str, current: str, base_href: str = None):
        # Annotate listing items with read state so the sidebar can show "new" markers.
        digests = _list_digests(digests_dir)
        try:
            read_digests = _read_digest_ids()
        except Exception:
            read_digests = set()
        for d in digests:
            d["unread"] = d["id"] not in read_digests
        unread_digest_count = sum(1 for d in digests if d["unread"])

        # Persistent banner above every page when neither auth path is
        # configured. Reading cached digests still works; only generation does,
        # so we warn rather than gate. The Setup page is the one place this is
        # hidden (the page itself IS the configuration UI).
        has_api_key = bool(os.environ.get("ANTHROPIC_API_KEY"))
        has_claude_code = (
            claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)
        )
        if current != "setup" and not (has_api_key or has_claude_code):
            banner = (
                '<div class="flash" style="border-left-color: #c00;">'
                '<strong>LLM auth not configured.</strong> '
                'Generating digests and panel discussions requires either an '
                'Anthropic API key or a Claude Code subscription login. '
                'Reading existing digests still works. '
                '<a href="/setup">Set it up →</a>'
                '</div>'
            )
            body = banner + body

        # Snapshot of library state so the sidebar's poll JS can detect a
        # new digest landing (auto-pipeline finishing, etc.) without a
        # manual refresh. We send count + max mtime to the page; the JS
        # polls /sidebar-status and shows a "new digest" banner if either
        # value drifts.
        max_mtime = max((d["mtime"] for d in digests), default=0.0)
        return render_template_string(
            SERVE_PAGE_TEMPLATE,
            body=body,
            title=title,
            current=current,
            base_href=base_href,
            digests=digests,
            channel_count=len(read_channels()),
            unread_digest_count=unread_digest_count,
            digests_count=len(digests),
            digests_max_mtime=max_mtime,
        )

    def _require_llm_or_redirect():
        """Helper for action endpoints: returns a redirect Response if neither
        an API key nor a logged-in Claude Code sandbox is configured. Use as:
        r = _require_llm_or_redirect(); if r is not None: return r."""
        from flask import redirect
        if os.environ.get("ANTHROPIC_API_KEY"):
            return None
        if claude_code_installed() and _claude_code_session_state.get("logged_in"):
            return None
        return redirect("/setup?msg=Auth+required+to+run+this+action.")

    @app.route("/")
    def home():
        from flask import redirect
        digests = _list_digests(digests_dir)
        channels = read_channels()

        # True first-run (no auth at all + nothing on disk yet): land directly
        # on the setup page so the user isn't asked to subscribe before they
        # can generate anything. Skip when there are existing digests — those
        # should still be readable even without a key.
        has_api_key = bool(os.environ.get("ANTHROPIC_API_KEY"))
        has_claude_code = (
            claude_code_installed()
            and _claude_code_session_state.get("logged_in", False)
        )
        if not digests and not channels and not (has_api_key or has_claude_code):
            return redirect("/setup")

        # Empty states first — show a real CTA, not a list of zero items.
        if not digests:
            if not channels:
                body = (
                    "<h1>Welcome to yt2md</h1>"
                    "<p>You haven't subscribed to any channels yet.</p>"
                    '<p><a class="cta" href="/channels">Add your first channel →</a></p>'
                )
            else:
                body = (
                    "<h1>Polling is set up</h1>"
                    f"<p>You're watching {len(channels)} channel(s). Your first digest "
                    "will appear after the next polling run (every few hours, or "
                    'fire one now from the <a href="/schedule">Schedule</a> page).</p>'
                )
            return page(body, title="yt2md", current="home")

        # Featured content: the latest digest's body, with an eyebrow link.
        featured = digests[0]
        featured_md = (digests_dir / featured["id"] / "digest.md").read_text()
        body = (
            f'<p class="featured-eyebrow">Latest digest · '
            f'<a href="/digests/{featured["id"]}/">{featured["title"]}</a></p>'
        )
        body += _render_markdown(featured_md)
        base_href = f"/digests/{featured['id']}/"

        # No "More digests" footer here — sidebar is the navigation surface;
        # showing the same list twice is just noise.

        return page(body, title="Home", current="home", base_href=base_href)

    @app.route("/channels", methods=["GET"])
    def channels_page():
        from flask import request
        channels = read_channels()
        digests = _list_digests(digests_dir)
        sched_state = _load_schedule_state()
        poll_has_run = bool((sched_state.get("poll") or {}).get("last_started_at"))
        flash = request.args.get("msg", "")

        body = "<h1>Subscriptions</h1>"
        if flash:
            body += f'<div class="flash">{flash}</div>'

        # "What's next" guidance — adapts to current state. Only shown until the
        # user has at least one digest, then disappears.
        next_step = None
        if not channels:
            next_step = (
                "Paste a YouTube channel URL below to get started. "
                "Each new video on this channel will be auto-digested."
            )
        elif not poll_has_run:
            next_step = (
                "You\'re subscribed but the scheduler hasn\'t fired its first poll yet. "
                'It runs every few hours — or trigger one now from the '
                '<a href="/schedule">Schedule page</a>.'
            )
        elif not digests:
            next_step = (
                "Polling fires every 6 hours. Your first digest will land after the next run "
                '— or fire one now from the <a href="/schedule">Schedule page</a>.'
            )
        if next_step:
            body += f'<div class="next-step"><strong>Next step:</strong> {next_step}</div>'

        body += (
            '<form method="post" action="/channels" class="add-form">'
            '<label for="channel-url" class="sr-only">YouTube channel URL</label>'
            '<input id="channel-url" type="text" name="url" '
            'placeholder="https://www.youtube.com/@channel/videos  (or @handle)" '
            'autofocus required>'
            '<button type="submit">Add</button>'
            '</form>'
        )
        if channels:
            body += '<ul class="channel-list">'
            for ch in channels:
                body += (
                    '<li>'
                    f'<span class="url">{ch}</span>'
                    '<form method="post" action="/channels/remove" style="margin:0;">'
                    f'<input type="hidden" name="url" value="{ch}">'
                    '<button type="submit">Remove</button>'
                    '</form>'
                    '</li>'
                )
            body += '</ul>'
        else:
            body += "<p class='empty-state'>No subscriptions yet. Paste a YouTube channel URL above.</p>"
        body += (
            "<p class='meta-info' style='margin-top:32px'>"
            "Stored in <code>~/yt2md/channels.txt</code>."
            "</p>"
        )

        # Per-channel Claude Project setup. Surfaces one link per channel
        # that already has at least one digest on disk (a Project on a
        # never-digested channel would have nothing to chat about).
        known = _known_channel_names(digests_dir)
        if known:
            from urllib.parse import quote
            body += (
                '<h2 style="margin-top:40px">Set up a Claude Project for this channel</h2>'
                '<p class="meta-info">'
                'Group a channel\'s digests into a <strong>Claude Project</strong> '
                'on claude.ai so follow-up questions chat against just that channel\'s '
                'library. Each link below renders the Custom Instructions you paste '
                'into a Project to scope it. Requires the '
                '<a href="https://github.com/jyouturner/youtube-to-markdown#talk-to-your-library-from-claude-mcp" target="_blank" rel="noopener">'
                'yt2md MCP server</a> connected in your Claude client.'
                '</p>'
                '<ul class="channel-list">'
            )
            for name in known:
                body += (
                    '<li>'
                    f'<span class="url">{name}</span>'
                    f'<a href="/channels/project-instructions?channel={quote(name)}">Setup →</a>'
                    '</li>'
                )
            body += '</ul>'

        return page(body, title="Subscriptions", current="channels")

    @app.route("/channels/project-instructions")
    def channels_project_instructions():
        """Render the per-channel Claude Project Custom Instructions in a
        copy-friendly textarea, with a 4-step setup walkthrough. The
        copyable text is exactly what `yt2md project-instructions
        --channel <name>` would print — single source of truth."""
        from flask import request, redirect
        from markupsafe import escape

        wanted = (request.args.get("channel") or "").strip()
        known = _known_channel_names(digests_dir)
        # Case-insensitive exact resolution so URL-encoded names round-trip
        # cleanly through casing differences.
        canonical = next(
            (n for n in known if n.lower() == wanted.lower()),
            None,
        )
        if not canonical:
            return redirect(
                "/channels?msg=No+digests+yet+for+that+channel"
            )

        instructions = render_claude_project_instructions(canonical)
        safe_name = escape(canonical)
        safe_instructions = escape(instructions)
        body = (
            f'<h1>Claude Project: {safe_name}</h1>'
            '<p>Paste the instructions below into a new Claude Project\'s '
            '<em>Custom Instructions</em> box. The Project will then scope '
            f'every chat to <strong>{safe_name}</strong>\'s digests via the '
            'yt2md MCP tools.</p>'
            '<ol class="setup-steps">'
            '<li>Open <a href="https://claude.ai/projects" target="_blank" rel="noopener">'
            'claude.ai/projects</a> and create a new Project. Name it after '
            f'the channel (e.g. <code>{safe_name}</code>).</li>'
            '<li>In the Project, open <strong>Custom Instructions</strong> '
            'and paste the text from the box below.</li>'
            '<li>Make sure the <a href="https://github.com/jyouturner/youtube-to-markdown#talk-to-your-library-from-claude-mcp" '
            'target="_blank" rel="noopener">yt2md MCP server</a> is connected '
            'in your Claude client. Without it the Project has no way to read '
            'your local library.</li>'
            '<li>Open a chat in the Project and ask away — try '
            f'<em>"what\'s the latest on {safe_name}?"</em> to confirm the '
            'agent is calling <code>list_digests</code> with the right channel filter.</li>'
            '</ol>'
            '<div style="margin-top:24px">'
            '<button type="button" data-copy-target="proj-instructions">'
            'Copy instructions</button>'
            '</div>'
            '<textarea id="proj-instructions" readonly '
            'style="width:100%; min-height:360px; margin-top:12px; '
            'font-family: ui-monospace, monospace; font-size: 13px; '
            'padding:12px; border:1px solid #ddd; border-radius:6px;">'
            f'{safe_instructions}'
            '</textarea>'
            '<p class="meta-info" style="margin-top:16px">'
            'Same text is available on the CLI: '
            f'<code>yt2md project-instructions --channel "{safe_name}"</code>.'
            '</p>'
        )
        return page(body, title=f"Claude Project — {canonical}",
                    current="channels")

    @app.route("/channels", methods=["POST"])
    def channels_add():
        from flask import request, redirect
        url = normalize_channel_url(request.form.get("url", ""))
        if not is_url(url):
            return redirect("/channels?msg=Not+a+valid+URL")
        channels = read_channels()
        if url in channels:
            return redirect(f"/channels?msg=Already+watching+{url}")
        channels.append(url)
        write_channels(channels)
        return redirect(f"/channels?msg=Added+{url}")

    @app.route("/channels/remove", methods=["POST"])
    def channels_remove():
        from flask import request, redirect
        url = request.form.get("url", "").strip()
        channels = [c for c in read_channels() if c != url]
        write_channels(channels)
        return redirect(f"/channels?msg=Removed+{url}")

    @app.route("/schedule")
    def schedule_page():
        from flask import request
        from html import escape as h
        import time as _t
        flash = request.args.get("msg", "")
        cfg = load_schedule_config()
        sched_state = _load_schedule_state()

        body = "<h1>Schedule</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'

        body += '<form method="post" action="/schedule/save" class="schedule-form">'
        body += '<div class="schedule-fields">'
        body += '<label>Polling interval'
        body += f'  <input type="number" name="poll_hours" value="{cfg["poll_interval_hours"]}" min="0.1" step="0.1" required>'
        body += '  <span class="suffix">hours</span>'
        body += '</label>'
        body += '</div>'  # schedule-fields

        body += '<button type="submit" class="primary">Save</button>'
        body += '</form>'

        body += (
            f'<p class="meta-info">Current schedule: {h(_format_schedule_summary(cfg))}. '
            'Scheduling runs inside this server — pauses while it\'s down, '
            'catches up on missed slots when you start it again.</p>'
        )

        for kind, friendly, desc in [
            ("poll", "Polling", "fires <code>yt2md watch run</code>"),
        ]:
            sentence, dot_class = _scheduler_status_summary(kind, sched_state)
            s = sched_state.get(kind) or {}
            next_at = _compute_next_poll(cfg, s.get("last_started_at"))

            body += f'<div class="job-block"><h3><span class="dot {dot_class}"></span>{friendly}</h3>'
            body += f'<p class="meta-info" style="margin: 0 0 12px;">{desc}</p>'
            body += f'<p class="job-summary">{sentence}</p>'
            body += (f'<p class="meta-info">Next run: {h(_format_next_run(next_at))} '
                     f'({h(_t.strftime("%Y-%m-%d %H:%M", _t.localtime(next_at)))}).</p>')
            running = kind in _scheduler_jobs
            disabled = " disabled" if running else ""
            running_label = " (already running)" if running else ""
            body += (
                f'<div class="job-actions">'
                f'<form method="post" action="/schedule/run/{kind}" style="margin:0;">'
                f'<button type="submit"{disabled}>Run now{running_label}</button>'
                f'</form></div>'
            )

            body += '<details><summary>Diagnostics</summary>'
            body += '<table class="status-table">'
            for k in ("last_started_at", "last_finished_at", "last_exit_code", "last_pid"):
                v = s.get(k)
                if v is None:
                    continue
                if k.endswith("_at"):
                    v = _t.strftime("%Y-%m-%d %H:%M:%S", _t.localtime(v))
                body += f'<tr><td>{k.replace("_", " ")}</td><td>{h(str(v))}</td></tr>'
            body += '</table></details>'

            log_path = data_dir / "logs" / f"{kind}.log"
            body += '<details style="margin-top: 8px;"><summary>Recent log (last 20 lines)</summary>'
            body += f'<div class="log-block">{h(_tail_log(log_path, 20))}</div>'
            body += '</details>'
            body += '</div>'

        body += '<p class="meta-info">Refresh the page to see updated status after a run.</p>'
        return page(body, title="Schedule", current="schedule")

    @app.route("/schedule/save", methods=["POST"])
    def schedule_save():
        from flask import redirect, request
        cfg = load_schedule_config()
        try:
            if request.form.get("poll_hours"):
                cfg["poll_interval_hours"] = float(request.form["poll_hours"])
        except Exception as e:
            return redirect(f"/schedule?msg=Invalid+input:+{e}")
        save_schedule_config(cfg)
        return redirect("/schedule?msg=Saved+(scheduler+picks+up+within+30s)")

    @app.route("/schedule/run/<job>", methods=["POST"])
    def schedule_run(job):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        if job != "poll":
            abort(404)
        with _schedule_lock():
            existing = _scheduler_jobs.get(job)
            if existing is not None and existing.poll() is None:
                return redirect(f"/schedule?msg={job}+is+already+running")
            proc = _fire_scheduled_job(job)
        if proc is None:
            return redirect(f"/schedule?msg=Failed+to+fire+{job}+(yt2md+not+on+PATH)")
        return redirect(
            f"/schedule?msg=Started+{job}+(pid+{proc.pid}).+Refresh+for+status."
        )

    @app.route("/settings", methods=["GET"])
    def settings_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")
        s = load_settings()
        # Show the EFFECTIVE value (settings.json with .env fallback) so the
        # form matches what the system is actually using. On Save we persist
        # whatever the user submits, which becomes the new canonical value.
        for key, env_name in (
            ("digest_model", "YT2MD_DIGEST_MODEL"),
            ("panel_model", "YT2MD_PANEL_MODEL"),
            ("whisper_model", "YT2MD_WHISPER_MODEL"),
            ("cookies_from_browser", "YT2MD_COOKIES_FROM_BROWSER"),
            ("digest_language", "YT2MD_DIGEST_LANGUAGE"),
        ):
            if not s.get(key) and os.environ.get(env_name):
                s[key] = os.environ[env_name]

        whisper_choices = ("tiny", "base", "small", "medium", "large-v2", "large-v3")
        cookie_choices = ("", "chrome", "firefox", "safari", "brave", "edge",
                          "chromium", "opera", "vivaldi")

        body = "<h1>Settings</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Stored in <code>~/yt2md/settings.json</code> '
            '(API key in <code>~/yt2md/.env</code>). '
            'New values take effect on the next one-off submit / scheduled poll / '
            '"Discuss with experts" click — no restart required.</p>'
        )

        body += '<form method="post" action="/settings/save" class="schedule-form">'
        body += '<div class="schedule-fields" style="grid-template-columns: 1fr;">'

        cur_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if cur_key:
            cur_state = (
                f'<span style="color: var(--accent);">set</span> '
                f'(<code>{h(cur_key[:7])}…{h(cur_key[-4:])}</code>)'
            )
        else:
            cur_state = '<span style="color: #c00;">not set</span>'
        body += (
            '<label>Anthropic API key'
            '  <input type="password" name="anthropic_api_key" '
            '    placeholder="sk-ant-... (leave blank to keep current)" '
            '    autocomplete="off">'
            '  <span class="suffix" style="display:block;">'
            f'    Current: {cur_state}. '
            f'    {API_KEY_COST_NOTE} '
            '    Get a key at '
            '    <a href="https://console.anthropic.com/settings/keys" target="_blank" '
            '       rel="noopener">console.anthropic.com/settings/keys</a>. '
            '    The key is validated on save with a 1-token test call.'
            '  </span>'
            '</label>'
        )

        body += (
            '<label>Digest model'
            f'  <input type="text" name="digest_model" value="{h(s["digest_model"])}" required>'
            '  <span class="suffix" style="display:block;">'
            'Anthropic model ID for the per-video digest. e.g. <code>claude-sonnet-4-6</code>, '
            '<code>claude-opus-4-7</code>, <code>claude-haiku-4-5-20251001</code>.'
            '</span>'
            '</label>'
        )

        body += (
            '<label>Panel-discussion model'
            f'  <input type="text" name="panel_model" value="{h(s["panel_model"])}" required>'
            '  <span class="suffix" style="display:block;">'
            'Used when you click "Discuss with experts" on a digest. Multi-perspective '
            'synthesis benefits from a stronger reasoning model — Opus is the default.'
            '</span>'
            '</label>'
        )

        body += '<label>Whisper model'
        body += '  <select name="whisper_model">'
        for w in whisper_choices:
            sel = ' selected' if s["whisper_model"] == w else ''
            body += f'    <option value="{w}"{sel}>{w}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Local STT fallback when YouTube has no captions. Larger = better quality, '
            'slower, bigger first-run download. <code>medium</code> is a good default.'
            '</span>'
            '</label>'
        )

        body += '<label>Digest language'
        body += '  <select name="digest_language">'
        for code, label in (
            ("auto", "auto — match the transcript's language"),
            ("en", "en — always English"),
        ):
            sel = ' selected' if s.get("digest_language", "auto") == code else ''
            body += f'    <option value="{code}"{sel}>{label}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Applies to both the per-video digest and the panel discussion. '
            '<code>auto</code> writes in the source language (e.g. Chinese for a Chinese-language video). '
            '<code>en</code> forces English regardless.'
            '</span>'
            '</label>'
        )

        body += '<label>Cookies from browser'
        body += '  <select name="cookies_from_browser">'
        for c in cookie_choices:
            sel = ' selected' if s.get("cookies_from_browser", "") == c else ''
            label = "(none)" if c == "" else c
            body += f'    <option value="{c}"{sel}>{label}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'YouTube increasingly requires logged-in cookies. Pick the browser you\'re '
            'signed into YouTube on; yt-dlp extracts cookies on each run. Leave as '
            '"(none)" if you only digest publicly-accessible videos.'
            '</span>'
            '</label>'
        )

        # LLM backend selector
        backend_choices = (
            ("auto", "auto — pick API when ANTHROPIC_API_KEY is set, else Claude Code"),
            ("api", "api — direct Anthropic API (requires ANTHROPIC_API_KEY)"),
            ("claude-code", "claude-code — bundled Claude Code via `claude -p` (will bill against the Agent SDK credit pool after Jun 15 2026)"),
            ("claude-code-pty", "claude-code-pty — primary local `claude` driven as interactive REPL (stays on Pro/Max plan; if ANTHROPIC_API_KEY is also set, vision calls auto-route to API so frame quality is preserved)"),
        )
        body += '<label>LLM backend'
        body += '  <select name="llm_backend">'
        for code, label in backend_choices:
            sel = ' selected' if s.get("llm_backend", "auto") == code else ''
            body += f'    <option value="{code}"{sel}>{h(label)}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Which auth path to use for digest / panel calls. The "auto" mode '
            'picks the cheapest available path. Switch from <a href="/setup">/setup</a>. '
            '<strong>claude-code-pty</strong> uses your machine\'s primary '
            '<code>claude</code> install (must be signed in to Pro/Max) and drives '
            'it as an interactive REPL — this is the only path that stays on '
            'subscription billing after Jun 15 2026, when Anthropic moves '
            '<code>claude -p</code> to a separate metered Agent SDK credit pool. '
            'When PTY is selected and an <code>ANTHROPIC_API_KEY</code> is also '
            'present, image-based calls (slide classifier and per-topic '
            'frame-picking) auto-route through the API so frame quality is '
            'preserved — the text-heavy 80% of the bill still goes to your plan.'
            '</span>'
            '</label>'
        )

        # Claude Code vision toggle
        cc_vision_checked = ' checked' if s.get("claude_code_vision") else ''
        body += (
            '<label>'
            f'  <input type="checkbox" name="claude_code_vision" value="1"{cc_vision_checked}> '
            'Enable vision frame-picking under Claude Code backend'
            '  <span class="suffix" style="display:block;">'
            'Off by default. The Claude Code CLI has no native image flag, so '
            'enabling this base64-embeds frames into prompts (token-heavy). '
            'Has no effect when the API backend is in use (which always uses '
            'native vision).'
            '  </span>'
            '</label>'
        )

        # TTS provider + per-provider settings for the 🎧 Listen feature
        # on viewer pages. macOS `say` is free + offline + lower quality;
        # ElevenLabs is paid + cloud + much higher quality.
        tts_provider_choices = (
            ("macos", "macOS `say` — free, offline, lower quality"),
            ("elevenlabs", "ElevenLabs — paid, cloud, much higher quality"),
        )
        body += '<label>TTS provider'
        body += '  <select name="tts_provider">'
        for code, label in tts_provider_choices:
            sel = ' selected' if s.get("tts_provider", "macos") == code else ''
            body += f'    <option value="{code}"{sel}>{h(label)}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Which backend renders the MP3 when you click 🎧 Listen on a '
            'digest / panel / takeaway. ElevenLabs sounds dramatically '
            'better but costs against your plan credits — set '
            '<code>ELEVENLABS_API_KEY</code> in <code>~/yt2md/.env</code> '
            '(get one at <a href="https://elevenlabs.io/app/settings/api-keys" '
            'target="_blank" rel="noopener">elevenlabs.io</a>).'
            '  </span>'
            '</label>'
        )

        # macOS-only fields.
        body += (
            '<label>macOS TTS voice'
            f'  <input type="text" name="tts_voice" '
            f'    value="{h(s.get("tts_voice") or "")}" '
            '    placeholder="(leave blank to use the system default voice)" '
            '    autocomplete="off">'
            '  <span class="suffix" style="display:block;">'
            'Only used when provider = macOS. For the best quality, leave '
            'this blank AND set a Siri voice as your system default in '
            '<strong>System Settings → Accessibility → Spoken Content → '
            'System Voice</strong> (download "Siri Voice 1"). For non-Siri '
            'options try <code>Fiona</code>, <code>Samantha (Enhanced)</code>, '
            'or run <code>say -v ?</code> in Terminal to see every voice '
            'installed on your machine.'
            '  </span>'
            '</label>'
        )
        body += (
            '<label>macOS TTS speaking rate'
            f'  <input type="text" name="tts_rate" '
            f'    value="{h(s.get("tts_rate") or "")}" '
            '    placeholder="(blank = system default, ~175 wpm)" '
            '    inputmode="numeric" autocomplete="off">'
            '  <span class="suffix" style="display:block;">'
            'Only used when provider = macOS. Words per minute — 150 reads '
            'slowly + relaxed, 200+ skims faster. Blank uses the system '
            'default (~175).'
            '  </span>'
            '</label>'
        )

        # ElevenLabs-only fields. Curated list = the "Default voices" set
        # ElevenLabs ships to every account (free tier included). Library
        # voices like the old "Rachel" require a paid plan via API, so
        # they're deliberately not in the dropdown — users on a paid plan
        # who want a library voice can fall back to the "Custom voice ID"
        # text field below.
        elevenlabs_voice_choices = (
            ("nPczCjzI2devNBz1zQrb", "Brian — US male, deep + calm (default)"),
            ("EXAVITQu4vr4xnSDxMaL", "Sarah — US female, soft"),
            ("9BWtsMINqrJLrRacOk9x", "Aria — US female, conversational"),
            ("cgSgspJ2msm6clMCkdW9", "Jessica — US female, expressive"),
            ("XrExE9yKIg1WjnnlVkGX", "Matilda — US female, friendly"),
            ("bIHbv24MWmeRgasZH58o", "Will — US male, friendly"),
            ("cjVigY5qzO86Huf0OWal", "Eric — US male, friendly"),
            ("iP95p4xoKVk53GoZ742B", "Chris — US male, casual"),
            ("pqHfZKP75CvOlQylNhV4", "Bill — US male, gruff/trustworthy"),
            ("TX3LPaxmHKxFdv7VOQHJ", "Liam — US male, articulate"),
            ("JBFqnCBsd6RMkjVDRZzb", "George — UK male, warm"),
            ("onwK4e9ZLuTAKqWW03F9", "Daniel — UK male, news-anchor"),
            ("Xb7hH8MSUJpSbSDYk0k2", "Alice — UK female, confident"),
            ("pFZP5JQG7iQjIQuC4Bku", "Lily — UK female, warm"),
        )
        current_voice = s.get("elevenlabs_voice_id") or "nPczCjzI2devNBz1zQrb"
        known_ids = {code for code, _ in elevenlabs_voice_choices}
        # If the saved voice ID isn't in the curated list (e.g. a custom
        # library voice on a paid plan), preserve it as "__custom__" so
        # the selector still reflects the user's actual configuration.
        is_custom = current_voice not in known_ids and current_voice != ""
        body += '<label>ElevenLabs voice'
        body += '  <select name="elevenlabs_voice_id">'
        for code, label in elevenlabs_voice_choices:
            sel = ' selected' if current_voice == code else ''
            body += f'    <option value="{code}"{sel}>{h(label)}</option>'
        custom_sel = ' selected' if is_custom else ''
        body += (
            f'    <option value="__custom__"{custom_sel}>'
            'Custom voice ID (use field below)</option>'
        )
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Only used when provider = ElevenLabs. These are the '
            '"Default voices" included with every plan (free tier '
            'included). Library voices from '
            '<a href="https://elevenlabs.io/app/voice-library" target="_blank" '
            'rel="noopener">elevenlabs.io/app/voice-library</a> require a '
            'paid plan — to use one, pick <em>Custom voice ID</em> above '
            'and paste the ID into the field below.'
            '  </span>'
            '</label>'
        )
        body += (
            '<label>Custom voice ID (optional)'
            f'  <input type="text" name="elevenlabs_voice_id_custom" '
            f'    value="{h(current_voice) if is_custom else ""}" '
            '    placeholder="paste a library voice ID here" '
            '    autocomplete="off">'
            '  <span class="suffix" style="display:block;">'
            'Only takes effect when "Custom voice ID" is selected above. '
            'Requires a paid ElevenLabs plan for library voices.'
            '  </span>'
            '</label>'
        )
        elevenlabs_model_choices = (
            ("eleven_multilingual_v2", "eleven_multilingual_v2 — best quality"),
            ("eleven_turbo_v2_5", "eleven_turbo_v2_5 — faster, slightly lower quality"),
            ("eleven_flash_v2_5", "eleven_flash_v2_5 — fastest, cheapest"),
        )
        body += '<label>ElevenLabs model'
        body += '  <select name="elevenlabs_model">'
        for code, label in elevenlabs_model_choices:
            cur = s.get("elevenlabs_model") or "eleven_multilingual_v2"
            sel = ' selected' if cur == code else ''
            body += f'    <option value="{code}"{sel}>{h(label)}</option>'
        body += '  </select>'
        body += (
            '  <span class="suffix" style="display:block;">'
            'Only used when provider = ElevenLabs. multilingual_v2 is the '
            'high-quality default; the turbo/flash options trade quality '
            'for speed + lower credit cost.'
            '  </span>'
            '</label>'
        )

        body += '</div>'  # schedule-fields
        body += '<button type="submit" class="primary">Save</button>'
        body += '</form>'
        return page(body, title="Settings", current="settings")

    @app.route("/settings/save", methods=["POST"])
    def settings_save():
        from flask import redirect, request
        from urllib.parse import quote_plus
        s = load_settings()
        for key in ("digest_model", "panel_model", "whisper_model",
                    "cookies_from_browser", "digest_language", "llm_backend",
                    "tts_provider", "tts_voice", "tts_rate",
                    "elevenlabs_voice_id", "elevenlabs_model"):
            v = request.form.get(key)
            if v is not None:
                s[key] = v.strip()
        # The voice selector ships the sentinel "__custom__" when the
        # user wants to paste a library voice ID. Swap in the custom
        # field's value so downstream code (which only reads
        # elevenlabs_voice_id) doesn't need to know about the sentinel.
        if s.get("elevenlabs_voice_id") == "__custom__":
            custom = (request.form.get("elevenlabs_voice_id_custom") or "").strip()
            s["elevenlabs_voice_id"] = custom or "nPczCjzI2devNBz1zQrb"
        # Checkboxes are absent from request.form when unchecked.
        s["claude_code_vision"] = bool(request.form.get("claude_code_vision"))
        save_settings(s)

        # API key is stored separately (.env, not settings.json) so non-secret
        # config can be checked into a shared settings file without leaking it.
        new_key = (request.form.get("anthropic_api_key") or "").strip()
        if new_key:
            err = validate_api_key(new_key)
            if err:
                return redirect(
                    f"/settings?msg=Settings+saved+but+API+key+rejected:+{quote_plus(err)}"
                )
            set_env_var("ANTHROPIC_API_KEY", new_key)
            return redirect("/settings?msg=Saved+(API+key+validated).")
        return redirect("/settings?msg=Saved.")

    @app.route("/setup", methods=["GET"])
    def setup_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")

        body = "<h1>Connect Claude</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Pick one of the two paths below. yt2md needs '
            'access to a Claude model to generate digests and panel discussions; '
            'either an Anthropic API key (per-call billing) or a Claude.ai '
            'subscription via the bundled Claude Code (no extra billing).</p>'
        )

        # --- Side-by-side dual-auth grid ---
        body += (
            '<div class="auth-grid" style="display: grid; '
            'grid-template-columns: repeat(auto-fit, minmax(360px, 1fr)); '
            'gap: 20px; margin-top: 20px;">'
        )

        # API key panel
        cur_key = os.environ.get("ANTHROPIC_API_KEY", "")
        api_already_set = bool(cur_key)
        api_panel = '<div class="schedule-form">'
        api_panel += '<h2 style="margin-top: 0;">Anthropic API key</h2>'
        if api_already_set:
            api_panel += (
                '<div class="flash" style="border-left-color: var(--accent); '
                'margin-top: 0;">'
                f'<strong>Configured</strong> '
                f'(<code>{h(cur_key[:7])}…{h(cur_key[-4:])}</code>). '
                '<a href="/">Go to library →</a>'
                '</div>'
            )
        api_panel += (
            '<p class="meta-info" style="margin-top: 0;">'
            'For developers / users who already have an Anthropic API account. '
            'Pay-per-call. Native vision support, prompt caching.'
            '</p>'
            '<ol style="margin: 8px 0 16px 20px; padding: 0; font-size: 14px;">'
            '<li>Get a key at <a href="https://console.anthropic.com/settings/keys" '
            'target="_blank" rel="noopener">console.anthropic.com</a>.</li>'
            '<li>Add a payment method (subscription does NOT cover API usage).</li>'
            '<li>Paste below.</li>'
            '</ol>'
            f'<p class="meta-info" style="font-size: 13px;">{API_KEY_COST_NOTE}</p>'
            '<form method="post" action="/setup/save-api-key">'
            '<label>API key'
            '  <input type="password" name="anthropic_api_key" '
            '    placeholder="sk-ant-..." autocomplete="off" '
            'style="width: 100%; box-sizing: border-box;">'
            '</label>'
            '<button type="submit" class="primary" style="margin-top: 12px;">'
            'Save and validate</button>'
            '</form>'
            '</div>'
        )

        # Claude Code panel — installation + login flow
        snap = claude_setup_snapshot()
        cc_panel = '<div class="schedule-form">'
        cc_panel += '<h2 style="margin-top: 0;">Claude.ai subscription</h2>'
        cc_panel += (
            '<p class="meta-info" style="margin-top: 0;">'
            'Uses your Pro/Max plan via a bundled, sandboxed copy of Claude Code. '
            'No extra billing. Vision off by default (toggle in Settings). '
            'Prompt caching disappears, so very long videos may use more tokens.'
            '</p>'
        )

        # Status & action area, populated by JS poll, with a server-rendered
        # initial state for users who hit the page without JS.
        cc_panel += '<div id="cc-status" style="margin-bottom: 12px;">'
        if snap["logged_in"]:
            cc_panel += (
                '<div class="flash" style="border-left-color: var(--accent); margin: 0;">'
                '<strong>Signed in.</strong> '
                '<a href="/">Go to library →</a>'
                '</div>'
            )
        elif not snap["node_ok"]:
            cc_panel += (
                '<div class="flash" style="border-left-color: #c00; margin: 0;">'
                f'<strong>Node.js {MIN_NODE_MAJOR}+ required.</strong> '
                'Install with <code>brew install node</code> (macOS) / '
                f'<code>winget install OpenJS.NodeJS</code> (Windows) or your '
                'package manager, then refresh this page.'
                '</div>'
            )
        elif not snap["installed"]:
            cc_panel += (
                '<p class="meta-info" style="margin: 0;">'
                f'Step 1: install Claude Code into <code>~/yt2md/claude-code/</code> '
                f'(~200&nbsp;MB; isolated from any system install).'
                '</p>'
            )
        else:
            cc_panel += (
                '<p class="meta-info" style="margin: 0;">'
                'Step 2: sign in. A new browser tab will open for Claude.ai OAuth. '
                'Complete sign-in there; this page will update automatically.'
                '</p>'
            )
        cc_panel += '</div>'

        # Action buttons (the JS toggles these based on status).
        install_disabled = (
            ' disabled' if not snap["node_ok"] or snap["install_running"]
            or snap["installed"] else ''
        )
        login_visible = snap["installed"] and not snap["logged_in"]
        cc_panel += '<div id="cc-actions" style="display: flex; gap: 8px; flex-wrap: wrap;">'
        cc_panel += (
            f'<form method="post" action="/setup/install-claude" style="display:inline;">'
            f'<button type="submit" class="primary"{install_disabled} '
            'id="cc-install-btn">Install Claude Code</button>'
            '</form>'
        )
        if login_visible:
            cc_panel += (
                '<form method="post" action="/setup/login-claude" style="display:inline;">'
                '<button type="submit" class="primary" id="cc-login-btn">'
                'Sign in with Claude.ai</button>'
                '</form>'
            )
        if snap["logged_in"]:
            cc_panel += (
                '<form method="post" action="/setup/logout-claude" style="display:inline;" '
                'onsubmit="return confirm(\'Sign out of Claude Code?\');">'
                '<button type="submit">Sign out</button>'
                '</form>'
            )
        cc_panel += '</div>'

        # Live log tail (hidden when both logs are empty).
        cc_panel += (
            '<details id="cc-log-details" style="margin-top: 12px;'
            + ('' if (snap["install_log_tail"] or snap["login_log_tail"]) else ' display: none;')
            + '">'
            '<summary>Recent log</summary>'
            '<pre id="cc-log-tail" class="log-block" style="max-height: 200px; '
            'overflow: auto; font-size: 11px;">'
            + h((snap["login_log_tail"] or snap["install_log_tail"] or "").strip())
            + '</pre>'
            '</details>'
        )
        cc_panel += '</div>'

        body += api_panel + cc_panel + "</div>"  # close auth-grid

        # Polling JS: refreshes the Claude Code panel every 2s while a job is
        # running, so the user sees install / login progress without refreshing.
        body += """
<script>
(function () {
  const statusEl = document.getElementById('cc-status');
  const actionsEl = document.getElementById('cc-actions');
  const logDetailsEl = document.getElementById('cc-log-details');
  const logTailEl = document.getElementById('cc-log-tail');
  if (!statusEl) return;
  let pollTimer = null;

  function escapeHtml(s) {
    const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML;
  }

  function renderStatus(s) {
    if (s.logged_in) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: var(--accent); margin: 0;">' +
        '<strong>Signed in.</strong> <a href="/">Go to library →</a></div>';
    } else if (!s.node_ok) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;">' +
        '<strong>Node.js ' + 18 + '+ required.</strong> ' +
        'Install with <code>brew install node</code> (macOS) / <code>winget install OpenJS.NodeJS</code> (Windows), then refresh this page.' +
        '</div>';
    } else if (s.install_running) {
      statusEl.innerHTML =
        '<p class="meta-info" style="margin: 0;">Installing Claude Code… this takes ~30s.</p>';
    } else if (s.login_running) {
      statusEl.innerHTML =
        '<p class="meta-info" style="margin: 0;">Waiting for Claude.ai OAuth to complete in the browser tab that just opened…</p>';
    } else if (s.install_error) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;"><strong>Install failed.</strong> ' +
        escapeHtml(s.install_error) + '</div>';
    } else if (s.login_error) {
      statusEl.innerHTML =
        '<div class="flash" style="border-left-color: #c00; margin: 0;"><strong>Login failed.</strong> ' +
        escapeHtml(s.login_error) + '</div>';
    } else if (!s.installed) {
      statusEl.innerHTML = '<p class="meta-info" style="margin: 0;">Step 1: install Claude Code into <code>~/yt2md/claude-code/</code>.</p>';
    } else {
      statusEl.innerHTML = '<p class="meta-info" style="margin: 0;">Step 2: sign in. A new browser tab will open for Claude.ai OAuth.</p>';
    }
  }

  function renderActions(s) {
    let html = '';
    const installDisabled = (!s.node_ok || s.install_running || s.installed) ? ' disabled' : '';
    html += '<form method="post" action="/setup/install-claude" style="display:inline;">' +
            '<button type="submit" class="primary"' + installDisabled + ' id="cc-install-btn">' +
            (s.install_running ? 'Installing…' : 'Install Claude Code') + '</button></form>';
    if (s.installed && !s.logged_in) {
      html += '<form method="post" action="/setup/login-claude" style="display:inline;">' +
              '<button type="submit" class="primary"' + (s.login_running ? ' disabled' : '') + ' id="cc-login-btn">' +
              (s.login_running ? 'Waiting for OAuth…' : 'Sign in with Claude.ai') + '</button></form>';
    }
    if (s.logged_in) {
      html += '<form method="post" action="/setup/logout-claude" style="display:inline;" ' +
              'onsubmit="return confirm(\\'Sign out of Claude Code?\\');">' +
              '<button type="submit">Sign out</button></form>';
    }
    actionsEl.innerHTML = html;
  }

  function renderLog(s) {
    const tail = (s.login_log_tail || s.install_log_tail || '').trim();
    if (!tail) {
      logDetailsEl.style.display = 'none';
      return;
    }
    logDetailsEl.style.display = '';
    logTailEl.textContent = tail;
  }

  async function poll() {
    try {
      const res = await fetch('/setup/claude-status');
      if (!res.ok) return;
      const s = await res.json();
      renderStatus(s);
      renderActions(s);
      renderLog(s);
      // Auto-redirect home once login lands.
      if (s.logged_in) {
        clearInterval(pollTimer);
        setTimeout(() => { window.location.href = '/?msg=Signed+in+via+Claude.ai+subscription.'; }, 1500);
        return;
      }
      // Stop polling if nothing is in flight (saves cycles).
      if (!s.install_running && !s.login_running) {
        clearInterval(pollTimer);
        pollTimer = null;
      }
    } catch (e) { /* network blip — try again next tick */ }
  }

  // Always do an immediate poll so transient state from a just-submitted
  // form is reflected without the 2s delay.
  poll();
  pollTimer = setInterval(poll, 2000);
})();
</script>
"""

        return page(body, title="Set up", current="setup")

    @app.route("/setup/save-api-key", methods=["POST"])
    def setup_save_api_key():
        from flask import redirect, request
        from urllib.parse import quote_plus
        new_key = (request.form.get("anthropic_api_key") or "").strip()
        if not new_key:
            return redirect("/setup?msg=Paste+a+key+first.")
        err = validate_api_key(new_key)
        if err:
            return redirect(f"/setup?msg=Key+rejected:+{quote_plus(err)}")
        set_env_var("ANTHROPIC_API_KEY", new_key)
        return redirect("/?msg=API+key+saved.+You+can+now+generate+digests.")

    # Legacy route name retained for back-compat with bookmarks/redirects.
    @app.route("/setup/save", methods=["POST"])
    def setup_save_legacy():
        return setup_save_api_key()

    @app.route("/setup/install-claude", methods=["POST"])
    def setup_install_claude():
        from flask import redirect
        err = start_install_job()
        if err:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg={quote_plus(err)}")
        return redirect("/setup")

    @app.route("/setup/login-claude", methods=["POST"])
    def setup_login_claude():
        from flask import redirect
        err = start_login_job()
        if err:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg={quote_plus(err)}")
        return redirect("/setup")

    @app.route("/setup/logout-claude", methods=["POST"])
    def setup_logout_claude():
        from flask import redirect
        rc, out = claude_logout()
        if rc != 0:
            from urllib.parse import quote_plus
            return redirect(f"/setup?msg=Logout+failed:+{quote_plus(out[:200])}")
        return redirect("/setup?msg=Signed+out+of+Claude+Code.")

    @app.route("/setup/claude-status")
    def setup_claude_status():
        from flask import jsonify
        return jsonify(claude_setup_snapshot())

    @app.route("/one-off", methods=["GET"])
    def one_off_page():
        from flask import request
        from html import escape as h
        flash = request.args.get("msg", "")
        active = _list_active_oneoff_jobs()  # also reaps exited subprocesses
        failures = _list_recent_oneoff_failures()

        body = "<h1>One-off digest</h1>"
        if flash:
            body += f'<div class="flash">{h(flash)}</div>'
        body += (
            '<p class="meta-info">Paste a YouTube video URL. The digest runs in the '
            'background and lands in your library when complete (1–25 min depending on '
            'video length and the vision pass). You can close this tab — it keeps running.</p>'
        )
        body += (
            '<form method="post" action="/one-off" class="add-form">'
            '<label for="video-url" class="sr-only">YouTube video URL</label>'
            '<input id="video-url" type="text" name="url" '
            'placeholder="https://youtu.be/... (or full watch URL)" '
            'autofocus required>'
            '<button type="submit">Digest</button>'
            '</form>'
        )

        import time as _t
        # Read the log once so per-job stage lookups don't hit disk repeatedly.
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""

        # Always render the section containers so the polling JS can target them
        # (display:none hides them when empty).
        active_hidden = "" if active else " style='display:none'"
        body += f"<section id='oneoff-active-section'{active_hidden}>"
        body += "<h2>In progress</h2>"
        body += "<ul id='oneoff-active-list' class='channel-list'>"
        for j in active:
            elapsed = int(_t.time() - j["started"])
            stage = _describe_job_stage(log_text, j["video_id"])
            body += (
                '<li>'
                f'<span class="url"><strong>{h(j["video_id"])}</strong> · '
                f'{h(j["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">'
                f'{elapsed//60}m {elapsed%60}s · {h(stage)}</span>'
                '</li>'
            )
        body += "</ul></section>"

        failures_hidden = "" if failures else " style='display:none'"
        body += f"<section id='oneoff-failures-section'{failures_hidden}>"
        body += "<h2>Recent failures</h2>"
        body += "<ul id='oneoff-failures-list' class='channel-list'>"
        for f in failures:
            ago = int(_t.time() - f["ended"])
            if ago < 60:
                when = f"{ago}s ago"
            elif ago < 3600:
                when = f"{ago // 60}m ago"
            else:
                when = f"{ago // 3600}h ago"
            err = f["error"] or f"exit code {f['exit_code']}"
            body += (
                '<li>'
                f'<span class="url"><strong>{h(f["video_id"])}</strong> · '
                f'{h(f["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">{when}</span>'
                f'<div style="color: var(--muted); font-size: 13px; margin-top: 4px;">{h(err)}</div>'
                '</li>'
            )
        body += "</ul>"
        body += (
            "<p class='meta-info'>"
            f"Full output: <code>{h(str(log_path))}</code>"
            "</p>"
        )
        body += "</section>"

        body += (
            "<p class='meta-info' style='margin-top: 32px;'>"
            "One-off digests share the same library as subscription mode. "
            "They appear in the sidebar's <strong>Digests</strong> section once ready."
            "</p>"
        )

        # Polling: refresh in-progress + failures every 2s without reloading the page.
        body += """
<script>
(function () {
  const activeSection  = document.getElementById('oneoff-active-section');
  const activeList     = document.getElementById('oneoff-active-list');
  const failSection    = document.getElementById('oneoff-failures-section');
  const failList       = document.getElementById('oneoff-failures-list');
  if (!activeSection || !failSection) return;

  const fmtElapsed = s => `${Math.floor(s/60)}m ${s%60}s`;
  const fmtAgo = s => s < 60 ? `${s}s ago` : (s < 3600 ? `${Math.floor(s/60)}m ago` : `${Math.floor(s/3600)}h ago`);
  const esc = s => { const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML; };

  function render(data) {
    const active = data.active || [];
    activeSection.style.display = active.length ? '' : 'none';
    activeList.innerHTML = active.map(j => `
      <li>
        <span class="url"><strong>${esc(j.video_id)}</strong> &middot; ${esc(j.url)}</span>
        <span style="color: var(--muted); font-size: 12px;">${esc(fmtElapsed(j.elapsed_secs))} &middot; ${esc(j.stage)}</span>
      </li>
    `).join('');

    const failures = data.failures || [];
    failSection.style.display = failures.length ? '' : 'none';
    failList.innerHTML = failures.map(f => `
      <li>
        <span class="url"><strong>${esc(f.video_id)}</strong> &middot; ${esc(f.url)}</span>
        <span style="color: var(--muted); font-size: 12px;">${esc(fmtAgo(f.ago_secs))}</span>
        <div style="color: var(--muted); font-size: 13px; margin-top: 4px;">${esc(f.error || ('exit code ' + f.exit_code))}</div>
      </li>
    `).join('');
  }

  async function poll() {
    try {
      const res = await fetch('/one-off/status', { cache: 'no-store' });
      if (res.ok) render(await res.json());
    } catch (_) { /* transient — try again next tick */ }
  }

  poll();
  setInterval(poll, 2000);
})();
</script>
"""
        return page(body, title="One-off digest", current="one-off")

    @app.route("/sidebar-status")
    def sidebar_status():
        """Lightweight poll target so the sidebar can detect when new digests
        land — auto-pipeline finish, scheduled poll producing a new digest,
        manual deletion, etc. Returns just count + max mtime; the page's JS
        compares to the values it was rendered with and shows a refresh
        banner if either has drifted.
        """
        from flask import jsonify
        digests = _list_digests(digests_dir)
        return jsonify({
            "digests_count": len(digests),
            "max_mtime": max((d["mtime"] for d in digests), default=0.0),
        })

    @app.route("/one-off/status")
    def one_off_status():
        from flask import jsonify
        import time as _t
        active = _list_active_oneoff_jobs()  # also reaps any just-exited subprocesses
        failures = _list_recent_oneoff_failures()
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""
        now = _t.time()
        return jsonify({
            "active": [
                {
                    "video_id": j["video_id"],
                    "url": j["url"],
                    "started": j["started"],
                    "elapsed_secs": int(now - j["started"]),
                    "stage": _describe_job_stage(log_text, j["video_id"]),
                }
                for j in active
            ],
            "failures": [
                {
                    "video_id": f["video_id"],
                    "url": f["url"],
                    "started": f["started"],
                    "ended": f["ended"],
                    "ago_secs": int(now - f["ended"]),
                    "exit_code": f["exit_code"],
                    "error": f["error"],
                }
                for f in failures
            ],
        })

    @app.route("/one-off", methods=["POST"])
    def one_off_submit():
        from flask import redirect, request
        import time as _t
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        url = request.form.get("url", "").strip()
        if not url:
            return redirect("/one-off?msg=URL+is+required")

        video_id = extract_video_id(url)
        if not video_id:
            return redirect(
                f"/one-off?msg=Couldn%27t+extract+a+YouTube+video+ID+from:+{url}"
            )

        # Already in library? Send them straight to the existing digest.
        existing = digests_dir / video_id / "digest.md"
        if existing.exists():
            return redirect(f"/digests/{video_id}/")

        # Already in progress? Show the page with a message rather than re-firing.
        for active in _list_active_oneoff_jobs():
            if active["video_id"] == video_id:
                return redirect(f"/one-off?msg=Already+digesting+{video_id}")

        # Fire and forget. start_new_session=True detaches the child so it survives
        # if the web server is killed. stdout/stderr go to oneoff.log.
        digest_path = digests_dir / video_id / "digest.md"
        digest_path.parent.mkdir(parents=True, exist_ok=True)
        log_path = data_dir / "logs" / "oneoff.log"
        log_path.parent.mkdir(parents=True, exist_ok=True)
        log_fd = open(log_path, "a")
        log_fd.write(f"\n===== {_t.strftime('%Y-%m-%d %H:%M:%S')} starting {video_id} ({url}) =====\n")
        log_fd.flush()

        yt2md_path = shutil.which("yt2md")
        if not yt2md_path:
            log_fd.close()
            return redirect("/one-off?msg=yt2md+not+on+PATH")

        try:
            proc = subprocess.Popen(
                [yt2md_path, url, "-o", str(digest_path),
                 "--source", "oneoff"],
                cwd=digest_path.parent,
                stdout=log_fd,
                stderr=subprocess.STDOUT,
                env={**os.environ, **_settings_to_env(load_settings())},
                **_DETACH_KWARGS,
            )
        finally:
            # Subprocess holds its own copy of the fd; safe for us to close.
            log_fd.close()

        _oneoff_jobs[proc.pid] = {
            "video_id": video_id,
            "started": _t.time(),
            "url": url,
            "proc": proc,
        }
        return redirect(
            f"/one-off?msg=Started+digesting+{video_id}+(check+sidebar+in+a+few+minutes)"
        )

    @app.route("/activity")
    def activity_page():
        from flask import request
        from html import escape as h
        import time as _t
        flt = request.args.get("status", "all")  # all | success | failed
        rows = _recent_runs(limit=200)
        if flt == "success":
            rows = [r for r in rows if r.get("success")]
        elif flt == "failed":
            rows = [r for r in rows if not r.get("success")]

        body = "<h1>Activity</h1>"
        body += '<p class="meta-info">Every completed one-off digest, success or failure. Persists across server restarts.</p>'

        # Budget gate status: this workspace's month-to-date spend vs the
        # warn/block thresholds. Defensive — never break Activity if the Admin
        # API is unreachable.
        try:
            _bs = budget_status()
            _mtd = _bs.get("month_to_date_usd")
            if _mtd is not None and (_bs.get("block_usd") or _bs.get("warn_usd")):
                _warn, _block = _bs.get("warn_usd") or 0, _bs.get("block_usd") or 0
                _color = "#c0392b" if (_block and _mtd >= _block) else (
                    "#e67e22" if (_warn and _mtd >= _warn) else "#27ae60")
                _src = "billing" if _bs.get("source") == "billing" else "local log (no admin key)"
                body += (
                    f'<p class="meta-info" style="border-left:3px solid {_color};'
                    f'padding-left:8px">Budget (this workspace, month-to-date): '
                    f'<strong style="color:{_color}">${_mtd:.2f}</strong> '
                    f'&middot; warn ${_warn:.0f} &middot; block ${_block:.0f} '
                    f'&middot; <span style="opacity:.7">source: {h(_src)}</span></p>')
        except Exception:
            pass

        # Cost summary across windows. Reads the central LLM usage log once
        # and bucket-totals — cheap, runs in a few ms even with thousands
        # of entries.
        usage_entries = read_llm_usage_log()
        now_ts = _t.time()
        WINDOWS = [
            ("today", 24 * 3600),
            ("7d", 7 * 24 * 3600),
            ("30d", 30 * 24 * 3600),
            ("all", None),
        ]
        window_totals: dict = {label: 0.0 for label, _ in WINDOWS}
        window_counts: dict = {label: 0 for label, _ in WINDOWS}
        backend_counts: dict = {}
        for e in usage_entries:
            cost = float(e.get("cost_usd", 0.0) or 0.0)
            ts = float(e.get("ts", 0) or 0)
            backend_counts[e.get("backend", "?")] = (
                backend_counts.get(e.get("backend", "?"), 0) + 1
            )
            for label, window in WINDOWS:
                if window is None or (now_ts - ts) <= window:
                    window_totals[label] += cost
                    window_counts[label] += 1
        backend_summary = ", ".join(
            f"{n} {b}" for b, n in sorted(backend_counts.items(), key=lambda kv: -kv[1])
        ) if backend_counts else ""
        body += (
            "<div class='schedule-form' style='display: grid; "
            "grid-template-columns: repeat(auto-fit, minmax(160px, 1fr)); "
            "gap: 12px 20px;'>"
        )
        for label, _ in WINDOWS:
            body += (
                f"<div><div style='color: var(--muted); font-size: 12px;'>"
                f"{h(label)}</div>"
                f"<div style='font-size: 22px; font-weight: 600;'>"
                f"${window_totals[label]:.2f}</div>"
                f"<div style='color: var(--muted); font-size: 11px;'>"
                f"{window_counts[label]} call(s)</div></div>"
            )
        body += "</div>"
        if backend_summary:
            body += (
                f"<p class='meta-info'>By backend: {h(backend_summary)}. "
                "Subscription (Claude Code) calls report $0 — billed via the "
                "user's plan. Pricing is estimated; treat as a guide.</p>"
            )

        # In-progress section — live one-off jobs the reaper hasn't logged yet.
        # Hidden by default; populated by the polling JS at the bottom of the page.
        active_now = _list_active_oneoff_jobs()
        log_path = data_dir / "logs" / "oneoff.log"
        try:
            log_text = log_path.read_text(errors="replace")
        except OSError:
            log_text = ""
        active_hidden = "" if active_now else " style='display:none'"
        body += f"<section id='activity-active-section'{active_hidden}>"
        body += "<h2>In progress</h2>"
        body += "<ul id='activity-active-list' class='channel-list'>"
        for j in active_now:
            elapsed = int(_t.time() - j["started"])
            stage = _describe_job_stage(log_text, j["video_id"])
            body += (
                '<li>'
                f'<span class="url"><strong>{h(j["video_id"])}</strong> · '
                f'{h(j["url"])}</span>'
                f'<span style="color: var(--muted); font-size: 12px;">'
                f'{elapsed//60}m {elapsed%60}s · {h(stage)}</span>'
                '</li>'
            )
        body += "</ul></section>"

        # Polling JS for the in-progress section — emitted here so both the
        # "no runs yet" early return and the full table path include it.
        active_poll_js = """
<script>
(function () {
  const section = document.getElementById('activity-active-section');
  const list = document.getElementById('activity-active-list');
  if (!section || !list) return;
  let prevCount = list.children.length;
  const fmtElapsed = s => `${Math.floor(s/60)}m ${s%60}s`;
  const esc = s => { const d = document.createElement('div'); d.textContent = String(s ?? ''); return d.innerHTML; };
  async function poll() {
    try {
      const res = await fetch('/one-off/status', { cache: 'no-store' });
      if (!res.ok) return;
      const data = await res.json();
      const active = data.active || [];
      section.style.display = active.length ? '' : 'none';
      list.innerHTML = active.map(j => `
        <li>
          <span class="url"><strong>${esc(j.video_id)}</strong> &middot; ${esc(j.url)}</span>
          <span style="color: var(--muted); font-size: 12px;">${esc(fmtElapsed(j.elapsed_secs))} &middot; ${esc(j.stage)}</span>
        </li>
      `).join('');
      if (active.length < prevCount) {
        window.location.reload();
        return;
      }
      prevCount = active.length;
    } catch (_) { /* try again */ }
  }
  poll();
  setInterval(poll, 2000);
})();
</script>
"""

        # Filter chips
        chip = lambda v, label, count: (
            f'<a href="/activity?status={v}" class="filter-chip'
            + (' active' if flt == v else '')
            + f'">{h(label)} <span class="filter-chip-count">({count})</span></a>'
        )
        all_runs = _recent_runs(limit=200)
        n_all = len(all_runs)
        n_ok = sum(1 for r in all_runs if r.get("success"))
        n_fail = n_all - n_ok
        body += "<div class='filter-row'>"
        body += chip("all", "All", n_all)
        body += chip("success", "Success", n_ok)
        body += chip("failed", "Failed", n_fail)
        body += "</div>"

        if not rows:
            body += "<p class='meta-info'>No runs recorded yet. Submit a one-off digest from the <a href='/one-off'>One-off digest</a> page.</p>"
            body += active_poll_js
            return page(body, title="Activity", current="activity")

        def _fmt_ago(secs: float) -> str:
            secs = int(secs)
            if secs < 60: return f"{secs}s ago"
            if secs < 3600: return f"{secs//60}m ago"
            if secs < 86400: return f"{secs//3600}h ago"
            return f"{secs//86400}d ago"

        def _fmt_dur(secs) -> str:
            if secs is None: return "—"
            secs = float(secs)
            if secs < 60: return f"{secs:.1f}s"
            return f"{int(secs)//60}m {int(secs)%60}s"

        def _fmt_int(n) -> str:
            if n is None: return "—"
            n = int(n)
            return f"{n:,}"

        now = _t.time()
        # Index usage by video_id for O(1) per-row cost lookup. Each row's
        # cost = sum of entries whose ts is within the row's [started_at,
        # ended_at] window AND video_id matches.
        usage_by_vid: dict = {}
        for e in usage_entries:
            usage_by_vid.setdefault(e.get("video_id", ""), []).append(e)

        body += "<table class='activity-table'>"
        body += (
            "<thead><tr>"
            "<th>When</th><th>Video</th><th>Outcome</th>"
            "<th>Duration</th><th>Stages</th><th>Tokens</th><th>Cost</th>"
            "</tr></thead><tbody>"
        )
        for r in rows:
            video_id = r.get("video_id") or ""
            url = r.get("url") or ""
            ago = _fmt_ago(now - (r.get("started_at") or now))
            dur = _fmt_dur(r.get("duration_secs"))
            success = r.get("success")
            if success:
                outcome = "<span class='ok'>✓ done</span>"
                if r.get("digest_path"):
                    title = f'<a href="/digests/{h(video_id)}/">{h(video_id)}</a>'
                else:
                    title = h(video_id)
            else:
                stage = r.get("stage_reached") or "?"
                outcome = f"<span class='fail'>✗ failed at {h(stage)}</span>"
                title = h(video_id)
            # Stage breakdown
            parts = []
            for label, key in [
                ("dl", "download_secs"),
                ("whisper", "whisper_secs"),
                ("frames", "frames_secs"),
                ("digest", "digest_secs"),
                ("vision", "vision_secs"),
            ]:
                v = r.get(key)
                if v is not None and v > 0.05:
                    parts.append(f"{label} {v:.1f}s")
            stages_cell = h(", ".join(parts)) if parts else "—"
            # Tokens (digest only)
            tin = r.get("digest_input_tokens")
            tout = r.get("digest_output_tokens")
            cache = r.get("digest_cache_read_tokens") or 0
            if tin or tout:
                tokens_cell = f"in {_fmt_int(tin)} · out {_fmt_int(tout)}"
                if cache:
                    tokens_cell += f" · cache {_fmt_int(cache)}"
                tokens_cell = h(tokens_cell)
            else:
                tokens_cell = "—"

            # Per-row cost: every LLM call recorded for this video_id whose
            # ts falls in the run's [started_at, ended_at] window.
            row_cost = 0.0
            row_backend = None
            started_at = r.get("started_at") or 0
            ended_at = r.get("ended_at") or now
            for e in usage_by_vid.get(video_id, []):
                ets = float(e.get("ts", 0) or 0)
                if started_at <= ets <= ended_at:
                    row_cost += float(e.get("cost_usd", 0.0) or 0.0)
                    row_backend = e.get("backend") or row_backend
            if row_cost > 0:
                cost_cell = f"${row_cost:.4f}"
            elif row_backend == "claude-code":
                cost_cell = "<span style='color: var(--muted);'>subscription</span>"
            else:
                cost_cell = "—"

            body += "<tr>"
            body += f"<td title='{h(url)}'>{h(ago)}</td>"
            body += f"<td>{title}"
            extras = []
            if r.get("source_lang"): extras.append(f"lang: {h(r['source_lang'])}")
            if r.get("used_whisper"):
                wm = r.get("whisper_model") or "?"
                extras.append(f"whisper: {h(wm)}")
            if extras:
                body += f"<div class='activity-meta'>{' · '.join(extras)}</div>"
            body += "</td>"
            body += f"<td>{outcome}"
            err = r.get("error")
            if err and not success:
                body += f"<div class='activity-meta activity-error'>{h(err)}</div>"
            body += "</td>"
            body += f"<td>{h(dur)}</td>"
            body += f"<td class='activity-stages'>{stages_cell}</td>"
            body += f"<td class='activity-tokens'>{tokens_cell}</td>"
            body += f"<td class='activity-cost'>{cost_cell}</td>"
            body += "</tr>"
        body += "</tbody></table>"
        body += (
            "<p class='meta-info' style='margin-top: 24px;'>"
            f"Raw log: <code>{h(str(_oneoff_log_path()))}</code><br>"
            f"JSONL: <code>{h(str(_runs_jsonl_path()))}</code>"
            "</p>"
        )
        body += active_poll_js
        return page(body, title="Activity", current="activity")

    @app.route("/digests/<video_id>/")
    def view_digest(video_id):
        from html import escape as h
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)
        try:
            _mark_digest_read(video_id)
        except Exception:
            pass  # never block reading on a DB error
        md_source = digest_md.read_text()
        rendered = _render_markdown(md_source)

        nav = _viewer_nav(video_id, "digest", digests_dir)
        any_running = _any_local_job_running(video_id)

        # Bottom-of-page actions: Copy markdown is the canonical "take this
        # elsewhere" affordance for any reader who scrolled to the end.
        # Delete digest sits below it in a small "danger zone" — only on
        # the digest viewer, not on takeaway/panel (must be on the main
        # anchor page).
        bottom_actions = (
            "<hr style='margin: 32px 0 16px; border: none; "
            "border-top: 1px solid var(--border);'>"
            "<div class='digest-actions' style='margin-bottom: 16px;'>"
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='page-md-source' "
            "title='Copy the markdown source — paste into your notes app, "
            "email, or another LLM'>Copy markdown</button>"
            # Refresh metadata: re-probe yt-dlp for the current title,
            # thumbnail, channel info. Useful when a YouTube creator
            # renames the video after upload (A/B testing titles is
            # common for data-driven creators).
            f"<form method='post' action='/digests/{h(video_id)}/refresh-metadata' "
            "style='display:inline;' "
            "onsubmit=\"this.querySelector('button').disabled=true;"
            "this.querySelector('button').textContent='Refreshing…';\">"
            "<button type='submit' class='discuss-btn-secondary' "
            "title='Re-probe YouTube and update title, thumbnail, channel "
            "info. Useful when a creator renames their video after upload.'>"
            "Refresh metadata</button></form>"
            "</div>"
            f"<textarea id='page-md-source' hidden aria-hidden='true'>"
            f"{h(md_source)}</textarea>"
            "<hr style='margin: 24px 0 12px; border: none; "
            "border-top: 1px solid var(--border);'>"
            "<div style='display:flex; justify-content: flex-end;'>"
            f"<form method='post' action='/digests/{h(video_id)}/delete' "
            "style='display:inline;' "
            "onsubmit=\"return confirm('Delete this digest? "
            "The rendered output, frames, and cached video will be wiped.');\">"
            "<button type='submit' class='delete-btn' "
            "title='Wipes digest.md, panel.md, takeaway.md, slides.pptx, "
            "frames, and the cached video. The video can be re-digested "
            "via One-off later.'>Delete digest</button></form></div>"
        )

        audio = _audio_section(video_id, "digest", digests_dir)
        poll_js = _JOB_POLL_JS if any_running else ""
        body = (
            nav
            + _COPY_BUTTON_JS
            + poll_js
            + "<hr style='margin: 16px 0 32px; border: none; "
            "border-top: 1px solid var(--border);'>"
            + audio
            + rendered
            + bottom_actions
        )
        return page(body, title=video_id, current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/")

    def _build_panel_for_digest(video_id: str) -> None:
        """Background worker for panel generation. Delegates to the
        module-level builder so the agent API can call the same path."""
        build_panel_for_video(video_id, digests_dir=digests_dir)

    def _build_takeaway_for_digest(video_id: str) -> None:
        """Background worker for takeaway generation. Delegates."""
        build_takeaway_for_video(video_id, digests_dir=digests_dir)

    @app.route("/digests/<video_id>/discuss", methods=["POST"])
    def generate_panel_route(video_id):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        if not (digests_dir / video_id / "digest.md").exists():
            abort(404)
        if (digests_dir / video_id / "panel.md").exists():
            return redirect(f"/digests/{video_id}/?msg=Panel+already+exists.")
        # Spawn in the background; the digest viewer's running-state
        # placeholder polls /job-status?kind=panel and reloads when done.
        start_local_job(f"{video_id}:panel", _build_panel_for_digest, video_id)
        return redirect(f"/digests/{video_id}/")

    @app.route("/digests/<video_id>/takeaway", methods=["POST"])
    def generate_takeaway_route(video_id):
        from flask import redirect
        gate = _require_llm_or_redirect()
        if gate is not None:
            return gate
        if not (digests_dir / video_id / "digest.md").exists():
            abort(404)
        if (digests_dir / video_id / "takeaway.md").exists():
            return redirect(f"/digests/{video_id}/?msg=Takeaway+already+exists.")
        start_local_job(
            f"{video_id}:takeaway", _build_takeaway_for_digest, video_id,
        )
        return redirect(f"/digests/{video_id}/")

    @app.route("/digests/<video_id>/takeaway/")
    def view_takeaway(video_id):
        from html import escape as h
        takeaway_md = digests_dir / video_id / "takeaway.md"
        if not takeaway_md.exists():
            abort(404)
        md_source = takeaway_md.read_text()
        rendered = _render_markdown(md_source)
        # Pre-assemble the chat handoff prompt server-side so the button can
        # just copy a hidden textarea (same pattern as Copy markdown). Loads
        # whichever artifacts exist (digest + panel + takeaway).
        nav = _viewer_nav(video_id, "takeaway", digests_dir)
        any_running = _any_local_job_running(video_id)
        # Bottom-of-page actions: Copy markdown for the takeaway itself,
        # plus Continue in chat (only on takeaway — that's the natural
        # "I want to ask a follow-up" endpoint after reading the synthesis).
        # Continue-in-chat now navigates to /handoff so the user can see
        # what's about to be sent, type their question, and pick the
        # destination (claude.ai paste flow or Claude Desktop / MCP).
        bottom_actions = (
            "<hr style='margin: 32px 0 16px; border: none; "
            "border-top: 1px solid var(--border);'>"
            "<div class='digest-actions' style='margin-bottom: 16px;'>"
            f"<a href='/digests/{h(video_id)}/handoff' "
            "class='discuss-btn' style='text-decoration:none;' "
            "title='Open the handoff page — review the bundle, type your "
            "question, then copy + open Claude.'>Continue in chat ▸</a>"
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='page-md-source' "
            "title='Copy the takeaway markdown'>Copy markdown</button>"
            "</div>"
            f"<textarea id='page-md-source' hidden aria-hidden='true'>"
            f"{h(md_source)}</textarea>"
        )
        audio = _audio_section(video_id, "takeaway", digests_dir)
        poll_js = _JOB_POLL_JS if any_running else ""
        body = (
            nav + _COPY_BUTTON_JS + poll_js
            + "<hr style='margin: 16px 0 32px; border: none; "
            "border-top: 1px solid var(--border);'>"
            + audio + rendered + bottom_actions
        )
        return page(body, title=f"Takeaway · {video_id}",
                    current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/takeaway/")

    @app.route("/digests/<video_id>/handoff")
    def handoff_page(video_id):
        """Visible 'Continue in chat' flow. Replaces the silent
        copy-then-open-tab dance with a page that shows what's about to
        be sent, lets the user type their question, and only THEN copies
        and opens claude.ai. Avoids the failure mode where the user
        clicks the button, lands on an empty chat, and doesn't realize
        the bundle is sitting in their clipboard waiting to be pasted.
        """
        from html import escape as h
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)

        # Title for the page heading — prefer metadata.json's video
        # title, fall back to the digest.md H1, then video_id.
        title = video_id
        meta_path = digests_dir / video_id / "metadata.json"
        if meta_path.exists():
            try:
                meta = json.loads(meta_path.read_text())
                if meta.get("title"):
                    title = meta["title"]
            except Exception:
                pass
        if title == video_id:
            try:
                for line in digest_md.read_text().splitlines():
                    if line.startswith("# "):
                        title = line[2:].strip()
                        break
            except Exception:
                pass

        # Bundle the artifacts. build_chat_handoff_prompt already ends
        # with "My question: " — the user types into a separate textarea
        # and the JS appends it to that trailer at copy time.
        bundle = build_chat_handoff_prompt(video_id, digests_dir)
        # Strip the trailing "My question: " so the JS can append
        # exactly "My question: <typed>" without doubling the label.
        bundle_no_q = bundle.rstrip()
        if bundle_no_q.endswith("My question:"):
            bundle_no_q = bundle_no_q[: -len("My question:")].rstrip()

        artifact_counts = []
        if (digests_dir / video_id / "digest.md").exists():
            artifact_counts.append("digest")
        if (digests_dir / video_id / "panel.md").exists():
            artifact_counts.append("panel")
        if (digests_dir / video_id / "takeaway.md").exists():
            artifact_counts.append("takeaway")
        artifact_str = " + ".join(artifact_counts) or "summary"

        bundle_size_kb = len(bundle.encode("utf-8")) / 1024
        nav = _viewer_nav(video_id, "takeaway", digests_dir)

        body = (
            nav
            + "<hr style='margin: 16px 0 32px; border: none; "
            "border-top: 1px solid var(--border);'>"
            f"<h1 style='margin-top:0;'>Continue in chat</h1>"
            f"<p class='meta-info' style='margin-bottom:24px;'>"
            f"Send the {artifact_str} for "
            f"<strong>{h(title)}</strong> "
            f"({bundle_size_kb:.1f} KB) to Claude with a follow-up "
            f"question. The bundle goes on your clipboard; Claude opens "
            f"in a new tab; you paste with ⌘V.</p>"

            "<label style='display:block; font-weight:600; "
            "margin-bottom:6px;'>Your follow-up question</label>"
            "<textarea id='user-question' rows='3' "
            "placeholder='e.g. How does this compare to the "
            "subscription-pricing model SAP just announced?' "
            "style='width:100%; padding:12px 14px; border-radius:4px; "
            "border:1px solid var(--border); font-family:inherit; "
            "font-size:15px; line-height:1.5; "
            "background:var(--bg);' autofocus></textarea>"

            "<div style='margin-top:14px; display:flex; gap:10px; "
            "flex-wrap:wrap; align-items:center;'>"
            "<button id='handoff-go' type='button' class='discuss-btn' "
            "title='Copy the bundle (with your question) to your "
            "clipboard and open claude.ai in a new tab.'>"
            "📋 Copy + open Claude</button>"
            "<a href='https://claude.ai/new' target='_blank' "
            "rel='noopener' class='discuss-btn-secondary' "
            "style='text-decoration:none;'>Open Claude only</a>"
            f"<a href='/digests/{h(video_id)}/takeaway/' "
            "class='discuss-btn-secondary' "
            "style='text-decoration:none;'>← Back</a>"
            "<span id='handoff-status' style='margin-left:8px; "
            "color:var(--muted); font-size:14px;'></span>"
            "</div>"

            # Hidden bundle — what gets concatenated with the user's
            # question at copy time. Stored as a textarea so we don't
            # have to worry about character escaping in attributes.
            f"<textarea id='handoff-bundle' hidden aria-hidden='true'>"
            f"{h(bundle_no_q)}</textarea>"

            # ---- Path B: Claude Desktop + MCP server ----
            # Same question, different destination. If the user has the
            # yt2md MCP server wired into Claude Desktop, this path is
            # strictly better — Claude pulls structured data via tools
            # instead of staring at a wall of pasted markdown, and can
            # navigate to specific topics / panel turns on demand.
            "<hr style='margin:32px 0 24px; border:none; "
            "border-top:1px dashed var(--border);'>"
            "<h2 style='margin-top:0; font-size:18px;'>"
            "Or: discuss in Claude Desktop "
            "<span style='font-size:12px; color:var(--muted); "
            "font-weight:normal;'>(if you have the yt2md MCP server "
            "set up)</span></h2>"

            "<p class='meta-info' style='margin-bottom:14px;'>"
            "Better than pasting a 60 KB blob — Claude Desktop pulls "
            "the digest through the <code>read_digest</code> MCP tool "
            "and can drill into specific topics / panel turns on demand. "
            "Setup is one config edit; see the "
            "<a href='https://github.com/jyouturner/youtube-to-markdown"
            "#talk-to-your-library-from-claude-mcp' target='_blank' "
            "rel='noopener'>MCP section in the README</a>.</p>"

            "<div style='display:flex; gap:10px; flex-wrap:wrap; "
            "align-items:center;'>"
            "<button id='handoff-mcp-go' type='button' "
            "class='discuss-btn-secondary' "
            "title='Copy a Claude Desktop prompt that calls the yt2md "
            "MCP read_digest tool for this video.'>"
            "📋 Copy Claude Desktop prompt</button>"
            "<span id='handoff-mcp-status' style='margin-left:8px; "
            "color:var(--muted); font-size:14px;'></span>"
            "</div>"

            # Hidden video ID + title for the MCP prompt builder.
            f"<input type='hidden' id='handoff-vid' value='{h(video_id)}'>"
            f"<input type='hidden' id='handoff-title' value='{h(title)}'>"

            # Collapsible preview so the user can see exactly what's
            # being sent (transparency without it dominating the page).
            "<details style='margin-top:32px;'>"
            "<summary style='cursor:pointer; color:var(--muted); "
            "font-size:14px;'>Preview what gets sent "
            f"({bundle_size_kb:.1f} KB)</summary>"
            "<pre style='margin-top:12px; padding:14px; "
            "background:var(--code-bg); border-radius:4px; "
            "max-height:400px; overflow:auto; font-size:12px; "
            "white-space:pre-wrap; word-break:break-word;'>"
            f"{h(bundle_no_q[:6000])}"
            + ("\n\n... (truncated in preview; full bundle is copied)"
               if len(bundle_no_q) > 6000 else "")
            + "</pre>"
            "</details>"

            # Inline JS — same-page, no globals. Builds the final
            # payload (bundle + "---" + "My question: " + typed), copies
            # via Clipboard API, opens new tab, shows clear "now paste"
            # confirmation. Falls back to a textarea-select trick for
            # browsers / contexts where the modern clipboard API is
            # restricted.
            "<script>"
            "(function(){"
            "const btn = document.getElementById('handoff-go');"
            "const q   = document.getElementById('user-question');"
            "const b   = document.getElementById('handoff-bundle');"
            "const st  = document.getElementById('handoff-status');"
            "function build() {"
            "  const ques = (q.value || '').trim();"
            "  const trail = ques ? ('\\n\\nMy question: ' + ques)"
            "                     : '\\n\\nMy question: ';"
            "  return b.value + trail;"
            "}"
            "async function copyAndOpen() {"
            "  const text = build();"
            "  let ok = false;"
            "  try {"
            "    await navigator.clipboard.writeText(text);"
            "    ok = true;"
            "  } catch(e) {"
            "    const ta = document.createElement('textarea');"
            "    ta.value = text; ta.style.position='fixed'; ta.style.left='-9999px';"
            "    document.body.appendChild(ta); ta.select();"
            "    try { ok = document.execCommand('copy'); } catch(_){}"
            "    document.body.removeChild(ta);"
            "  }"
            "  if (!ok) { st.textContent = '✗ Copy failed — try Open Claude only and paste manually.'; return; }"
            "  st.innerHTML = '✓ Copied. Claude opens in a new tab — paste with <kbd>⌘V</kbd> (or <kbd>Ctrl+V</kbd>).';"
            "  st.style.color = '#0a7d2c';"
            "  st.style.fontWeight = '600';"
            "  setTimeout(() => window.open('https://claude.ai/new', '_blank', 'noopener'), 250);"
            "}"
            "btn.addEventListener('click', copyAndOpen);"
            "q.addEventListener('keydown', (e) => {"
            "  if ((e.metaKey || e.ctrlKey) && e.key === 'Enter') copyAndOpen();"
            "});"

            # ---- MCP path button ----
            "const mcpBtn = document.getElementById('handoff-mcp-go');"
            "const mcpSt  = document.getElementById('handoff-mcp-status');"
            "const vid    = document.getElementById('handoff-vid').value;"
            "const ttl    = document.getElementById('handoff-title').value;"
            "function buildMcpPrompt() {"
            "  const ques = (q.value || '').trim() || '[type your question above first]';"
            "  return ('Using the yt2md MCP server, call read_digest with '"
            "    + 'digest_id=\"' + vid + '\" and section=\"full\" to get the '"
            "    + 'context for \"' + ttl + '\". Also pull '"
            "    + 'section=\"panel\" if the panel critique is relevant. '"
            "    + 'Then let\\u2019s discuss:\\n\\n' + ques + '\\n\\n'"
            "    + 'Reference specific topics (section=\"topic:N\") or panel '"
            "    + 'turns (section=\"panel:turn:N\") as needed.');"
            "}"
            "async function copyMcp() {"
            "  const text = buildMcpPrompt();"
            "  let ok = false;"
            "  try { await navigator.clipboard.writeText(text); ok = true; }"
            "  catch(e) {"
            "    const ta = document.createElement('textarea');"
            "    ta.value = text; ta.style.position='fixed'; ta.style.left='-9999px';"
            "    document.body.appendChild(ta); ta.select();"
            "    try { ok = document.execCommand('copy'); } catch(_){}"
            "    document.body.removeChild(ta);"
            "  }"
            "  if (!ok) { mcpSt.textContent = '✗ Copy failed.'; return; }"
            "  mcpSt.innerHTML = '✓ Copied. Paste into Claude Desktop chat.';"
            "  mcpSt.style.color = '#0a7d2c';"
            "  mcpSt.style.fontWeight = '600';"
            "}"
            "mcpBtn.addEventListener('click', copyMcp);"
            "})();"
            "</script>"
        )
        return page(body, title=f"Continue in chat · {h(title)}",
                    current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/")

    @app.route("/digests/<video_id>/panel/")
    def view_panel(video_id):
        from html import escape as h
        panel_md = digests_dir / video_id / "panel.md"
        if not panel_md.exists():
            abort(404)
        md_source = panel_md.read_text()
        rendered = _render_markdown(md_source)
        nav = _viewer_nav(video_id, "panel", digests_dir)
        any_running = _any_local_job_running(video_id)
        bottom_actions = (
            "<hr style='margin: 32px 0 16px; border: none; "
            "border-top: 1px solid var(--border);'>"
            "<div class='digest-actions' style='margin-bottom: 16px;'>"
            "<button type='button' class='discuss-btn-secondary' "
            "data-copy-target='page-md-source' "
            "title='Copy the panel discussion markdown'>Copy markdown</button>"
            "</div>"
            f"<textarea id='page-md-source' hidden aria-hidden='true'>"
            f"{h(md_source)}</textarea>"
        )
        audio = _audio_section(video_id, "panel", digests_dir)
        poll_js = _JOB_POLL_JS if any_running else ""
        body = (
            nav + _COPY_BUTTON_JS + poll_js
            + "<hr style='margin: 16px 0 32px; border: none; "
            "border-top: 1px solid var(--border);'>"
            + audio + rendered + bottom_actions
        )
        return page(body, title=f"Panel · {video_id}",
                    current=f"digest:{video_id}",
                    base_href=f"/digests/{video_id}/panel/")

    @app.route("/digests/<video_id>/refresh-metadata", methods=["POST"])
    def refresh_metadata_route(video_id):
        """Re-probe yt-dlp for current title, thumbnail, channel info — and
        update the digest's H1 + metadata.json + thumbnail files in place.
        Useful when a YouTube creator renames their video after upload
        (Nate B Jones, for instance, A/B tests titles regularly). The
        digest content itself is not regenerated; only the metadata."""
        from flask import redirect
        from urllib.parse import quote_plus
        digest_md = digests_dir / video_id / "digest.md"
        if not digest_md.exists():
            abort(404)

        # Probe via yt-dlp with the user's configured cookies (paywalled
        # videos won't resolve otherwise).
        try:
            import yt_dlp
        except ImportError:
            return redirect(f"/digests/{video_id}/?msg=yt-dlp+not+available")
        s = load_settings()
        cookies = s.get("cookies_from_browser") or os.environ.get(
            "YT2MD_COOKIES_FROM_BROWSER")
        # Permissive format selector — yt-dlp validates the requested
        # format even with skip_download=True, and the default selector
        # ("bestvideo*+bestaudio") errors on videos with non-standard
        # streams. Same string fetch_youtube uses on its probe pass.
        ydl_opts: dict = {
            "quiet": True, "no_warnings": True, "skip_download": True,
            "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best",
            "ignore_no_formats_error": True,
        }
        if cookies:
            ydl_opts["cookiesfrombrowser"] = (cookies,)
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(
                    f"https://www.youtube.com/watch?v={video_id}",
                    download=False,
                )
        except Exception as e:
            return redirect(
                f"/digests/{video_id}/?msg=Refresh+failed:+"
                f"{quote_plus(str(e)[:200])}"
            )

        new_title = info.get("title") or video_id
        new_url = info.get("webpage_url") or f"https://www.youtube.com/watch?v={video_id}"
        new_upload = info.get("upload_date")
        new_thumb_url = info.get("thumbnail")
        new_channel_id = info.get("channel_id") or info.get("uploader_id") or ""
        new_channel_name = info.get("channel") or info.get("uploader") or ""
        new_channel_url = (info.get("channel_url")
                           or info.get("uploader_url") or "")

        # Update the H1 heading in digest.md in place (preserves the rest
        # of the file — Watch on YouTube link, topics, slide images).
        text = digest_md.read_text()
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if line.startswith("# "):
                lines[i] = f"# {new_title}"
                break
        digest_md.write_text("\n".join(lines))

        # Re-download the video thumbnail if the URL changed (or to
        # backfill a digest that was generated before thumbnails landed).
        has_thumb = False
        if new_thumb_url:
            has_thumb = download_image(
                new_thumb_url, digests_dir / video_id / "thumbnail.jpg",
            )

        # Channel avatar: download if missing for this channel.
        has_channel_thumb = False
        if new_channel_id:
            ch_path = (get_data_dir() / "channel_thumbnails"
                       / f"{new_channel_id}.jpg")
            if ch_path.exists():
                has_channel_thumb = True
            elif new_channel_url:
                ch_url = probe_channel_thumbnail_url(
                    new_channel_url, cookies_from_browser=cookies,
                )
                if ch_url:
                    has_channel_thumb = download_image(ch_url, ch_path)

        # Update metadata.json — merge with existing fields if any.
        meta_path = digests_dir / video_id / "metadata.json"
        metadata = {}
        if meta_path.exists():
            try:
                metadata = json.loads(meta_path.read_text())
            except Exception:
                pass
        metadata.update({
            "video_id": video_id,
            "title": new_title,
            "url": new_url,
            "upload_date": new_upload,
            "channel_id": new_channel_id or None,
            "channel_name": new_channel_name or None,
            "channel_url": new_channel_url or None,
            "has_thumbnail": has_thumb,
            "has_channel_thumbnail": has_channel_thumb,
        })
        meta_path.write_text(json.dumps(metadata, indent=2) + "\n")

        return redirect(
            f"/digests/{video_id}/?msg=Metadata+refreshed+(title:+"
            f"{quote_plus(new_title[:60])})"
        )

    @app.route("/digests/<video_id>/delete", methods=["POST"])
    def delete_digest(video_id):
        from flask import redirect
        target = digests_dir / video_id
        if not target.exists():
            abort(404)
        # Wipe artifacts + cached source media. shutil handles missing children.
        shutil.rmtree(target, ignore_errors=True)
        # Clear read-state row.
        try:
            with _library_connect() as conn:
                conn.execute("DELETE FROM digest_reads WHERE digest_id = ?", (video_id,))
        except Exception:
            pass
        return redirect("/?msg=Deleted+" + video_id)

    @app.route("/digests/<video_id>/digest_images/<path:filename>")
    def digest_image(video_id, filename):
        return send_from_directory(digests_dir / video_id / "digest_images", filename)

    @app.route("/digests/<video_id>/thumbnail.jpg")
    def video_thumbnail(video_id):
        thumb = digests_dir / video_id / "thumbnail.jpg"
        if not thumb.exists():
            abort(404)
        return send_from_directory(digests_dir / video_id, "thumbnail.jpg")

    @app.route("/channel-thumbnails/<channel_id>.jpg")
    def channel_thumbnail(channel_id):
        # Channel avatars are shared across all digests from the same channel.
        # Disallow any path traversal in channel_id (yt-dlp channel IDs are
        # always alphanumeric + dashes/underscores).
        if not re.match(r"^[\w\-]+$", channel_id):
            abort(404)
        ch_dir = get_data_dir() / "channel_thumbnails"
        target = ch_dir / f"{channel_id}.jpg"
        if not target.exists():
            abort(404)
        return send_from_directory(ch_dir, f"{channel_id}.jpg")

    @app.route("/digests/<video_id>/slides.pptx")
    def download_slides(video_id):
        slides_path = digests_dir / video_id / "slides.pptx"
        if not slides_path.exists():
            abort(404)
        return send_from_directory(
            digests_dir / video_id, "slides.pptx",
            as_attachment=True, download_name=f"{video_id}.pptx",
        )

    def _build_slides_for_digest(video_id: str) -> None:
        """Background worker for slides generation. Delegates to the
        module-level builder so the agent API shares this code path."""
        build_slides_for_video(video_id, digests_dir=digests_dir)

    @app.route("/digests/<video_id>/slides", methods=["POST"])
    def generate_slides_route(video_id):
        from flask import redirect
        # Slides already on disk — no-op, just bounce back.
        if (digests_dir / video_id / "slides.pptx").exists():
            return redirect(f"/digests/{video_id}/?msg=Slides+already+exist.")
        # Spawn the job and return immediately. The UI polls
        # /job-status?kind=slides for progress and reloads when done.
        start_local_job(f"{video_id}:slides", _build_slides_for_digest, video_id)
        return redirect(f"/digests/{video_id}/")

    def _build_audio_for_artifact(video_id: str, kind: str) -> None:
        """Background worker — delegates to module-level builder so the
        agent API can share this code path."""
        build_audio_for_artifact(video_id, kind, digests_dir=digests_dir)

    @app.route("/digests/<video_id>/audio/<kind>", methods=["POST"])
    def generate_audio_route(video_id, kind):
        from flask import redirect
        if kind not in AUDIO_SOURCE_BY_KIND:
            abort(404)
        mp3_path = digests_dir / video_id / f"{kind}.mp3"
        if mp3_path.exists():
            return redirect(f"/digests/{video_id}/?msg=Audio+already+exists.")
        src_name = AUDIO_SOURCE_BY_KIND[kind]
        if not (digests_dir / video_id / src_name).exists():
            from urllib.parse import quote_plus
            return redirect(
                f"/digests/{video_id}/?msg="
                f"{quote_plus(f'{src_name} missing — generate it first.')}"
            )
        start_local_job(
            f"{video_id}:audio_{kind}",
            _build_audio_for_artifact, video_id, kind,
        )
        # Bounce back to the kind-specific viewer the user clicked from
        # (digest / panel / takeaway). Going to the wrong tab would
        # leave them staring at a different page with no visible sign
        # that anything is happening.
        viewer_path = {
            "digest":   f"/digests/{video_id}/",
            "panel":    f"/digests/{video_id}/panel/",
            "takeaway": f"/digests/{video_id}/takeaway/",
        }[kind]
        return redirect(viewer_path)

    @app.route("/digests/<video_id>/audio/<kind>.mp3")
    def serve_audio(video_id, kind):
        if kind not in AUDIO_SOURCE_BY_KIND:
            abort(404)
        mp3_path = digests_dir / video_id / f"{kind}.mp3"
        if not mp3_path.exists():
            abort(404)
        return send_from_directory(
            digests_dir / video_id, f"{kind}.mp3",
            mimetype="audio/mpeg",
        )

    @app.route("/listen")
    def listen_page():
        """Friction-free "subscribe in Podcasts" page. Computes the right
        URL based on bind address + your hostname, generates a QR you
        scan with the iPhone camera (one tap → Podcasts opens), and
        offers a copy-to-clipboard backup."""
        from flask import request
        from html import escape as h
        import socket as _socket
        import io as _io

        # Compose the feed URL the phone should hit. Prefer the request's
        # host (so a user already browsing via .local sees their LAN URL),
        # falling back to socket.gethostname for the localhost case.
        request_host = request.host  # "<host>:<port>" as seen by client
        host_only = request_host.split(":")[0]
        is_localhost = host_only in ("127.0.0.1", "localhost")
        if is_localhost:
            # User opened /listen via localhost — the URL we expose
            # must be reachable from the phone, so swap in the .local
            # hostname. Server may not be bound to 0.0.0.0 yet; warn.
            hostname = _socket.gethostname()
            if not hostname.endswith(".local"):
                hostname = f"{hostname}.local"
            display_host = f"{hostname}:{request.host.split(':', 1)[1]}"
        else:
            display_host = request_host

        feed_url = f"http://{display_host}/podcast.xml"
        # podcast:// is Apple Podcasts' deep-link scheme. Tapping a
        # podcast:// URL on iOS opens Podcasts with a Subscribe prompt.
        deep_link = f"podcast://{display_host}/podcast.xml"

        # Render the deep-link as an SVG QR — no PIL needed for SVG path.
        import qrcode as _qrcode
        from qrcode.image.svg import SvgPathImage as _SvgPath
        qr = _qrcode.QRCode(box_size=10, border=2)
        qr.add_data(deep_link)
        qr.make(fit=True)
        buf = _io.BytesIO()
        qr.make_image(image_factory=_SvgPath).save(buf)
        qr_svg = buf.getvalue().decode()
        # The library writes a full XML doc; strip the prelude so we can
        # inline the <svg> directly.
        if "<svg" in qr_svg:
            qr_svg = qr_svg[qr_svg.index("<svg"):]

        # Bind-host check: warn if the server is loopback-only. The
        # browser tab the user is reading this page in might be on
        # localhost regardless (they're on the Mac), so we can't infer
        # bind state from the request — read it from app.config where
        # cmd_serve stashed it.
        bind_host = app.config.get("YT2MD_BIND_HOST", "127.0.0.1")
        lan_unreachable = bind_host in ("127.0.0.1", "localhost")

        warning_html = ""
        if lan_unreachable:
            warning_html = (
                '<div style="background:#fff3cd; border:1px solid #ffc107; '
                'border-left-width:4px; padding:14px 18px; margin:0 0 24px; '
                'border-radius:4px;">'
                '<strong>⚠ Server isn\'t LAN-reachable yet.</strong><br>'
                'Restart the server with <code>yt2md serve --host 0.0.0.0</code> '
                'before subscribing, otherwise your phone can\'t reach the feed.'
                '</div>'
            )

        body = (
            warning_html
            + '<h1>📱 Listen on phone</h1>'
            '<p class="meta-info" style="margin-bottom: 28px;">Subscribe '
            'in Apple Podcasts (or Overcast, Pocket Casts, etc.) to '
            'auto-download new takeaways as they\'re generated. '
            'Plays offline, resumes mid-episode.</p>'

            '<div style="display:flex; gap:32px; flex-wrap:wrap; '
            'align-items:flex-start;">'

            # Left column: QR code
            '<div style="flex:0 0 auto;">'
            '<h2 style="margin-top:0;">Easy: scan with your phone</h2>'
            '<div style="background:#fff; padding:16px; border-radius:8px; '
            'display:inline-block; border:1px solid var(--border);">'
            + qr_svg.replace(
                "<svg",
                '<svg width="260" height="260" style="display:block;"'
            )
            + '</div>'
            '<ol style="font-size:14px; color:var(--muted); margin-top:14px; '
            'padding-left:20px; max-width:260px;">'
            '<li>Open Camera app on your iPhone.</li>'
            '<li>Point it at the code above.</li>'
            '<li>Tap the notification — Apple Podcasts opens.</li>'
            '<li>Tap <strong>Subscribe</strong>.</li>'
            '</ol>'
            '</div>'

            # Right column: manual subscribe details
            '<div style="flex:1 1 320px; min-width:280px;">'
            '<h2 style="margin-top:0;">Or paste this URL</h2>'
            f'<input type="text" readonly value="{h(feed_url)}" '
            'id="feed-url" style="width:100%; font-family:monospace; '
            'font-size:13px; padding:10px 12px; border:1px solid var(--border); '
            'border-radius:4px; background:var(--code-bg);" '
            'onclick="this.select();">'
            '<div style="margin-top:10px; display:flex; gap:8px;">'
            '<button class="primary" '
            'data-copy-target="feed-url-raw">📋 Copy URL</button>'
            f'<a href="{h(deep_link)}" class="discuss-btn-secondary" '
            'style="text-decoration:none;">Open in Podcasts ▸</a>'
            '</div>'
            f'<textarea id="feed-url-raw" style="position:absolute; '
            'left:-9999px;">{h(feed_url)}</textarea>'

            '<h3 style="margin-top:24px;">Apple Podcasts</h3>'
            '<ol style="font-size:14px; line-height:1.6; color:var(--muted);">'
            '<li>Open <strong>Podcasts</strong> on your iPhone.</li>'
            '<li>Tap <strong>Library</strong> (bottom right).</li>'
            '<li>Tap the <strong>⋯</strong> menu → '
            '<strong>Add a Show by URL</strong>.</li>'
            '<li>Paste the URL above and tap <strong>Subscribe</strong>.</li>'
            '</ol>'

            '<h3>Overcast / Pocket Casts</h3>'
            '<p style="font-size:14px; line-height:1.6; color:var(--muted);">'
            'Look for "Add URL" or "Add by URL" — same flow.</p>'

            '</div>'  # right column
            '</div>'  # flex container

            f'<p class="meta-info" style="margin-top:32px;"><strong>Feed URL:</strong> '
            f'<a href="{h(feed_url)}" target="_blank">{h(feed_url)}</a> · '
            f'<strong>Deep link:</strong> <code>{h(deep_link)}</code></p>'
        )
        return page(body, title="Listen on phone", current="listen")

    @app.route("/podcast.xml")
    def podcast_feed():
        """RSS 2.0 podcast feed listing every *.mp3 in the library as an
        episode. Subscribe in Apple Podcasts / Overcast / etc. with the
        URL of this endpoint — once subscribed, the app auto-downloads
        new takeaways, plays offline, and resumes mid-episode.

        Audio URLs are derived from request.host_url so subscribing from
        a phone (over LAN with --host 0.0.0.0) bakes phone-reachable
        URLs into the feed."""
        from flask import request, Response
        from html import escape as h
        from email.utils import format_datetime
        import datetime as _dt

        base = request.host_url.rstrip("/")
        items_xml: list = []
        latest_pub = None

        # Walk all digest dirs, pull each MP3 as a separate episode so
        # podcast apps queue digest/panel/takeaway independently. Sort
        # by audio mtime desc so the newest renderings appear first
        # (podcast apps care about pubDate, not feed position, but
        # ordering helps anyone eyeballing the raw XML).
        episodes: list = []
        for d in digests_dir.iterdir() if digests_dir.exists() else []:
            if not d.is_dir():
                continue
            for kind in AUDIO_SOURCE_BY_KIND:  # digest / panel / takeaway
                mp3 = d / f"{kind}.mp3"
                if not mp3.exists():
                    continue
                episodes.append((mp3.stat().st_mtime, d.name, kind, mp3))
        episodes.sort(reverse=True)

        for mtime, vid, kind, mp3 in episodes:
            # Pull title + overview from digest.json (cheap; cached).
            try:
                dj = load_digest_json(vid, digests_dir=digests_dir)
                title = dj["video"].get("title") or vid
                overview = dj.get("overview") or ""
            except Exception:
                title = vid
                overview = ""
            kind_label = kind.capitalize()
            ep_title = f"{kind_label}: {title}"
            ep_guid = f"yt2md:{vid}:{kind}"
            ep_pub = format_datetime(
                _dt.datetime.fromtimestamp(mtime, tz=_dt.timezone.utc)
            )
            if latest_pub is None or mtime > latest_pub:
                latest_pub = mtime
            ep_url = f"{base}/digests/{vid}/audio/{kind}.mp3"
            try:
                ep_size = mp3.stat().st_size
            except OSError:
                ep_size = 0
            # Description: kind hint + the digest's overview, so the
            # podcast app's episode notes give the listener context.
            desc_body = f"({kind_label} narration) {overview}".strip()
            items_xml.append(
                f"  <item>\n"
                f"    <title>{h(ep_title)}</title>\n"
                f"    <description>{h(desc_body)}</description>\n"
                f"    <enclosure url=\"{h(ep_url)}\" "
                f"length=\"{ep_size}\" type=\"audio/mpeg\" />\n"
                f"    <guid isPermaLink=\"false\">{h(ep_guid)}</guid>\n"
                f"    <pubDate>{h(ep_pub)}</pubDate>\n"
                f"    <link>{h(base)}/digests/{h(vid)}/</link>\n"
                f"  </item>"
            )

        last_build = format_datetime(
            _dt.datetime.fromtimestamp(
                latest_pub or _dt.datetime.now().timestamp(),
                tz=_dt.timezone.utc,
            )
        )
        body = (
            '<?xml version="1.0" encoding="UTF-8"?>\n'
            '<rss version="2.0" '
            'xmlns:itunes="http://www.itunes.com/dtds/podcast-1.0.dtd">\n'
            '<channel>\n'
            '  <title>yt2md library</title>\n'
            f'  <link>{h(base)}/</link>\n'
            '  <description>Audio renditions of your distilled YouTube '
            'digests, panels, and takeaways — generated locally by yt2md.</description>\n'
            '  <language>en</language>\n'
            f'  <lastBuildDate>{h(last_build)}</lastBuildDate>\n'
            '  <itunes:author>yt2md</itunes:author>\n'
            '  <itunes:summary>Audio renditions of your distilled YouTube digests.</itunes:summary>\n'
            '  <itunes:explicit>no</itunes:explicit>\n'
            '  <itunes:category text="Technology"/>\n'
            + "\n".join(items_xml)
            + "\n</channel>\n</rss>\n"
        )
        return Response(body, mimetype="application/rss+xml")

    @app.route("/digests/<video_id>/job-status")
    def job_status_route(video_id):
        from flask import jsonify, request
        kind = request.args.get("kind", "slides")
        snap = local_job_status(f"{video_id}:{kind}")
        # Surface artifact presence so the polling UI can decide whether to
        # reload to the success state or show "missing" after a clean exit.
        artifact_paths = {
            "slides": digests_dir / video_id / "slides.pptx",
            "panel": digests_dir / video_id / "panel.md",
            "takeaway": digests_dir / video_id / "takeaway.md",
            "audio_digest":   digests_dir / video_id / "digest.mp3",
            "audio_panel":    digests_dir / video_id / "panel.mp3",
            "audio_takeaway": digests_dir / video_id / "takeaway.mp3",
        }
        target = artifact_paths.get(kind)
        if target is not None:
            snap["artifact_exists"] = target.exists()
        return jsonify(snap)

    @app.errorhandler(404)
    def not_found(e):
        return page(
            "<h1>Not found</h1><p>That digest doesn't exist yet.</p>",
            title="404", current="home",
        ), 404

    host = getattr(args, "host", None) or "127.0.0.1"
    # Localhost URL is what we open in the browser locally, regardless of
    # bind address — the user is still on this Mac. LAN URL is what
    # other devices (phone, tablet) should hit when binding to 0.0.0.0.
    url = f"http://127.0.0.1:{args.port}/"
    print(f"yt2md reader: {url}")
    if host == "0.0.0.0":
        import socket as _socket
        hostname = _socket.gethostname()
        if not hostname.endswith(".local"):
            hostname = f"{hostname}.local"
        lan_url = f"http://{hostname}:{args.port}/"
        print(f"LAN access:   {lan_url}  (reachable from phone/tablet on same Wi-Fi)")
        print(f"📱 Listen on phone: {lan_url}listen  "
              "(QR code + Apple Podcasts deep link)")
        print(f"  ⚠ Library is now readable to anyone on your network. "
              "Use --host 127.0.0.1 (default) if that's not what you want.")
    else:
        # Friendly nudge so users don't miss the phone-playback path
        # just because they don't know about the --host flag.
        print(f"  (tip: restart with `yt2md serve --host 0.0.0.0` to "
              f"listen to MP3s on a phone — see {url}listen)")
    print(f"Data dir: {data_dir}")
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if api_key:
        print(f"API key: set ({api_key[:7]}…{api_key[-4:]})")
    else:
        print(f"API key: (not set)")
    # Cheap sentinel-based probe to populate the session-state cache so the
    # banner / gate know whether Claude Code is logged in across restarts.
    claude_probe_login_state()
    cc_status = (
        "logged in" if _claude_code_session_state.get("logged_in")
        else ("installed (run /setup to log in)" if claude_code_installed()
              else "not installed")
    )
    print(f"Claude Code: {cc_status}")
    if not api_key and not _claude_code_session_state.get("logged_in"):
        print(f"  → first-run setup at {url}setup")
    # Surface the YouTube prerequisites at startup so the user notices missing
    # cookies / JS runtime before a one-off digest fails 30s later.
    cookies_browser = os.environ.get("YT2MD_COOKIES_FROM_BROWSER")
    print(f"Cookies: {cookies_browser or '(none — set YT2MD_COOKIES_FROM_BROWSER for paywalled YouTube)'}")
    js_rt = _ensure_js_runtime_available()
    print(f"JS runtime: {js_rt or '(not found — yt-dlp n-challenge will fail; install deno or node)'}")
    _cleanup_legacy_launchd()
    start_scheduler()
    print("Press Ctrl-C to stop.\n")

    if not args.no_browser:
        import webbrowser
        import threading
        threading.Timer(0.5, lambda: webbrowser.open(url)).start()

    # Stash the bind host so /listen can tell the user whether the phone
    # can actually reach the server (independent of which URL they used
    # to load the page locally).
    app.config["YT2MD_BIND_HOST"] = host
    app.run(host=host, port=args.port, debug=False, use_reloader=False)
    return 0


def cmd_doctor(args) -> int:
    """Diagnose prerequisites and config. Prints check/X per item with fix hints.

    Exits 0 if everything looks usable, 1 if any blocking issue was found.
    Designed to be the first thing a new user runs after install — gives a
    concrete punch list instead of failing mid-pipeline 30s into the first
    digest.
    """
    OK, WARN, FAIL = "\033[32m✓\033[0m", "\033[33m!\033[0m", "\033[31m✗\033[0m"
    blocking: list = []
    advisory: list = []

    def ok(msg: str) -> None:
        print(f"  {OK} {msg}")

    def warn(msg: str, hint: Optional[str] = None) -> None:
        print(f"  {WARN} {msg}")
        if hint:
            print(f"      → {hint}")
        advisory.append(msg)

    def fail(msg: str, hint: Optional[str] = None) -> None:
        print(f"  {FAIL} {msg}")
        if hint:
            print(f"      → {hint}")
        blocking.append(msg)

    print("\nyt2md doctor — checking prerequisites and config\n")

    print("System tools:")
    for tool in ("ffmpeg", "ffprobe"):
        if shutil.which(tool):
            ok(f"{tool} on PATH")
        else:
            fail(f"{tool} not found",
                 "macOS: `brew install ffmpeg` · "
                 "Windows: `winget install Gyan.FFmpeg` · "
                 "Linux: `apt install ffmpeg`")

    if shutil.which("uv"):
        try:
            v = subprocess.run(["uv", "--version"], capture_output=True, text=True,
                               timeout=5).stdout.strip()
            ok(f"uv ({v})")
        except Exception:
            ok("uv on PATH")
    else:
        warn("uv not on PATH",
             "Recommended: install via `curl -LsSf https://astral.sh/uv/install.sh | sh`")

    rt = _ensure_js_runtime_available()
    if rt:
        rt_path = shutil.which(rt) or "(unknown)"
        try:
            ver = subprocess.run([rt, "--version"], capture_output=True, text=True,
                                 timeout=5).stdout.strip()
        except Exception:
            ver = "?"
        ok(f"JS runtime: {rt} {ver} ({rt_path})")
    else:
        fail("No JS runtime found (needed for yt-dlp's n-challenge solver)",
             "`brew install deno` (simplest) or `nvm install 20`")

    print("\nPython packages:")
    try:
        import yt_dlp  # type: ignore
        ok(f"yt-dlp {yt_dlp.version.__version__}")
    except ImportError:
        fail("yt-dlp not installed", "Run `uv sync` from the project directory")
    try:
        import faster_whisper  # type: ignore  # noqa: F401
        ok("faster-whisper installed")
    except ImportError:
        warn("faster-whisper not installed",
             "Whisper fallback won't work for captionless videos. `uv sync` to install.")
    try:
        import anthropic  # type: ignore  # noqa: F401
        ok("anthropic SDK installed")
    except ImportError:
        fail("anthropic SDK not installed", "Run `uv sync`")

    print("\nAPI / auth:")
    load_env_files()
    key = os.environ.get("ANTHROPIC_API_KEY", "")
    if key:
        ok(f"ANTHROPIC_API_KEY set ({key[:10]}…{key[-4:]})")
    else:
        fail("ANTHROPIC_API_KEY not set",
             f"Get a key at https://console.anthropic.com/settings/keys; first run "
             f"of yt2md will prompt and save it to {get_data_dir() / '.env'}")

    settings = load_settings()
    cookies = settings.get("cookies_from_browser") or os.environ.get("YT2MD_COOKIES_FROM_BROWSER", "")
    if cookies:
        ok(f"YouTube cookies: from browser '{cookies}'")
    else:
        warn("YouTube cookies not configured",
             "Many videos now require login. Set in /settings or as "
             "YT2MD_COOKIES_FROM_BROWSER=firefox in ~/yt2md/.env")

    # Cost controls — this app spends real money per digest. Surface the
    # safety setup (Console spend cap + the optional in-app budget gate)
    # so a new user configures the "right way" before the first run.
    print("\nCost controls:")
    print("  reminder: set a monthly spend cap in the Anthropic Console "
          "(https://console.anthropic.com/settings/limits) — the hard backstop "
          "nothing can exceed.")
    admin = os.environ.get("ANTHROPIC_ADMIN_KEY", "")
    if admin:
        ok(f"ANTHROPIC_ADMIN_KEY set ({admin[:14]}…) — authoritative billing + price calibration")
    else:
        warn("ANTHROPIC_ADMIN_KEY not set (optional)",
             "Without it the budget gate uses the local usage log and prices come "
             "from the built-in table. Add an org admin key to ~/yt2md/.env for real "
             "billed month-to-date spend and `yt2md refresh-pricing`.")
    block_usd = settings.get("budget_block_usd")
    warn_usd = settings.get("budget_warn_usd")
    if block_usd:
        ok(f"budget gate: warn ${float(warn_usd or 0):.0f} / block new digests "
           f"${float(block_usd):.0f} (month-to-date)")
    else:
        warn("in-app budget gate off",
             "Set budget_warn_usd / budget_block_usd in ~/yt2md/settings.json to "
             "refuse new digests past a monthly threshold.")
    if _pricing_cache_path().exists():
        ok("prices: calibrated from real billing (pricing_cache.json)")
    else:
        print("  prices: built-in table (run `yt2md refresh-pricing` to "
              "calibrate from real billing — needs admin key)")

    print("\nConfig:")
    data_dir = get_data_dir()
    print(f"  data dir: {data_dir}")
    if _settings_file().exists():
        ok(f"settings.json exists")
    else:
        warn("settings.json not yet created — defaults in effect",
             f"Open http://localhost:7682/settings to configure (after `yt2md serve`)")
    print(f"  digest model: {settings.get('digest_model')}")
    print(f"  panel model:  {settings.get('panel_model')}")
    print(f"  whisper model: {settings.get('whisper_model')}")
    print(f"  digest language: {settings.get('digest_language')}")

    cfg = load_schedule_config()
    print(f"  schedule: {_format_schedule_summary(cfg)}")

    digests_dir = data_dir / "digests"
    n_digests = len([d for d in digests_dir.iterdir() if (d / "digest.md").exists()]) if digests_dir.exists() else 0
    n_channels = len(read_channels())
    print(f"  library: {n_digests} digest(s), {n_channels} channel(s)")

    print()
    if blocking:
        print(f"{FAIL} {len(blocking)} blocking issue{'s' if len(blocking) != 1 else ''}; "
              "fix the items marked above.")
        if advisory:
            print(f"{WARN} {len(advisory)} advisory item{'s' if len(advisory) != 1 else ''} "
                  "(non-blocking but worth setting up).")
        return 1
    if advisory:
        n = len(advisory)
        print(f"{OK} Core requirements met. {n} advisory item"
              f"{'s' if n != 1 else ''} {'are' if n != 1 else 'is'} optional.")
    else:
        print(f"{OK} All checks passed. Run `yt2md serve` to start the reader.")
    return 0


# ---- Library CLI subcommands -------------------------------------------
#
# Thin wrappers over the agent API for shell / scheduler use. Designed
# to be pipe-friendly: --json flips text output (human-readable) to a
# JSON stream that jq or a Claude Desktop scheduled task can consume.

def _maybe_print_json(obj, *, as_json: bool) -> None:
    """Pretty-print as JSON or pass to the human-readable formatter."""
    if as_json:
        print(json.dumps(obj, indent=2, ensure_ascii=False))
        return
    _print_human(obj)


def _print_human(obj) -> None:
    """Minimal human-readable view for the CLI subcommands. JSON is
    available via --json for any caller that needs full fidelity."""
    if isinstance(obj, list):
        if not obj:
            print("(empty)")
            return
        for item in obj:
            if isinstance(item, dict) and "id" in item and "title" in item:
                # list_digests entry
                flags = "".join([
                    "P" if item.get("has_panel") else "-",
                    "T" if item.get("has_takeaway") else "-",
                    "S" if item.get("has_slides") else "-",
                ])
                date = item.get("published_at") or "    -    "
                read = "·" if item.get("read") else " "
                print(f"  {read} {item['id']:14s} [{flags}] {date}  {item['title']}")
            elif isinstance(item, dict) and "digest_id" in item and "section" in item:
                # search_library hit
                print(f"  [{item['score']:>2d}] {item['digest_id']:14s} "
                      f"{item['section']:20s} {item.get('snippet', '')}")
            elif isinstance(item, dict) and "url" in item:
                print(f"  {item['url']}")
            else:
                print(f"  {item}")
        return
    if isinstance(obj, dict):
        for k, v in obj.items():
            print(f"  {k}: {v}")
        return
    print(obj)


def cmd_list(args) -> int:
    """yt2md list — browse the local digest library."""
    entries = list_digests(
        channel=args.channel or None,
        since=args.since or None,
        unread=args.unread,
        q=args.q or None,
        topic=args.topic or None,
        source=args.source or None,
        saved=True if args.saved else None,
        dismissed=True if args.dismissed else False,
        limit=args.limit,
    )
    _maybe_print_json(entries, as_json=args.json)
    return 0


def cmd_read(args) -> int:
    """yt2md read <id> — read a section of a digest."""
    try:
        result = read_digest(args.video_id, section=args.section)
    except (ValueError, NotImplementedError, FileNotFoundError) as e:
        sys.exit(f"read failed: {e}")
    _maybe_print_json(result, as_json=args.json)
    return 0


def cmd_search(args) -> int:
    """yt2md search <query> — substring search across the library."""
    hits = search_library(args.query, k=args.k)
    _maybe_print_json(hits, as_json=args.json)
    return 0


def cmd_digest(args) -> int:
    """yt2md digest <url> — kick off the ingestion pipeline for a URL.
    Non-blocking by default; --wait blocks until done."""
    result = digest_video(args.url, blocking=args.wait, source=args.source)
    _maybe_print_json(result, as_json=args.json)
    return 0


def cmd_topics(args) -> int:
    """yt2md topics — list the topic taxonomy across the library."""
    topics = list_topics(
        min_digests=args.min_digests, limit=args.limit,
        source=args.source or None,
    )
    if args.json:
        print(json.dumps(topics, indent=2, ensure_ascii=False))
        return 0
    if not topics:
        print("(no topics yet — tag a digest with `yt2md retrofit-topics`)")
        return 0
    for t in topics:
        srcs = []
        if t["sources"]["llm"]: srcs.append("llm")
        if t["sources"]["user"]: srcs.append("user")
        print(f"  {t['n_digests']:>3d}  {t['topic']:30s}  ({'+'.join(srcs)})")
    return 0


def cmd_retrofit_topics(args) -> int:
    """yt2md retrofit-topics — tag any digest that doesn't already have
    LLM topics. Resumable. Use --dry-run first to see what would happen."""
    result = retrofit_topics(
        since=args.since or None,
        limit=args.limit,
        dry_run=args.dry_run,
    )
    if args.dry_run:
        print(f"\nWould tag {result['count']} digest(s):")
        for vid in result["would_tag"][:20]:
            print(f"  {vid}")
        if result["count"] > 20:
            print(f"  ... and {result['count'] - 20} more")
        return 0
    print(f"\nTagged {result['tagged']} of {result['candidates']}; "
          f"errors: {len(result['errors'])}")
    return 0 if not result["errors"] else 1


# ---- Claude Project per-channel instructions ---------------------------
#
# A Claude Project on claude.ai can be pinned to a single channel by giving
# it Custom Instructions that tell the agent to scope every yt2md MCP call
# to that channel's name. This pair (render_claude_project_instructions +
# cmd_project_instructions) emits that text. Surfaced in the web UI on the
# Subscriptions page as a "Set up Claude Project" link, and on the CLI as
# `yt2md project-instructions [--channel <name>]`.


def _known_channel_names(digests_dir: Optional[Path] = None) -> List[str]:
    """Distinct channel_names that have at least one digest on disk.

    Read from per-digest metadata.json, sorted case-insensitively. Used
    by both the CLI lister (no --channel arg) and the web Subscriptions
    page to decide which channels can offer a Claude Project setup link.
    """
    if digests_dir is None:
        digests_dir = get_data_dir() / "digests"
    if not digests_dir.exists():
        return []
    seen: dict = {}
    for d in digests_dir.iterdir():
        if not d.is_dir():
            continue
        meta_path = d / "metadata.json"
        if not meta_path.exists():
            continue
        try:
            meta = json.loads(meta_path.read_text())
        except Exception:
            continue
        name = (meta.get("channel_name") or "").strip()
        if name and name.lower() not in seen:
            seen[name.lower()] = name
    return sorted(seen.values(), key=lambda s: s.lower())


def render_claude_project_instructions(channel_name: str) -> str:
    """Return Markdown text to paste into a Claude Project's Custom
    Instructions. Pins the agent to `channel_name` via the yt2md MCP tools.

    The text is deliberately the *contents* of the Custom Instructions
    box — no surrounding setup prose (that lives on the page that hosts
    this text).
    """
    return f"""You are a reading assistant scoped to a single YouTube channel: **{channel_name}**.

Your knowledge of this channel comes from the yt2md MCP server running on the user's machine. It reads a local library of per-video artifacts (digest, panel discussion, takeaway synthesis, slides) that yt2md has already generated. Always pass `channel="{channel_name}"` to scope your queries — never query the global library.

How to work:

1. **Searching.** Default to `search_library(q="…")` first; hits include the `digest_id` you can pass to `read_digest`. The search is library-wide, so check that each hit's `channel` matches "{channel_name}" before quoting it. If search returns nothing, try `list_digests(channel="{channel_name}", q="<title keyword>")` and pick by title.
2. **No question yet.** When the user opens a fresh chat without a specific ask, call `list_digests(channel="{channel_name}", unread=True, limit=5)`. Offer to walk through the most recent unread digest. If nothing is unread, fall back to `list_digests(channel="{channel_name}", limit=5)`.
3. **Reading efficiently.** For a deep question, read the relevant `takeaway` first (it's the synthesis), then `panel` (critique), then specific `topics`. Don't read `section="full"` unless you actually need everything — it burns tokens.
4. **Citing.** Every factual claim should reference the digest's title and, where possible, the YouTube timestamp link from the digest body. `read_digest` returns a `video` block with the canonical `url`.

What you do NOT have:
- The raw video. You cannot watch, transcribe, or re-extract frames — the digest text is the source of truth. If something isn't in the digest/panel/takeaway, say so rather than guess.
- Other channels. If the user asks about a video that's not on **{channel_name}**, tell them — they should open a different Project (one per channel) for that.
- Live ingestion expectations. You *can* call `digest_video(url)` if the user pastes a new URL, but the pipeline takes 5–10 minutes; let them know it'll land in the library after, and offer to come back to it.
"""


def cmd_project_instructions(args) -> int:
    """`yt2md project-instructions [--channel <name>]`.

    With no --channel: list channel_names that have at least one digest
    on disk (the ones a Project would actually have something to chat
    about). With --channel: case-insensitive substring match against
    those known names; render the template if exactly one matches.
    """
    known = _known_channel_names()
    if not args.channel:
        if not known:
            print("No digested channels yet. Subscribe and let the scheduler "
                  "poll, or run `yt2md digest <url>` first.")
            return 0
        print("Channels in your library (pass one to --channel):")
        for name in known:
            print(f"  {name}")
        return 0

    needle = args.channel.strip().lower()
    matches = [n for n in known if needle in n.lower()]
    if not matches:
        print(f"No channel in the library matches {args.channel!r}.",
              file=sys.stderr)
        if known:
            print("Known channels:", file=sys.stderr)
            for n in known:
                print(f"  {n}", file=sys.stderr)
        return 1
    if len(matches) > 1:
        print(f"{args.channel!r} matches multiple channels — be more specific:",
              file=sys.stderr)
        for n in matches:
            print(f"  {n}", file=sys.stderr)
        return 1

    print(render_claude_project_instructions(matches[0]))
    return 0


# ---- MCP server subcommand ---------------------------------------------
#
# Exposes the Phase A agent API (read_digest, search_library, etc.) over
# the Model Context Protocol so Claude Desktop / Claude Code / any MCP
# client can talk to the local library. Tools intentionally mirror the
# Python API one-to-one — the server is a thin transport adapter, not a
# new abstraction layer. Logic lives in the existing functions.

def cmd_mcp(args) -> int:
    """Run the yt2md MCP server (stdio transport).

    Wire it into Claude Desktop by adding to
    ~/Library/Application Support/Claude/claude_desktop_config.json:

        {
          "mcpServers": {
            "yt2md": {"command": "yt2md", "args": ["mcp"]}
          }
        }

    Wire it into Claude Code with:
        claude mcp add yt2md -- yt2md mcp
    """
    try:
        from mcp.server.fastmcp import FastMCP
    except ImportError:
        sys.exit(
            "MCP server requires the `mcp` package. Reinstall with:\n"
            "  uv tool install --force --python 3.11 .\n"
            "or:\n"
            "  pip install 'mcp>=1.2'"
        )

    mcp = FastMCP("yt2md")

    # ---- Library navigation -------------------------------------------

    @mcp.tool()
    def list_digests(
        channel: str = "",
        since: str = "",
        unread: bool = False,
        q: str = "",
        topic: str = "",
        source: str = "",
        saved: bool = False,
        include_dismissed: bool = False,
        only_dismissed: bool = False,
        limit: int = 20,
    ) -> list:
        """List digests in the local yt2md library, most recent first.

        Use to browse the library or to find videos before drilling in
        with read_digest. Filters compose (AND).

        Args:
            channel: case-insensitive substring match against channel name.
            since: ISO date "YYYY-MM-DD"; only digests at or after this date.
            unread: true → only digests not yet marked read.
            q: case-insensitive substring match against title.
            topic: exact-match on a topic tag (LLM or user-assigned).
                Use list_topics to discover available tags.
            source: "subscription" | "oneoff" | "meta" (provenance filter).
            saved: true → only user-saved digests.
            include_dismissed: false (default) hides user-dismissed digests.
                true → show all including dismissed.
            only_dismissed: true → only dismissed digests (overrides
                include_dismissed).
            limit: max entries (default 20, hard cap 200).

        Returns: list of {id, title, url, channel, channel_url,
            published_at, mtime, read, topics, topics_split, source,
            user_saved, user_dismissed, has_panel, has_takeaway,
            has_slides, has_audio}.
        """
        dismissed_filter: Optional[bool]
        if only_dismissed:
            dismissed_filter = True
        elif include_dismissed:
            dismissed_filter = None  # don't filter
        else:
            dismissed_filter = False
        return globals()["list_digests"](
            channel=channel or None, since=since or None, unread=unread,
            q=q or None, topic=topic or None, source=source or None,
            saved=True if saved else None, dismissed=dismissed_filter,
            limit=min(limit, 200),
        )

    @mcp.tool()
    def read_digest(digest_id: str, section: str = "full") -> dict:
        """Read a section of a digest as structured data.

        section options:
            "meta"           - just video metadata
            "overview"       - the opening overview paragraph
            "topics"         - all topics (titles + bodies + bullets)
            "topic:N"        - the Nth topic (1-indexed)
            "topic:<slug>"   - first topic whose title contains <slug>
            "panel"          - panelists + all turns
            "panel:panelists"- just the panelist bios
            "panel:turn:N"   - the Nth panel turn
            "takeaway"       - paragraphs + inline citations
            "full"           - everything (digest only; panel/takeaway
                               are separate)

        Returns: {section, video, content}. Always includes a `video`
        block so the caller can deep-link back to YouTube. Raises
        ValueError for unknown sections or out-of-range indices.
        """
        return globals()["read_digest"](digest_id, section=section)

    @mcp.tool()
    def search_library(q: str, k: int = 10) -> list:
        """Substring search across the local library (digest + panel +
        takeaway). Case-insensitive whole-token AND match.

        Use this to find which digests mention a topic, then pipe a
        hit's `section` directly into read_digest for the full content.

        Returns: list of {digest_id, title, section, snippet, score,
            url, video} sorted by score desc (title hits weight highest).
        """
        return globals()["search_library"](q, k=min(k, 50))

    # ---- Ingestion + generation ---------------------------------------

    @mcp.tool()
    def digest_video(url: str) -> dict:
        """Ingest a YouTube video into the library. Non-blocking: spawns
        the digest pipeline as a detached child process and returns the
        job handle. The pipeline takes 5-10 minutes for a typical video.

        Skip-if-exists: if the video is already digested, returns the
        existing digest immediately without re-running.

        Returns one of:
            {video_id, status: "exists", digest: <full digest JSON>}
            {video_id, job_id, status: "started", log_path}
        """
        return globals()["digest_video"](url, blocking=False)

    @mcp.tool()
    def generate_panel(video_id: str) -> dict:
        """Generate the panel discussion for a digested video. Non-
        blocking: runs in a background thread. Poll with job_status
        or just read_digest(id, "panel") after a couple of minutes.

        Returns: {video_id, kind: "panel", status, job_key}.
        """
        return globals()["generate_panel"](video_id, blocking=False)

    @mcp.tool()
    def generate_takeaway(video_id: str) -> dict:
        """Generate the takeaway synthesis for a digested video. Best
        run after the panel exists so it can integrate panel critique.
        Non-blocking; see generate_panel for shape."""
        return globals()["generate_takeaway"](video_id, blocking=False)

    @mcp.tool()
    def generate_slides(video_id: str) -> dict:
        """Generate the slide deck (.pptx) from cached frames for a
        digested video. Non-blocking; see generate_panel for shape."""
        return globals()["generate_slides"](video_id, blocking=False)

    @mcp.tool()
    def generate_audio(video_id: str, kind: str) -> dict:
        """Render a digest/panel/takeaway to MP3 using the configured
        TTS provider (macOS `say` or ElevenLabs). Non-blocking; poll
        via job_status with kind=f"audio_{kind}".

        kind: "digest" | "panel" | "takeaway".
        Skip-if-exists: returns immediately if the mp3 is already on disk."""
        return globals()["generate_audio"](video_id, kind, blocking=False)

    @mcp.tool()
    def mark_digest_read(digest_id: str) -> dict:
        """Mark a digest as read. Use after surfacing a digest to the
        user so future `list_digests(unread=True)` calls skip it.
        Idempotent."""
        return globals()["mark_digest_read"](digest_id)

    @mcp.tool()
    def mark_digest_unread(digest_id: str) -> dict:
        """Mark a digest as unread. Idempotent."""
        return globals()["mark_digest_unread"](digest_id)

    # ---- Topic taxonomy + curation -----------------------------------

    @mcp.tool()
    def list_topics(min_digests: int = 1, limit: int = 50,
                    source: str = "") -> list:
        """List the topic taxonomy across the library.

        Returns [{topic, n_digests, last_seen, sources: {llm, user}}],
        sorted by digest count desc. Use to answer "what have I been
        reading about?" and to pick a tag for list_digests(topic=...).

        source: '' (union) | 'llm' | 'user'. Filter to one provenance
        when you want to see only LLM-assigned vs only your manual tags.
        """
        return globals()["list_topics"](
            min_digests=min_digests, limit=limit,
            source=source or None,
        )

    @mcp.tool()
    def tag_digest(digest_id: str, tags: list) -> dict:
        """Add user tags to a digest. Idempotent. Kept separate from
        LLM-assigned tags, so this never overwrites the auto-tagging.
        Tags are normalized to lowercase-hyphen-separated."""
        return globals()["tag_digest"](digest_id, tags)

    @mcp.tool()
    def untag_digest(digest_id: str, tags: list) -> dict:
        """Remove user tags from a digest. LLM tags are not touched."""
        return globals()["untag_digest"](digest_id, tags)

    @mcp.tool()
    def save_digest(digest_id: str) -> dict:
        """Mark a digest as 'saved' (worth keeping / returning to).
        Surfaces in list_digests(saved=True)."""
        return globals()["save_digest"](digest_id)

    @mcp.tool()
    def unsave_digest(digest_id: str) -> dict:
        """Remove the 'saved' flag."""
        return globals()["unsave_digest"](digest_id)

    @mcp.tool()
    def dismiss_digest(digest_id: str) -> dict:
        """Mark a digest as 'dismissed' so it's hidden from briefings.
        list_digests(dismissed=False) excludes it by default."""
        return globals()["dismiss_digest"](digest_id)

    @mcp.tool()
    def undismiss_digest(digest_id: str) -> dict:
        """Remove the 'dismissed' flag."""
        return globals()["undismiss_digest"](digest_id)

    @mcp.tool()
    def retag_digest(digest_id: str) -> dict:
        """Re-run the LLM tagging step. Useful when the taxonomy has
        grown — a digest tagged early can pick up tags introduced later."""
        return globals()["retag_digest"](digest_id)

    @mcp.tool()
    def job_status(video_id: str, kind: str) -> dict:
        """Check the status of a background generation job.

        kind: "panel" | "takeaway" | "slides" | "audio_digest" |
              "audio_panel" | "audio_takeaway"

        Returns: {phase: "idle"|"running"|"done"|"error", elapsed?, error?}
        """
        return local_job_status(f"{video_id}:{kind}")

    # ---- Subscriptions -------------------------------------------------

    @mcp.tool()
    def list_subscriptions() -> list:
        """List subscribed YouTube channels."""
        return globals()["list_subscriptions"]()

    @mcp.tool()
    def add_subscription(channel_url: str) -> dict:
        """Subscribe to a YouTube channel. Accepts a full URL, a
        youtube.com/@handle path, or a bare @handle. Idempotent."""
        return globals()["add_subscription"](channel_url)

    @mcp.tool()
    def remove_subscription(channel_url: str) -> dict:
        """Unsubscribe from a YouTube channel. Idempotent."""
        return globals()["remove_subscription"](channel_url)

    mcp.run(transport="stdio")
    return 0


# ---- subcommand dispatcher ----

def _subcommand_main(argv: List[str]) -> int:
    """Handle yt2md {watch,serve} ..."""
    ap = argparse.ArgumentParser(prog="yt2md", description="yt2md subcommands")
    sub = ap.add_subparsers(dest="cmd", required=True)

    watch = sub.add_parser("watch", help="Manage watched channels and run polling")
    watch_sub = watch.add_subparsers(dest="watch_cmd", required=True)
    p = watch_sub.add_parser("add", help="Add a channel URL"); p.add_argument("url")
    p.set_defaults(func=cmd_watch_add)
    p = watch_sub.add_parser("list", help="List watched channels"); p.set_defaults(func=cmd_watch_list)
    p = watch_sub.add_parser("remove", help="Remove a channel URL"); p.add_argument("url")
    p.set_defaults(func=cmd_watch_remove)
    p = watch_sub.add_parser("run", help="Poll all channels and digest new videos")
    p.set_defaults(func=cmd_watch_run)

    serve = sub.add_parser("serve", help="Start a local web reader (also runs the in-process scheduler)")
    serve.add_argument("--port", type=int, default=7682, help="Port (default: 7682)")
    serve.add_argument("--host", default="127.0.0.1",
                       help="Bind address (default: 127.0.0.1, localhost only). "
                            "Pass 0.0.0.0 to make the library reachable from "
                            "other devices on the same Wi-Fi (e.g. listen to "
                            "MP3s on a phone). Exposes the whole library to "
                            "anyone on your network — only do this on trusted Wi-Fi.")
    serve.add_argument("--no-browser", action="store_true",
                       help="Don't auto-open a browser tab on start")
    serve.set_defaults(func=cmd_serve)

    doctor = sub.add_parser("doctor", help="Check prerequisites and config; print a punch list")
    doctor.set_defaults(func=cmd_doctor)

    mcp_parser = sub.add_parser(
        "mcp",
        help="Run the MCP server (stdio). Wire into Claude Desktop / Claude Code "
             "to give an agent access to the library.",
    )
    mcp_parser.set_defaults(func=cmd_mcp)

    # ---- Library query/ingest CLI (same surface as the MCP tools) ----
    list_p = sub.add_parser("list", help="List digests in the library")
    list_p.add_argument("--channel", default="", help="Substring match on channel name")
    list_p.add_argument("--since", default="", help='ISO date "YYYY-MM-DD"')
    list_p.add_argument("--unread", action="store_true", help="Only unread digests")
    list_p.add_argument("-q", default="", help="Substring match on title")
    list_p.add_argument("--topic", default="", help="Exact-match on a topic tag")
    list_p.add_argument("--source", default="",
                        choices=("", "subscription", "oneoff", "meta"),
                        help="Provenance filter")
    list_p.add_argument("--saved", action="store_true",
                        help="Only digests marked saved")
    list_p.add_argument("--dismissed", action="store_true",
                        help="Only digests marked dismissed (hidden by default)")
    list_p.add_argument("--limit", type=int, default=20, help="Max results (default: 20)")
    list_p.add_argument("--json", action="store_true", help="Emit JSON (default: human-readable)")
    list_p.set_defaults(func=cmd_list)

    read_p = sub.add_parser("read", help="Read a section of a digest as structured data")
    read_p.add_argument("video_id", help="YouTube video ID (digest dir name)")
    read_p.add_argument("--section", default="full",
                        help='full | meta | overview | topics | topic:N | '
                             'topic:<slug> | panel | panel:turn:N | '
                             'panel:panelists | takeaway')
    read_p.add_argument("--json", action="store_true",
                        help="Emit JSON (default: pretty-printed)")
    read_p.set_defaults(func=cmd_read)

    search_p = sub.add_parser("search", help="Substring search across the library")
    search_p.add_argument("query", help="Search query (case-insensitive, AND across tokens)")
    search_p.add_argument("-k", type=int, default=10, help="Max hits (default: 10)")
    search_p.add_argument("--json", action="store_true", help="Emit JSON")
    search_p.set_defaults(func=cmd_search)

    digest_p = sub.add_parser("digest", help="Ingest a YouTube URL")
    digest_p.add_argument("url", help="YouTube URL")
    digest_p.add_argument("--wait", action="store_true",
                          help="Block until the pipeline finishes (default: detached)")
    digest_p.add_argument("--source", default="oneoff",
                          choices=("oneoff", "subscription", "meta"),
                          help="Provenance stamp (default: oneoff)")
    digest_p.add_argument("--json", action="store_true", help="Emit JSON")
    digest_p.set_defaults(func=cmd_digest)

    topics_p = sub.add_parser("topics", help="Show the topic taxonomy")
    topics_p.add_argument("--min-digests", type=int, default=1,
                          help="Only show tags with at least N digests")
    topics_p.add_argument("--limit", type=int, default=50, help="Max tags")
    topics_p.add_argument("--source", default="",
                          choices=("", "llm", "user"),
                          help="Filter by tag provenance")
    topics_p.add_argument("--json", action="store_true", help="Emit JSON")
    topics_p.set_defaults(func=cmd_topics)

    proj_p = sub.add_parser(
        "project-instructions",
        help="Render Custom Instructions for a per-channel Claude Project "
             "scoped to one channel's library via the yt2md MCP tools",
    )
    proj_p.add_argument("--channel", default="",
                        help="Channel name (case-insensitive substring, "
                             "must match exactly one channel that has at "
                             "least one digest). Omit to list known channels.")
    proj_p.set_defaults(func=cmd_project_instructions)

    retro_p = sub.add_parser(
        "retrofit-topics",
        help="Tag any digest in the library that doesn't have LLM topics yet "
             "(resumable)",
    )
    retro_p.add_argument("--since", default="",
                         help='Only tag digests on or after this date (ISO "YYYY-MM-DD")')
    retro_p.add_argument("--limit", type=int, default=None,
                         help="Max digests to tag this run")
    retro_p.add_argument("--dry-run", action="store_true",
                         help="Print what would be tagged without calling the LLM")
    retro_p.set_defaults(func=cmd_retrofit_topics)

    refresh_p = sub.add_parser(
        "refresh-pricing",
        help="Recalibrate the model price table from your actual Anthropic "
             "billing (Admin API) and cache it",
    )
    refresh_p.add_argument("--lookback", type=int, default=14, metavar="DAYS",
                           help="Billing window to derive rates from (default: 14)")
    refresh_p.set_defaults(func=cmd_refresh_pricing)

    args = ap.parse_args(argv)
    return args.func(args)


# ---------- Main ----------

def main():
    # Subcommand dispatch — short-circuit the single-video flow when the user
    # invokes yt2md watch / serve / doctor.
    if len(sys.argv) > 1 and sys.argv[1] in (
        "watch", "serve", "doctor", "mcp",
        "list", "read", "search", "digest", "topics", "retrofit-topics",
        "project-instructions", "refresh-pricing",
    ):
        load_env_files()
        sys.exit(_subcommand_main(sys.argv[1:]))

    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("video", help="Input MP4 file path OR a YouTube URL (auto-downloads mp4 + SRT)")
    ap.add_argument("srt", nargs="?", default=None,
                    help="SRT transcript file (required if 'video' is a local path; "
                         "ignored for URLs — fetched automatically)")
    ap.add_argument("-o", "--output", type=Path, default=None, metavar="PATH",
                    help="Digest output path (default: <video-name>_digest.md)")
    ap.add_argument("--slides", action="store_true",
                    help="Build slides.pptx alongside the digest. Off by default. "
                         "Runs frame extraction, pHash dedup, vision classification, "
                         "and deck generation.")
    ap.add_argument("--no-slide-classification", action="store_true",
                    help="When building slides, skip the vision-LLM classifier that filters "
                         "raw frames down to actual deck slides. Falls back to pHash dedup "
                         "only — cheaper but typically produces 2-3x more (mostly redundant) "
                         "slides for talks where the speaker has a deck.")
    ap.add_argument("--slide-classifier-model",
                    default=os.environ.get("YT2MD_SLIDE_CLASSIFIER_MODEL"),
                    help="Vision model for the slide-classification step (default: "
                         "claude-haiku-4-5-20251001 — cheap and accurate enough for the task).")
    ap.add_argument("--deck", nargs="?", const="__default__", default=None, metavar="PATH",
                    help="Override the slides output path (default: <digest_dir>/slides.pptx).")
    ap.add_argument("--deck-only", action="store_true",
                    help="Skip the digest entirely — only build the deck. No API key needed.")
    ap.add_argument("--no-vision", action="store_true",
                    help="Disable vision-based frame picking for the digest. Cheaper but the "
                         "frames may be less illustrative.")
    ap.add_argument("--no-panel", action="store_true",
                    help="Skip the panel-of-experts discussion (saves ~1 Opus call per video). "
                         "Distillation will still run but without panel-informed confidence tags.")
    ap.add_argument("--no-takeaway", action="store_true",
                    help="Skip the takeaway step (synthesis prose) appended to digest.md.")
    ap.add_argument("--no-tagging", action="store_true",
                    help="Skip the LLM topic-tagging step. Tags drive the agent API's "
                         "topic filter; skip when iterating on the pipeline locally.")
    ap.add_argument("--source", choices=("oneoff", "subscription", "meta"),
                    default="oneoff",
                    help="Provenance tag stamped into digest_meta. Defaults to 'oneoff' "
                         "(direct CLI / one-off web route). The watch run loop passes "
                         "'subscription'; meta-digest runs pass 'meta'.")
    ap.add_argument("--digest-model",
                    default=os.environ.get("YT2MD_DIGEST_MODEL") or "claude-sonnet-4-6",
                    help="Claude model for the digest (default: claude-sonnet-4-6). "
                         "Use claude-opus-4-7 for the highest-quality summarization.")
    ap.add_argument("--panel-model",
                    default=os.environ.get("YT2MD_PANEL_MODEL") or DEFAULT_PANEL_MODEL,
                    help=f"Claude model for the panel discussion (default: {DEFAULT_PANEL_MODEL}). "
                         "Multi-perspective synthesis benefits from a stronger model.")
    ap.add_argument("--takeaway-model",
                    default=os.environ.get("YT2MD_TAKEAWAY_MODEL") or DEFAULT_TAKEAWAY_MODEL,
                    help=f"Claude model for the takeaway step (default: "
                         f"{DEFAULT_TAKEAWAY_MODEL}).")
    ap.add_argument("--scene-threshold", type=float, default=0.2,
                    help="Scene-detection sensitivity, 0.1=lots of frames, 0.5=only major changes (default: 0.2)")
    ap.add_argument("--interval", type=float, default=20.0,
                    help="Also sample one frame every N seconds (0 to disable). Useful for "
                         "screen recordings where gradual changes don't trip scene detection. (default: 20)")
    ap.add_argument("--hash-distance", type=int, default=4,
                    help="Perceptual hash dedup threshold; lower = stricter. Compared only against "
                         "the previous kept frame, so recurring views are preserved. (default: 4)")
    ap.add_argument("--keep-frames", action="store_true",
                    help="Keep extracted frames in ./frames_<videoname>/ instead of cleaning up")
    ap.add_argument("--downloads-dir", type=Path, default=Path("downloads"),
                    help="Where to cache YouTube downloads (default: ./downloads)")
    ap.add_argument("--source-lang", default=None, metavar="CODE",
                    help="BCP-47 language code of a local SRT (e.g. 'zh-Hans'). "
                         "Drives the digest's output language when --digest-language=auto. "
                         "Ignored for URLs — yt-dlp / Whisper picks the track and the "
                         "lang is read from the file.")
    ap.add_argument("--digest-language",
                    default=os.environ.get("YT2MD_DIGEST_LANGUAGE") or "auto",
                    choices=["auto", "en"],
                    help="Output language for the digest + panel discussion. "
                         "'auto' (default) writes in the transcript's language. "
                         "'en' forces English regardless of source.")
    ap.add_argument("--whisper-model",
                    default=os.environ.get("YT2MD_WHISPER_MODEL") or DEFAULT_WHISPER_MODEL,
                    choices=["tiny", "base", "small", "medium", "large-v2", "large-v3"],
                    help=f"faster-whisper model used as fallback when a YouTube "
                         f"video has no captions (default: {DEFAULT_WHISPER_MODEL}).")
    ap.add_argument("--no-whisper", action="store_true",
                    help="Disable Whisper fallback; fail when a video has no captions.")
    ap.add_argument("--cookies-from-browser", default=os.environ.get("YT2MD_COOKIES_FROM_BROWSER"),
                    choices=["chrome", "firefox", "safari", "brave", "edge",
                             "chromium", "opera", "vivaldi"],
                    metavar="BROWSER",
                    help="Pass cookies from this browser to yt-dlp. Required when "
                         "YouTube returns 'Sign in to confirm you're not a bot'. "
                         "Defaults to $YT2MD_COOKIES_FROM_BROWSER if set.")
    args = ap.parse_args()

    load_env_files()

    for tool in ("ffmpeg", "ffprobe"):
        if not shutil.which(tool):
            sys.exit(
                f"{tool} not found on PATH. Install with "
                "`brew install ffmpeg` (macOS), "
                "`winget install Gyan.FFmpeg` (Windows), "
                "or your distro's package manager (Linux)."
            )

    do_digest = not args.deck_only
    if do_digest:
        ensure_api_key()
        # Budget gate (worker side): covers direct `yt2md <url>` and the
        # subprocesses spawned by digest_video / the scheduler. Checked once at
        # the start of the pipeline, never mid-run, so a started digest always
        # finishes. --deck-only makes no LLM calls, so it's exempt.
        _blocked = check_budget(action="run this digest")
        if _blocked:
            sys.exit(_blocked)

    import time as _time
    timings: dict = {}
    fetch_meta: dict = {
        "used_whisper": False, "whisper_model": None,
    }
    video_title: Optional[str] = None
    video_url: Optional[str] = None
    upload_date: Optional[str] = None
    thumbnail_url: Optional[str] = None
    channel_id: str = ""
    channel_name: str = ""
    channel_url: str = ""
    # Run-start timestamp lets the summary block aggregate cost from the
    # usage log without needing to thread a list through every call site
    # (slide_classifier in particular records from inside its function).
    _run_start_ts = _time.time()

    if is_url(args.video):
        print(f"[0/5] Fetching YouTube video: {args.video}")
        result = fetch_youtube(
            args.video,
            args.downloads_dir,
            whisper_model=args.whisper_model,
            allow_whisper=not args.no_whisper,
            cookies_from_browser=args.cookies_from_browser,
        )
        video_path = result["mp4"]
        srt_path = result["srt"]
        source_lang = result["lang"]
        video_title = result.get("title")
        video_url = result.get("webpage_url")
        upload_date = result.get("upload_date")
        thumbnail_url = result.get("thumbnail_url")
        channel_id = result.get("channel_id") or ""
        channel_name = result.get("channel_name") or ""
        channel_url = result.get("channel_url") or ""
        timings["download"] = round(result["download_secs"], 3)
        timings["whisper"] = round(result["whisper_secs"], 3)
        fetch_meta["used_whisper"] = result["used_whisper"]
        fetch_meta["whisper_model"] = result["whisper_model"]
        print(f"      mp4: {video_path}")
        print(f"      srt: {srt_path} (lang: {source_lang})")
    else:
        video_path = Path(args.video)
        if not video_path.exists():
            sys.exit(f"Video not found: {video_path}")
        if args.srt is None:
            sys.exit("SRT path is required when 'video' is a local file.")
        srt_path = Path(args.srt)
        if not srt_path.exists():
            sys.exit(f"SRT not found: {srt_path}")
        # Local-file path: caller didn't tell us the language; assume English.
        # Override with --source-lang if you're digesting a non-English local SRT.
        source_lang = args.source_lang or "en"

    base = video_path.stem
    digest_path = args.output if args.output is not None else Path(f"{base}_digest.md")
    # Slides off by default. --slides opts in; --deck path overrides;
    # --deck-only forces it on even without --slides.
    if args.deck and args.deck != "__default__":
        deck_path: Optional[Path] = Path(args.deck)
    elif args.deck_only or args.slides:
        deck_path = digest_path.parent / "slides.pptx"
    else:
        deck_path = None

    # Frame extraction is needed only when building slides OR when vision
    # picking is enabled (to supply the frame pool for vision_pick_frames).
    # Compute vision capability now so we can decide before spawning ffmpeg.
    _vision_backend_for_gate = select_backend(for_vision=True)
    _vision_capable = getattr(_vision_backend_for_gate, "vision_supported", False)
    need_frames = deck_path is not None or (not args.no_vision and _vision_capable)

    workdir = Path(tempfile.mkdtemp(prefix="v2d_"))
    scene_dir = workdir / "scene"
    interval_dir = workdir / "interval"

    try:
        duration = get_video_duration(video_path)

        # Single-pass extraction: scene detection + interval sampling share
        # one ffmpeg decode via filter_complex split. Cuts wall time roughly
        # in half on slide-heavy talks vs. running two parallel ffmpegs that
        # both decode the same MP4 and contend for CPU + disk read bandwidth.
        if need_frames:
            print(f"[1/5] Extracting frames "
                  f"(scene threshold={args.scene_threshold}, interval={args.interval}s, "
                  f"single-pass)...")
            _frames_t0 = _time.monotonic()
            scene_frames, interval_frames = extract_scene_and_interval_frames(
                video_path, scene_dir, interval_dir,
                scene_threshold=args.scene_threshold,
                interval=args.interval, duration=duration,
            )
            timings["frames_extract"] = round(_time.monotonic() - _frames_t0, 3)
            frames = merge_frames(scene_frames, interval_frames)
            cap_note = (
                f" (capped from >={SCENE_FRAME_HARD_CAP})"
                if len(scene_frames) == SCENE_FRAME_HARD_CAP else ""
            )
            print(f"      {len(scene_frames)} scene{cap_note} + "
                  f"{len(interval_frames)} interval = {len(frames)} candidate frames "
                  f"({timings['frames_extract']}s)")

            print(f"[2/5] Deduping consecutive near-identical frames "
                  f"(hash distance <= {args.hash_distance})...")
            _dedupe_t0 = _time.monotonic()
            frames = dedupe_frames(frames, args.hash_distance)
            timings["frames_dedupe"] = round(_time.monotonic() - _dedupe_t0, 3)
            print(f"      {len(frames)} unique frames ({timings['frames_dedupe']}s)")
            timings["frames"] = round(_time.monotonic() - _frames_t0, 3)
        else:
            frames = []
            print("[1-2/5] Frame extraction skipped (no --slides, no vision)")

        print(f"[3/5] Parsing SRT: {srt_path.name}")
        segments = parse_srt(srt_path)
        print(f"      {len(segments)} transcript segments")

        print("[4/5] Aligning transcript to frames...")

        if deck_path is not None:
            # Build a slides-specific frame set: pHash cluster globally to
            # collapse "speaker → slide → speaker → same slide" into one
            # representative, then (optionally) ask a vision LLM to filter
            # to actual deck slides only. The digest's vision-pick step
            # still uses the richer `frames` pool — the two consumers want
            # different shapes of the same data.
            deck_frames = global_phash_cluster(frames)
            print(f"      slides: {len(frames)} candidates → "
                  f"{len(deck_frames)} after global pHash dedup")
            settings = load_settings()
            slide_classify_enabled = (
                settings.get("slide_classification", True)
                and not args.no_slide_classification
            )
            if slide_classify_enabled and len(deck_frames) > _GRID_CELLS:
                try:
                    classifier_backend = select_backend(for_vision=True)
                    if getattr(classifier_backend, "vision_supported", False):
                        _classify_t0 = _time.monotonic()
                        deck_frames = classify_slides_via_grids(
                            deck_frames,
                            backend=classifier_backend,
                            model=(args.slide_classifier_model
                                   or settings.get("slide_classifier_model")
                                   or "claude-haiku-4-5-20251001"),
                            workdir=workdir,
                            log_video_id=video_path.stem,
                        )
                        timings["slide_classify"] = round(
                            _time.monotonic() - _classify_t0, 3,
                        )
                        print(f"      slides: vision-classified → "
                              f"{len(deck_frames)} kept "
                              f"({timings['slide_classify']}s)")
                except Exception as e:
                    print(f"      slides: classification skipped "
                          f"({type(e).__name__}: {e}); using pHash dedup only.")
            slides_data = assign_transcript_to_frames(
                deck_frames, segments, duration,
            )
            print(f"[5/5] Building slides ({len(slides_data)} slides) -> {deck_path}")
            build_deck(slides_data, deck_path, video_title or video_path.stem)
        else:
            print("[5/5] Slides skipped (no --slides)")

        usage = None
        if do_digest:
            digest_path.parent.mkdir(parents=True, exist_ok=True)
            images_dir = digest_path.parent / f"{digest_path.stem}_images"
            backend = select_backend()
            log_video_id = video_path.stem
            print(f"[+] Generating digest with {args.digest_model} via {backend.name} backend -> {digest_path}")
            _digest_t0 = _time.monotonic()
            digest, usage = generate_digest(
                segments, video_title or video_path.stem, args.digest_model,
                source_lang=source_lang,
                output_language=args.digest_language,
                backend=backend,
            )
            digest_log_entry = record_llm_usage(
                video_id=log_video_id, kind="digest", model=args.digest_model,
                backend_name=backend.name, usage=usage,
            )
            timings["digest"] = round(_time.monotonic() - _digest_t0, 3)
            print(f"      {len(digest.topics)} topics  |  "
                  f"input: {usage.input_tokens} tokens "
                  f"(cache read: {getattr(usage, 'cache_read_input_tokens', 0)}, "
                  f"cache write: {getattr(usage, 'cache_creation_input_tokens', 0)})  |  "
                  f"output: {usage.output_tokens} tokens  |  "
                  f"cost: ${digest_log_entry['cost_usd']:.4f}")

            vision_picks = None
            if not args.no_vision:
                # Vision routing: even if the text/parse backend is PTY (no
                # vision), select_backend(for_vision=True) will return the API
                # backend if a key is available — hybrid mode without a separate
                # user setting. Falls back to timestamp-based picks otherwise.
                vision_backend = select_backend(for_vision=True)
                if not getattr(vision_backend, "vision_supported", False):
                    print(f"[+] Vision skipped — {vision_backend.name} backend has vision disabled "
                          "(timestamp-based picks will be used).")
                else:
                    print(f"[+] Vision-picking frames with {args.digest_model} via {vision_backend.name}...")
                    _vision_t0 = _time.monotonic()
                    try:
                        vision_picks, v_usage = vision_pick_frames(
                            digest, frames, duration, args.digest_model,
                            segments=segments, backend=vision_backend,
                        )
                        v_log_entry = record_llm_usage(
                            video_id=log_video_id, kind="vision_pick",
                            model=args.digest_model,
                            backend_name=vision_backend.name, usage=v_usage,
                        )
                        timings["vision"] = round(_time.monotonic() - _vision_t0, 3)
                        print(f"      vision-selected {len(vision_picks)}/{len(digest.topics)} topics  |  "
                              f"input: {v_usage.input_tokens} tokens  |  "
                              f"output: {v_usage.output_tokens} tokens  |  "
                              f"cost: ${v_log_entry['cost_usd']:.4f}")
                    except Exception as e:
                        # Anything from the vision call — VisionUnsupported,
                        # auth errors, rate limits, network blips — should
                        # NOT kill the pipeline. Digest is the load-bearing
                        # artifact; timestamp-based frame picks are a fine
                        # fallback and slides.pptx / panel / takeaway still
                        # generate cleanly.
                        print(f"      vision skipped ({type(e).__name__}: {e}); "
                              "falling back to timestamp picks.")

            write_markdown_digest(
                digest, frames, duration, digest_path, images_dir, vision_picks,
                video_title=video_title, video_url=video_url,
            )
            print(f"      Digest written. Images in {images_dir}/")

            # Persist thumbnails + metadata sidecar. These power the sidebar's
            # visual scan — small video poster next to each title, channel
            # avatar as a secondary anchor. Best-effort: thumbnails are
            # nice-to-have, not load-bearing, so we don't fail the pipeline
            # on a missed download.
            digest_dir = digest_path.parent
            thumbnail_local: Optional[Path] = None
            channel_thumbnail_local: Optional[Path] = None
            if thumbnail_url:
                cand = digest_dir / "thumbnail.jpg"
                if download_image(thumbnail_url, cand):
                    thumbnail_local = cand
                    print(f"      thumbnail: {cand.name}")
            if channel_id and channel_url:
                # Shared across all digests from the same channel — avoids
                # downloading the same avatar 50× for a 50-video subscription.
                ch_dir = get_data_dir() / "channel_thumbnails"
                ch_cand = ch_dir / f"{channel_id}.jpg"
                if not ch_cand.exists():
                    ch_url = probe_channel_thumbnail_url(
                        channel_url,
                        cookies_from_browser=args.cookies_from_browser,
                    )
                    if ch_url and download_image(ch_url, ch_cand):
                        channel_thumbnail_local = ch_cand
                        print(f"      channel avatar: {ch_cand.name}")
                else:
                    channel_thumbnail_local = ch_cand
            import datetime as _dt
            added_at_iso = _dt.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
            metadata = {
                "video_id": video_path.stem,
                "title": video_title,
                "url": video_url,
                "upload_date": upload_date,
                "channel_id": channel_id or None,
                "channel_name": channel_name or None,
                "channel_url": channel_url or None,
                "has_thumbnail": thumbnail_local is not None,
                "has_channel_thumbnail": channel_thumbnail_local is not None,
                # New: provenance + curation fields. Topics are filled
                # in by the tagging step below (best-effort).
                "source": {
                    "kind": args.source,
                    "added_at": added_at_iso,
                },
                "topics": [],
                "user_tags": [],
                "user_dismissed": False,
                "user_saved": False,
            }
            try:
                (digest_dir / "metadata.json").write_text(
                    json.dumps(metadata, indent=2) + "\n"
                )
            except OSError:
                pass
            # Mirror source provenance into digest_meta so list_digests
            # filters by `source` are indexed (avoid scanning all JSONs).
            try:
                _record_digest_added(video_path.stem, source_kind=args.source)
            except Exception as _e:
                print(f"      [warn] couldn't record digest_meta: {_e}")

            # Topic tagging (Haiku, ~$0.001). Best-effort: a failure here
            # just leaves `topics` empty in metadata.json; can be
            # backfilled later via `yt2md retrofit-topics`.
            if not args.no_tagging:
                print("[+] Tagging topics with Haiku...")
                try:
                    tag_result = tag_digest_via_llm(video_path.stem)
                    tags_str = ", ".join(tag_result["tags"]) or "(none)"
                    new_str = (
                        f"  (new: {', '.join(tag_result['new_tags'])})"
                        if tag_result["new_tags"] else ""
                    )
                    print(f"      tags: {tags_str}{new_str}")
                except Exception as _e:
                    print(f"      [warn] tagging failed: {_e}")

            # Render the digest to markdown text once for the panel + takeaway
            # prompts (saves a re-read on each step). For the panel/takeaway we
            # want the source-of-truth digest the user will see, so read back
            # the just-written file rather than reconstructing from `digest`.
            digest_md_text = digest_path.read_text()

            panel_md_text: Optional[str] = None
            panel_path = digest_path.parent / "panel.md"
            if not args.no_panel:
                print(f"[+] Generating panel discussion with {args.panel_model}...")
                _panel_t0 = _time.monotonic()
                try:
                    panel_md_text, p_usage = generate_panel_discussion(
                        digest_md_text, segments, model=args.panel_model,
                        source_lang=source_lang,
                        output_language=args.digest_language,
                        backend=backend,
                    )
                    p_log_entry = record_llm_usage(
                        video_id=log_video_id, kind="panel",
                        model=args.panel_model,
                        backend_name=backend.name, usage=p_usage,
                    )
                    panel_path.write_text(panel_md_text)
                    timings["panel"] = round(_time.monotonic() - _panel_t0, 3)
                    print(f"      Panel written -> {panel_path}  |  "
                          f"input: {p_usage.input_tokens} tokens  |  "
                          f"output: {p_usage.output_tokens} tokens  |  "
                          f"cost: ${p_log_entry['cost_usd']:.4f}")
                except Exception as e:
                    # Takeaway can still run without a panel — just with less
                    # explicit pushback to weave in. Don't take the whole
                    # pipeline down for one downstream step's failure.
                    print(f"      Panel generation failed ({type(e).__name__}: {e}); "
                          "continuing without panel.")

            takeaway_path = digest_path.parent / "takeaway.md"
            if not args.no_takeaway:
                print(f"[+] Generating takeaway with {args.takeaway_model}...")
                _take_t0 = _time.monotonic()
                try:
                    takeaway_text, t_usage = generate_takeaway_prose(
                        digest_md_text, panel_md_text, segments,
                        model=args.takeaway_model,
                        publish_date=upload_date,
                        source_lang=source_lang,
                        output_language=args.digest_language,
                        backend=backend,
                    )
                    t_log_entry = record_llm_usage(
                        video_id=log_video_id, kind="takeaway",
                        model=args.takeaway_model,
                        backend_name=backend.name, usage=t_usage,
                    )
                    body = render_takeaway_markdown(
                        takeaway_text, video_url=video_url,
                    )
                    takeaway_path.write_text(body)
                    timings["takeaway"] = round(_time.monotonic() - _take_t0, 3)
                    print(f"      Takeaway written -> {takeaway_path}  |  "
                          f"input: {t_usage.input_tokens} tokens  |  "
                          f"output: {t_usage.output_tokens} tokens  |  "
                          f"cost: ${t_log_entry['cost_usd']:.4f}")
                except Exception as e:
                    print(f"      Takeaway failed ({type(e).__name__}: {e}); "
                          "digest is still complete without it.")

        if args.keep_frames:
            dest = Path.cwd() / f"frames_{video_path.stem}"
            dest.mkdir(parents=True, exist_ok=True)
            for d in (scene_dir, interval_dir):
                if d.exists():
                    shutil.copytree(d, dest, dirs_exist_ok=True)
            print(f"      frames saved to {dest}")

        outputs = []
        if do_digest:
            outputs.append(str(digest_path))
        if deck_path is not None:
            outputs.append(str(deck_path))
        print(f"\nDone. Wrote: {', '.join(outputs)}")

        # Cost audit: aggregate every LLM call recorded during this run.
        # Reads from the canonical usage log so it picks up records made
        # inside helper functions (e.g. slide_classifier) too.
        try:
            log_video_id = video_path.stem
            run_log_entries = [
                e for e in read_llm_usage_log()
                if e.get("ts", 0) >= _run_start_ts
                and e.get("video_id") == log_video_id
            ]
            costs_by_kind: dict = {}
            for e in run_log_entries:
                costs_by_kind[e["kind"]] = (
                    costs_by_kind.get(e["kind"], 0.0) + float(e.get("cost_usd", 0.0))
                )
            total_cost = round(sum(costs_by_kind.values()), 4)
            costs_by_kind = {k: round(v, 4) for k, v in costs_by_kind.items()}
            backend_used = (
                run_log_entries[0]["backend"] if run_log_entries else None
            )
        except Exception:
            costs_by_kind, total_cost, backend_used = {}, 0.0, None

        # Print a human-readable cost summary above the structured line.
        if costs_by_kind:
            print(f"\nCost summary (backend: {backend_used}):")
            for kind in ("digest", "vision_pick", "panel", "takeaway",
                         "slide_classifier", "validation"):
                if kind in costs_by_kind:
                    print(f"  {kind:18s} ${costs_by_kind[kind]:.4f}")
            note = " (subscription — no per-call billing)" if backend_used == "claude-code" else ""
            print(f"  {'TOTAL':18s} ${total_cost:.4f}{note}")

        # Structured one-line summary parsed by the web reaper. Keep this on
        # one line and as the LAST thing printed on success.
        summary = {
            "source_lang": source_lang,
            "used_whisper": fetch_meta["used_whisper"],
            "whisper_model": fetch_meta["whisper_model"],
            "timings": timings,
            "tokens": {
                "input": getattr(usage, "input_tokens", None) if usage else None,
                "output": getattr(usage, "output_tokens", None) if usage else None,
                "cache_read": getattr(usage, "cache_read_input_tokens", None) if usage else None,
                "cache_creation": getattr(usage, "cache_creation_input_tokens", None) if usage else None,
            } if usage else None,
            "cost": {
                "total_usd": total_cost,
                "by_kind": costs_by_kind,
                "backend": backend_used,
            },
            "digest_path": str(digest_path) if do_digest else None,
        }
        print("[summary] " + json.dumps(summary))
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    main()
